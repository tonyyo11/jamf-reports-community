#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""jamf-reports-community.py — Community Jamf Pro reporting tool.

Generates Excel reports from jamf-cli output and optional CSV exports.
Config-driven: no org-specific values are hardcoded.

Requirements: Python 3.9+, xlsxwriter, pandas, pyyaml, matplotlib (optional)
Usage:
    python3 jamf-reports-community.py generate [--config config.yaml] [--csv path/to/export.csv]
                                               [--historical-csv-dir path/to/snapshots/]
    python3 jamf-reports-community.py collect [--config config.yaml] [--csv path/to/export.csv]
                                              [--historical-csv-dir path/to/snapshots/]
    python3 jamf-reports-community.py inventory-csv [--config config.yaml]
                                                    [--out-file path/to/inventory.csv]
    python3 jamf-reports-community.py launchagent-setup [--config config.yaml]
                                                        [--mode csv-assisted]
                                                        [--schedule weekdays]
                                                        [--time-of-day 07:00]
    python3 jamf-reports-community.py scaffold [--csv path/to/export.csv] [--out config.yaml]
    python3 jamf-reports-community.py check [--csv path/to/export.csv]
"""

import argparse
import copy
import hashlib
import json
import math
import os
import plistlib
import re
import shlex
import shutil
import subprocess
import sys
import tempfile
import unicodedata
import urllib.error
import urllib.parse
import urllib.request
from collections import Counter
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timedelta, timezone
from fnmatch import fnmatch
from pathlib import Path
from typing import Any, Optional

import pandas as pd
import xlsxwriter
import yaml

plt: Any = None
mdates: Any = None
HAS_MATPLOTLIB: Optional[bool] = None

pptx_Presentation: Any = None   # pptx.Presentation class, set by _load_pptx()
pptx_Inches: Any = None          # pptx.util.Inches
pptx_Pt: Any = None              # pptx.util.Pt
HAS_PPTX: Optional[bool] = None

LAUNCHAGENT_LABEL_PREFIX = "com.github.tonyyo11.jamf-reports-community"
DEFAULT_LAUNCHD_PATH = "/opt/homebrew/bin:/usr/local/bin:/usr/bin:/bin:/usr/sbin:/sbin"
DEFAULT_AUTOMATION_MODE = "jamf-cli-only"
DEFAULT_AUTOMATION_SCHEDULE = "weekdays"
DEFAULT_AUTOMATION_TIME_OF_DAY = "07:00"
DEFAULT_CSV_FRESHNESS_DAYS = 14
AUTOMATION_MODE_DESCRIPTIONS: dict[str, str] = {
    "snapshot-only": "Refresh jamf-cli snapshots and archive per-family CSV history without writing a workbook.",
    "jamf-cli-only": "Generate a workbook from jamf-cli live data and/or cached JSON snapshots.",
    "jamf-cli-full": "Build a jamf-cli baseline CSV, refresh snapshots, and generate a workbook.",
    "csv-assisted": "Prefer manifest-selected CSV input when available, then fall back to inbox CSV, plus jamf-cli data.",
}
AUTOMATION_SCHEDULE_DESCRIPTIONS: dict[str, str] = {
    "daily": "Every day at the chosen time",
    "weekdays": "Monday through Friday at the chosen time",
    "weekly": "One weekday each week at the chosen time",
    "monthly": "One day of the month at the chosen time",
}
REPORT_FAMILY_NAMES = ("computers", "mobile", "compliance")
WEEKDAY_NAME_TO_VALUE: dict[str, int] = {
    "sun": 0,
    "sunday": 0,
    "mon": 1,
    "monday": 1,
    "tue": 2,
    "tues": 2,
    "tuesday": 2,
    "wed": 3,
    "wednesday": 3,
    "thu": 4,
    "thur": 4,
    "thurs": 4,
    "thursday": 4,
    "fri": 5,
    "friday": 5,
    "sat": 6,
    "saturday": 6,
}
WEEKDAY_VALUE_TO_NAME: dict[int, str] = {
    0: "Sunday",
    1: "Monday",
    2: "Tuesday",
    3: "Wednesday",
    4: "Thursday",
    5: "Friday",
    6: "Saturday",
}


# ---------------------------------------------------------------------------
# Defaults
# ---------------------------------------------------------------------------

DEFAULT_CONFIG: dict[str, Any] = {
    "columns": {
        "computer_name": "",
        "serial_number": "",
        "operating_system": "",
        "last_checkin": "",
        "department": "",
        "manager": "",
        "email": "",
        "filevault": "",
        "sip": "",
        "firewall": "",
        "gatekeeper": "",
        "secure_boot": "",
        "bootstrap_token": "",
        "disk_percent_full": "",
        "architecture": "",
        "model": "",
        "last_enrollment": "",
        "mdm_expiry": "",
    },
    "mobile_columns": {
        "device_name": "",
        "serial_number": "",
        "operating_system": "",
        "last_checkin": "",
        "email": "",
        "model": "",
        "device_family": "",
        "managed": "",
        "supervised": "",
    },
    "security_agents": [],
    "jamf_cli": {
        "enabled": True,
        "data_dir": "jamf-cli-data",
        "profile": "",
        "use_cached_data": True,
        "allow_live_overview": True,
    },
    "protect": {
        "enabled": False,
    },
    "platform": {
        "enabled": False,
        "compliance_benchmarks": [],
    },
    "compliance": {
        "enabled": False,
        "failures_count_column": "",
        "failures_list_column": "",
        "baseline_label": "mSCP Compliance",
    },
    "custom_eas": [],
    "report_families": {
        "computers": {
            "enabled": False,
            "current_dir": "",
            "historical_dir": "",
            "include_globs": [],
            "exclude_globs": [],
            "prefer_name_contains": [],
        },
        "mobile": {
            "enabled": False,
            "current_dir": "",
            "historical_dir": "",
            "include_globs": [],
            "exclude_globs": [],
            "prefer_name_contains": [],
        },
        "compliance": {
            "enabled": False,
            "current_dir": "",
            "historical_dir": "",
            "include_globs": [],
            "exclude_globs": [],
            "prefer_name_contains": [],
        },
    },
    "thresholds": {
        "stale_device_days": 30,
        "checkin_overdue_days": 7,
        "critical_disk_percent": 90,
        "warning_disk_percent": 80,
        "cert_warning_days": 90,
        "profile_error_critical": 50,
        "profile_error_warning": 10,
    },
    "output": {
        "output_dir": "Generated Reports",
        "timestamp_outputs": True,
        "archive_enabled": True,
        "archive_dir": "",
        "keep_latest_runs": 10,
        "export_pptx": False,
    },
    "charts": {
        "enabled": True,
        "save_png": True,
        "embed_in_xlsx": True,
        "historical_csv_dir": "",
        "archive_current_csv": True,
        "os_adoption": {
            "enabled": True,
            "per_major_charts": True,
        },
        "compliance_trend": {
            "enabled": True,
            "bands": [
                {"label": "Pass", "min_failures": 0, "max_failures": 0, "color": "#4472C4"},
                {"label": "Low (1-10)", "min_failures": 1, "max_failures": 10, "color": "#2E9E7D"},
                {"label": "Med-Low (11-30)", "min_failures": 11, "max_failures": 30,
                 "color": "#FFCA30"},
                {"label": "Medium (31-50)", "min_failures": 31, "max_failures": 50,
                 "color": "#F07C21"},
                {"label": "High (>50)", "min_failures": 51, "max_failures": 9999,
                 "color": "#C0392B"},
            ],
        },
        "device_state_trend": {
            "enabled": True,
        },
    },
    "branding": {
        "org_name": "",
        "logo_path": "",
        "accent_color": "#2D5EA2",
        "accent_dark": "#004165",
    },
}

# Fuzzy-match candidates for scaffold auto-detection
COLUMN_HINTS: dict[str, list[str]] = {
    "computer_name": ["computer name", "device name", "hostname", "name"],
    "serial_number": ["serial number", "serial", "serialnumber"],
    "operating_system": ["operating system version", "operating system", "macos version"],
    "last_checkin": [
        "last check-in",
        "last checkin",
        "last contact",
        "last inventory update",
        "checkin",
    ],
    "department": ["department", "dept"],
    "manager": ["manager", "managed by", "direct manager"],
    "email": ["email address", "email", "e-mail"],
    "filevault": ["filevault 2 status", "filevault 2 - status", "filevault status", "filevault"],
    "sip": ["system integrity protection", "sip"],
    "firewall": ["firewall", "firewall enabled", "fw"],
    "gatekeeper": ["gatekeeper"],
    "secure_boot": ["secure boot level", "secure boot"],
    "bootstrap_token": ["bootstrap token escrowed", "bootstrap token is escrowed"],
    "disk_percent_full": ["boot drive percentage full", "percentage full", "disk percent full"],
    "architecture": ["architecture", "arch", "cpu type"],
    "model": ["model", "hardware model", "device model"],
    "last_enrollment": ["last enrollment", "enrollment date", "enrolled"],
    "mdm_expiry": ["mdm profile expiration date", "mdm expiry", "profile expiration date"],
}

COLUMN_EXCLUDES: dict[str, list[str]] = {
    "manager": ["managed", "unmanaged"],
    "secure_boot": ["external boot"],
    "bootstrap_token": ["allowed"],
    "disk_percent_full": ["available mb", "capacity mb", "free mb"],
}


# ---------------------------------------------------------------------------
# Module-level helpers
# ---------------------------------------------------------------------------


def _safe_write(
    worksheet: xlsxwriter.workbook.Worksheet,
    row: int,
    col: int,
    value: Any,
    fmt: Optional[xlsxwriter.workbook.Format] = None,
) -> None:
    """Write a value to a worksheet cell with sanitization.

    Args:
        worksheet: Target xlsxwriter worksheet.
        row: Zero-based row index.
        col: Zero-based column index.
        value: Value to write; sanitized before writing.
        fmt: Optional xlsxwriter format object.
    """
    if value is None:
        worksheet.write_blank(row, col, None, fmt)
        return

    if isinstance(value, float):
        if not (value == value) or value in (float("inf"), float("-inf")):  # noqa: PLR0124
            value = 0

    if isinstance(value, str):
        value = "".join(
            ch for ch in value if unicodedata.category(ch)[0] != "C" or ch in ("\n", "\t")
        )
        value = value[:32000]
        if value.lstrip() and value.lstrip()[0] in ("=", "+", "-", "@"):
            if fmt:
                worksheet.write_string(row, col, value, fmt)
            else:
                worksheet.write_string(row, col, value)
            return

    if fmt:
        worksheet.write(row, col, value, fmt)
    else:
        worksheet.write(row, col, value)


def _parse_manager(raw_value: Any) -> str:
    """Extract a human-readable name from an AD DN or plain string.

    Args:
        raw_value: Raw manager field value from CSV.

    Returns:
        Plain display name, or empty string if blank/invalid.
    """
    if raw_value is None:
        return ""
    if isinstance(raw_value, float) and (raw_value != raw_value):
        return ""
    s = str(raw_value).strip()
    if not s:
        return ""
    # AD DN: CN=SMITH\, JOHN,OU=...  (backslash-comma escapes commas inside the CN value)
    match = re.match(r"CN=((?:[^,\\]|\\,)+)", s, re.IGNORECASE)
    if match:
        cn = match.group(1).replace("\\,", ",").strip()
        return cn.title()
    return s


def _now_ts() -> str:
    """Return current UTC timestamp as a compact string."""
    return datetime.now(timezone.utc).strftime("%Y-%m-%dT%H%M%S%f")


def _file_stamp() -> str:
    """Return a local timestamp string suitable for output filenames."""
    return datetime.now().strftime("%Y-%m-%d_%H%M%S")


def _sha256_file(path: Path) -> str:
    """Return the SHA-256 digest for a file."""
    digest = hashlib.sha256()
    with open(path, "rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _path_has_timestamp(path: Path) -> bool:
    """Return True when the path stem already contains a date/time stamp."""
    return bool(re.search(r"\d{4}-\d{2}-\d{2}(?:[_T]\d{4,6})?", path.stem))


def _timestamped_output_path(path: Path, stamp: str, enabled: bool) -> Path:
    """Append a timestamp to a path unless timestamping is disabled or already present."""
    if not enabled or _path_has_timestamp(path):
        return path
    if path.suffix:
        return path.with_name(f"{path.stem}_{stamp}{path.suffix}")
    return path.with_name(f"{path.name}_{stamp}")


def _filename_component(text: str) -> str:
    """Return a filesystem-friendly label component."""
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", str(text).strip())
    cleaned = re.sub(r"_+", "_", cleaned).strip("._")
    return cleaned or "jamf_report"


def _strip_timestamp_suffix(stem: str) -> str:
    """Remove a trailing date/time stamp from a filename stem when present."""
    return re.sub(r"[_-]\d{4}-\d{2}-\d{2}(?:[_T]\d{4,6})?$", "", stem)


def _run_group_prefix(stem: str, family_base: str) -> str:
    """Return a per-run grouping prefix for a timestamped output stem."""
    escaped = re.escape(family_base)
    pattern = rf"^({escaped}_\d{{4}}-\d{{2}}-\d{{2}}(?:[_T]\d{{4,6}})?)"
    match = re.match(pattern, stem)
    if match:
        return match.group(1)
    return stem


def _archive_old_output_runs(
    directory: Path,
    family_base: str,
    suffixes: set[str],
    keep_latest_runs: int,
    archive_dir: Path,
) -> list[Path]:
    """Archive older output runs for a report family and return moved paths."""
    if keep_latest_runs < 1 or not directory.is_dir():
        return []

    grouped: dict[str, list[Path]] = {}
    for path in directory.iterdir():
        if not path.is_file() or path.suffix.lower() not in suffixes:
            continue
        if not path.stem.startswith(family_base):
            continue
        group_key = _run_group_prefix(path.stem, family_base)
        grouped.setdefault(group_key, []).append(path)

    if len(grouped) <= keep_latest_runs:
        return []

    group_order = sorted(
        grouped.items(),
        key=lambda item: max(candidate.stat().st_mtime for candidate in item[1]),
        reverse=True,
    )
    archive_dir.mkdir(parents=True, exist_ok=True)

    moved: list[Path] = []
    for _, paths in group_order[keep_latest_runs:]:
        for path in paths:
            family_archive_dir = archive_dir / family_base
            family_archive_dir.mkdir(parents=True, exist_ok=True)
            dest = family_archive_dir / path.name
            if dest.exists():
                dest = family_archive_dir / f"{path.stem}_{_file_stamp()}{path.suffix}"
            shutil.move(str(path), str(dest))
            moved.append(dest)

    return moved


def _cli_path(path_value: Optional[str]) -> Optional[Path]:
    """Return an expanded Path for a CLI-supplied path string."""
    if path_value is None:
        return None
    text = str(path_value).strip()
    if not text:
        return None
    return Path(text).expanduser()


def _cli_input_candidates(
    path_value: Optional[str],
    config: Optional["Config"] = None,
) -> list[Path]:
    """Return candidate locations for a CLI-supplied input path."""
    path = _cli_path(path_value)
    if path is None:
        return []

    candidates = [path]
    if config is not None and not path.is_absolute():
        config_relative = config.base_dir / path
        if config_relative not in candidates:
            candidates.append(config_relative)
    return candidates


def _resolve_cli_input_path(
    path_value: Optional[str],
    config: Optional["Config"] = None,
) -> Optional[Path]:
    """Resolve a CLI-supplied input path from existing candidate locations."""
    candidates = _cli_input_candidates(path_value, config)
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return candidates[0] if candidates else None


def _describe_cli_input_candidates(
    path_value: Optional[str],
    config: Optional["Config"] = None,
) -> str:
    """Return a human-readable list of input paths checked for a CLI argument."""
    candidates = _cli_input_candidates(path_value, config)
    if not candidates:
        return ""

    seen: set[str] = set()
    rendered: list[str] = []
    for candidate in candidates:
        text = str(candidate)
        if text not in seen:
            seen.add(text)
            rendered.append(text)
    return " or ".join(rendered)


def _expand_setup_path(path_value: str, base_dir: Path) -> Path:
    """Return an absolute setup path, resolving relative values from base_dir."""
    path = Path(path_value).expanduser()
    if not path.is_absolute():
        path = base_dir / path
    return path.resolve()


def _prompt_text(prompt: str, default: Optional[str] = None) -> str:
    """Prompt for a text value, returning the default on empty input."""
    if not sys.stdin.isatty():
        raise SystemExit(f"Error: {prompt} requires interactive input.")

    suffix = f" [{default}]" if default not in (None, "") else ""
    while True:
        try:
            raw = input(f"{prompt}{suffix}: ").strip()
        except (EOFError, KeyboardInterrupt):
            raise SystemExit("Interactive setup aborted.") from None
        if raw:
            return raw
        if default is not None:
            return default


def _prompt_yes_no(prompt: str, default: bool = True) -> bool:
    """Prompt for a yes/no answer."""
    if not sys.stdin.isatty():
        raise SystemExit(f"Error: {prompt} requires interactive input.")

    default_hint = "Y/n" if default else "y/N"
    while True:
        try:
            raw = input(f"{prompt} ({default_hint}): ").strip().lower()
        except (EOFError, KeyboardInterrupt):
            raise SystemExit("Interactive setup aborted.") from None
        if not raw:
            return default
        if raw in {"y", "yes"}:
            return True
        if raw in {"n", "no"}:
            return False
        print("  Please answer yes or no.")


def _prompt_choice(
    title: str,
    options: list[tuple[str, str]],
    default: str,
) -> str:
    """Prompt the user to choose a key from a small list of options."""
    if not sys.stdin.isatty():
        raise SystemExit(f"Error: {title} requires interactive input.")

    print(f"\n{title}")
    for idx, (_, label) in enumerate(options, 1):
        marker = " (default)" if options[idx - 1][0] == default else ""
        print(f"  {idx}. {label}{marker}")

    while True:
        try:
            raw = input("  Choice: ").strip()
        except (EOFError, KeyboardInterrupt):
            raise SystemExit("Interactive setup aborted.") from None
        if not raw:
            return default
        if raw.isdigit():
            choice = int(raw)
            if 1 <= choice <= len(options):
                return options[choice - 1][0]
        print(f"  Enter a number between 1 and {len(options)}, or press Enter for default.")


def _parse_time_of_day(value: str) -> tuple[int, int]:
    """Parse a local HH:MM time string."""
    text = str(value).strip()
    match = re.fullmatch(r"(\d{1,2}):(\d{2})", text)
    if not match:
        raise ValueError("Time must be in HH:MM format, for example 07:00.")
    hour = int(match.group(1))
    minute = int(match.group(2))
    if hour > 23 or minute > 59:
        raise ValueError("Time must be within 00:00 to 23:59.")
    return hour, minute


def _parse_weekday(value: str) -> tuple[int, str]:
    """Parse a weekday name for launchd weekly schedules."""
    text = str(value).strip().lower()
    if text not in WEEKDAY_NAME_TO_VALUE:
        raise ValueError("Weekday must be Sunday through Saturday.")
    day_value = WEEKDAY_NAME_TO_VALUE[text]
    return day_value, WEEKDAY_VALUE_TO_NAME[day_value]


def _parse_day_of_month(value: Any) -> int:
    """Parse a monthly schedule day, restricting to 1-28 for predictability."""
    try:
        day_value = int(value)
    except (TypeError, ValueError) as exc:
        raise ValueError("Day of month must be an integer between 1 and 28.") from exc
    if not 1 <= day_value <= 28:
        raise ValueError("Day of month must be between 1 and 28.")
    return day_value


def _default_launchagent_label(config: "Config") -> str:
    """Return the default LaunchAgent label for a config/profile combination."""
    profile_name = str(config.jamf_cli.get("profile", "") or "").strip()
    slug_source = profile_name or config.path.stem or "default"
    return f"{LAUNCHAGENT_LABEL_PREFIX}.{_filename_component(slug_source)}"


def _launchagent_schedule_items(
    schedule: str,
    hour: int,
    minute: int,
    weekday: Optional[int] = None,
    day_of_month: Optional[int] = None,
) -> list[dict[str, int]]:
    """Return launchd StartCalendarInterval entries for a schedule preset."""
    if schedule == "daily":
        return [{"Hour": hour, "Minute": minute}]
    if schedule == "weekdays":
        return [{"Weekday": day, "Hour": hour, "Minute": minute} for day in range(1, 6)]
    if schedule == "weekly":
        if weekday is None:
            raise ValueError("Weekly schedules require a weekday.")
        return [{"Weekday": weekday, "Hour": hour, "Minute": minute}]
    if schedule == "monthly":
        if day_of_month is None:
            raise ValueError("Monthly schedules require a day of month.")
        return [{"Day": day_of_month, "Hour": hour, "Minute": minute}]
    raise ValueError(f"Unsupported schedule: {schedule}")


def _launchagent_schedule_summary(
    schedule: str,
    hour: int,
    minute: int,
    weekday_name: Optional[str] = None,
    day_of_month: Optional[int] = None,
) -> str:
    """Return a human-readable summary for a schedule preset."""
    time_str = f"{hour:02d}:{minute:02d}"
    if schedule == "daily":
        return f"Daily at {time_str}"
    if schedule == "weekdays":
        return f"Weekdays at {time_str}"
    if schedule == "weekly":
        return f"Weekly on {weekday_name or 'Monday'} at {time_str}"
    if schedule == "monthly":
        return f"Monthly on day {day_of_month or 1} at {time_str}"
    return f"{schedule} at {time_str}"


def _latest_csv_inbox_file(
    csv_inbox_dir: Optional[str],
    freshness_days: int,
) -> tuple[Optional[Path], str]:
    """Return the newest CSV in an inbox folder, enforcing an optional age limit."""
    if not csv_inbox_dir:
        return None, "No CSV inbox configured."

    inbox_dir = Path(csv_inbox_dir).expanduser()
    if not inbox_dir.is_dir():
        return None, f"CSV inbox not found: {inbox_dir}"

    candidates = sorted(
        (
            path for path in inbox_dir.rglob("*.csv")
            if path.is_file() and not path.is_symlink() and not path.name.startswith(".")
        ),
        key=lambda path: path.stat().st_mtime,
        reverse=True,
    )
    if not candidates:
        return None, f"No CSV files found in {inbox_dir}"

    newest = candidates[0]
    age = datetime.now() - datetime.fromtimestamp(newest.stat().st_mtime)
    if freshness_days > 0 and age > timedelta(days=freshness_days):
        age_days = int(age.total_seconds() // 86400)
        return None, (
            f"Newest CSV is stale ({age_days} day(s) old): {newest.name}"
        )
    return newest, f"Using newest CSV from inbox: {newest.name}"


def _select_automation_csv(
    config: "Config",
    csv_inbox_dir: Optional[str],
    freshness_days: int,
) -> tuple[Optional[Path], Optional[str], str, str]:
    """Select the best automation CSV, preferring report families over inbox files."""
    manifest_csv, family_name, manifest_note = _default_generate_csv(config)
    if manifest_csv is not None:
        note = manifest_note or f"Using report_families.{family_name}: {manifest_csv.name}"
        return manifest_csv, family_name, f"report_families.{family_name}", note

    inbox_csv, inbox_note = _latest_csv_inbox_file(csv_inbox_dir, freshness_days)
    if inbox_csv is not None:
        return inbox_csv, None, "csv_inbox", inbox_note

    note_parts = [part for part in [manifest_note, inbox_note] if part]
    return None, None, "", " | ".join(note_parts)


def _list_of_strings(value: Any) -> list[str]:
    """Return a config value normalized to a list of non-empty strings."""
    if isinstance(value, list):
        return [str(item).strip() for item in value if str(item).strip()]
    if value is None:
        return []
    text = str(value).strip()
    return [text] if text else []


def _report_family_config(config: "Config", family_name: str) -> dict[str, Any]:
    """Return one report-family config block."""
    families = config.report_families
    value = families.get(family_name, {})
    return value if isinstance(value, dict) else {}


def _report_family_label(family_name: str) -> str:
    """Return a user-facing label for a report family key."""
    return family_name.replace("_", " ").title()


def _report_family_current_dir(config: "Config", family_name: str) -> Optional[Path]:
    """Return the current/raw CSV directory for a report family."""
    family = _report_family_config(config, family_name)
    return config.resolve_path_value(family.get("current_dir", ""))


def _report_family_historical_dir(config: "Config", family_name: str) -> Optional[Path]:
    """Return the historical snapshot directory for a report family."""
    family = _report_family_config(config, family_name)
    return config.resolve_path_value(family.get("historical_dir", ""))


def _report_family_matches(path: Path, base_dir: Path, patterns: list[str]) -> bool:
    """Return True when a path matches any configured family glob pattern."""
    if not patterns:
        return False
    relative = str(path.relative_to(base_dir))
    name = path.name
    return any(fnmatch(name, pattern) or fnmatch(relative, pattern) for pattern in patterns)


def _report_family_header_score(
    config: "Config",
    family_name: str,
    header_cols: list[str],
) -> tuple[int, int]:
    """Return a header-based specificity score for a report-family candidate."""
    normalized_headers = {_normalized_text(item) for item in header_cols}
    configured_core = [
        str(config.columns.get(field, "") or "").strip()
        for field in ("computer_name", "serial_number", "operating_system", "last_checkin")
    ]
    configured_all = [
        str(value or "").strip()
        for value in config.columns.values()
        if str(value or "").strip()
    ]

    if family_name == "computers":
        core_matches = sum(
            1 for item in configured_core if item and _normalized_text(item) in normalized_headers
        )
        total_matches = sum(
            1 for item in configured_all if item and _normalized_text(item) in normalized_headers
        )
        return core_matches, total_matches

    if family_name == "mobile":
        configured_mobile = [
            str(value or "").strip()
            for value in config.mobile_columns.values()
            if str(value or "").strip()
        ]
        expected = configured_mobile or [
            "Display Name",
            "Serial Number",
            "OS Version",
            "Model",
            "Last Inventory Update",
            "Email Address",
        ]
        total_matches = sum(
            1 for item in expected if _normalized_text(item) in normalized_headers
        )
        return total_matches, total_matches

    if family_name == "compliance":
        expected = [
            str(config.compliance.get("failures_count_column", "") or "").strip(),
            str(config.compliance.get("failures_list_column", "") or "").strip(),
            str(config.columns.get("computer_name", "") or "").strip(),
            str(config.columns.get("serial_number", "") or "").strip(),
        ]
        total_matches = sum(
            1 for item in expected if item and _normalized_text(item) in normalized_headers
        )
        return total_matches, total_matches

    return 0, 0


def _report_family_candidates(
    config: "Config",
    family_name: str,
) -> tuple[list[Path], str]:
    """Return matching CSV candidates for a configured report family."""
    family = _report_family_config(config, family_name)
    if family.get("enabled") is not True:
        return [], f"report_families.{family_name} is disabled."

    current_dir = _report_family_current_dir(config, family_name)
    if current_dir is None:
        return [], f"report_families.{family_name}.current_dir is not configured."
    if not current_dir.is_dir():
        return [], f"{_report_family_label(family_name)} current_dir not found: {current_dir}"

    include_globs = _list_of_strings(family.get("include_globs", []))
    exclude_globs = _list_of_strings(family.get("exclude_globs", []))
    candidates = [
        path for path in current_dir.rglob("*.csv")
        if path.is_file()
        and not path.is_symlink()
        and not path.name.startswith(".")
        and (not include_globs or _report_family_matches(path, current_dir, include_globs))
        and (not exclude_globs or not _report_family_matches(path, current_dir, exclude_globs))
    ]
    if not candidates:
        return [], f"No CSV files matched report_families.{family_name} in {current_dir}"
    return candidates, f"Found {len(candidates)} candidate CSV(s) for {family_name}"


def _latest_report_family_file(
    config: "Config",
    family_name: str,
) -> tuple[Optional[Path], str]:
    """Return the best current CSV candidate for a report family."""
    candidates, note = _report_family_candidates(config, family_name)
    if not candidates:
        return None, note

    family = _report_family_config(config, family_name)
    preferred_terms = [
        item.casefold() for item in _list_of_strings(family.get("prefer_name_contains", []))
    ]
    scored: list[tuple[tuple[int, int, int, float, str], Path]] = []
    for path in candidates:
        try:
            header_df = pd.read_csv(path, nrows=0, encoding="utf-8-sig")
            header_cols = header_df.columns.tolist()
        except Exception:
            header_cols = []
        preferred_hits = sum(1 for term in preferred_terms if term in path.name.casefold())
        header_primary, header_secondary = _report_family_header_score(
            config, family_name, header_cols,
        )
        key = (
            header_primary,
            preferred_hits,
            header_secondary,
            path.stat().st_mtime,
            path.name,
        )
        scored.append((key, path))

    best = max(scored, key=lambda item: item[0])[1]
    return best, f"Using report_families.{family_name}: {best.name}"


def _family_for_csv_path(config: "Config", csv_path: Path) -> Optional[str]:
    """Return the matching report family for a CSV path, if any."""
    for family_name in REPORT_FAMILY_NAMES:
        family = _report_family_config(config, family_name)
        if family.get("enabled") is not True:
            continue
        current_dir = _report_family_current_dir(config, family_name)
        if current_dir is None:
            continue
        try:
            csv_path.relative_to(current_dir)
        except ValueError:
            continue
        include_globs = _list_of_strings(family.get("include_globs", []))
        exclude_globs = _list_of_strings(family.get("exclude_globs", []))
        if include_globs and not _report_family_matches(csv_path, current_dir, include_globs):
            continue
        if exclude_globs and _report_family_matches(csv_path, current_dir, exclude_globs):
            continue
        return family_name
    return None


def _default_generate_csv(config: "Config") -> tuple[Optional[Path], Optional[str], Optional[str]]:
    """Return the default primary CSV and family for generate/check workflows."""
    notes: list[str] = []
    for family_name in ("computers", "mobile"):
        path, note = _latest_report_family_file(config, family_name)
        if path is not None:
            return path, family_name, note
        if note:
            notes.append(f"{family_name}: {note}")
    return None, None, " | ".join(notes)


def _guess_report_family_from_headers(
    config: "Config",
    header_cols: list[str],
) -> Optional[str]:
    """Infer the most likely report family from CSV headers."""
    best_family: Optional[str] = None
    best_score = (0, 0)
    for family_name in REPORT_FAMILY_NAMES:
        score = _report_family_header_score(config, family_name, header_cols)
        if score > best_score:
            best_family = family_name
            best_score = score
    return best_family if best_score > (0, 0) else None


def _default_historical_dir(
    config: "Config",
    family_name: Optional[str],
    fallback_value: Optional[str] = None,
) -> Optional[Path]:
    """Return the historical CSV directory for the current workflow."""
    explicit = _cli_path(fallback_value)
    if explicit is not None:
        return explicit
    if family_name:
        family_hist = _report_family_historical_dir(config, family_name)
        if family_hist is not None:
            return family_hist
    return config.resolve_path("charts", "historical_csv_dir")


def _write_status_file(path_value: Optional[str], status: dict[str, Any]) -> None:
    """Persist an automation run status JSON file when a path is configured."""
    if not path_value:
        return
    path = Path(path_value).expanduser()
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w") as fh:
        json.dump(status, fh, indent=2, sort_keys=True)
        fh.write("\n")


def _require_existing_config_path(path_value: str) -> Path:
    """Return a resolved config path, raising when the file does not exist."""
    path = Path(path_value).expanduser()
    if not path.is_absolute():
        path = (Path.cwd() / path).resolve()
    else:
        path = path.resolve()
    if not path.exists():
        raise SystemExit(
            f"Error: config file not found: {path}\n"
            "Create it first with scaffold or copy config.example.yaml into place."
        )
    return path


def _seed_config_template_path() -> Path:
    """Return the repository config.example.yaml path."""
    return Path(__file__).resolve().with_name("config.example.yaml")


def _load_workspace_seed_config(seed_config_path: Optional[str]) -> tuple["Config", str]:
    """Load the seed config for workspace bootstrapping."""
    seed_path_obj = _cli_path(seed_config_path)
    if seed_path_obj is not None:
        if not seed_path_obj.exists():
            raise SystemExit(f"Error: seed config not found: {seed_path_obj}")
        return Config(str(seed_path_obj)), str(seed_path_obj.resolve())

    template_path = _seed_config_template_path()
    if template_path.exists():
        return Config(str(template_path)), str(template_path)

    return Config("__workspace_init_defaults__.yaml"), "DEFAULT_CONFIG"


def _find_jamf_cli_binary() -> Optional[str]:
    """Return the best available jamf-cli binary path."""
    candidates = [
        os.environ.get("JAMFCLI_PATH", ""),
        shutil.which("jamf-cli") or "",
        "/opt/homebrew/bin/jamf-cli",
        "/usr/local/bin/jamf-cli",
    ]
    for path in candidates:
        if path and Path(path).is_file() and os.access(path, os.X_OK):
            return path
    return None


def _jamf_cli_enabled(config: "Config") -> bool:
    """Return True when jamf-cli integration is enabled in config."""
    return config.jamf_cli.get("enabled", True) is not False


def _profile_isolation_guidance(config: "Config") -> list[str]:
    """Return advisory notes for multi-profile path isolation."""
    if not _jamf_cli_enabled(config):
        return []
    profile_name = str(config.jamf_cli.get("profile", "") or "").strip()
    if not profile_name:
        return []

    profile_component = _filename_component(profile_name)
    guidance = [
        "Active jamf_cli.profile is set. Keep one config/workspace per tenant,"
        " or make snapshot/output paths profile-specific. Consider workspace-init"
        " for a per-profile workspace skeleton.",
    ]

    data_dir_raw = str(config.get("jamf_cli", "data_dir", default="jamf-cli-data") or "").strip()
    if _normalized_text(data_dir_raw) == "jamf-cli-data":
        guidance.append(
            "jamf_cli.data_dir is the shared default."
            f" Consider jamf-cli-data/{profile_component}."
        )

    hist_dir_raw = str(config.get("charts", "historical_csv_dir", default="") or "").strip()
    if not hist_dir_raw or _normalized_text(hist_dir_raw) == "snapshots":
        guidance.append(
            "charts.historical_csv_dir is blank or generic."
            f" Consider snapshots/{profile_component} when multiple tenants share a parent workspace."
        )

    output_dir_raw = str(
        config.get("output", "output_dir", default="Generated Reports") or ""
    ).strip()
    if _normalized_text(output_dir_raw) == "generated reports":
        guidance.append(
            "output.output_dir is the shared default."
            f" Consider Generated Reports/{profile_component} if multiple tenants share a workspace."
        )

    return guidance


def _days_since(date_str: str) -> Optional[int]:
    """Parse a date string and return days elapsed since then.

    Args:
        date_str: Date string in common formats.

    Returns:
        Integer days since that date, or None if unparseable.
    """
    text = str(date_str).strip()
    if not text:
        return None

    iso_text = text[:-1] + "+00:00" if text.endswith("Z") else text
    try:
        parsed = datetime.fromisoformat(iso_text)
        if parsed.tzinfo is not None:
            return (datetime.now(timezone.utc) - parsed.astimezone(timezone.utc)).days
        return (datetime.now() - parsed).days
    except ValueError:
        pass

    for fmt in ("%Y-%m-%d %H:%M:%S", "%Y-%m-%dT%H:%M:%S", "%m/%d/%Y", "%Y-%m-%d"):
        try:
            dt = datetime.strptime(text, fmt)
            return (datetime.now() - dt).days
        except ValueError:
            continue
    return None


def _display_value(value: Any, empty_label: str = "Unknown / Not Reported") -> str:
    """Return a report-friendly label for a value, substituting blanks."""
    text = str(value or "").strip()
    return text if text else empty_label


def _load_matplotlib() -> bool:
    """Import matplotlib lazily so non-chart commands avoid startup overhead."""
    global HAS_MATPLOTLIB, mdates, plt
    if HAS_MATPLOTLIB is not None:
        return HAS_MATPLOTLIB

    mpl_config_dir = os.environ.get("MPLCONFIGDIR", "").strip()
    if not mpl_config_dir:
        mpl_dir = Path(tempfile.gettempdir()) / "jamf-reports-community-mpl"
        mpl_dir.mkdir(parents=True, exist_ok=True)
        os.environ["MPLCONFIGDIR"] = str(mpl_dir)

    if not os.environ.get("XDG_CACHE_HOME", "").strip():
        cache_dir = Path(tempfile.gettempdir()) / "jamf-reports-community-cache"
        cache_dir.mkdir(parents=True, exist_ok=True)
        os.environ["XDG_CACHE_HOME"] = str(cache_dir)

    try:
        import matplotlib

        matplotlib.use("Agg")  # non-interactive backend — must be set before pyplot import
        import matplotlib.dates as matplotlib_dates
        import matplotlib.pyplot as pyplot

        plt = pyplot
        mdates = matplotlib_dates
        HAS_MATPLOTLIB = True
    except ImportError:
        HAS_MATPLOTLIB = False
    return HAS_MATPLOTLIB


def _load_pptx() -> bool:
    """Import python-pptx lazily; sets HAS_PPTX and module-level pptx_* globals.

    Returns:
        True if python-pptx is available, False otherwise.
    """
    global HAS_PPTX, pptx_Presentation, pptx_Inches, pptx_Pt
    if HAS_PPTX is not None:
        return HAS_PPTX
    try:
        from pptx import Presentation as _Presentation  # type: ignore[import-untyped]
        from pptx.util import Inches as _Inches, Pt as _Pt  # type: ignore[import-untyped]

        pptx_Presentation = _Presentation
        pptx_Inches = _Inches
        pptx_Pt = _Pt
        HAS_PPTX = True
    except ImportError:
        HAS_PPTX = False
    return HAS_PPTX


def _normalized_text(value: Any) -> str:
    """Return a lowercase, single-spaced representation of a header or cell value."""
    return re.sub(r"\s+", " ", str(value).strip().lower())


def _header_tokens(value: Any) -> set[str]:
    """Return lowercase alphanumeric tokens from a header or cell value."""
    return set(re.findall(r"[a-z0-9]+", _normalized_text(value)))


def _column_match_score(header: str, logical: str) -> int:
    """Return a heuristic scaffold score for a header/logical field pairing."""
    normalized = _normalized_text(header)
    score = 0
    for candidate in COLUMN_HINTS.get(logical, []):
        candidate_norm = _normalized_text(candidate)
        if normalized == candidate_norm:
            score = max(score, 100 + len(candidate_norm))
        elif normalized.startswith(candidate_norm):
            score = max(score, 80 + len(candidate_norm))
        elif candidate_norm in normalized:
            score = max(score, 60 + len(candidate_norm))
        elif len(normalized) >= 8 and normalized in candidate_norm:
            score = max(score, 40 + len(normalized))

    for blocked in COLUMN_EXCLUDES.get(logical, []):
        if _normalized_text(blocked) in normalized:
            score -= 80

    return score if score >= 60 else 0


def _best_header_match(
    headers: list[str],
    logical: str,
    used_headers: Optional[set[str]] = None,
) -> tuple[Optional[str], int]:
    """Return the best-scoring header for a logical field."""
    best_header: Optional[str] = None
    best_score = 0
    reserved = used_headers or set()
    for header in headers:
        if header in reserved:
            continue
        score = _column_match_score(header, logical)
        if score > best_score:
            best_header = header
            best_score = score
    return best_header, best_score


def _best_hint_match(
    headers: list[str],
    hints: list[str],
    used_headers: Optional[set[str]] = None,
) -> Optional[str]:
    """Return the best header matching one of the provided hint substrings."""
    normalized_hints = [_normalized_text(hint) for hint in hints if str(hint).strip()]
    reserved = used_headers or set()
    best_header: Optional[str] = None
    best_score = 0
    for header in headers:
        if header in reserved:
            continue
        normalized = _normalized_text(header)
        score = 0
        for hint in normalized_hints:
            if normalized == hint:
                score = max(score, 100)
            elif hint in normalized:
                score = max(score, 90)
        if score > best_score:
            best_header = header
            best_score = score
    return best_header if best_score else None


def _contains_case_insensitive(series: Any, needle: str) -> Any:
    """Return a boolean Series for a case-insensitive substring match."""
    if not needle:
        return series.str.strip() != ""
    return series.str.contains(needle, case=False, regex=False, na=False)


def _split_multi_value_cell(value: Any) -> list[str]:
    """Split a compliance/list cell on common delimiters and return non-empty entries."""
    return [item.strip() for item in re.split(r"[|\n\r]+", str(value)) if item.strip()]


def _to_int(value: Any, default: int = 0) -> int:
    """Coerce a string/number-like value to int, returning default on failure."""
    try:
        return int(float(str(value).strip()))
    except (AttributeError, TypeError, ValueError):
        return default


def _parse_percent(value: Any) -> Optional[float]:
    """Return a 0.0-1.0 ratio from a percent-like value, or None if unavailable."""
    if value in (None, ""):
        return None
    text = str(value).strip()
    has_percent = text.endswith("%")
    try:
        numeric = float(text.rstrip("%").strip())
    except (AttributeError, TypeError, ValueError):
        return None
    if has_percent or numeric > 1:
        numeric /= 100.0
    if numeric < 0:
        return None
    return numeric


def _to_bool(value: Any) -> bool:
    """Coerce common boolean-like values to a Python bool."""
    if isinstance(value, bool):
        return value
    return _normalized_text(value) in {"true", "1", "yes", "y"}


MOBILE_INVENTORY_FIELD_CANDIDATES: dict[str, list[str]] = {
    "id": [
        "id",
        "mobileDeviceId",
        "deviceId",
        "general.id",
        "general.mobileDeviceId",
    ],
    "name": [
        "displayName",
        "name",
        "general.displayName",
        "general.name",
    ],
    "serial": [
        "serialNumber",
        "serial",
        "general.serialNumber",
        "hardware.serialNumber",
    ],
    "model": [
        "model",
        "modelIdentifier",
        "general.model",
        "hardware.model",
        "hardware.modelIdentifier",
    ],
    "os_version": [
        "osVersion",
        "operatingSystemVersion",
        "general.osVersion",
        "hardware.osVersion",
    ],
    "managed": [
        "managed",
        "isManaged",
        "general.managed",
        "general.isManaged",
    ],
    "supervised": [
        "supervised",
        "general.supervised",
        "security.supervised",
    ],
    "shared_ipad": [
        "sharedIpad",
        "general.sharedIpad",
    ],
    "username": [
        "username",
        "userAndLocation.username",
        "location.username",
    ],
    "email": [
        "emailAddress",
        "email",
        "userAndLocation.emailAddress",
        "location.emailAddress",
    ],
    "department": [
        "department",
        "userAndLocation.department",
        "location.department",
    ],
    "building": [
        "building",
        "userAndLocation.building",
        "location.building",
    ],
    "last_inventory": [
        "lastInventoryUpdateDate",
        "general.lastInventoryUpdateDate",
    ],
    "activation_lock": [
        "activationLockEnabled",
        "security.activationLockEnabled",
    ],
    "passcode_compliant": [
        "passcodeCompliant",
        "passcodeCompliantWithProfile",
        "security.passcodeCompliant",
        "security.passcodeCompliantWithProfile",
    ],
    "data_protection": [
        "dataProtection",
        "security.dataProtection",
    ],
    "jailbreak_status": [
        "jailbreakStatus",
        "security.jailbreakStatus",
    ],
    "ownership": [
        "deviceOwnershipType",
        "general.deviceOwnershipType",
    ],
}

MOBILE_PROFILE_FIELD_CANDIDATES: dict[str, list[str]] = {
    "id": ["id", "general.id"],
    "name": ["name", "general.name"],
    "category": ["category.name", "categoryName", "category"],
    "site": ["site.name", "siteName", "site"],
    "description": ["description", "general.description"],
}

PROTECT_COMPUTER_FIELD_CANDIDATES: dict[str, list[str]] = {
    "name": [
        "hostname",
        "hostName",
        "computerName",
        "displayName",
        "name",
    ],
    "serial": [
        "serialNumber",
        "serial",
        "deviceSerialNumber",
    ],
    "plan_name": [
        "plan.name",
        "planName",
        "assignedPlanName",
        "protectPlanName",
        "protect_plan_name",
        "plan",
    ],
    "plan_id": [
        "plan.id",
        "planId",
        "assignedPlanId",
        "protectPlanId",
        "protect_plan_id",
    ],
    "status": [
        "status",
        "state",
        "agentStatus",
        "computerStatus",
        "protectStatus",
    ],
    "last_seen": [
        "lastSeen",
        "lastSeenAt",
        "updated",
        "updatedAt",
        "lastCheckin",
        "lastCheckIn",
    ],
}

PROTECT_PLAN_FIELD_CANDIDATES: dict[str, list[str]] = {
    "name": ["name", "planName", "displayName"],
    "id": ["id", "planId"],
    "enabled": ["enabled", "active", "isEnabled"],
}

PROTECT_ANALYTIC_FIELD_CANDIDATES: dict[str, list[str]] = {
    "name": ["name", "analyticName", "displayName", "title"],
    "severity": ["severity", "level"],
    "enabled": ["enabled", "active", "isEnabled"],
}


def _extract_items(raw: Any) -> list[Any]:
    """Return a list payload from list- or envelope-shaped jamf-cli output."""
    if isinstance(raw, list):
        return raw
    if not isinstance(raw, dict):
        return []

    for key in (
        "results",
        "items",
        "rows",
        "devices",
        "mobileDevices",
        "data",
    ):
        value = raw.get(key)
        if isinstance(value, list):
            return value

    for value in raw.values():
        if isinstance(value, list):
            return value
    return [raw] if raw else []


def _flatten_record(data: Any, prefix: str = "") -> dict[str, Any]:
    """Flatten nested dictionaries to dot-separated keys."""
    if not isinstance(data, dict):
        return {}

    flattened: dict[str, Any] = {}
    for key, value in data.items():
        full_key = f"{prefix}.{key}" if prefix else str(key)
        if isinstance(value, dict):
            flattened.update(_flatten_record(value, full_key))
        else:
            flattened[full_key] = value
    return flattened


def _first_value(mapping: dict[str, Any], candidates: list[str], default: Any = "") -> Any:
    """Return the first non-empty value present in mapping for candidate keys."""
    for key in candidates:
        value = mapping.get(key)
        if value not in (None, "", []):
            return value
    return default


def _optional_bool(value: Any) -> Optional[bool]:
    """Parse a boolean-ish value, returning None when the value is unknown."""
    if isinstance(value, bool):
        return value

    normalized = _normalized_text(value)
    if normalized in {
        "true",
        "1",
        "yes",
        "y",
        "enabled",
        "on",
        "managed",
        "supervised",
        "shared",
        "compliant",
        "present",
        "activated",
    }:
        return True
    if normalized in {
        "false",
        "0",
        "no",
        "n",
        "disabled",
        "off",
        "unmanaged",
        "unsupervised",
        "not compliant",
        "absent",
        "deactivated",
    }:
        return False
    return None


def _yes_no_unknown(value: Any) -> str:
    """Return a display-friendly boolean label."""
    parsed = _optional_bool(value)
    if parsed is None:
        return "Unknown"
    return "Yes" if parsed else "No"


def _days_since_timestamp(value: Any) -> Optional[int]:
    """Return whole days since a timestamp-like value, or None when unavailable."""
    if value in (None, ""):
        return None
    parsed = pd.to_datetime(value, utc=True, errors="coerce")
    if pd.isna(parsed):
        return None
    return max(0, int((datetime.now(timezone.utc) - parsed.to_pydatetime()).days))


def _mobile_device_family(model: Any, name: Any) -> str:
    """Infer a friendly device-family label from mobile model/name fields."""
    text = f"{model} {name}".strip().lower()
    if "ipad" in text:
        return "iPad"
    if "iphone" in text:
        return "iPhone"
    if "ipod" in text:
        return "iPod"
    if "appletv" in text or "apple tv" in text:
        return "Apple TV"
    if "vision" in text:
        return "Vision"
    return "Mobile"


def _normalize_mobile_inventory_row(item: Any) -> dict[str, Any]:
    """Normalize a mobile-device record from jamf-cli into report columns."""
    flat = _flatten_record(item if isinstance(item, dict) else {})
    model = _first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["model"])
    name = _first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["name"])
    last_inventory = _first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["last_inventory"])

    return {
        "Jamf Pro ID": str(_first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["id"])).strip(),
        "Device Name": str(name or "").strip(),
        "Serial Number": str(_first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["serial"])).strip(),
        "Device Family": _mobile_device_family(model, name),
        "Managed": _yes_no_unknown(_first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["managed"])),
        "Supervised": _yes_no_unknown(
            _first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["supervised"])
        ),
        "Shared iPad": _yes_no_unknown(
            _first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["shared_ipad"])
        ),
        "Model": str(model or "").strip(),
        "OS Version": str(
            _first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["os_version"])
        ).strip(),
        "Username": str(
            _first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["username"])
        ).strip(),
        "Email": str(_first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["email"])).strip(),
        "Department": str(
            _first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["department"])
        ).strip(),
        "Building": str(_first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["building"])).strip(),
        "Last Inventory Update": str(last_inventory or "").strip(),
        "Days Since Inventory": _days_since_timestamp(last_inventory),
        "Activation Lock": _yes_no_unknown(
            _first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["activation_lock"])
        ),
        "Passcode Compliant": _yes_no_unknown(
            _first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["passcode_compliant"])
        ),
        "Data Protection": str(
            _first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["data_protection"])
        ).strip(),
        "Jailbreak Status": str(
            _first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["jailbreak_status"])
        ).strip(),
        "Ownership": str(_first_value(flat, MOBILE_INVENTORY_FIELD_CANDIDATES["ownership"])).strip(),
    }


def _normalize_mobile_profile_row(item: Any) -> dict[str, Any]:
    """Normalize a mobile configuration profile row from jamf-cli."""
    flat = _flatten_record(item if isinstance(item, dict) else {})
    profile_id = _first_value(flat, MOBILE_PROFILE_FIELD_CANDIDATES["id"])
    profile_name = _first_value(flat, MOBILE_PROFILE_FIELD_CANDIDATES["name"])
    return {
        "Profile ID": str(profile_id or "").strip(),
        "Profile Name": str(profile_name or "").strip(),
        "Category": str(_first_value(flat, MOBILE_PROFILE_FIELD_CANDIDATES["category"])).strip(),
        "Site": str(_first_value(flat, MOBILE_PROFILE_FIELD_CANDIDATES["site"])).strip(),
        "Description": str(_first_value(flat, MOBILE_PROFILE_FIELD_CANDIDATES["description"])).strip(),
    }


def _summarize_mobile_inventory(rows: list[dict[str, Any]]) -> dict[str, Any]:
    """Aggregate device-family, model, and status counts for mobile inventory rows."""
    summary = {
        "total": len(rows),
        "managed": 0,
        "unmanaged": 0,
        "supervised": 0,
        "shared_ipad": 0,
        "assigned": 0,
        "activation_lock": 0,
        "passcode_compliant": 0,
        "inventory_age_known": 0,
        "families": Counter(),
        "os_versions": Counter(),
        "models": Counter(),
    }
    for row in rows:
        if row.get("Managed") == "Yes":
            summary["managed"] += 1
        elif row.get("Managed") == "No":
            summary["unmanaged"] += 1
        if row.get("Supervised") == "Yes":
            summary["supervised"] += 1
        if row.get("Shared iPad") == "Yes":
            summary["shared_ipad"] += 1
        if row.get("Activation Lock") == "Yes":
            summary["activation_lock"] += 1
        if row.get("Passcode Compliant") == "Yes":
            summary["passcode_compliant"] += 1
        if row.get("Username") or row.get("Email"):
            summary["assigned"] += 1
        if isinstance(row.get("Days Since Inventory"), int):
            summary["inventory_age_known"] += 1

        family = str(row.get("Device Family", "")).strip()
        if family:
            summary["families"][family] += 1
        os_version = str(row.get("OS Version", "")).strip()
        if os_version:
            summary["os_versions"][os_version] += 1
        model = str(row.get("Model", "")).strip()
        if model:
            summary["models"][model] += 1
    return summary


def _security_control_is_compliant(logical: str, value: Any) -> bool:
    """Return True when a CSV value represents a compliant security control state."""
    normalized = _normalized_text(value)
    if not normalized:
        return False

    if logical == "filevault":
        ratio = re.fullmatch(r"(\d+)/(\d+)", normalized)
        if ratio:
            encrypted, total = (int(part) for part in ratio.groups())
            return total > 0 and encrypted == total
        return normalized in {
            "encrypted",
            "all partitions encrypted",
            "boot partitions encrypted",
            "yes",
            "true",
            "enabled",
            "on",
        }

    if logical == "secure_boot":
        return normalized in {"full security", "medium security"}

    if logical == "bootstrap_token":
        return normalized in {"escrowed", "yes", "true", "enabled"}

    if logical == "gatekeeper":
        return normalized in {
            "enabled",
            "yes",
            "true",
            "1",
            "on",
            "active",
            "running",
            "connected",
            "app_store",
            "app store",
            "app_store_and_identified_developers",
            "app store and identified developers",
            "mac_app_store",
            "mac app store",
            "mac_app_store_and_identified_developers",
            "mac app store and identified developers",
        }

    return normalized in {"enabled", "yes", "true", "1", "on", "active", "running", "connected"}


def _compliance_label(comp_cfg: dict[str, Any]) -> str:
    """Return the configured compliance label."""
    return str(comp_cfg.get("baseline_label", "Compliance")).strip() or "Compliance"


def _semantic_warnings(config: "Config", df: pd.DataFrame) -> list[str]:
    """Return warnings for columns that exist but are semantically suspicious."""
    warnings: list[str] = []
    columns = config.columns

    manager_col = columns.get("manager", "")
    if manager_col:
        manager_header = _normalized_text(manager_col)
        if "managed" in manager_header and "manager" not in manager_header:
            warnings.append(
                "columns.manager points to a management-state column."
                " Use a real manager EA or leave it blank."
            )
        elif manager_col in df.columns:
            sample = {_normalized_text(v) for v in df[manager_col].dropna().astype(str).head(10)}
            if sample and sample.issubset({"managed", "unmanaged"}):
                warnings.append(
                    "columns.manager sample values look like Jamf management"
                    " status, not a person's manager."
                )

    disk_col = columns.get("disk_percent_full", "")
    if disk_col:
        disk_header = _normalized_text(disk_col)
        if any(token in disk_header for token in ("available mb", "capacity mb", "free mb")):
            warnings.append(
                "columns.disk_percent_full should point to a percentage-used column,"
                " not MB available/capacity."
            )

    secure_boot_col = columns.get("secure_boot", "")
    if secure_boot_col and "external boot" in _normalized_text(secure_boot_col):
        warnings.append(
            "columns.secure_boot points to External Boot Level. Use Secure Boot Level instead."
        )

    bootstrap_col = columns.get("bootstrap_token", "")
    if bootstrap_col and "allowed" in _normalized_text(bootstrap_col):
        warnings.append(
            "columns.bootstrap_token points to Bootstrap Token Allowed."
            " Use Bootstrap Token Escrowed for compliance tracking."
        )

    filevault_col = columns.get("filevault", "")
    if filevault_col and filevault_col in df.columns:
        sample_vals = {_normalized_text(v) for v in df[filevault_col].dropna().astype(str).head(10)}
        if any(re.fullmatch(r"\d+/\d+", val) for val in sample_vals):
            warnings.append(
                "columns.filevault sample values look like counts"
                " (for example 1/1). FileVault 2 Status is usually a better"
                " compliance column."
            )

    # Heuristic: date columns with values more than 20 years in the future
    # indicate clock drift, test data, or a wrong column mapping.
    far_future = datetime.now() + timedelta(days=365 * 20)
    date_logical_fields = ["mdm_expiry", "last_enrollment", "last_checkin"]
    for field in date_logical_fields:
        col_name = columns.get(field, "")
        if not col_name or col_name not in df.columns:
            continue
        sample = df[col_name].dropna().head(30)
        parsed_sample = pd.to_datetime(sample.astype(str), errors="coerce", utc=True)
        far_ts = pd.Timestamp(far_future, tz="UTC")
        future_count = int((parsed_sample > far_ts).sum())
        if future_count > 0:
            warnings.append(
                f"columns.{field} ({col_name!r}) has {future_count} value(s) more than"
                " 20 years in the future — verify the column mapping or check for"
                " data quality issues (clock drift, test devices)."
            )

    return warnings


def _archive_csv_snapshot(csv_path: str, historical_dir: str) -> tuple[Optional[Path], bool]:
    """Copy the current CSV into the historical snapshot directory for future trend runs."""
    source = Path(csv_path)
    if not source.is_file():
        return None, False

    out_dir = Path(historical_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    source_size = source.stat().st_size
    source_digest = _sha256_file(source)
    # Compare against the 20 most-recent snapshots only. A full-corpus scan across
    # hundreds of large CSVs would read gigabytes on each archive call.
    recent = sorted(
        (p for p in out_dir.rglob("*.csv") if p.is_file() and not p.is_symlink()),
        key=lambda p: p.stat().st_mtime,
        reverse=True,
    )[:20]
    for existing in recent:
        try:
            if existing.stat().st_size == source_size and _sha256_file(existing) == source_digest:
                return existing, False
        except OSError:
            continue
    prefix = re.sub(r"[^a-z0-9]+", "_", source.stem.lower()).strip("_") or "inventory"
    dest = out_dir / f"{prefix}_{_file_stamp()}.csv"
    shutil.copy2(source, dest)
    return dest, True


def _age_label_from_seconds(total_seconds: int) -> str:
    """Return a short human-readable age label from a second count."""
    if total_seconds < 120:
        return "just now"
    if total_seconds < 3600:
        return f"{total_seconds // 60}m ago"
    if total_seconds < 86400:
        return f"{total_seconds // 3600}h ago"
    return f"{total_seconds // 86400}d ago"


def _path_timestamp_label(path: Optional[Path]) -> str:
    """Return an mtime timestamp for a file path, or an empty string."""
    if path is None or not path.exists():
        return ""
    return datetime.fromtimestamp(path.stat().st_mtime).strftime("%Y-%m-%d %H:%M")


def _path_age_label(path: Optional[Path]) -> str:
    """Return a short age label for a file path, or an empty string."""
    if path is None or not path.exists():
        return ""
    delta = datetime.now() - datetime.fromtimestamp(path.stat().st_mtime)
    return _age_label_from_seconds(max(0, int(delta.total_seconds())))


def _source_family_for_report_type(report_type: str) -> str:
    """Return a broad family label for a jamf-cli report type."""
    mobile_tokens = ("mobile", "ios")
    compliance_tokens = ("compliance",)
    report_type_lc = str(report_type or "").strip().casefold()
    if any(token in report_type_lc for token in mobile_tokens):
        return "mobile"
    if any(token in report_type_lc for token in compliance_tokens):
        return "compliance"
    return "computers"


def _source_kind_from_mode(mode: str) -> str:
    """Return a normalized source-kind label for a jamf-cli source mode."""
    if mode == "live":
        return "jamf_cli_live"
    if mode == "cached-fallback":
        return "jamf_cli_cached_fallback"
    return "jamf_cli_cached"


def _site_name(site_value: Any) -> str:
    """Return a friendly site name from a site object or plain string."""
    if isinstance(site_value, dict):
        return str(site_value.get("name", "")).strip()
    return str(site_value or "").strip()


def _inventory_export_row(computer: dict[str, Any]) -> dict[str, Any]:
    """Return a wide CSV row built from jamf-cli `computers list` output."""
    location = computer.get("location", {})
    if not isinstance(location, dict):
        location = {}

    name = str(computer.get("name", "")).strip()
    if not name:
        name = f"Computer {computer.get('id', '')}".strip()

    return {
        "Jamf Pro ID": str(computer.get("id", "")).strip(),
        "Computer Name": name,
        "Serial Number": str(computer.get("serialNumber", "") or "").strip(),
        "Managed": "Managed" if _to_bool(computer.get("isManaged")) else "Unmanaged",
        "Operating System": str(computer.get("operatingSystemVersion", "") or "").strip(),
        "OS Build": str(computer.get("operatingSystemBuild", "") or "").strip(),
        "OS Rapid Security Response": str(
            computer.get("operatingSystemRapidSecurityResponse", "") or ""
        ).strip(),
        "Model": str(computer.get("modelIdentifier", "") or "").strip(),
        "Asset Tag": str(computer.get("assetTag", "") or "").strip(),
        "IP Address": str(computer.get("ipAddress", "") or "").strip(),
        "Last Check-in": str(computer.get("lastContactDate", "") or "").strip(),
        "Last Report": str(computer.get("lastReportDate", "") or "").strip(),
        "Last Enrollment": str(computer.get("lastEnrolledDate", "") or "").strip(),
        "Username": str(location.get("username", "") or "").strip(),
        "Real Name": str(location.get("realName", "") or "").strip(),
        "Email Address": str(location.get("emailAddress", "") or "").strip(),
        "Position": str(location.get("position", "") or "").strip(),
        "Department": str(location.get("department", "") or "").strip(),
        "Building": str(location.get("building", "") or "").strip(),
        "Room": str(location.get("room", "") or "").strip(),
        "Site": _site_name(computer.get("site")),
        "UDID": str(computer.get("udid", "") or "").strip(),
        "Management ID": str(computer.get("managementId", "") or "").strip(),
    }


def _inventory_lookup_key(value: Any) -> str:
    """Return a normalized inventory join key."""
    text = str(value or "").strip()
    if not text:
        return ""
    return re.sub(r"\s+", " ", text).casefold()


def _inventory_row_lookup_keys(row: dict[str, Any]) -> list[str]:
    """Return join keys for an exported inventory row."""
    keys: list[str] = []
    for field in ("Jamf Pro ID", "Serial Number", "UDID", "Management ID", "Computer Name"):
        key = _inventory_lookup_key(row.get(field, ""))
        if key and key not in keys:
            keys.append(key)
    return keys


def _inventory_build_row_index(rows: list[dict[str, Any]]) -> dict[str, list[dict[str, Any]]]:
    """Index inventory rows by the strongest available join keys."""
    row_index: dict[str, list[dict[str, Any]]] = {}
    for row in rows:
        for key in _inventory_row_lookup_keys(row):
            row_index.setdefault(key, []).append(row)
    return row_index


def _inventory_resolve_row(
    row_index: dict[str, list[dict[str, Any]]],
    values: list[Any],
) -> Optional[dict[str, Any]]:
    """Return the first uniquely matched inventory row for the supplied values."""
    for value in values:
        key = _inventory_lookup_key(value)
        if not key:
            continue
        matches = row_index.get(key, [])
        if len(matches) == 1:
            return matches[0]
    return None


def _inventory_detail_lookup_values(computer: dict[str, Any]) -> list[Any]:
    """Return the best identifiers for a per-device jamf-cli detail lookup."""
    return [
        computer.get("id", ""),
        computer.get("serialNumber", ""),
        computer.get("udid", ""),
        computer.get("managementId", ""),
        computer.get("name", ""),
    ]


def _inventory_ea_lookup_values(item: dict[str, Any]) -> list[Any]:
    """Return the best identifiers for matching an EA row to an inventory row."""
    return [
        item.get("device_id", ""),
        item.get("serial", ""),
        item.get("serial_number", ""),
        item.get("computer_id", ""),
        item.get("id", ""),
        item.get("udid", ""),
        item.get("management_id", ""),
        item.get("device", ""),
        item.get("device_name", ""),
        item.get("name", ""),
    ]


def _inventory_detail_identifier(computer: dict[str, Any]) -> str:
    """Return the best jamf-cli identifier for per-device detail lookups."""
    for key in ("id", "serialNumber", "name"):
        value = str(computer.get(key, "") or "").strip()
        if value:
            return value
    return ""


def _inventory_security_detail_fields(detail_rows: Any) -> dict[str, str]:
    """Extract generic security posture fields from `jamf-cli pro device` output."""
    extracted = {column: "" for column in INVENTORY_SECURITY_DETAIL_COLUMNS}
    rows = detail_rows if isinstance(detail_rows, list) else []
    for item in rows:
        if not isinstance(item, dict):
            continue
        if str(item.get("section", "") or "").strip() != "Security":
            continue
        resource = str(item.get("resource", "") or "").strip()
        column = INVENTORY_SECURITY_RESOURCE_MAP.get(resource)
        if not column:
            continue
        extracted[column] = str(item.get("value", "") or "").strip()
    return extracted


def _enrich_inventory_rows_with_security_details(
    bridge: "JamfCLIBridge",
    computers: list[dict[str, Any]],
    row_index: dict[str, list[dict[str, Any]]],
) -> tuple[int, int, int]:
    """Merge per-device security posture values into inventory export rows."""
    targets: list[tuple[dict[str, Any], str]] = []
    unresolved = 0
    for computer in computers:
        if not isinstance(computer, dict):
            continue
        identifier = _inventory_detail_identifier(computer)
        row = _inventory_resolve_row(row_index, _inventory_detail_lookup_values(computer))
        if row is None or not identifier:
            unresolved += 1
            continue
        targets.append((row, identifier))

    if not targets:
        return 0, 0, unresolved

    enriched = 0
    failures = 0
    max_workers = min(8, len(targets))
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_map = {
            executor.submit(bridge.device_detail, identifier): row
            for row, identifier in targets
        }
        for future in as_completed(future_map):
            row = future_map[future]
            try:
                detail_fields = _inventory_security_detail_fields(future.result())
            except RuntimeError:
                failures += 1
                continue
            row.update(detail_fields)
            if any(detail_fields.values()):
                enriched += 1
    return enriched, failures, unresolved


# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------


class Config:
    """Loads and validates the YAML configuration file.

    Args:
        path: Path to config.yaml. Defaults to ./config.yaml.
    """

    _WORKSPACE_INIT_DEFAULTS_NAME = "__workspace_init_defaults__.yaml"

    def __init__(self, path: str = "config.yaml") -> None:
        self._path = Path(path).expanduser()
        self._data: dict[str, Any] = {}
        self._load()

    def _load(self) -> None:
        if self._path.name == self._WORKSPACE_INIT_DEFAULTS_NAME:
            self._data = copy.deepcopy(DEFAULT_CONFIG)
            return
        if self._path.exists():
            with open(self._path, encoding="utf-8") as fh:
                try:
                    loaded = yaml.safe_load(fh) or {}
                except yaml.YAMLError as exc:
                    raise SystemExit(
                        f"Error: config file '{self._path}' has invalid YAML syntax:\n{exc}"
                    ) from None
            if not isinstance(loaded, dict):
                raise SystemExit(
                    f"Error: config file '{self._path}' must contain a top-level mapping,"
                    f" not {type(loaded).__name__}."
                )
            self._data = self._merge(DEFAULT_CONFIG, loaded)
            return
        raise SystemExit(
            f"Error: config file not found: {self.path}\n"
            "Create it first with scaffold or copy config.example.yaml into place."
        )

    def _merge(self, base: dict, override: dict) -> dict:
        result = copy.deepcopy(base)
        for k, v in override.items():
            if isinstance(v, dict) and isinstance(result.get(k), dict):
                result[k] = self._merge(result[k], v)
            else:
                result[k] = v
        return result

    def get(self, *keys: str, default: Any = None) -> Any:
        """Retrieve a nested config value by key path."""
        node = self._data
        for k in keys:
            if not isinstance(node, dict):
                return default
            node = node.get(k, default)
        return node

    def to_dict(self) -> dict[str, Any]:
        """Return a deep copy of the loaded config data."""
        return copy.deepcopy(self._data)

    @property
    def base_dir(self) -> Path:
        """Return the directory relative config-managed paths should use."""
        return self.path.parent

    @property
    def path(self) -> Path:
        """Return the resolved config file path."""
        config_path = self._path
        if not config_path.is_absolute():
            config_path = (Path.cwd() / config_path).resolve()
        else:
            config_path = config_path.resolve()
        return config_path

    def resolve_path_value(self, value: Any) -> Optional[Path]:
        """Resolve a config-managed path value relative to the config location."""
        if value is None:
            return None
        text = str(value).strip()
        if not text:
            return None
        path = Path(text).expanduser()
        if path.is_absolute():
            return path
        return self.base_dir / path

    def resolve_path(self, *keys: str, default: Any = None) -> Optional[Path]:
        """Resolve a nested config path value relative to the config location."""
        return self.resolve_path_value(self.get(*keys, default=default))

    @property
    def columns(self) -> dict[str, str]:
        value = self._data.get("columns", {})
        return value if isinstance(value, dict) else {}

    @property
    def mobile_columns(self) -> dict[str, str]:
        value = self._data.get("mobile_columns", {})
        return value if isinstance(value, dict) else {}

    @property
    def security_agents(self) -> list[dict]:
        value = self._data.get("security_agents", [])
        return value if isinstance(value, list) else []

    @property
    def custom_eas(self) -> list[dict]:
        value = self._data.get("custom_eas", [])
        return value if isinstance(value, list) else []

    @property
    def report_families(self) -> dict:
        value = self._data.get("report_families") or {}
        return value if isinstance(value, dict) else {}

    @property
    def compliance(self) -> dict:
        value = self._data.get("compliance") or {}
        return value if isinstance(value, dict) else {}

    @property
    def jamf_cli(self) -> dict:
        value = self._data.get("jamf_cli") or {}
        return value if isinstance(value, dict) else {}

    @property
    def protect(self) -> dict:
        value = self._data.get("protect") or {}
        return value if isinstance(value, dict) else {}

    @property
    def platform(self) -> dict:
        value = self._data.get("platform") or {}
        return value if isinstance(value, dict) else {}

    @property
    def thresholds(self) -> dict:
        value = self._data.get("thresholds") or {}
        return value if isinstance(value, dict) else {}

    @property
    def output(self) -> dict:
        value = self._data.get("output") or {}
        return value if isinstance(value, dict) else {}


# ---------------------------------------------------------------------------
# ColumnMapper
# ---------------------------------------------------------------------------


class ColumnMapper:
    """Resolves logical field names to actual CSV column names from config.

    Args:
        config: Loaded Config instance.
        section: Config section containing the logical-field mapping.
    """

    def __init__(self, config: Config, section: str = "columns") -> None:
        self._config = config
        self._section = section

    def _mapping(self) -> dict[str, str]:
        if self._section == "mobile_columns":
            return self._config.mobile_columns
        return self._config.columns

    def get(self, logical: str) -> Optional[str]:
        """Return the configured column name for a logical field, or None."""
        col = self._mapping().get(logical, "")
        return col if col else None

    def extract(self, row: Any, logical: str) -> str:
        """Extract a cell value from a DataFrame row by logical field name.

        Args:
            row: A pandas Series (DataFrame row).
            logical: Logical field name as defined in config columns.

        Returns:
            String value of the cell, or empty string if column not found.
        """
        col = self.get(logical)
        if col is None or col not in row.index:
            return ""
        val = row[col]
        if val is None:
            return ""
        if isinstance(val, float) and val != val:
            return ""
        return str(val).strip()


# ---------------------------------------------------------------------------
# JamfCLIBridge
# ---------------------------------------------------------------------------


class JamfCLIBridge:
    """Thin subprocess wrapper around jamf-cli.

    Args:
        save_output: If True, persist JSON output to jamf-cli-data/.
        data_dir: Directory used to save or read cached jamf-cli JSON snapshots.
        profile: Optional jamf-cli profile name passed as -p to every command.
        use_cached_data: If True, fall back to saved snapshots when live commands fail.
    """

    def __init__(
        self,
        save_output: bool = True,
        data_dir: str = "jamf-cli-data",
        profile: str = "",
        use_cached_data: bool = True,
    ) -> None:
        self._binary = self._find_binary()
        self._save = save_output
        self._data_dir = Path(data_dir).expanduser()
        self._profile = str(profile).strip()
        self._use_cached_data = use_cached_data
        self._report_commands_cache: Optional[set[str]] = None
        self._protect_commands_cache: Optional[set[str]] = None
        self._last_source_info: dict[str, dict[str, Any]] = {}

    def _find_binary(self) -> Optional[str]:
        return _find_jamf_cli_binary()

    def is_available(self) -> bool:
        """Return True if jamf-cli binary is found and executable."""
        return self._binary is not None

    def has_cached_data(
        self,
        include_protect: bool = True,
        include_platform: bool = False,
        platform_benchmarks: Optional[list[str]] = None,
    ) -> bool:
        """Return True when the configured data directory contains cached JSON snapshots."""
        report_names = [
            "overview",
            "security",
            "patch-status",
            "patch_status",
            "policy-status",
            "policy_status",
            "inventory-summary",
            "inventory_summary",
            "device-compliance",
            "device_compliance",
            "ea-results",
            "ea_results",
            "software-installs",
            "software_installs",
            "app-status",
            "app_status",
            "update-status",
            "update_status",
            "update-device-failures",
            "update_device_failures",
            "patch-device-failures",
            "patch_device_failures",
            "computer-extension-attributes",
            "computer_extension_attributes",
            "classic-macos-profiles",
            "classic_macos_profiles",
            "mobile-devices-list",
            "mobile_devices_list",
            "mobile-device-inventory-details",
            "mobile_device_inventory_details",
            "classic-ios-profiles",
            "classic_ios_profiles",
            "checkin-status",
            "checkin_status",
            "hardware-models",
            "hardware_models",
            "env-stats",
            "env_stats",
        ]
        if include_protect:
            report_names.extend(
                [
                    "protect-overview",
                    "protect_overview",
                    "protect-computers",
                    "protect_computers",
                    "protect-analytics",
                    "protect_analytics",
                    "protect-plans",
                    "protect_plans",
                ]
            )
        if include_platform:
            report_names.extend(
                [
                    "blueprint-status",
                    "blueprint_status",
                    "ddm-status",
                    "ddm_status",
                    "compliance-rules",
                    "compliance_rules",
                    "compliance-devices",
                    "compliance_devices",
                ]
            )
            for benchmark in platform_benchmarks or []:
                benchmark_name = str(benchmark or "").strip()
                if not benchmark_name:
                    continue
                slug = _benchmark_slug(benchmark_name)
                legacy_slug = _legacy_benchmark_slug(benchmark_name)
                report_names.extend(
                    [
                        f"compliance-rules-{slug}",
                        f"compliance-devices-{slug}",
                    ]
                )
                if legacy_slug != slug:
                    report_names.extend(
                        [
                            f"compliance-rules-{legacy_slug}",
                            f"compliance-devices-{legacy_slug}",
                        ]
                    )
        return self._latest_cached_json(report_names) is not None

    def _report_commands(self) -> set[str]:
        """Return the installed jamf-cli report subcommands, if discoverable."""
        if self._report_commands_cache is not None:
            return self._report_commands_cache
        if not self._binary:
            self._report_commands_cache = set()
            return self._report_commands_cache

        try:
            cmd = [self._binary]
            if self._profile:
                cmd.extend(["-p", self._profile])
            cmd.extend(["pro", "report", "--help"])
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                check=True,
                timeout=30,
                stdin=subprocess.DEVNULL,
            )
        except (subprocess.SubprocessError, PermissionError):
            self._report_commands_cache = set()
            return self._report_commands_cache

        commands: set[str] = set()
        in_available_section = False
        section_headers = {
            "Core Commands:",
            "Security Configuration:",
            "Endpoints:",
            "Organization:",
            "Access & Identity:",
            "Available Commands:",
        }
        for line in result.stdout.splitlines():
            stripped = line.strip()
            if stripped in section_headers:
                in_available_section = True
                continue
            if not in_available_section:
                continue
            if (
                not stripped
                or stripped.startswith("Flags:")
                or stripped.startswith("Global Flags:")
            ):
                if stripped.startswith("Flags:") or stripped.startswith("Global Flags:"):
                    break
                in_available_section = False
                continue
            commands.add(stripped.split()[0])

        self._report_commands_cache = commands
        return commands

    def _require_report_command(
        self,
        command_name: str,
        cache_names: Optional[list[str]] = None,
    ) -> None:
        """Raise when the installed jamf-cli does not support a report subcommand."""
        commands = self._report_commands()
        if commands and command_name not in commands:
            if self._use_cached_data and cache_names and self._latest_cached_json(cache_names):
                return
            raise RuntimeError(
                f"jamf-cli report '{command_name}' is not available in the"
                " installed jamf-cli build."
            )

    def _protect_commands(self) -> set[str]:
        """Return the installed jamf-cli protect subcommands, if discoverable."""
        if self._protect_commands_cache is not None:
            return self._protect_commands_cache
        if not self._binary:
            self._protect_commands_cache = set()
            return self._protect_commands_cache

        try:
            cmd = [self._binary]
            if self._profile:
                cmd.extend(["-p", self._profile])
            cmd.extend(["protect", "--help"])
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                check=True,
                timeout=30,
                stdin=subprocess.DEVNULL,
            )
        except (subprocess.SubprocessError, PermissionError):
            self._protect_commands_cache = set()
            return self._protect_commands_cache

        commands: set[str] = set()
        in_available_section = False
        section_headers = {
            "Core Commands:",
            "Security Configuration:",
            "Endpoints:",
            "Organization:",
            "Access & Identity:",
            "Available Commands:",
        }
        for line in result.stdout.splitlines():
            stripped = line.strip()
            if stripped in section_headers:
                in_available_section = True
                continue
            if not in_available_section:
                continue
            if (
                not stripped
                or stripped.startswith("Flags:")
                or stripped.startswith("Global Flags:")
            ):
                if stripped.startswith("Flags:") or stripped.startswith("Global Flags:"):
                    break
                in_available_section = False
                continue
            commands.add(stripped.split()[0])

        self._protect_commands_cache = commands
        return commands

    def _require_protect_command(
        self,
        command_name: str,
        cache_names: Optional[list[str]] = None,
    ) -> None:
        """Raise when the installed jamf-cli does not support a protect subcommand."""
        commands = self._protect_commands()
        if commands and command_name not in commands:
            if self._use_cached_data and cache_names and self._latest_cached_json(cache_names):
                return
            raise RuntimeError(
                f"jamf-cli protect '{command_name}' is not available in the"
                " installed jamf-cli build."
            )

    @staticmethod
    def _parse_json_output(raw_output: str) -> Any:
        """Return parsed JSON from output that may include banners or prefixes."""
        text = raw_output.lstrip("\ufeff").strip()
        if not text:
            raise json.JSONDecodeError("empty output", raw_output, 0)

        try:
            return json.loads(text)
        except json.JSONDecodeError:
            decoder = json.JSONDecoder()
            for match in re.finditer(r"(?m)^[ \t]*[\[{]", text):
                try:
                    parsed, _ = decoder.raw_decode(text[match.start():].lstrip())
                except json.JSONDecodeError:
                    continue
                if isinstance(parsed, (dict, list)):
                    return parsed
            for i, ch in enumerate(text):
                if ch not in "[{":
                    continue
                try:
                    parsed, _ = decoder.raw_decode(text[i:])
                except json.JSONDecodeError:
                    continue
                if isinstance(parsed, (dict, list)):
                    return parsed
            raise

    def _set_source_info(
        self,
        report_type: str,
        source_mode: str,
        cached_path: Optional[Path] = None,
    ) -> None:
        """Record the most recent source provenance for a report."""
        if not report_type:
            return
        info: dict[str, Any] = {"mode": source_mode}
        if cached_path is not None:
            info["cached_path"] = cached_path
        self._last_source_info[report_type] = info

    def _run(self, args: list[str]) -> Any:
        """Run jamf-cli with args and return parsed JSON.

        Args:
            args: List of command arguments (excluding the binary itself).

        Returns:
            Parsed JSON object.

        Raises:
            RuntimeError: If binary not found or command fails.
        """
        if not self._binary:
            raise RuntimeError(
                "jamf-cli binary not found. Set JAMFCLI_PATH or install via Homebrew."
            )
        cmd = [self._binary, "--output", "json", "--no-input"]
        if self._profile:
            cmd.extend(["-p", self._profile])
        cmd.extend(args)
        try:
            result = subprocess.run(
                cmd, capture_output=True, text=True, check=True,
                timeout=120, stdin=subprocess.DEVNULL,
            )
        except subprocess.TimeoutExpired as e:
            raise RuntimeError(f"jamf-cli timed out after 120s: {e}") from e
        except PermissionError:
            raise RuntimeError("jamf-cli is not executable. Check file permissions.")
        except subprocess.CalledProcessError as exc:
            detail = (exc.stderr or exc.stdout).strip()
            raise RuntimeError(f"jamf-cli failed ({exc.returncode}): {detail}") from exc
        stdout = result.stdout or ""
        stderr = result.stderr or ""
        candidates = [stdout, "\n".join(part for part in [stdout, stderr] if part)]
        if stderr and stdout != stderr:
            candidates.append("\n".join(part for part in [stderr, stdout] if part))

        last_exc: Optional[json.JSONDecodeError] = None
        for candidate in candidates:
            if not candidate.strip():
                continue
            try:
                return self._parse_json_output(candidate)
            except json.JSONDecodeError as exc:
                last_exc = exc

        detail = "\n".join(part for part in [stdout.strip(), stderr.strip()] if part)
        raise RuntimeError(
            f"jamf-cli returned non-JSON output: {detail[:1000]}"
        ) from last_exc

    def _latest_cached_json(self, report_names: list[str]) -> Optional[Path]:
        """Return the newest cached JSON snapshot for any of the supplied report names."""
        candidates: list[Path] = []
        for report_name in report_names:
            report_dir = self._data_dir / report_name
            if report_dir.is_dir():
                candidates.extend(
                    path for path in report_dir.rglob("*.json")
                    if path.is_file() and not path.is_symlink() and ".partial" not in path.name
                )
            elif self._data_dir.is_dir():
                pattern = f"{report_name}_*.json"
                candidates.extend(
                    path for path in self._data_dir.rglob(pattern)
                    if path.is_file() and not path.is_symlink() and ".partial" not in path.name
                )

        if not candidates:
            return None
        return max(candidates, key=lambda path: path.stat().st_mtime)

    def _load_cached_json(
        self,
        report_names: list[str],
        report_type: str = "",
        source_mode: str = "cached",
    ) -> Any:
        """Load and return the newest cached JSON snapshot for the supplied report names."""
        cached_path = self._latest_cached_json(report_names)
        if cached_path is None:
            raise RuntimeError("no cached jamf-cli snapshot is available")
        try:
            with open(cached_path, encoding="utf-8") as fh:
                data = json.load(fh)
        except json.JSONDecodeError as exc:
            raise RuntimeError(
                f"Cached snapshot is malformed and cannot be parsed: {cached_path}\n"
                f"  Delete it and re-run 'collect' to refresh. Detail: {exc}"
            ) from exc
        except OSError as exc:
            raise RuntimeError(f"Could not read cached snapshot {cached_path}: {exc}") from exc
        print(f"  [cache] {cached_path}")
        self._set_source_info(report_type, source_mode, cached_path)
        return data

    def snapshot_age_label(self, report_names: list[str]) -> str:
        """Return a human-readable age string for the newest cached snapshot.

        Args:
            report_names: Cache directory/file name candidates to search.

        Returns:
            A string like "snapshot 2026-04-04 14:32 (18h ago)", or "" if none found.
        """
        path = self._latest_cached_json(report_names)
        if path is None:
            return ""
        mtime = datetime.fromtimestamp(path.stat().st_mtime)
        delta = datetime.now() - mtime
        total_seconds = int(delta.total_seconds())
        if total_seconds < 120:
            age = "just now"
        elif total_seconds < 3600:
            age = f"{total_seconds // 60}m ago"
        elif total_seconds < 86400:
            age = f"{total_seconds // 3600}h ago"
        else:
            age = f"{total_seconds // 86400}d ago"
        return f"snapshot {mtime.strftime('%Y-%m-%d %H:%M')} ({age})"

    def overview_source_label(self) -> str:
        """Return a human-readable provenance label for the latest overview fetch."""
        info = self._last_source_info.get("overview", {})
        mode = str(info.get("mode", "") or "").strip()
        age_label = self.snapshot_age_label(["overview"])
        if mode == "live":
            return "jamf-cli pro overview"
        if mode == "cached-fallback":
            if age_label:
                return f"cached jamf-cli pro overview (live fallback; {age_label})"
            return "cached jamf-cli pro overview (live fallback)"
        if age_label:
            return f"cached jamf-cli pro overview ({age_label})"
        return "cached jamf-cli pro overview"

    def source_info(self, report_type: str) -> dict[str, Any]:
        """Return recorded provenance info for one report type."""
        return dict(self._last_source_info.get(report_type, {}))

    def all_source_info(self) -> dict[str, dict[str, Any]]:
        """Return recorded provenance info for all fetched report types."""
        return {
            report_type: dict(info)
            for report_type, info in self._last_source_info.items()
        }

    def _run_and_save(
        self,
        report_type: str,
        args: list[str],
        cache_names: Optional[list[str]] = None,
    ) -> Any:
        """Run a command and optionally save output to jamf-cli-data/.

        Args:
            report_type: Subdirectory name under jamf-cli-data/.
            args: Arguments passed to _run.
            cache_names: Cache directory/file name candidates for fallback lookup.

        Returns:
            Parsed JSON result.
        """
        cache_candidates = cache_names or [report_type]
        try:
            data = self._run(args)
        except RuntimeError as exc:
            if not self._use_cached_data:
                raise
            try:
                return self._load_cached_json(
                    cache_candidates,
                    report_type=report_type,
                    source_mode="cached-fallback",
                )
            except RuntimeError as cache_exc:
                raise RuntimeError(f"{exc} | cache fallback: {cache_exc}") from exc

        self._set_source_info(report_type, "live")
        if self._save:
            out_dir = self._data_dir / report_type
            try:
                out_dir.mkdir(parents=True, exist_ok=True)
                final_path = out_dir / f"{report_type}_{_now_ts()}.json"
                tmp_path = final_path.with_suffix(".partial")
                with open(tmp_path, "w", encoding="utf-8") as fh:
                    json.dump(data, fh, indent=2)
                tmp_path.rename(final_path)
            except OSError as exc:
                print(f"  [warn] Could not save snapshot for '{report_type}': {exc}")
        return data

    def overview(self, cached_only: bool = False) -> Any:
        """Fetch fleet overview from jamf-cli pro overview."""
        if cached_only:
            return self._load_cached_json(
                ["overview"],
                report_type="overview",
                source_mode="cached",
            )
        return self._run_and_save("overview", ["pro", "overview"], ["overview"])

    def security_report(self) -> Any:
        """Fetch security posture report from jamf-cli pro report security."""
        self._require_report_command("security", ["security"])
        return self._run_and_save("security", ["pro", "report", "security"], ["security"])

    def policy_status(self, scan_failures: bool = False) -> Any:
        """Fetch policy health report from jamf-cli pro report policy-status."""
        self._require_report_command("policy-status", ["policy-status", "policy_status"])
        args = ["pro", "report", "policy-status"]
        if scan_failures:
            args.append("--scan-failures")
        return self._run_and_save("policy-status", args, ["policy-status", "policy_status"])

    def profile_status(self) -> Any:
        """Fetch profile status report from jamf-cli pro report profile-status."""
        self._require_report_command("profile-status", ["profile-status", "profile_status"])
        return self._run_and_save(
            "profile-status",
            ["pro", "report", "profile-status"],
            ["profile-status", "profile_status"],
        )

    def patch_status(self) -> Any:
        """Fetch patch compliance report from jamf-cli pro report patch-status."""
        return self._run_and_save(
            "patch-status",
            ["pro", "report", "patch-status"],
            ["patch-status", "patch_status"],
        )

    def patch_device_failures(self) -> Any:
        """Fetch per-device patch failures via pro report patch-status --scan-failures.

        Requires jamf-cli v1.4.0+. Returns one row per failing device per patch policy,
        enriched with inventory data and the last action taken from the patch log.
        JSON shape:
          [{"policy":"Firefox 130.0","policy_id":"42","device":"MacBook-001",
            "device_id":"123","status_date":"2026-04-01","attempt":3,
            "last_action":"Retrying","serial":"ABC123",
            "os_version":"15.7.3","username":"jdoe"}, ...]
        """
        return self._run_and_save(
            "patch-device-failures",
            ["pro", "report", "patch-status", "--scan-failures"],
            ["patch-device-failures", "patch_device_failures"],
        )

    def app_status(self) -> Any:
        """Fetch managed app deployment report from jamf-cli pro report app-status."""
        self._require_report_command("app-status", ["app-status", "app_status"])
        return self._run_and_save(
            "app-status",
            ["pro", "report", "app-status"],
            ["app-status", "app_status"],
        )

    def update_status(self) -> Any:
        """Fetch managed software update report from jamf-cli pro report update-status."""
        self._require_report_command("update-status", ["update-status", "update_status"])
        try:
            return self._run_and_save(
                "update-status",
                ["pro", "report", "update-status"],
                ["update-status", "update_status"],
            )
        except RuntimeError as exc:
            detail = str(exc)
            if (
                "No managed software update data found." in detail
                or "Managed Software Update Plans toggle is off." in detail
            ):
                return {
                    "message": "No managed software update data found.",
                    "summary": {},
                    "ErrorDevices": [],
                }
            raise

    def update_device_failures(self) -> Any:
        """Fetch per-device update failures via pro report update-status --scan-failures.

        Requires jamf-cli v1.6.0+. Enriches error devices and failed plans with
        inventory details (name, serial, OS, username) and per-plan last events.
        JSON shape (v1.6):
          [{"total": N, "status_summary": [{"status": "...", "count": N}],
            "error_devices": [{"name": "...", "serial": "...", "device_type": "...",
                               "os_version": "...", "username": "...", "status": "...",
                               "product_key": "...", "updated": "..."}],
            "plan_total": N,
            "plan_state_summary": [{"state": "...", "count": N}],
            "failed_plans": [{"name": "...", "serial": "...", "device_type": "...",
                              "os_version": "...", "username": "...", "state": "...",
                              "action": "...", "version": "...", "error": "...",
                              "last_event": "..."}]}]
        """
        self._require_report_command("update-status", ["update-status", "update_status"])
        try:
            return self._run_and_save(
                "update-device-failures",
                ["pro", "report", "update-status", "--scan-failures"],
                ["update-device-failures", "update_device_failures"],
            )
        except RuntimeError as exc:
            detail = str(exc)
            if (
                "No managed software update data found." in detail
                or "Managed Software Update Plans toggle is off." in detail
            ):
                return [{
                    "message": "No managed software update data found.",
                    "total": 0,
                    "status_summary": [],
                    "error_devices": [],
                    "plan_total": 0,
                    "plan_state_summary": [],
                    "failed_plans": [],
                }]
            raise

    def checkin_status(self, threshold_days: int = 7) -> Any:
        """Fetch check-in health summary.

        Tries jamf-cli pro report checkin-status (requires dashboard-era jamf-cli).
        Falls back to device-compliance with threshold_days when that command is absent.

        Native JSON shape:
          {"computers": {"total": N, "overdue": N, "threshold_days": N},
           "mobile": {"total": N, "overdue": N, "threshold_days": N}}

        Fallback shape: same list as device_compliance() — sheet writer detects and
        aggregates accordingly.

        Args:
            threshold_days: Overdue threshold used for the device-compliance fallback.
        """
        try:
            self._require_report_command("checkin-status", ["checkin-status", "checkin_status"])
            return self._run_and_save(
                "checkin-status",
                ["pro", "report", "checkin-status"],
                ["checkin-status", "checkin_status"],
            )
        except RuntimeError:
            pass
        return self._run_and_save(
            "device-compliance",
            ["pro", "report", "device-compliance", "--days-since-checkin", str(threshold_days)],
            ["device-compliance", "device_compliance"],
        )

    def hardware_models(self) -> Any:
        """Fetch hardware model distribution by count.

        Tries jamf-cli pro report hardware-models (requires dashboard-era jamf-cli).
        Falls back to inventory-summary, which the sheet writer aggregates by model.

        Native JSON shape:
          {"computers": [{"model": "MacBook Pro 14-inch", "count": N}, ...],
           "mobile": [{"model": "iPhone 15", "count": N}, ...]}

        Fallback shape: same list as inventory_summary() — sheet writer detects and
        aggregates accordingly.
        """
        try:
            self._require_report_command("hardware-models", ["hardware-models", "hardware_models"])
            return self._run_and_save(
                "hardware-models",
                ["pro", "report", "hardware-models"],
                ["hardware-models", "hardware_models"],
            )
        except RuntimeError:
            pass
        return self._run_and_save(
            "inventory-summary",
            ["pro", "report", "inventory-summary"],
            ["inventory-summary", "inventory_summary"],
        )

    def env_stats(self) -> Any:
        """Fetch environment object counts (policies, profiles, scripts, packages, etc.).

        Requires jamf-cli pro report env-stats (available in dashboard-era builds).
        JSON shape:
          {"policies": N, "config_profiles": N, "scripts": N, "packages": N,
           "smart_groups_computer": N, "smart_groups_mobile": N,
           "extension_attributes": N, "categories": N}
        """
        self._require_report_command("env-stats", ["env-stats", "env_stats"])
        return self._run_and_save(
            "env-stats",
            ["pro", "report", "env-stats"],
            ["env-stats", "env_stats"],
        )

    def blueprint_status(self) -> Any:
        """Fetch blueprint deployment status from jamf-cli pro report blueprint-status."""
        self._require_report_command("blueprint-status", ["blueprint-status", "blueprint_status"])
        return self._run_and_save(
            "blueprint-status",
            ["pro", "report", "blueprint-status"],
            ["blueprint-status", "blueprint_status"],
        )

    def compliance_rules(self, benchmark_title: str) -> Any:
        """Fetch benchmark rule compliance from jamf-cli pro report compliance-rules."""
        benchmark = str(benchmark_title).strip()
        if not benchmark:
            raise RuntimeError("platform compliance rules require a benchmark title or ID")
        self._require_report_command("compliance-rules", ["compliance-rules", "compliance_rules"])
        slug = _benchmark_slug(benchmark)
        legacy_slug = _legacy_benchmark_slug(benchmark)
        cache_key = f"compliance-rules-{slug}"
        cache_names = [cache_key, "compliance-rules", "compliance_rules"]
        legacy_cache_key = f"compliance-rules-{legacy_slug}"
        if legacy_cache_key != cache_key:
            cache_names.insert(1, legacy_cache_key)
        return self._run_and_save(
            cache_key,
            ["pro", "report", "compliance-rules", benchmark],
            cache_names,
        )

    def compliance_devices(self, benchmark_title: str) -> Any:
        """Fetch failing benchmark devices from jamf-cli pro report compliance-devices."""
        benchmark = str(benchmark_title).strip()
        if not benchmark:
            raise RuntimeError("platform compliance devices require a benchmark title or ID")
        self._require_report_command(
            "compliance-devices",
            ["compliance-devices", "compliance_devices"],
        )
        slug = _benchmark_slug(benchmark)
        legacy_slug = _legacy_benchmark_slug(benchmark)
        cache_key = f"compliance-devices-{slug}"
        cache_names = [cache_key, "compliance-devices", "compliance_devices"]
        legacy_cache_key = f"compliance-devices-{legacy_slug}"
        if legacy_cache_key != cache_key:
            cache_names.insert(1, legacy_cache_key)
        return self._run_and_save(
            cache_key,
            ["pro", "report", "compliance-devices", benchmark],
            cache_names,
        )

    def ddm_status(self) -> Any:
        """Fetch declaration health from jamf-cli pro report ddm-status."""
        self._require_report_command("ddm-status", ["ddm-status", "ddm_status"])
        return self._run_and_save(
            "ddm-status",
            ["pro", "report", "ddm-status"],
            ["ddm-status", "ddm_status"],
        )

    def protect_overview(self) -> Any:
        """Fetch a Jamf Protect instance summary from jamf-cli protect overview."""
        self._require_protect_command("overview", ["protect-overview", "protect_overview"])
        return self._run_and_save(
            "protect-overview",
            ["protect", "overview"],
            ["protect-overview", "protect_overview"],
        )

    def protect_computers_list(self) -> Any:
        """Fetch Jamf Protect computer rows from jamf-cli protect computers list."""
        self._require_protect_command("computers", ["protect-computers", "protect_computers"])
        return self._run_and_save(
            "protect-computers",
            ["protect", "computers", "list"],
            ["protect-computers", "protect_computers"],
        )

    def protect_analytics(self) -> Any:
        """Fetch Jamf Protect analytics from jamf-cli protect analytics list."""
        self._require_protect_command("analytics", ["protect-analytics", "protect_analytics"])
        return self._run_and_save(
            "protect-analytics",
            ["protect", "analytics", "list"],
            ["protect-analytics", "protect_analytics"],
        )

    def protect_plans(self) -> Any:
        """Fetch Jamf Protect plans from jamf-cli protect plans list."""
        self._require_protect_command("plans", ["protect-plans", "protect_plans"])
        return self._run_and_save(
            "protect-plans",
            ["protect", "plans", "list"],
            ["protect-plans", "protect_plans"],
        )

    def device_detail(self, identifier: str) -> Any:
        """Fetch the aggregated device detail view for one computer."""
        ident = str(identifier).strip()
        if not ident:
            raise RuntimeError("device detail lookup requires a non-empty identifier")
        return self._run(["pro", "device", ident])

    def computers_list(self) -> Any:
        """Fetch the lightweight computer inventory index from jamf-cli pro computers list."""
        return self._run_and_save(
            "computers-list",
            ["pro", "computers", "list"],
            ["computers-list", "computers_list"],
        )

    def ea_results(self, name_filter: str = "", include_all: bool = True) -> Any:
        """Fetch computer extension attribute values from jamf-cli pro report ea-results."""
        self._require_report_command("ea-results")
        args = ["pro", "report", "ea-results"]
        if name_filter:
            args.extend(["--name", name_filter])
        if include_all:
            args.append("--all")
        return self._run(args)

    def ea_results_report(self, name_filter: str = "", include_all: bool = True) -> Any:
        """Fetch and cache computer extension attribute values from jamf-cli."""
        self._require_report_command("ea-results", ["ea-results", "ea_results"])
        args = ["pro", "report", "ea-results"]
        if name_filter:
            args.extend(["--name", name_filter])
        if include_all:
            args.append("--all")
        return self._run_and_save("ea-results", args, ["ea-results", "ea_results"])

    def inventory_summary(self) -> Any:
        """Fetch hardware model and OS breakdown from jamf-cli pro report inventory-summary."""
        self._require_report_command(
            "inventory-summary",
            ["inventory-summary", "inventory_summary"],
        )
        return self._run_and_save(
            "inventory-summary",
            ["pro", "report", "inventory-summary"],
            ["inventory-summary", "inventory_summary"],
        )

    def device_compliance(self, days_since_checkin: Optional[int] = None) -> Any:
        """Fetch device compliance rows from jamf-cli pro report device-compliance."""
        self._require_report_command(
            "device-compliance",
            ["device-compliance", "device_compliance"],
        )
        args = ["pro", "report", "device-compliance"]
        if days_since_checkin is not None:
            args.extend(["--days-since-checkin", str(days_since_checkin)])
        return self._run_and_save(
            "device-compliance",
            args,
            ["device-compliance", "device_compliance"],
        )

    def computer_extension_attributes(self) -> Any:
        """Fetch computer extension attribute definitions from jamf-cli."""
        return self._run_and_save(
            "computer-extension-attributes",
            ["pro", "computer-extension-attributes", "list"],
            ["computer-extension-attributes", "computer_extension_attributes"],
        )

    def software_installs(
        self,
        title_filter: str = "",
        include_system: bool = False,
    ) -> Any:
        """Fetch installed software distribution from jamf-cli pro report software-installs."""
        self._require_report_command(
            "software-installs",
            ["software-installs", "software_installs"],
        )
        args = ["pro", "report", "software-installs"]
        if include_system:
            args.append("--include-system")
        if title_filter:
            args.extend(["--title", title_filter])
        return self._run_and_save(
            "software-installs",
            args,
            ["software-installs", "software_installs"],
        )

    def groups(self) -> Any:
        """Fetch smart and static group inventory from jamf-cli pro groups.

        Returns a list of group objects. Confirmed available in jamf-cli v1.7.0.
        Exact JSON shape must be validated against a live or test instance before
        implementing the parser in CoreDashboard._write_smart_groups.

        Expected shape (to confirm against fixture):
            [
              {
                "id": "123",
                "name": "All Managed Macs",
                "type": "computer",           # "computer" | "mobile_device"
                "is_smart": true,
                "member_count": 142,
                "criteria": [...]             # may be absent for static groups
              },
              ...
            ]

        TODO: Run `jamf-cli pro groups --output json` against a test instance,
              confirm the JSON shape, commit the result as
              tests/fixtures/jamf-cli-data/groups/groups.json,
              then implement the bridge method body below.

        Returns:
            Parsed JSON list of group objects, or raises RuntimeError if unavailable.
        """
        # TODO: implement using _run_and_save once the command name is confirmed.
        # Pattern to follow:
        #   return self._run_and_save(
        #       "groups",
        #       ["pro", "groups", "--output", "json"],
        #       ["groups"],
        #   )
        raise NotImplementedError(
            "JamfCLIBridge.groups() requires fixture validation. "
            "Run `jamf-cli pro groups --output json` and commit the result to "
            "tests/fixtures/jamf-cli-data/groups/groups.json before implementing."
        )

    def device_lookup(self, device_id: str) -> Any:
        """Fetch a per-device detail view from jamf-cli pro device.

        Args:
            device_id: Jamf Pro computer ID or serial number to look up.

        Returns:
            Parsed JSON response from jamf-cli.

        Raises:
            RuntimeError: If jamf-cli is unavailable or the request fails.
        """
        return self._run(["pro", "device", device_id])

    def computers_inventory_patch(
        self,
        serial: str,
        field_values: dict[str, str],
    ) -> Any:
        """Patch writable fields on a computer identified by serial number.

        Requires jamf-cli v1.6.0+ (computers-inventory patch subcommand).
        Uses --serial to identify the device; --set for each field/value pair.

        Args:
            serial: Device serial number.
            field_values: Mapping of field paths to values, e.g.
                {"general.managed": "true"}.

        Returns:
            Parsed JSON response (updated computer object).

        Raises:
            RuntimeError: If jamf-cli is unavailable, the subcommand is not
                supported by the installed version, or the API call fails.
        """
        args = ["pro", "computers-inventory", "patch", "--serial", serial]
        for field, value in field_values.items():
            args.extend(["--set", f"{field}={value}"])
        return self._run(args)

    # ── List-endpoint methods used by the HTML report ────────────────────────

    def classic_policies_list(self) -> Any:
        """Fetch the classic policy list from jamf-cli pro classic-policies list."""
        return self._run_and_save(
            "classic-policies",
            ["pro", "classic-policies", "list"],
            ["classic-policies", "classic_policies"],
        )

    def macos_profiles_list(self) -> Any:
        """Fetch the macOS config profile list from jamf-cli."""
        return self._run_and_save(
            "classic-macos-profiles",
            ["pro", "classic-macos-config-profiles", "list"],
            ["classic-macos-profiles", "classic_macos_profiles"],
        )

    def ios_profiles_list(self) -> Any:
        """Fetch the iOS/mobile config profile list from jamf-cli."""
        return self._run_and_save(
            "classic-ios-profiles",
            ["pro", "classic-mobile-config-profiles", "list"],
            ["classic-ios-profiles", "classic_ios_profiles"],
        )

    def mobile_devices_list(self) -> Any:
        """Fetch the mobile-device list from jamf-cli."""
        return self._run_and_save(
            "mobile-devices-list",
            ["pro", "mobile-devices", "list"],
            ["mobile-devices-list", "mobile_devices_list"],
        )

    def mobile_device_inventory_details(self) -> Any:
        """Fetch paginated mobile-device inventory details from jamf-cli."""
        return self._run_and_save(
            "mobile-device-inventory-details",
            ["pro", "mobile-device-inventory-details", "list"],
            ["mobile-device-inventory-details", "mobile_device_inventory_details"],
        )

    def smart_groups_list(self) -> Any:
        """Fetch the smart computer group list from jamf-cli."""
        return self._run_and_save(
            "smart-computer-groups",
            ["pro", "smart-computer-groups", "list"],
            ["smart-computer-groups", "smart_computer_groups"],
        )

    def scripts_list(self) -> Any:
        """Fetch the script list from jamf-cli pro scripts list."""
        return self._run_and_save(
            "scripts",
            ["pro", "scripts", "list"],
            ["scripts"],
        )

    def packages_list(self) -> Any:
        """Fetch the package list from jamf-cli pro packages list."""
        return self._run_and_save(
            "packages",
            ["pro", "packages", "list"],
            ["packages"],
        )

    def categories_list(self) -> Any:
        """Fetch the category list from jamf-cli pro categories list."""
        return self._run_and_save(
            "categories",
            ["pro", "categories", "list"],
            ["categories"],
        )

    def device_enrollments_list(self) -> Any:
        """Fetch the ADE/device enrollment list from jamf-cli."""
        cache_names = [
            "device-enrollment-instances",
            "device_enrollment_instances",
            "device-enrollments",
            "device_enrollments",
        ]
        try:
            return self._run_and_save(
                "device-enrollment-instances",
                ["pro", "device-enrollment-instances", "list"],
                cache_names,
            )
        except RuntimeError as exc:
            detail = str(exc).lower()
            if "unknown command" not in detail and "device-enrollment-instances" not in detail:
                raise
        return self._run_and_save(
            "device-enrollments",
            ["pro", "device-enrollments", "list"],
            cache_names,
        )

    def sites_list(self) -> Any:
        """Fetch the site list from jamf-cli pro sites list."""
        return self._run_and_save(
            "sites",
            ["pro", "sites", "list"],
            ["sites"],
        )

    def buildings_list(self) -> Any:
        """Fetch the building list from jamf-cli pro buildings list."""
        return self._run_and_save(
            "buildings",
            ["pro", "buildings", "list"],
            ["buildings"],
        )

    def departments_list(self) -> Any:
        """Fetch the department list from jamf-cli pro departments list."""
        return self._run_and_save(
            "departments",
            ["pro", "departments", "list"],
            ["departments"],
        )


# ---------------------------------------------------------------------------
# jamf-cli helpers
# ---------------------------------------------------------------------------


def _platform_benchmark_titles(config: Config) -> list[str]:
    """Return normalized platform benchmark titles from config."""
    raw = config.get("platform", "compliance_benchmarks", default=[]) or []
    if isinstance(raw, str):
        raw = [raw]
    return [str(title).strip() for title in raw if str(title).strip()]


def _build_jamf_cli_bridge(
    config: Config,
    *,
    save_output: bool,
    use_cached_data: Optional[bool] = None,
) -> JamfCLIBridge:
    """Construct a JamfCLIBridge from config with consistent defaults."""
    jamf_cli_cfg = config.jamf_cli
    jamf_cli_dir = config.resolve_path("jamf_cli", "data_dir", default="jamf-cli-data")
    if use_cached_data is None:
        use_cached = jamf_cli_cfg.get("use_cached_data", True) is not False
    else:
        use_cached = use_cached_data
    return JamfCLIBridge(
        save_output=save_output,
        data_dir=str(jamf_cli_dir or Path("jamf-cli-data")),
        profile=str(jamf_cli_cfg.get("profile", "") or "").strip(),
        use_cached_data=use_cached,
    )


# ---------------------------------------------------------------------------
# Workbook format helpers
# ---------------------------------------------------------------------------


def _build_formats(wb: xlsxwriter.Workbook, accent_color: str = "#2D5EA2") -> dict[str, Any]:
    """Create and return a dict of named xlsxwriter formats.

    Args:
        wb: Active xlsxwriter Workbook.

    Returns:
        Dict mapping name -> Format object.
    """
    return {
        "title": wb.add_format({"bold": True, "font_size": 14}),
        "subtitle": wb.add_format({"italic": True, "font_color": "#595959", "font_size": 10}),
        "header": wb.add_format(
            {"bold": True, "bg_color": accent_color, "font_color": "white", "border": 1}
        ),
        "cell": wb.add_format({"border": 1}),
        "green": wb.add_format({"bg_color": "#C6EFCE", "border": 1}),
        "yellow": wb.add_format({"bg_color": "#FFEB9C", "border": 1}),
        "red": wb.add_format({"bg_color": "#FFC7CE", "border": 1}),
        "pct": wb.add_format({"num_format": "0.0%", "border": 1}),
        "pct_green": wb.add_format({"num_format": "0.0%", "bg_color": "#C6EFCE", "border": 1}),
        "pct_yellow": wb.add_format({"num_format": "0.0%", "bg_color": "#FFEB9C", "border": 1}),
        "pct_red": wb.add_format({"num_format": "0.0%", "bg_color": "#FFC7CE", "border": 1}),
        "date": wb.add_format({"num_format": "yyyy-mm-dd", "border": 1}),
        "int": wb.add_format({"num_format": "0", "border": 1}),
    }


def _write_sheet_header(
    ws: xlsxwriter.workbook.Worksheet,
    title: str,
    subtitle: str,
    fmts: dict,
    ncols: int = 8,
) -> int:
    """Write title/subtitle rows and return the next data row index.

    Args:
        ws: Target worksheet.
        title: Bold title text for row 0.
        subtitle: Italic subtitle text for row 1.
        fmts: Format dict from _build_formats.
        ncols: Number of columns to merge for title.

    Returns:
        Row index where data should begin (typically 3).
    """
    ws.merge_range(0, 0, 0, ncols - 1, title, fmts["title"])
    ws.merge_range(1, 0, 1, ncols - 1, subtitle, fmts["subtitle"])
    return 3


def _org_title(org_name: str, base: str) -> str:
    """Prefix a sheet title with the org name when configured.

    Args:
        org_name: Organisation name from branding config, or empty string.
        base: The base sheet title (e.g. "Fleet Overview").

    Returns:
        "{org_name} \u2014 {base}" when org_name is set, otherwise base unchanged.
    """
    return f"{org_name} \u2014 {base}" if org_name else base


def _pct_format(fmts: dict, pct: float) -> Any:
    """Return a color-coded percentage format based on value.

    Args:
        fmts: Format dict from _build_formats.
        pct: Compliance percentage (0.0 to 1.0).

    Returns:
        xlsxwriter Format object.
    """
    if pct >= 0.95:
        return fmts["pct_green"]
    if pct >= 0.80:
        return fmts["pct_yellow"]
    return fmts["pct_red"]


def _legacy_benchmark_slug(title: str) -> str:
    """Return the legacy benchmark slug used by earlier cache layouts."""
    import re as _re

    return _re.sub(r"[^a-z0-9]+", "-", title.lower()).strip("-")[:48]


def _benchmark_slug(title: str) -> str:
    """Return a collision-resistant filesystem-safe slug from a benchmark title."""
    normalized = _legacy_benchmark_slug(title) or "benchmark"
    digest = hashlib.sha1(str(title).encode("utf-8")).hexdigest()[:8]
    stem = normalized[:39].rstrip("-") or "benchmark"
    return f"{stem}-{digest}"


def _excel_sheet_name(
    base: str,
    suffix: str,
    max_len: int = 31,
    existing_names: Optional[set[str]] = None,
) -> str:
    """Return an Excel-safe, collision-resistant sheet name.

    Args:
        base: The primary label (for example a benchmark title or EA name).
        suffix: Short suffix appended after truncation room is reserved.
        max_len: Excel sheet name character limit (default 31).
        existing_names: Optional set of existing sheet names for de-duplication.

    Returns:
        A sanitized string guaranteed to be <= max_len characters and unique
        against existing_names when supplied.
    """
    cleaned_base = re.sub(r"[\[\]:*?/\\\\]", " ", str(base or ""))
    cleaned_base = re.sub(r"\s+", " ", cleaned_base).strip().strip("'")
    cleaned_suffix = re.sub(r"[\[\]:*?/\\\\]", " ", str(suffix or ""))
    cleaned_suffix = re.sub(r"\s+", " ", cleaned_suffix).strip().strip("'")
    if not cleaned_base:
        cleaned_base = "Sheet"

    room = max(1, max_len - len(cleaned_suffix))
    candidate = (cleaned_base[:room] + cleaned_suffix).strip().strip("'")
    if not candidate:
        candidate = "Sheet"
    candidate = candidate[:max_len]

    if not existing_names:
        return candidate

    existing_lower = {name.lower() for name in existing_names}
    if candidate.lower() not in existing_lower:
        return candidate

    index = 2
    while True:
        marker = f" ({index})"
        room = max(1, max_len - len(cleaned_suffix) - len(marker))
        candidate = (cleaned_base[:room] + cleaned_suffix + marker).strip().strip("'")
        candidate = candidate[:max_len]
        if candidate.lower() not in existing_lower:
            return candidate
        index += 1


def _write_report_sources_sheet(
    wb: xlsxwriter.Workbook,
    fmts: dict[str, Any],
    generated_at: str,
    config: Config,
    csv_path: Optional[str],
    csv_family: str,
    csv_origin: str,
    hist_dir: Optional[str],
    jamf_cli_dir: Optional[str],
    jamf_cli_profile: str,
    live_overview_allowed: bool,
    overview_source_label: str,
    jamf_cli_sheets: list[str],
    csv_sheets: list[str],
    chart_source: str,
    org_name: str = "",
    source_details: Optional[list[dict[str, str]]] = None,
) -> None:
    """Write a workbook sheet describing the data sources used for the report."""
    ws = wb.add_worksheet("Report Sources")
    row = _write_sheet_header(
        ws,
        _org_title(org_name, "Report Sources"),
        f"Generated: {generated_at}",
        fmts,
        ncols=8,
    )
    ws.set_column(0, 0, 24)
    ws.set_column(1, 1, 80)
    ws.set_column(2, 7, 18)

    logo_path = config.resolve_path("branding", "logo_path")
    if logo_path and logo_path.exists():
        try:
            ws.insert_image(
                0, 5, str(logo_path),
                {"x_scale": 0.25, "y_scale": 0.25, "object_position": 1},
            )
        except Exception:
            pass  # logo insertion is best-effort; never block report generation

    if jamf_cli_sheets and csv_sheets:
        report_mode = "Mixed (jamf-cli + CSV)"
    elif jamf_cli_sheets:
        report_mode = "jamf-cli only"
    else:
        report_mode = "CSV only"

    summary_rows = [
        ("Report Mode", report_mode),
        ("Config Base Dir", str(config.base_dir)),
        ("CSV Input", csv_path or ""),
        ("CSV Family", csv_family),
        ("CSV Origin", csv_origin),
        ("Historical CSV Dir", hist_dir or ""),
        ("jamf-cli Data Dir", jamf_cli_dir or ""),
        ("jamf-cli Profile", jamf_cli_profile),
        ("Live Overview Setting", "Enabled" if live_overview_allowed else "Cached only"),
        ("Overview Source", overview_source_label),
    ]
    for label, value in summary_rows:
        _safe_write(ws, row, 0, label, fmts["cell"])
        _safe_write(ws, row, 1, value, fmts["cell"])
        row += 1

    if source_details:
        row += 1
        detail_headers = ["Scope", "Kind", "Origin", "Family", "Path", "Timestamp", "Age", "Notes"]
        for col_i, header in enumerate(detail_headers):
            _safe_write(ws, row, col_i, header, fmts["header"])
        row += 1
        for detail in source_details:
            values = [
                detail.get("scope", ""),
                detail.get("kind", ""),
                detail.get("origin", ""),
                detail.get("family", ""),
                detail.get("path", ""),
                detail.get("timestamp", ""),
                detail.get("age", ""),
                detail.get("notes", ""),
            ]
            for col_i, value in enumerate(values):
                _safe_write(ws, row, col_i, value, fmts["cell"])
            row += 1

    row += 1
    headers = ["Sheet", "Source", "Included"]
    for col_i, header in enumerate(headers):
        _safe_write(ws, row, col_i, header, fmts["header"])
    row += 1

    sheet_rows: list[tuple[str, str, str]] = []
    sheet_rows.extend((sheet, "jamf-cli", "Yes") for sheet in jamf_cli_sheets)
    sheet_rows.extend((sheet, "CSV", "Yes") for sheet in csv_sheets)
    if chart_source:
        sheet_rows.append(("Charts", chart_source, "Yes"))

    for sheet, source, included in sheet_rows:
        _safe_write(ws, row, 0, sheet, fmts["cell"])
        _safe_write(ws, row, 1, source, fmts["cell"])
        _safe_write(ws, row, 2, included, fmts["cell"])
        row += 1


# ---------------------------------------------------------------------------
# CoreDashboard helpers
# ---------------------------------------------------------------------------


def _extract_envelope(raw: Any) -> dict:
    """Unwrap a single-element list response and return a dict, or empty dict.

    Several jamf-cli report commands wrap their result in a one-element list.
    This helper handles both the list and bare-dict shapes consistently.

    Args:
        raw: The parsed JSON value returned by a jamf-cli command.

    Returns:
        The inner dict, or an empty dict if the response is absent/invalid.
    """
    node = raw[0] if isinstance(raw, list) and raw else raw
    return node if isinstance(node, dict) and node else {}


def _protect_overview_has_data(raw: Any) -> bool:
    """Return True when a Protect overview response contains a non-placeholder value."""
    for item in _extract_items(raw):
        if not isinstance(item, dict):
            continue
        value = str(item.get("value", "") or "").strip()
        if value and _normalized_text(value) not in {"n/a", "na"}:
            return True

    if isinstance(raw, dict):
        for value in _flatten_record(raw).values():
            text = str(value or "").strip()
            if text and _normalized_text(text) not in {"n/a", "na"}:
                return True
    return False


# ---------------------------------------------------------------------------
# CoreDashboard
# ---------------------------------------------------------------------------


class CoreDashboard:
    """Generates Excel sheets from jamf-cli data only (no CSV required).

    Args:
        config: Loaded Config instance.
        bridge: JamfCLIBridge instance.
        workbook: Active xlsxwriter Workbook.
        fmts: Format dict from _build_formats.
    """

    def __init__(
        self,
        config: Config,
        bridge: JamfCLIBridge,
        workbook: xlsxwriter.Workbook,
        fmts: dict,
    ) -> None:
        self._config = config
        self._bridge = bridge
        self._wb = workbook
        self._fmts = fmts
        self._overview_rows_cache: Optional[list[dict[str, Any]]] = None
        self._overview_source_label: str = ""
        self._mobile_inventory_cache: Optional[tuple[list[dict[str, Any]], str]] = None
        self._mobile_profile_cache: Optional[list[dict[str, Any]]] = None

    @property
    def _org_name(self) -> str:
        """Return the configured org name, or empty string."""
        return (self._config.get("branding", "org_name") or "").strip()

    def _t(self, base: str) -> str:
        """Return sheet title prefixed with org name when configured."""
        return _org_title(self._org_name, base)

    def _severity_fmt(self, value: int, warn: int, crit: int) -> Any:
        """Return red/yellow/normal cell format based on value vs thresholds.

        Args:
            value: The integer value to test.
            warn: Threshold at or above which the yellow format is returned.
            crit: Threshold at or above which the red format is returned.

        Returns:
            An xlsxwriter format object from self._fmts.
        """
        if value >= crit:
            return self._fmts["red"]
        if value >= warn:
            return self._fmts["yellow"]
        return self._fmts["cell"]

    @staticmethod
    def _overview_value(rows: list[dict[str, Any]], resource: str) -> Any:
        """Return an overview value by resource label."""
        for row in rows:
            if row.get("resource") == resource:
                return row.get("value", "")
        return ""

    def _overview_rows_cached(self) -> list[dict[str, Any]]:
        """Return normalized overview rows, caching the result for this workbook."""
        if self._overview_rows_cache is not None:
            return self._overview_rows_cache

        live_overview_allowed = self._config.jamf_cli.get("allow_live_overview", True) is True
        data = self._bridge.overview(cached_only=not live_overview_allowed)
        self._overview_rows_cache = self._overview_rows(data)
        self._overview_source_label = self._bridge.overview_source_label()
        return self._overview_rows_cache

    def _mobile_inventory_rows(self) -> tuple[list[dict[str, Any]], str]:
        """Return normalized mobile inventory rows and the source command used."""
        if self._mobile_inventory_cache is not None:
            return self._mobile_inventory_cache

        attempts = [
            ("jamf-cli pro mobile-device-inventory-details list",
             self._bridge.mobile_device_inventory_details),
            ("jamf-cli pro mobile-devices list", self._bridge.mobile_devices_list),
        ]
        errors: list[str] = []
        for source, fetcher in attempts:
            try:
                items = _extract_items(fetcher())
            except RuntimeError as exc:
                errors.append(str(exc))
                continue

            rows = [
                _normalize_mobile_inventory_row(item)
                for item in items
                if isinstance(item, dict)
            ]
            rows = [
                row
                for row in rows
                if any(
                    str(row.get(key, "")).strip()
                    for key in ("Device Name", "Serial Number", "Model", "OS Version")
                )
            ]
            if rows:
                self._mobile_inventory_cache = (rows, source)
                return self._mobile_inventory_cache

        self._mobile_inventory_cache = ([], "")
        if errors:
            raise RuntimeError(errors[0])
        raise RuntimeError("jamf-cli returned no mobile device inventory rows")

    def _mobile_profile_rows(self) -> list[dict[str, Any]]:
        """Return normalized mobile configuration profile rows."""
        if self._mobile_profile_cache is not None:
            return self._mobile_profile_cache

        raw = self._bridge.ios_profiles_list()
        rows = [
            _normalize_mobile_profile_row(item)
            for item in _extract_items(raw)
            if isinstance(item, dict)
        ]
        rows = [row for row in rows if row.get("Profile Name") or row.get("Profile ID")]
        if not rows:
            raise RuntimeError("jamf-cli returned no mobile configuration profiles")
        self._mobile_profile_cache = rows
        return rows

    def _write_counter_block(
        self,
        ws: Any,
        row: int,
        title: str,
        left_header: str,
        counter: Counter,
        max_rows: int = 15,
    ) -> int:
        """Write a title plus a two-column frequency table and return the last row used."""
        if not counter:
            return row

        row += 1
        _safe_write(ws, row, 0, title, self._fmts["header"])
        row += 1
        _safe_write(ws, row, 0, left_header, self._fmts["header"])
        _safe_write(ws, row, 1, "Count", self._fmts["header"])
        row += 1
        for label, count in counter.most_common(max_rows):
            _safe_write(ws, row, 0, label, self._fmts["cell"])
            _safe_write(ws, row, 1, count, self._fmts["cell"])
            row += 1
        return row - 1

    def _write_table_block(
        self,
        ws: Any,
        row: int,
        title: str,
        headers: list[str],
        items: list[dict[str, Any]],
        max_rows: int = 25,
    ) -> int:
        """Write a simple table block and return the last row used."""
        if not items:
            return row

        row += 2
        _safe_write(ws, row, 0, title, self._fmts["header"])
        row += 1
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1
        for item in items[:max_rows]:
            for col_i, header in enumerate(headers):
                _safe_write(ws, row, col_i, item.get(header, ""), self._fmts["cell"])
            row += 1
        return row - 1

    def write_all(self) -> list[str]:
        """Write all core sheets. Returns list of sheet names written."""
        written = []
        sheets = [("Fleet Overview", self._write_overview)]
        if self._protect_enabled():
            sheets.append(("Protect Overview", self._write_protect_overview))
        if self._platform_enabled():
            sheets.append(("Platform Blueprints", self._write_platform_blueprints))
            for bench in self._platform_benchmark_titles():
                rules_label = _excel_sheet_name(bench, " Rules")
                devices_label = _excel_sheet_name(bench, " Devices")
                sheets.append((
                    rules_label,
                    lambda b=bench: self._write_platform_compliance_rules(b),
                ))
                sheets.append((
                    devices_label,
                    lambda b=bench: self._write_platform_compliance_devices(b),
                ))
            sheets.append(("Platform DDM Status", self._write_platform_ddm_status))
        sheets.extend(
            [
                ("Mobile Fleet Summary", self._write_mobile_fleet_summary),
                ("Inventory Summary", self._write_inventory_summary),
                ("Hardware Models", self._write_hardware_models),
                ("Mobile Inventory", self._write_mobile_inventory),
                ("Security Posture", self._write_security),
                ("Device Compliance", self._write_device_compliance),
                ("Check-in Health", self._write_checkin_health),
                ("EA Coverage", self._write_ea_coverage),
                ("Environment Stats", self._write_env_stats),
                ("EA Definitions", self._write_ea_definitions),
                ("Software Installs", self._write_software_installs),
                ("Policy Health", self._write_policy),
                ("Profile Status", self._write_profile_status),
                ("Mobile Config Profiles", self._write_mobile_config_profiles),
                ("App Status", self._write_app_status),
                ("Patch Compliance", self._write_patch),
                ("Patch Failures", self._write_patch_failures),
                ("Update Status", self._write_update_status),
                ("Update Failures", self._write_update_failures),
                # Smart Groups: wired but skipped until JamfCLIBridge.groups() is implemented.
                ("Smart Groups", self._write_smart_groups),
            ]
        )
        for name, fn in sheets:
            try:
                fn()
                written.append(name)
                print(f"  [ok] {name}")
            except RuntimeError as exc:
                print(f"  [skip] {name}: {exc}")
            except Exception as exc:  # unexpected shape or type from jamf-cli JSON
                print(f"  [skip] {name}: unexpected error — {type(exc).__name__}: {exc}")
        return written

    def _write_overview(self) -> None:
        ws = self._wb.add_worksheet("Fleet Overview")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        rows = self._overview_rows_cached()
        has_status = any(row["status"] for row in rows)
        source_name = self._overview_source_label or self._bridge.overview_source_label()
        row = _write_sheet_header(
            ws,
            self._t("Fleet Overview"),
            f"Source: {source_name} | Generated: {ts}",
            self._fmts,
            ncols=4 if has_status else 3,
        )
        ws.set_column(0, 0, 24)
        ws.set_column(1, 1, 42)
        ws.set_column(2, 2, 24)
        headers = ["Section", "Resource", "Value"] + (["Status"] if has_status else [])
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1
        for item in rows:
            _safe_write(ws, row, 0, item["section"], self._fmts["cell"])
            _safe_write(ws, row, 1, item["resource"], self._fmts["cell"])
            _safe_write(ws, row, 2, item["value"], self._fmts["cell"])
            if has_status:
                status_fmt = self._fmts["red"] if item["status"] == "red" else self._fmts["cell"]
                if item["status"] == "yellow":
                    status_fmt = self._fmts["yellow"]
                _safe_write(ws, row, 3, item["status"], status_fmt)
            row += 1

    def _protect_enabled(self) -> bool:
        """Return True when experimental Jamf Protect reporting is enabled."""
        return self._config.get("protect", "enabled", default=False) is True

    def _platform_enabled(self) -> bool:
        """Return True when preview Platform API reporting is enabled."""
        return self._config.get("platform", "enabled", default=False) is True

    def _platform_benchmark_titles(self) -> list[str]:
        """Return all configured Platform compliance benchmark titles.

        Reads from platform.compliance_benchmarks (list). Returns an empty list
        when the key is absent or empty, which causes compliance sheets to be skipped.
        """
        return _platform_benchmark_titles(self._config)

    @staticmethod
    def _platform_rows(raw: Any) -> list[dict[str, Any]]:
        """Normalize Platform report JSON into a list of row dictionaries."""
        return [item for item in _extract_items(raw) if isinstance(item, dict)]

    def _platform_status_fmt(self, value: Any) -> Any:
        """Return a warning/error format for platform state-like labels."""
        normalized = _normalized_text(value)
        if any(token in normalized for token in ("fail", "error", "unsuccessful")):
            return self._fmts["red"]
        if any(
            token in normalized
            for token in ("pending", "draft", "partial", "progress", "scheduled")
        ):
            return self._fmts["yellow"]
        return self._fmts["cell"]

    def _protect_text(self, flat: dict[str, Any], candidates: list[str]) -> str:
        """Return the first non-empty flattened value for a Protect field."""
        value = _first_value(flat, candidates, "")
        if value in (None, "", []):
            return ""
        if isinstance(value, bool):
            return "Yes" if value else "No"
        return str(value).strip()

    def _protect_has_plan(self, value: str) -> bool:
        """Return True when a Protect plan value appears assigned."""
        return _normalized_text(value) not in {"", "n/a", "na", "none", "null", "unassigned"}

    def _protect_overview_rows(self, data: Any) -> list[dict[str, Any]]:
        """Normalize Protect overview output into section/resource/value rows."""
        items = _extract_items(data)
        structured_rows = [
            {
                "Section": str(item.get("section", "") or ""),
                "Resource": str(item.get("resource", "") or ""),
                "Value": item.get("value", ""),
            }
            for item in items
            if isinstance(item, dict) and "resource" in item
        ]
        if structured_rows:
            return structured_rows

        if isinstance(data, dict):
            return [
                {"Section": "", "Resource": key, "Value": value}
                for key, value in sorted(_flatten_record(data).items())
            ]

        rows: list[dict[str, Any]] = []
        multiple = len(items) > 1
        for idx, item in enumerate(items, start=1):
            if not isinstance(item, dict):
                continue
            prefix = str(
                item.get("section") or item.get("group") or item.get("name") or f"item_{idx}"
            ).strip()
            for key, value in _flatten_record(item).items():
                resource = key
                if multiple and prefix and key not in {"section", "group", "name"}:
                    if key != prefix and not key.startswith(f"{prefix}."):
                        resource = f"{prefix}.{key}"
                rows.append({"Section": prefix if multiple else "", "Resource": resource, "Value": value})
        return rows

    def _protect_metric(self, flat: dict[str, Any]) -> tuple[str, Optional[int]]:
        """Return the most likely analytic hit-count field and value."""
        preferred_keys = [
            "hitCount",
            "hit_count",
            "matchCount",
            "match_count",
            "eventCount",
            "event_count",
            "detectionCount",
            "detection_count",
            "count",
        ]
        for key in preferred_keys:
            if key in flat:
                value = _to_int(flat.get(key), default=-1)
                if value >= 0:
                    return key, value

        best_key = ""
        best_value: Optional[int] = None
        best_score = -1
        for key, value in flat.items():
            key_norm = key.lower()
            if key_norm.endswith("id") or value in (None, "", True, False):
                continue
            try:
                numeric = int(float(str(value).strip()))
            except (TypeError, ValueError):
                continue
            score = sum(
                token in key_norm for token in ("hit", "match", "detection", "event", "count")
            )
            if score > best_score:
                best_key = key
                best_value = numeric
                best_score = score
        return best_key, best_value

    def _protect_computer_rows(self, raw: Any) -> list[dict[str, Any]]:
        """Normalize Protect computer rows across likely field shapes."""
        rows: list[dict[str, Any]] = []
        for item in _extract_items(raw):
            if not isinstance(item, dict):
                continue
            flat = _flatten_record(item)
            name = self._protect_text(flat, PROTECT_COMPUTER_FIELD_CANDIDATES["name"])
            serial = self._protect_text(flat, PROTECT_COMPUTER_FIELD_CANDIDATES["serial"])
            plan = self._protect_text(flat, PROTECT_COMPUTER_FIELD_CANDIDATES["plan_name"])
            if not plan:
                plan = self._protect_text(flat, PROTECT_COMPUTER_FIELD_CANDIDATES["plan_id"])
            row = {
                "Computer": name or serial or self._protect_text(flat, ["id"]),
                "Serial Number": serial,
                "Plan": plan,
                "Status": self._protect_text(flat, PROTECT_COMPUTER_FIELD_CANDIDATES["status"]),
                "Last Seen": self._protect_text(flat, PROTECT_COMPUTER_FIELD_CANDIDATES["last_seen"]),
            }
            if any(str(value).strip() for value in row.values()):
                rows.append(row)
        return rows

    def _protect_plan_rows(self, raw: Any) -> list[dict[str, Any]]:
        """Normalize Protect plan rows across likely field shapes."""
        rows: list[dict[str, Any]] = []
        for item in _extract_items(raw):
            if not isinstance(item, dict):
                continue
            flat = _flatten_record(item)
            name = self._protect_text(flat, PROTECT_PLAN_FIELD_CANDIDATES["name"])
            plan_id = self._protect_text(flat, PROTECT_PLAN_FIELD_CANDIDATES["id"])
            row = {
                "Plan Name": name or plan_id,
                "Plan ID": plan_id,
                "Enabled": self._protect_text(flat, PROTECT_PLAN_FIELD_CANDIDATES["enabled"]),
            }
            if any(str(value).strip() for value in row.values()):
                rows.append(row)
        return rows

    def _protect_analytic_rows(self, raw: Any) -> list[dict[str, Any]]:
        """Normalize Protect analytics rows across likely field shapes."""
        rows: list[dict[str, Any]] = []
        for item in _extract_items(raw):
            if not isinstance(item, dict):
                continue
            flat = _flatten_record(item)
            metric_key, metric_value = self._protect_metric(flat)
            metric_label = metric_key.split(".")[-1].replace("_", " ").title() if metric_key else ""
            name = self._protect_text(flat, PROTECT_ANALYTIC_FIELD_CANDIDATES["name"])
            row = {
                "Analytic": name or self._protect_text(flat, ["id"]),
                "Severity": self._protect_text(flat, PROTECT_ANALYTIC_FIELD_CANDIDATES["severity"]),
                "Enabled": self._protect_text(flat, PROTECT_ANALYTIC_FIELD_CANDIDATES["enabled"]),
                "Metric": metric_label,
                "Value": metric_value if metric_value is not None else "",
            }
            if any(str(value).strip() for value in row.values()):
                rows.append(row)
        return rows

    def _write_protect_overview(self) -> None:
        """Write an experimental Jamf Protect summary sheet from jamf-cli data."""
        if not self._protect_enabled():
            raise RuntimeError("disabled in config (set protect.enabled: true to opt in)")

        overview_raw = self._bridge.protect_overview()
        overview_rows = self._protect_overview_rows(overview_raw)
        if not overview_rows:
            raise RuntimeError("jamf-cli protect overview returned no usable rows")
        overview_has_data = _protect_overview_has_data(overview_raw)

        source_parts = ["jamf-cli protect overview"]
        optional_errors: list[str] = []
        computers: list[dict[str, Any]] = []
        plans: list[dict[str, Any]] = []
        analytics: list[dict[str, Any]] = []
        optional_sources = [
            ("jamf-cli protect computers list", self._bridge.protect_computers_list, self._protect_computer_rows),
            ("jamf-cli protect plans list", self._bridge.protect_plans, self._protect_plan_rows),
            ("jamf-cli protect analytics list", self._bridge.protect_analytics, self._protect_analytic_rows),
        ]
        for source_label, fetcher, normalizer in optional_sources:
            try:
                rows = normalizer(fetcher())
            except RuntimeError as exc:
                optional_errors.append(f"{source_label}: {exc}")
                continue
            if rows:
                source_parts.append(source_label)
            if "computers list" in source_label:
                computers = rows
            elif "plans list" in source_label:
                plans = rows
            else:
                analytics = rows

        ws = self._wb.add_worksheet("Protect Overview")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Protect Overview"),
            f"Source: {' + '.join(source_parts)} | Generated: {ts}",
            self._fmts,
            ncols=5,
        )
        ws.set_column(0, 0, 34)
        ws.set_column(1, 1, 42)
        ws.set_column(2, 4, 20)

        _safe_write(ws, row, 0, "Support Status", self._fmts["cell"])
        _safe_write(ws, row, 1, "Experimental", self._fmts["yellow"])
        row += 1
        _safe_write(ws, row, 0, "Validation Note", self._fmts["cell"])
        _safe_write(
            ws,
            row,
            1,
            "Built from jamf-cli 1.6 Protect commands; not fully tested against a live tenant.",
            self._fmts["yellow"],
        )
        row += 1
        if not overview_has_data:
            _safe_write(ws, row, 0, "Protect Auth Signal", self._fmts["cell"])
            _safe_write(
                ws,
                row,
                1,
                "Overview returned placeholder values only (N/A). Protect auth may be incomplete.",
                self._fmts["yellow"],
            )
            row += 1

        if computers:
            assigned = sum(1 for item in computers if self._protect_has_plan(str(item.get("Plan", ""))))
            missing = len(computers) - assigned
            _safe_write(ws, row, 0, "Protect Computers", self._fmts["cell"])
            _safe_write(ws, row, 1, len(computers), self._fmts["cell"])
            row += 1
            _safe_write(ws, row, 0, "Plan Assignment Rate", self._fmts["cell"])
            _safe_write(
                ws,
                row,
                1,
                assigned / len(computers) if computers else 0.0,
                self._fmts["pct"],
            )
            row += 1
            _safe_write(ws, row, 0, "Computers Missing a Plan", self._fmts["cell"])
            _safe_write(ws, row, 1, missing, self._severity_fmt(missing, 1, 5))
            row += 1
        if plans:
            _safe_write(ws, row, 0, "Protect Plans Listed", self._fmts["cell"])
            _safe_write(ws, row, 1, len(plans), self._fmts["cell"])
            row += 1
        if analytics:
            metric_rows = [item for item in analytics if item.get("Value", "") != ""]
            _safe_write(ws, row, 0, "Protect Analytics Listed", self._fmts["cell"])
            _safe_write(ws, row, 1, len(analytics), self._fmts["cell"])
            row += 1
            if metric_rows:
                _safe_write(ws, row, 0, "Analytics with Count Metrics", self._fmts["cell"])
                _safe_write(ws, row, 1, len(metric_rows), self._fmts["cell"])
                row += 1
        if optional_errors:
            _safe_write(ws, row, 0, "Optional Source Gaps", self._fmts["cell"])
            _safe_write(ws, row, 1, "See notes below", self._fmts["yellow"])
            row += 1

        row += 1
        _safe_write(ws, row, 0, "Section", self._fmts["header"])
        _safe_write(ws, row, 1, "Resource", self._fmts["header"])
        _safe_write(ws, row, 2, "Value", self._fmts["header"])
        row += 1
        for item in overview_rows:
            _safe_write(ws, row, 0, item.get("Section", ""), self._fmts["cell"])
            _safe_write(ws, row, 1, item.get("Resource", ""), self._fmts["cell"])
            _safe_write(ws, row, 2, item.get("Value", ""), self._fmts["cell"])
            row += 1

        if computers:
            plan_counts = Counter(
                item["Plan"] if self._protect_has_plan(str(item.get("Plan", ""))) else "Unassigned"
                for item in computers
            )
            row = self._write_counter_block(ws, row, "Plan Distribution", "Plan", plan_counts)

        row = self._write_table_block(
            ws,
            row,
            "Available Protect Plans",
            ["Plan Name", "Plan ID", "Enabled"],
            plans,
        )

        metric_rows = [
            item for item in analytics
            if item.get("Value", "") != ""
        ]
        metric_rows.sort(key=lambda item: -_to_int(item.get("Value", 0)))
        row = self._write_table_block(
            ws,
            row,
            "Top Analytics by Count",
            ["Analytic", "Severity", "Enabled", "Metric", "Value"],
            metric_rows,
            max_rows=20,
        )

        missing_plan_rows = [
            item for item in computers if not self._protect_has_plan(str(item.get("Plan", "")))
        ]
        missing_plan_rows.sort(
            key=lambda item: (str(item.get("Computer", "")), str(item.get("Serial Number", "")))
        )
        row = self._write_table_block(
            ws,
            row,
            "Computers Missing a Plan",
            ["Computer", "Serial Number", "Status", "Last Seen"],
            missing_plan_rows,
            max_rows=50,
        )

        if optional_errors:
            notes = [{"Note": message} for message in optional_errors[:5]]
            self._write_table_block(ws, row, "Optional Source Notes", ["Note"], notes, max_rows=5)

    def _write_platform_blueprints(self) -> None:
        """Write a blueprint deployment summary from Platform API report data."""
        if not self._platform_enabled():
            raise RuntimeError("disabled in config (set platform.enabled: true to opt in)")

        rows = self._platform_rows(self._bridge.blueprint_status())
        if not rows:
            raise RuntimeError("jamf-cli blueprint-status returned no rows")

        deployed = sum(1 for item in rows if _normalized_text(item.get("state", "")) == "deployed")
        with_failures = sum(1 for item in rows if _to_int(item.get("failed", 0)) > 0)
        with_pending = sum(1 for item in rows if _to_int(item.get("pending", 0)) > 0)

        ws = self._wb.add_worksheet("Platform Blueprints")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            "Platform Blueprints",
            "Source: jamf-cli pro report blueprint-status"
            f" | Generated: {ts}",
            self._fmts,
            ncols=7,
        )
        ws.set_column(0, 0, 34)
        ws.set_column(1, 6, 16)
        for label, value in [
            ("Total Blueprints", len(rows)),
            ("Deployed Blueprints", deployed),
            ("Blueprints with Failures", with_failures),
            ("Blueprints with Pending Devices", with_pending),
        ]:
            _safe_write(ws, row, 0, label, self._fmts["cell"])
            _safe_write(ws, row, 1, value, self._fmts["cell"])
            row += 1

        row += 1
        headers = ["Blueprint", "State", "Scope Groups", "Steps", "Succeeded", "Failed", "Pending"]
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1

        for item in sorted(
            rows,
            key=lambda current: (
                -_to_int(current.get("failed", 0)),
                -_to_int(current.get("pending", 0)),
                str(current.get("name", "")).lower(),
            ),
        ):
            failed = _to_int(item.get("failed", 0))
            pending = _to_int(item.get("pending", 0))
            _safe_write(ws, row, 0, item.get("name", ""), self._fmts["cell"])
            _safe_write(ws, row, 1, item.get("state", ""), self._platform_status_fmt(item.get("state")))
            _safe_write(ws, row, 2, item.get("scope", ""), self._fmts["cell"])
            _safe_write(ws, row, 3, item.get("steps", ""), self._fmts["cell"])
            _safe_write(ws, row, 4, item.get("succeeded", ""), self._fmts["cell"])
            _safe_write(ws, row, 5, item.get("failed", ""), self._fmts["red"] if failed else self._fmts["cell"])
            _safe_write(
                ws,
                row,
                6,
                item.get("pending", ""),
                self._fmts["yellow"] if pending else self._fmts["cell"],
            )
            row += 1

    def _write_platform_compliance_rules(self, benchmark: str) -> None:
        """Write per-rule Platform compliance results for a single benchmark.

        Args:
            benchmark: Benchmark title passed to jamf-cli compliance-rules.
        """
        rows = self._platform_rows(self._bridge.compliance_rules(benchmark))
        if not rows:
            raise RuntimeError("jamf-cli compliance-rules returned no rows")

        avg_pass_rate = [
            ratio for ratio in (_parse_percent(item.get("passRate", "")) for item in rows)
            if ratio is not None
        ]
        rules_with_failures = sum(1 for item in rows if _to_int(item.get("failed", 0)) > 0)
        rules_with_unknown = sum(1 for item in rows if _to_int(item.get("unknown", 0)) > 0)

        sheet_title = "Compliance Rules"
        sheet_name = _excel_sheet_name(
            benchmark,
            " Rules",
            existing_names=set(self._wb.sheetnames),
        )
        ws = self._wb.add_worksheet(sheet_name)
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            sheet_title,
            f"Source: jamf-cli pro report compliance-rules {benchmark} | Generated: {ts}",
            self._fmts,
            ncols=6,
        )
        ws.set_column(0, 0, 44)
        ws.set_column(1, 5, 16)
        for label, value in [
            ("Benchmark", benchmark),
            ("Rules Returned", len(rows)),
            ("Rules with Failures", rules_with_failures),
            ("Rules with Unknown Results", rules_with_unknown),
            (
                "Average Pass Rate",
                (sum(avg_pass_rate) / len(avg_pass_rate)) if avg_pass_rate else "",
            ),
        ]:
            _safe_write(ws, row, 0, label, self._fmts["cell"])
            fmt = self._fmts["pct"] if label == "Average Pass Rate" and value != "" else self._fmts["cell"]
            _safe_write(ws, row, 1, value, fmt)
            row += 1

        row += 1
        headers = ["Rule", "Passed", "Failed", "Unknown", "Devices", "Pass Rate"]
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1

        for item in rows:
            failed = _to_int(item.get("failed", 0))
            unknown = _to_int(item.get("unknown", 0))
            pass_rate = _parse_percent(item.get("passRate", ""))
            base_fmt = self._fmts["red"] if failed else self._fmts["yellow"] if unknown else self._fmts["cell"]
            _safe_write(ws, row, 0, item.get("rule", ""), base_fmt)
            _safe_write(ws, row, 1, _to_int(item.get("passed", 0)), self._fmts["cell"])
            _safe_write(ws, row, 2, failed, self._fmts["red"] if failed else self._fmts["cell"])
            _safe_write(
                ws,
                row,
                3,
                unknown,
                self._fmts["yellow"] if unknown else self._fmts["cell"],
            )
            _safe_write(ws, row, 4, _to_int(item.get("devices", 0)), self._fmts["cell"])
            if pass_rate is None:
                _safe_write(ws, row, 5, item.get("passRate", ""), self._fmts["cell"])
            else:
                _safe_write(ws, row, 5, pass_rate, _pct_format(self._fmts, pass_rate))
            row += 1

    def _write_platform_compliance_devices(self, benchmark: str) -> None:
        """Write failing Platform compliance devices for a single benchmark.

        Args:
            benchmark: Benchmark title passed to jamf-cli compliance-devices.
        """
        rows = self._platform_rows(self._bridge.compliance_devices(benchmark))
        if not rows:
            raise RuntimeError("jamf-cli compliance-devices returned no rows")

        compliance_values = [
            ratio for ratio in (_parse_percent(item.get("compliance", "")) for item in rows)
            if ratio is not None
        ]

        sheet_title = "Compliance Devices"
        sheet_name = _excel_sheet_name(
            benchmark,
            " Devices",
            existing_names=set(self._wb.sheetnames),
        )
        ws = self._wb.add_worksheet(sheet_name)
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            sheet_title,
            f"Source: jamf-cli pro report compliance-devices {benchmark} | Generated: {ts}",
            self._fmts,
            ncols=5,
        )
        ws.set_column(0, 0, 34)
        ws.set_column(1, 1, 18)
        ws.set_column(2, 4, 16)
        for label, value in [
            ("Benchmark", benchmark),
            ("Devices Returned", len(rows)),
            ("Devices with Failing Rules", sum(1 for item in rows if _to_int(item.get("rulesFailed", 0)) > 0)),
            (
                "Average Compliance",
                (sum(compliance_values) / len(compliance_values)) if compliance_values else "",
            ),
        ]:
            _safe_write(ws, row, 0, label, self._fmts["cell"])
            fmt = self._fmts["pct"] if label == "Average Compliance" and value != "" else self._fmts["cell"]
            _safe_write(ws, row, 1, value, fmt)
            row += 1

        row += 1
        headers = ["Device", "Device ID", "Rules Failed", "Rules Passed", "Compliance"]
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1

        for item in sorted(
            rows,
            key=lambda current: (
                -_to_int(current.get("rulesFailed", 0)),
                str(current.get("device", "")).lower(),
            ),
        ):
            rules_failed = _to_int(item.get("rulesFailed", 0))
            compliance = _parse_percent(item.get("compliance", ""))
            _safe_write(ws, row, 0, item.get("device", ""), self._fmts["cell"])
            _safe_write(ws, row, 1, item.get("deviceId", ""), self._fmts["cell"])
            _safe_write(
                ws,
                row,
                2,
                rules_failed,
                self._fmts["red"] if rules_failed else self._fmts["cell"],
            )
            _safe_write(ws, row, 3, _to_int(item.get("rulesPassed", 0)), self._fmts["cell"])
            if compliance is None:
                _safe_write(ws, row, 4, item.get("compliance", ""), self._fmts["cell"])
            else:
                _safe_write(ws, row, 4, compliance, _pct_format(self._fmts, compliance))
            row += 1

    def _write_platform_ddm_status(self) -> None:
        """Write DDM declaration health from Platform API report data."""
        if not self._platform_enabled():
            raise RuntimeError("disabled in config (set platform.enabled: true to opt in)")

        rows = self._platform_rows(self._bridge.ddm_status())
        if not rows:
            raise RuntimeError("jamf-cli ddm-status returned no rows")

        sources_with_issues = sum(1 for item in rows if _to_int(item.get("unsuccessful", 0)) > 0)
        unsuccessful_total = sum(_to_int(item.get("unsuccessful", 0)) for item in rows)

        ws = self._wb.add_worksheet("Platform DDM Status")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            "Platform DDM Status",
            "Source: jamf-cli pro report ddm-status"
            f" | Generated: {ts}",
            self._fmts,
            ncols=6,
        )
        ws.set_column(0, 1, 28)
        ws.set_column(2, 5, 16)
        for label, value in [
            ("Sources Returned", len(rows)),
            ("Sources with Unsuccessful Declarations", sources_with_issues),
            ("Total Unsuccessful Declarations", unsuccessful_total),
        ]:
            _safe_write(ws, row, 0, label, self._fmts["cell"])
            _safe_write(ws, row, 1, value, self._fmts["cell"])
            row += 1

        row += 1
        headers = ["Type", "Source", "Devices", "Declarations", "Successful", "Unsuccessful"]
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1

        for item in sorted(
            rows,
            key=lambda current: (
                -_to_int(current.get("unsuccessful", 0)),
                str(current.get("source", "")).lower(),
            ),
        ):
            unsuccessful = _to_int(item.get("unsuccessful", 0))
            _safe_write(ws, row, 0, item.get("type", ""), self._fmts["cell"])
            _safe_write(ws, row, 1, item.get("source", ""), self._fmts["cell"])
            _safe_write(ws, row, 2, _to_int(item.get("devices", 0)), self._fmts["cell"])
            _safe_write(ws, row, 3, _to_int(item.get("declarations", 0)), self._fmts["cell"])
            _safe_write(ws, row, 4, _to_int(item.get("successful", 0)), self._fmts["cell"])
            _safe_write(
                ws,
                row,
                5,
                unsuccessful,
                self._fmts["red"] if unsuccessful else self._fmts["cell"],
            )
            row += 1

    def _write_mobile_fleet_summary(self) -> None:
        """Write a mobile-focused fleet summary using overview and mobile inventory sources."""
        overview_rows: list[dict[str, Any]] = []
        mobile_rows: list[dict[str, Any]] = []
        profile_rows: list[dict[str, Any]] = []
        source_parts: list[str] = []

        try:
            overview_rows = self._overview_rows_cached()
            source_parts.append("jamf-cli pro overview")
        except RuntimeError:
            overview_rows = []

        try:
            mobile_rows, mobile_source = self._mobile_inventory_rows()
            if mobile_source:
                source_parts.append(mobile_source)
        except RuntimeError:
            mobile_rows = []

        try:
            profile_rows = self._mobile_profile_rows()
            source_parts.append("jamf-cli pro classic-mobile-config-profiles list")
        except RuntimeError:
            profile_rows = []

        if not overview_rows and not mobile_rows and not profile_rows:
            raise RuntimeError("no mobile device sources were available")

        summary = _summarize_mobile_inventory(mobile_rows) if mobile_rows else {}
        stale_days = int(self._config.thresholds.get("stale_device_days", 30))
        inventory_stale = sum(
            1
            for row_data in mobile_rows
            if isinstance(row_data.get("Days Since Inventory"), int)
            and row_data["Days Since Inventory"] > stale_days
        )

        ws = self._wb.add_worksheet("Mobile Fleet Summary")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        source_name = " + ".join(dict.fromkeys(source_parts)) or "jamf-cli mobile device sources"
        row = _write_sheet_header(
            ws,
            self._t("Mobile Fleet Summary"),
            f"Source: {source_name} | Generated: {ts}",
            self._fmts,
            ncols=3,
        )
        ws.set_column(0, 0, 30)
        ws.set_column(1, 1, 20)
        ws.set_column(2, 2, 20)

        summary_rows: list[tuple[str, Any]] = []
        for resource, label in [
            ("Managed Devices", "Managed Mobile Devices"),
            ("Unmanaged Devices", "Unmanaged Mobile Devices"),
            ("Mobile Device Prestages", "Mobile Device Prestages"),
            ("Mobile Device Smart Groups", "Mobile Device Smart Groups"),
            ("MDM Auto Renew (Mobile)", "MDM Auto Renew (Mobile)"),
            ("iOS Config Profiles", "iOS Config Profiles"),
        ]:
            value = self._overview_value(overview_rows, resource)
            if value not in (None, ""):
                summary_rows.append((label, value))

        if mobile_rows:
            summary_rows.extend(
                [
                    ("Inventory Rows Returned", summary.get("total", 0)),
                    ("Managed Rows", summary.get("managed", 0)),
                    ("Unmanaged Rows", summary.get("unmanaged", 0)),
                    ("Supervised Devices", summary.get("supervised", 0)),
                    ("Shared iPad Devices", summary.get("shared_ipad", 0)),
                    ("Assigned Users", summary.get("assigned", 0)),
                    ("Activation Lock Enabled", summary.get("activation_lock", 0)),
                    ("Passcode Compliant", summary.get("passcode_compliant", 0)),
                    (f"Inventory Older Than {stale_days} Days", inventory_stale),
                ]
            )
        if profile_rows:
            summary_rows.append(("Mobile Config Profiles (List)", len(profile_rows)))

        for label, value in summary_rows:
            _safe_write(ws, row, 0, label, self._fmts["cell"])
            _safe_write(ws, row, 1, value, self._fmts["cell"])
            row += 1

        if mobile_rows:
            row = self._write_counter_block(
                ws,
                row,
                "Device Family Distribution",
                "Device Family",
                summary["families"],
            )
            row = self._write_counter_block(
                ws,
                row,
                "OS Version Distribution",
                "OS Version",
                summary["os_versions"],
            )
            self._write_counter_block(ws, row, "Top Models", "Model", summary["models"], max_rows=10)

    def _write_mobile_inventory(self) -> None:
        """Write normalized mobile device inventory rows from jamf-cli."""
        rows, source_name = self._mobile_inventory_rows()
        if not rows:
            raise RuntimeError("jamf-cli returned no mobile device inventory rows")

        summary = _summarize_mobile_inventory(rows)
        stale_days = int(self._config.thresholds.get("stale_device_days", 30))
        inventory_stale = sum(
            1
            for row_data in rows
            if isinstance(row_data.get("Days Since Inventory"), int)
            and row_data["Days Since Inventory"] > stale_days
        )

        ws = self._wb.add_worksheet("Mobile Inventory")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Mobile Inventory"),
            f"Source: {source_name} | Generated: {ts}",
            self._fmts,
            ncols=20,
        )
        summary_rows = [
            ("Total Mobile Devices", summary["total"]),
            ("Managed", summary["managed"]),
            ("Unmanaged", summary["unmanaged"]),
            ("Supervised", summary["supervised"]),
            ("Shared iPad", summary["shared_ipad"]),
            ("Assigned Users", summary["assigned"]),
            ("Inventory Older Than Threshold", inventory_stale),
        ]
        for label, value in summary_rows:
            _safe_write(ws, row, 0, label, self._fmts["cell"])
            _safe_write(ws, row, 1, value, self._fmts["cell"])
            row += 1

        row += 1
        headers = list(rows[0].keys())
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1

        widths = [12, 26, 18, 14, 11, 11, 11, 24, 12, 18, 24, 18, 18, 22, 18, 16, 18, 18, 18, 14]
        for col_i, width in enumerate(widths):
            ws.set_column(col_i, col_i, width)

        for item in sorted(
            rows,
            key=lambda current: (
                str(current.get("Device Family", "")),
                str(current.get("Device Name", "")),
                str(current.get("Serial Number", "")),
            ),
        ):
            for col_i, header in enumerate(headers):
                value = item.get(header, "")
                if header == "Days Since Inventory" and isinstance(value, int):
                    fmt = self._severity_fmt(value, stale_days, stale_days * 2)
                    _safe_write(ws, row, col_i, value, fmt)
                else:
                    _safe_write(ws, row, col_i, value, self._fmts["cell"])
            row += 1

    def _write_mobile_config_profiles(self) -> None:
        """Write mobile configuration profile visibility from jamf-cli."""
        rows = self._mobile_profile_rows()
        if not rows:
            raise RuntimeError("jamf-cli returned no mobile configuration profiles")

        category_counts = Counter(row.get("Category") or "Uncategorized" for row in rows)
        ws = self._wb.add_worksheet("Mobile Config Profiles")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Mobile Config Profiles"),
            "Source: jamf-cli pro classic-mobile-config-profiles list"
            f" | Generated: {ts}",
            self._fmts,
            ncols=5,
        )
        ws.set_column(0, 0, 34)
        ws.set_column(1, 1, 14)
        ws.set_column(2, 2, 24)
        ws.set_column(3, 3, 20)
        ws.set_column(4, 4, 44)

        for label, value in [
            ("Total Profiles", len(rows)),
            ("Unique Categories", len(category_counts)),
            ("Uncategorized Profiles", category_counts.get("Uncategorized", 0)),
        ]:
            _safe_write(ws, row, 0, label, self._fmts["cell"])
            _safe_write(ws, row, 1, value, self._fmts["cell"])
            row += 1

        row = self._write_counter_block(
            ws,
            row,
            "Profiles by Category",
            "Category",
            category_counts,
            max_rows=15,
        )
        row += 2

        headers = ["Profile Name", "Profile ID", "Category", "Site", "Description"]
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1
        for item in sorted(
            rows,
            key=lambda current: (
                str(current.get("Category", "")),
                str(current.get("Profile Name", "")),
            ),
        ):
            for col_i, header in enumerate(headers):
                _safe_write(ws, row, col_i, item.get(header, ""), self._fmts["cell"])
            row += 1

    def _flatten_dict(self, d: Any, prefix: str = "") -> dict[str, Any]:
        """Recursively flatten a nested dict into dot-separated keys."""
        out: dict[str, Any] = {}
        if not isinstance(d, dict):
            return {prefix: d} if prefix else {}
        for k, v in d.items():
            full_key = f"{prefix}.{k}" if prefix else k
            if isinstance(v, dict):
                out.update(self._flatten_dict(v, full_key))
            elif isinstance(v, list):
                out[full_key] = f"[{len(v)} items]"
            else:
                out[full_key] = v
        return out

    def _overview_rows(self, data: Any) -> list[dict[str, Any]]:
        """Normalize jamf-cli overview data across older and newer JSON shapes."""
        if isinstance(data, dict):
            return [
                {"section": "", "resource": key, "value": value, "status": ""}
                for key, value in self._flatten_dict(data).items()
            ]

        rows: list[dict[str, Any]] = []
        for item in data if isinstance(data, list) else []:
            if not isinstance(item, dict):
                continue
            rows.append(
                {
                    "section": str(item.get("section", "")),
                    "resource": str(item.get("resource", "")),
                    "value": item.get("value", ""),
                    "status": str(item.get("status", "")),
                }
            )
        return rows

    def _write_security(self) -> None:
        # jamf-cli pro report security --output json returns a flat mixed list:
        #   [{"section":"summary","data":{total_devices,filevault_encrypted,...}},
        #    {"section":"device",...},
        #    {"section":"os_version","os_version":"15.7.3","count":N,"pct":"N%"}]
        raw = self._bridge.security_report()
        items = raw if isinstance(raw, list) else []
        summary = next(
            (i.get("data", i) for i in items if i.get("section") == "summary"), {}
        )
        total = int(summary.get("total_devices", 0))
        os_rows = [i for i in items if i.get("section") == "os_version"]

        ws = self._wb.add_worksheet("Security Posture")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Security Posture"),
            f"Source: jamf-cli pro report security | Generated: {ts}",
            self._fmts,
            ncols=5,
        )
        ws.set_column(0, 0, 35)
        ws.set_column(1, 4, 18)
        headers = ["Control", "Compliant", "Non-Compliant", "Total", "Compliance %"]
        for c, h in enumerate(headers):
            _safe_write(ws, row, c, h, self._fmts["header"])
        row += 1

        controls = [
            ("FileVault", "filevault_encrypted"),
            ("Gatekeeper", "gatekeeper_enabled"),
            ("SIP", "sip_enabled"),
            ("Firewall", "firewall_enabled"),
        ]
        for label, key in controls:
            compliant = int(summary.get(key, 0))
            non_compliant = total - compliant
            pct = compliant / total if total > 0 else 0.0
            _safe_write(ws, row, 0, label, self._fmts["cell"])
            _safe_write(ws, row, 1, compliant, self._fmts["cell"])
            _safe_write(ws, row, 2, non_compliant, self._fmts["cell"])
            _safe_write(ws, row, 3, total, self._fmts["cell"])
            _safe_write(ws, row, 4, pct, _pct_format(self._fmts, pct))
            row += 1

        if os_rows:
            row += 1
            _safe_write(ws, row, 0, "OS Version Distribution", self._fmts["header"])
            _safe_write(ws, row, 1, "Device Count", self._fmts["header"])
            _safe_write(ws, row, 2, "% of Fleet", self._fmts["header"])
            row += 1
            for item in os_rows:
                _safe_write(ws, row, 0, item.get("os_version", ""), self._fmts["cell"])
                _safe_write(ws, row, 1, item.get("count", 0), self._fmts["cell"])
                _safe_write(ws, row, 2, item.get("pct", ""), self._fmts["cell"])
                row += 1

    def _write_inventory_summary(self) -> None:
        """Write a flat inventory breakdown by model and OS version."""
        raw = self._bridge.inventory_summary()
        rows = raw if isinstance(raw, list) else []
        if not rows:
            raise RuntimeError("jamf-cli inventory-summary returned no rows")

        ws = self._wb.add_worksheet("Inventory Summary")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Inventory Summary"),
            f"Source: jamf-cli pro report inventory-summary | Generated: {ts}",
            self._fmts,
            ncols=3,
        )
        ws.set_column(0, 0, 36)
        ws.set_column(1, 1, 18)
        ws.set_column(2, 2, 14)
        headers = ["Model", "OS Version", "Device Count"]
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1

        for item in sorted(
            rows,
            key=lambda current: (
                -_to_int(current.get("count", 0)),
                str(current.get("model", "")),
                str(current.get("os_version", "")),
            ),
        ):
            _safe_write(ws, row, 0, item.get("model", "Unknown"), self._fmts["cell"])
            _safe_write(ws, row, 1, item.get("os_version", "Unknown"), self._fmts["cell"])
            _safe_write(ws, row, 2, _to_int(item.get("count", 0)), self._fmts["cell"])
            row += 1

    def _write_hardware_models(self) -> None:
        """Write hardware model distribution from jamf-cli.

        Consumes either the native hardware-models report shape (dashboard-era jamf-cli)
        or inventory-summary rows (current jamf-cli), aggregating by model.

        Native shape: {"computers": [{"model": ..., "count": N}], "mobile": [...]}
        Fallback shape: [{"model": ..., "os_version": ..., "count": N}]
        """
        raw = self._bridge.hardware_models()

        # Detect native hardware-models shape vs inventory-summary fallback.
        if isinstance(raw, dict) and ("computers" in raw or "mobile" in raw):
            computer_rows = [
                {"model": r.get("model", "Unknown"), "count": _to_int(r.get("count", 0))}
                for r in (raw.get("computers") or [])
                if isinstance(r, dict)
            ]
            mobile_rows = [
                {"model": r.get("model", "Unknown"), "count": _to_int(r.get("count", 0))}
                for r in (raw.get("mobile") or [])
                if isinstance(r, dict)
            ]
        else:
            # inventory-summary: aggregate counts by model (sum across OS versions).
            model_counts: dict[str, int] = {}
            for item in (raw if isinstance(raw, list) else []):
                if not isinstance(item, dict):
                    continue
                model = str(item.get("model") or "Unknown").strip()
                model_counts[model] = model_counts.get(model, 0) + _to_int(item.get("count", 0))
            computer_rows = sorted(
                [{"model": m, "count": c} for m, c in model_counts.items()],
                key=lambda r: -r["count"],
            )
            mobile_rows = []

        if not computer_rows and not mobile_rows:
            raise RuntimeError("hardware-models returned no rows")

        ws = self._wb.add_worksheet("Hardware Models")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Hardware Models"),
            f"Source: jamf-cli pro report hardware-models | Generated: {ts}",
            self._fmts,
            ncols=2,
        )
        ws.set_column(0, 0, 40)
        ws.set_column(1, 1, 14)

        top_n = 20
        for section_label, section_rows in [
            ("Computer Models", computer_rows),
            ("Mobile Models", mobile_rows),
        ]:
            if not section_rows:
                continue
            _safe_write(ws, row, 0, section_label, self._fmts["header"])
            _safe_write(ws, row, 1, "Count", self._fmts["header"])
            row += 1
            for item in section_rows[:top_n]:
                _safe_write(ws, row, 0, item["model"], self._fmts["cell"])
                _safe_write(ws, row, 1, item["count"], self._fmts["cell"])
                row += 1
            row += 1

    def _write_device_compliance(self) -> None:
        """Write a stale check-in summary and device detail table from jamf-cli."""
        stale_days = int(self._config.thresholds.get("stale_device_days", 30))
        raw = self._bridge.device_compliance(stale_days)
        rows = raw if isinstance(raw, list) else []
        if not rows:
            raise RuntimeError("jamf-cli device-compliance returned no rows")

        total = len(rows)
        stale_count = sum(1 for item in rows if _to_bool(item.get("stale")))
        managed_count = sum(1 for item in rows if _to_bool(item.get("managed")))

        ws = self._wb.add_worksheet("Device Compliance")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Device Compliance"),
            (
                "Source: jamf-cli pro report device-compliance"
                f" (--days-since-checkin {stale_days}) | Generated: {ts}"
            ),
            self._fmts,
            ncols=7,
        )
        summary_rows = [
            ("Total Devices", total),
            ("Managed Devices", managed_count),
            ("Unmanaged Devices", total - managed_count),
            (f"Stale Devices (>{stale_days} days)", stale_count),
            ("Current Devices", total - stale_count),
        ]
        for label, value in summary_rows:
            _safe_write(ws, row, 0, label, self._fmts["cell"])
            _safe_write(ws, row, 1, value, self._fmts["cell"])
            row += 1

        row += 1
        headers = [
            "Computer Name",
            "Serial Number",
            "Managed",
            "OS Version",
            "Last Contact",
            "Days Since Contact",
            "Stale",
        ]
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1
        ws.set_column(0, 0, 30)
        ws.set_column(1, 1, 18)
        ws.set_column(2, 2, 12)
        ws.set_column(3, 6, 18)

        for item in sorted(
            rows,
            key=lambda current: (
                not _to_bool(current.get("stale")),
                -_to_int(current.get("days_since_contact", -1), -1),
                str(current.get("name", "")),
            ),
        ):
            stale = _to_bool(item.get("stale"))
            managed = _to_bool(item.get("managed"))
            row_fmt = self._fmts["red"] if stale else self._fmts["cell"]
            _safe_write(ws, row, 0, item.get("name", ""), row_fmt)
            _safe_write(ws, row, 1, item.get("serial", ""), row_fmt)
            _safe_write(ws, row, 2, "Yes" if managed else "No", row_fmt)
            _safe_write(ws, row, 3, item.get("os_version", ""), row_fmt)
            _safe_write(ws, row, 4, item.get("last_contact", ""), row_fmt)
            _safe_write(ws, row, 5, item.get("days_since_contact", ""), row_fmt)
            _safe_write(ws, row, 6, "Yes" if stale else "No", row_fmt)
            row += 1

    def _write_checkin_health(self) -> None:
        """Write a check-in health summary from jamf-cli.

        Uses the native checkin-status report when available (dashboard-era jamf-cli),
        otherwise derives the summary from device-compliance rows using the configured
        checkin_overdue_days threshold.

        Native shape: {"computers": {"total": N, "overdue": N, "threshold_days": N},
                       "mobile": {"total": N, "overdue": N, "threshold_days": N}}
        Fallback shape: same list as device_compliance() output.
        """
        threshold = int(self._config.thresholds.get("checkin_overdue_days", 7))
        raw = self._bridge.checkin_status(threshold)

        ws = self._wb.add_worksheet("Check-in Health")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Check-in Health"),
            (
                f"Source: jamf-cli pro report checkin-status"
                f" (threshold: {threshold} days) | Generated: {ts}"
            ),
            self._fmts,
            ncols=4,
        )
        ws.set_column(0, 0, 30)
        ws.set_column(1, 3, 18)

        def _write_checkin_section(
            section_label: str,
            total: int,
            overdue: int,
        ) -> None:
            nonlocal row
            current = total - overdue
            pct_current = current / total if total > 0 else 0.0
            _safe_write(ws, row, 0, section_label, self._fmts["header"])
            _safe_write(ws, row, 1, "Count", self._fmts["header"])
            _safe_write(ws, row, 2, "% of Total", self._fmts["header"])
            row += 1
            _safe_write(ws, row, 0, "Total Devices", self._fmts["cell"])
            _safe_write(ws, row, 1, total, self._fmts["cell"])
            _safe_write(ws, row, 2, "", self._fmts["cell"])
            row += 1
            _safe_write(ws, row, 0, f"Checked In (within {threshold} days)", self._fmts["cell"])
            _safe_write(ws, row, 1, current, self._fmts["cell"])
            _safe_write(ws, row, 2, pct_current, _pct_format(self._fmts, pct_current))
            row += 1
            overdue_pct = overdue / total if total > 0 else 0.0
            overdue_fmt = self._fmts["red"] if overdue > 0 else self._fmts["cell"]
            _safe_write(ws, row, 0, f"Overdue (>{threshold} days)", overdue_fmt)
            _safe_write(ws, row, 1, overdue, overdue_fmt)
            _safe_write(ws, row, 2, overdue_pct, overdue_fmt)
            row += 2

        # Native shape: dict with computers/mobile sub-dicts.
        if isinstance(raw, dict) and ("computers" in raw or "mobile" in raw):
            for section_key, label in (("computers", "Computers"), ("mobile", "Mobile Devices")):
                sub = raw.get(section_key) or {}
                if not isinstance(sub, dict):
                    continue
                total = _to_int(sub.get("total", 0))
                overdue = _to_int(sub.get("overdue", 0))
                if total == 0:
                    continue
                _write_checkin_section(label, total, overdue)
            return

        # Fallback: device-compliance list — compute summary from stale flags.
        items = raw if isinstance(raw, list) else []
        if not items:
            raise RuntimeError("checkin-status returned no data")
        total = len(items)
        overdue = sum(1 for i in items if _to_bool(i.get("stale")))
        _write_checkin_section("Computers", total, overdue)

    def _write_ea_coverage(self) -> None:
        """Write a fleet-wide extension attribute coverage summary from jamf-cli."""
        raw = self._bridge.ea_results_report(include_all=True)
        rows = raw if isinstance(raw, list) else []
        if not rows:
            raise RuntimeError("jamf-cli ea-results returned no rows")

        definitions_by_id: dict[str, dict[str, Any]] = {}
        definitions_by_name: dict[str, dict[str, Any]] = {}
        try:
            definitions_raw = self._bridge.computer_extension_attributes()
            definitions = definitions_raw if isinstance(definitions_raw, list) else []
        except RuntimeError as exc:
            print(f"  [warn] EA definitions unavailable; EA Coverage will lack type/input"
                  f" metadata: {exc}")
            definitions = []

        for item in definitions:
            if not isinstance(item, dict):
                continue
            item_id = str(item.get("id", "") or "").strip()
            item_name = str(item.get("name", "") or "").strip()
            if item_id:
                definitions_by_id[item_id] = item
            if item_name:
                definitions_by_name[item_name] = item

        coverage: dict[str, dict[str, Any]] = {}
        for item in rows:
            if not isinstance(item, dict):
                continue
            ea_name = str(item.get("ea_name", "") or "").strip()
            definition_id = str(item.get("definition_id", "") or "").strip()
            if not ea_name and not definition_id:
                continue
            key = ea_name or definition_id
            info = coverage.setdefault(
                key,
                {
                    "ea_name": ea_name,
                    "definition_id": definition_id,
                    "total_devices": 0,
                    "populated_devices": 0,
                    "value_counts": Counter(),
                },
            )
            value = str(item.get("value", "") or "").strip()
            info["total_devices"] += 1
            if value:
                info["populated_devices"] += 1
                info["value_counts"][value] += 1

        if not coverage:
            raise RuntimeError("jamf-cli ea-results returned no usable rows")

        fleet_size = max((item["total_devices"] for item in coverage.values()), default=0)
        eas_with_data = sum(1 for item in coverage.values() if item["populated_devices"] > 0)

        ws = self._wb.add_worksheet("EA Coverage")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        subtitle = "Source: jamf-cli pro report ea-results --all"
        if definitions:
            subtitle += " + computer-extension-attributes list"
        row = _write_sheet_header(
            ws,
            self._t("Extension Attribute Coverage"),
            f"{subtitle} | Generated: {ts}",
            self._fmts,
            ncols=11,
        )
        summary_rows = [
            ("Total Extension Attributes", len(coverage)),
            ("EAs With Populated Values", eas_with_data),
            ("EAs With No Populated Values", len(coverage) - eas_with_data),
            ("Inferred Fleet Size", fleet_size),
        ]
        for label, value in summary_rows:
            _safe_write(ws, row, 0, label, self._fmts["cell"])
            _safe_write(ws, row, 1, value, self._fmts["cell"])
            row += 1

        row += 1
        headers = [
            "EA Name",
            "Definition ID",
            "Data Type",
            "Input Type",
            "Enabled",
            "Populated Devices",
            "Empty Devices",
            "Total Devices",
            "Coverage %",
            "Unique Non-Empty Values",
            "Top Values",
        ]
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1
        ws.set_column(0, 0, 34)
        ws.set_column(1, 4, 16)
        ws.set_column(5, 8, 16)
        ws.set_column(9, 9, 20)
        ws.set_column(10, 10, 56)

        for item in sorted(coverage.values(), key=lambda current: str(current["ea_name"]).lower()):
            definition = (
                definitions_by_id.get(item["definition_id"])
                or definitions_by_name.get(item["ea_name"])
                or {}
            )
            total_devices = item["total_devices"]
            populated_devices = item["populated_devices"]
            empty_devices = max(total_devices - populated_devices, 0)
            coverage_pct = populated_devices / total_devices if total_devices > 0 else 0.0
            top_values = ", ".join(
                f"{value} ({count})"
                for value, count in item["value_counts"].most_common(3)
            )
            _safe_write(ws, row, 0, item["ea_name"], self._fmts["cell"])
            _safe_write(ws, row, 1, item["definition_id"], self._fmts["cell"])
            _safe_write(ws, row, 2, definition.get("dataType", ""), self._fmts["cell"])
            _safe_write(ws, row, 3, definition.get("inputType", ""), self._fmts["cell"])
            enabled_raw = definition.get("enabled")
            enabled_value = ""
            if enabled_raw is True:
                enabled_value = "Yes"
            elif enabled_raw is False:
                enabled_value = "No"
            _safe_write(ws, row, 4, enabled_value, self._fmts["cell"])
            _safe_write(ws, row, 5, populated_devices, self._fmts["cell"])
            _safe_write(ws, row, 6, empty_devices, self._fmts["cell"])
            _safe_write(ws, row, 7, total_devices, self._fmts["cell"])
            _safe_write(ws, row, 8, coverage_pct, _pct_format(self._fmts, coverage_pct))
            _safe_write(ws, row, 9, len(item["value_counts"]), self._fmts["cell"])
            _safe_write(ws, row, 10, top_values, self._fmts["cell"])
            row += 1

    def _write_env_stats(self) -> None:
        """Write environment object counts from jamf-cli pro report env-stats.

        Requires dashboard-era jamf-cli. Sheet is skipped when the command is absent.
        JSON shape: {"policies": N, "config_profiles": N, "scripts": N, "packages": N,
                     "smart_groups_computer": N, "smart_groups_mobile": N,
                     "extension_attributes": N, "categories": N}
        """
        raw = self._bridge.env_stats()
        if not isinstance(raw, dict) or not raw:
            raise RuntimeError("env-stats returned no data")

        display_fields = [
            ("policies", "Policies"),
            ("config_profiles", "Configuration Profiles"),
            ("scripts", "Scripts"),
            ("packages", "Packages"),
            ("smart_groups_computer", "Smart Groups — Computer"),
            ("smart_groups_mobile", "Smart Groups — Mobile"),
            ("extension_attributes", "Extension Attributes"),
            ("categories", "Categories"),
        ]
        rows = [(label, raw[key]) for key, label in display_fields if key in raw]
        if not rows:
            raise RuntimeError("env-stats contained no recognised fields")

        ws = self._wb.add_worksheet("Environment Stats")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Environment Stats"),
            f"Source: jamf-cli pro report env-stats | Generated: {ts}",
            self._fmts,
            ncols=2,
        )
        ws.set_column(0, 0, 36)
        ws.set_column(1, 1, 14)
        _safe_write(ws, row, 0, "Object Type", self._fmts["header"])
        _safe_write(ws, row, 1, "Count", self._fmts["header"])
        row += 1
        for label, value in rows:
            _safe_write(ws, row, 0, label, self._fmts["cell"])
            _safe_write(ws, row, 1, _to_int(value), self._fmts["cell"])
            row += 1

    def _write_ea_definitions(self) -> None:
        """Write computer extension attribute definitions from jamf-cli."""
        raw = self._bridge.computer_extension_attributes()
        rows = raw if isinstance(raw, list) else []
        if not rows:
            raise RuntimeError("jamf-cli computer-extension-attributes returned no rows")

        ws = self._wb.add_worksheet("EA Definitions")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Extension Attribute Definitions"),
            "Source: jamf-cli pro computer-extension-attributes list"
            f" | Generated: {ts}",
            self._fmts,
            ncols=7,
        )
        headers = [
            "EA Name",
            "Definition ID",
            "Data Type",
            "Input Type",
            "Enabled",
            "Inventory Display",
            "Description",
        ]
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1
        ws.set_column(0, 0, 34)
        ws.set_column(1, 4, 16)
        ws.set_column(5, 5, 22)
        ws.set_column(6, 6, 60)

        for item in sorted(rows, key=lambda current: str(current.get("name", "")).lower()):
            enabled_raw = item.get("enabled")
            enabled_value = ""
            if enabled_raw is True:
                enabled_value = "Yes"
            elif enabled_raw is False:
                enabled_value = "No"
            _safe_write(ws, row, 0, item.get("name", ""), self._fmts["cell"])
            _safe_write(ws, row, 1, item.get("id", ""), self._fmts["cell"])
            _safe_write(ws, row, 2, item.get("dataType", ""), self._fmts["cell"])
            _safe_write(ws, row, 3, item.get("inputType", ""), self._fmts["cell"])
            _safe_write(ws, row, 4, enabled_value, self._fmts["cell"])
            _safe_write(ws, row, 5, item.get("inventoryDisplayType", ""), self._fmts["cell"])
            _safe_write(ws, row, 6, item.get("description", ""), self._fmts["cell"])
            row += 1

    def _write_software_installs(self) -> None:
        """Write installed software version distribution from jamf-cli."""
        raw = self._bridge.software_installs()
        rows = raw if isinstance(raw, list) else []
        if not rows:
            raise RuntimeError("jamf-cli software-installs returned no rows")

        ws = self._wb.add_worksheet("Software Installs")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Software Installs"),
            "Source: jamf-cli pro report software-installs"
            f" | Generated: {ts}",
            self._fmts,
            ncols=3,
        )
        title_count = len({str(item.get('title', '') or '').strip() for item in rows})
        _safe_write(ws, row, 0, "Distinct Titles", self._fmts["cell"])
        _safe_write(ws, row, 1, title_count, self._fmts["cell"])
        row += 1
        _safe_write(ws, row, 0, "Title-Version Rows", self._fmts["cell"])
        _safe_write(ws, row, 1, len(rows), self._fmts["cell"])
        row += 2

        headers = ["Application Title", "Version", "Device Count"]
        for col_i, header in enumerate(headers):
            _safe_write(ws, row, col_i, header, self._fmts["header"])
        row += 1
        ws.set_column(0, 0, 38)
        ws.set_column(1, 1, 20)
        ws.set_column(2, 2, 16)

        for item in sorted(
            rows,
            key=lambda current: (
                -_to_int(current.get("device_count", 0)),
                str(current.get("title", "")),
                str(current.get("version", "")),
            ),
        ):
            _safe_write(ws, row, 0, item.get("title", ""), self._fmts["cell"])
            _safe_write(ws, row, 1, item.get("version", ""), self._fmts["cell"])
            _safe_write(ws, row, 2, _to_int(item.get("device_count", 0)), self._fmts["cell"])
            row += 1

    def _write_policy(self) -> None:
        # jamf-cli pro report policy-status --output json returns:
        #   [{"summary":{total_policies,enabled,disabled,config_findings,warnings,info},
        #     "config_findings":[{severity,policy,policy_id,check,detail},...]}]
        raw = self._bridge.policy_status()
        envelope = _extract_envelope(raw)
        if not envelope:
            raise RuntimeError("policy-status returned no data")
        summary = envelope.get("summary", {}) if isinstance(envelope, dict) else {}
        findings = envelope.get("config_findings", []) if isinstance(envelope, dict) else []

        ws = self._wb.add_worksheet("Policy Health")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Policy Health"),
            f"Source: jamf-cli pro report policy-status | Generated: {ts}",
            self._fmts,
            ncols=4,
        )
        ws.set_column(0, 0, 40)
        ws.set_column(1, 3, 20)

        summary_fields = [
            ("total_policies", "Total Policies"),
            ("enabled", "Enabled"),
            ("disabled", "Disabled"),
            ("config_findings", "Config Findings"),
            ("warnings", "Warnings"),
            ("info", "Info"),
        ]
        for key, label in summary_fields:
            val = summary.get(key)
            if val is not None:
                _safe_write(ws, row, 0, label, self._fmts["cell"])
                _safe_write(ws, row, 1, val, self._fmts["cell"])
                row += 1

        if findings:
            row += 1
            headers = ["Severity", "Policy", "Check", "Detail"]
            for c, h in enumerate(headers):
                _safe_write(ws, row, c, h, self._fmts["header"])
            row += 1
            for finding in findings[:100]:
                sev = finding.get("severity", "")
                fmt = self._fmts["yellow"] if sev == "warning" else self._fmts["cell"]
                _safe_write(ws, row, 0, sev, fmt)
                _safe_write(ws, row, 1, finding.get("policy", ""), self._fmts["cell"])
                _safe_write(ws, row, 2, finding.get("check", ""), self._fmts["cell"])
                _safe_write(ws, row, 3, finding.get("detail", ""), self._fmts["cell"])
                row += 1

    def _write_profile_status(self) -> None:
        # jamf-cli pro report profile-status --output json returns:
        #   [{"summary": {"total_errors": N, "unique_profiles": N, "unique_devices": N, "days": N,
        #                 "devices_high_failure": N, "devices_high_pending": N},
        #     "failures": [{"device_type":"...","name":"...","id":"...","errors":N,
        #                   "devices":N,"last_error":"...","top_error":"..."},...],
        #     "device_failures": [...],
        #     "device_pending": [...]}]
        raw = self._bridge.profile_status()
        envelope = _extract_envelope(raw)
        if not envelope:
            raise RuntimeError("profile-status returned no data")
        summary = envelope.get("summary", {})
        failures = envelope.get("failures", [])
        device_failures = envelope.get("device_failures", [])

        ws = self._wb.add_worksheet("Profile Status")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Profile Status"),
            f"Source: jamf-cli pro report profile-status | Generated: {ts}",
            self._fmts,
            ncols=7,
        )
        ws.set_column(0, 0, 40)
        ws.set_column(1, 6, 18)

        summary_fields = [
            ("total_errors", "Total Errors"),
            ("unique_profiles", "Unique Profiles"),
            ("unique_devices", "Unique Devices"),
            ("days", "Days Lookback"),
            ("devices_high_failure", "Devices w/ High Failure"),
            ("devices_high_pending", "Devices w/ High Pending"),
        ]
        for key, label in summary_fields:
            val = summary.get(key)
            if val is not None:
                _safe_write(ws, row, 0, label, self._fmts["cell"])
                _safe_write(ws, row, 1, val, self._fmts["cell"])
                row += 1

        if failures:
            err_crit = int(self._config.thresholds.get("profile_error_critical", 50))
            err_warn = int(self._config.thresholds.get("profile_error_warning", 10))
            row += 1
            headers = ["Device Type", "Profile Name", "Profile ID", "Errors", "Devices",
                       "Last Error", "Top Error"]
            for c, h in enumerate(headers):
                _safe_write(ws, row, c, h, self._fmts["header"])
            row += 1
            for item in failures[:200]:
                errors = _to_int(item.get("errors", 0))
                fmt = self._severity_fmt(errors, err_warn, err_crit)
                _safe_write(ws, row, 0, item.get("device_type", ""), fmt)
                _safe_write(ws, row, 1, item.get("name", ""), fmt)
                _safe_write(ws, row, 2, item.get("id", ""), self._fmts["cell"])
                _safe_write(ws, row, 3, errors, fmt)
                _safe_write(ws, row, 4, _to_int(item.get("devices", 0)), fmt)
                _safe_write(ws, row, 5, item.get("last_error", ""), self._fmts["cell"])
                _safe_write(ws, row, 6, item.get("top_error", ""), self._fmts["cell"])
                row += 1

        if device_failures:
            row += 1
            _safe_write(ws, row, 0, "Devices with High Failure Count (>5 errors)",
                        self._fmts["header"])
            row += 1
            dev_headers = ["Name", "Serial", "Device Type", "OS Version", "Username", "Count"]
            for c, h in enumerate(dev_headers):
                _safe_write(ws, row, c, h, self._fmts["header"])
            row += 1
            for item in device_failures[:50]:
                _safe_write(ws, row, 0, item.get("name", ""), self._fmts["cell"])
                _safe_write(ws, row, 1, item.get("serial", ""), self._fmts["cell"])
                _safe_write(ws, row, 2, item.get("device_type", ""), self._fmts["cell"])
                _safe_write(ws, row, 3, item.get("os_version", ""), self._fmts["cell"])
                _safe_write(ws, row, 4, item.get("username", ""), self._fmts["cell"])
                _safe_write(ws, row, 5, _to_int(item.get("count", 0)), self._fmts["cell"])
                row += 1

    def _write_patch(self) -> None:
        # jamf-cli pro report patch-status --output json returns a flat list:
        #   [{"title":"Firefox","id":"123","on_latest":100,"on_other":20,
        #     "total":120,"latest":"130.0","compliance_pct":"83%"}, ...]
        # Newer builds may instead return:
        #   [{"title":"Firefox","id":"123","installed":100,
        #     "total":120,"latest":"130.0","compliance_pct":"83%"}, ...]
        raw = self._bridge.patch_status()
        titles = raw if isinstance(raw, list) else []
        uses_installed_shape = any("installed" in item for item in titles if isinstance(item, dict))

        ws = self._wb.add_worksheet("Patch Compliance")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Patch Compliance"),
            f"Source: jamf-cli pro report patch-status | Generated: {ts}",
            self._fmts,
            ncols=6,
        )
        ws.set_column(0, 0, 40)
        ws.set_column(1, 5, 18)
        headers = (
            ["Software Title", "Installed", "Not Installed", "Total", "Latest Version",
             "Compliance %"]
            if uses_installed_shape
            else ["Software Title", "On Latest", "On Other", "Total", "Latest Version",
                  "Compliance %"]
        )
        for c, h in enumerate(headers):
            _safe_write(ws, row, c, h, self._fmts["header"])
        row += 1
        for item in titles:
            total = _to_int(item.get("total", 0))
            if uses_installed_shape:
                primary = _to_int(item.get("installed", 0))
                secondary = max(total - primary, 0)
            else:
                primary = _to_int(item.get("on_latest", 0))
                secondary = _to_int(item.get("on_other", 0))
                if total == 0:
                    total = primary + secondary

            pct_raw = str(item.get("compliance_pct", "")).strip()
            pct_match = re.fullmatch(r"(\d+(?:\.\d+)?)%", pct_raw)
            pct_value = float(pct_match.group(1)) / 100 if pct_match else None
            _safe_write(ws, row, 0, item.get("title", ""), self._fmts["cell"])
            _safe_write(ws, row, 1, primary, self._fmts["cell"])
            _safe_write(ws, row, 2, secondary, self._fmts["cell"])
            _safe_write(ws, row, 3, total, self._fmts["cell"])
            _safe_write(ws, row, 4, item.get("latest", ""), self._fmts["cell"])
            if pct_value is not None:
                _safe_write(ws, row, 5, pct_value, _pct_format(self._fmts, pct_value))
            else:
                _safe_write(ws, row, 5, pct_raw or "N/A", self._fmts["cell"])
            row += 1

    def _write_patch_failures(self) -> None:
        # jamf-cli pro report patch-status --scan-failures --output json (v1.4.0+) returns:
        #   [{"policy":"Firefox 130.0","policy_id":"42","device":"MacBook-001",
        #     "device_id":"123","status_date":"2026-04-01","attempt":3,
        #     "last_action":"Retrying","serial":"ABC123",
        #     "os_version":"15.7.3","username":"jdoe"}, ...]
        # Each row is one failing device × one patch policy.
        raw = self._bridge.patch_device_failures()
        rows = raw if isinstance(raw, list) else []

        ws = self._wb.add_worksheet("Patch Failures")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Patch Failures"),
            f"Source: jamf-cli pro report patch-status --scan-failures | Generated: {ts}",
            self._fmts,
            ncols=8,
        )
        ws.set_column(0, 0, 30)  # Device
        ws.set_column(1, 1, 16)  # Serial
        ws.set_column(2, 2, 14)  # OS Version
        ws.set_column(3, 3, 22)  # Username
        ws.set_column(4, 4, 42)  # Policy
        ws.set_column(5, 5, 18)  # Status Date
        ws.set_column(6, 6, 10)  # Attempts
        ws.set_column(7, 7, 42)  # Last Action

        headers = [
            "Device", "Serial", "OS Version", "Username",
            "Policy", "Status Date", "Attempts", "Last Action",
        ]
        for c, h in enumerate(headers):
            _safe_write(ws, row, c, h, self._fmts["header"])
        row += 1

        if not rows:
            _safe_write(ws, row, 0, "No patch device failures found.", self._fmts["cell"])
            return

        for item in rows:
            _safe_write(ws, row, 0, item.get("device", ""), self._fmts["cell"])
            _safe_write(ws, row, 1, item.get("serial", ""), self._fmts["cell"])
            _safe_write(ws, row, 2, item.get("os_version", ""), self._fmts["cell"])
            _safe_write(ws, row, 3, item.get("username", ""), self._fmts["cell"])
            _safe_write(ws, row, 4, item.get("policy", ""), self._fmts["cell"])
            _safe_write(ws, row, 5, item.get("status_date", ""), self._fmts["cell"])
            _safe_write(ws, row, 6, _to_int(item.get("attempt", 0)), self._fmts["cell"])
            _safe_write(ws, row, 7, item.get("last_action", ""), self._fmts["cell"])
            row += 1

    def _write_app_status(self) -> None:
        # jamf-cli pro report app-status --output json shares the profile-status envelope:
        #   [{"summary": {"total_errors": N, "unique_profiles": N, "unique_devices": N,
        #                 "days": N, "devices_high_failure": N, "devices_high_pending": N},
        #     "failures": [{"device_type":"...","name":"...","id":"...","errors":N,
        #                   "devices":N,"last_error":"...","top_error":"..."},...],
        #     "device_failures": [...]}]
        raw = self._bridge.app_status()
        envelope = _extract_envelope(raw)
        if not envelope:
            raise RuntimeError("app-status returned no data")
        summary = envelope.get("summary", {})
        failures = envelope.get("failures", [])
        device_failures = envelope.get("device_failures", [])

        ws = self._wb.add_worksheet("App Status")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("App Status"),
            f"Source: jamf-cli pro report app-status | Generated: {ts}",
            self._fmts,
            ncols=7,
        )
        ws.set_column(0, 0, 40)
        ws.set_column(1, 6, 18)

        summary_fields = [
            ("total_errors", "Total Errors"),
            ("unique_profiles", "Unique Apps"),
            ("unique_devices", "Unique Devices"),
            ("days", "Days Lookback"),
            ("devices_high_failure", "Devices w/ High Failure"),
            ("devices_high_pending", "Devices w/ High Pending"),
        ]
        for key, label in summary_fields:
            val = summary.get(key)
            if val is not None:
                _safe_write(ws, row, 0, label, self._fmts["cell"])
                _safe_write(ws, row, 1, val, self._fmts["cell"])
                row += 1

        if failures:
            err_crit = int(self._config.thresholds.get("profile_error_critical", 50))
            err_warn = int(self._config.thresholds.get("profile_error_warning", 10))
            row += 1
            headers = ["Device Type", "App Name", "App ID", "Errors", "Devices",
                       "Last Error", "Top Error"]
            for c, h in enumerate(headers):
                _safe_write(ws, row, c, h, self._fmts["header"])
            row += 1
            for item in failures[:200]:
                errors = _to_int(item.get("errors", 0))
                fmt = self._severity_fmt(errors, err_warn, err_crit)
                _safe_write(ws, row, 0, item.get("device_type", ""), fmt)
                _safe_write(ws, row, 1, item.get("name", ""), fmt)
                _safe_write(ws, row, 2, item.get("id", ""), self._fmts["cell"])
                _safe_write(ws, row, 3, errors, fmt)
                _safe_write(ws, row, 4, _to_int(item.get("devices", 0)), fmt)
                _safe_write(ws, row, 5, item.get("last_error", ""), self._fmts["cell"])
                _safe_write(ws, row, 6, item.get("top_error", ""), self._fmts["cell"])
                row += 1

        if device_failures:
            row += 1
            _safe_write(ws, row, 0, "Devices with High Failure Count (>5 errors)",
                        self._fmts["header"])
            row += 1
            dev_headers = ["Name", "Serial", "Device Type", "OS Version", "Username", "Count"]
            for c, h in enumerate(dev_headers):
                _safe_write(ws, row, c, h, self._fmts["header"])
            row += 1
            for item in device_failures[:50]:
                _safe_write(ws, row, 0, item.get("name", ""), self._fmts["cell"])
                _safe_write(ws, row, 1, item.get("serial", ""), self._fmts["cell"])
                _safe_write(ws, row, 2, item.get("device_type", ""), self._fmts["cell"])
                _safe_write(ws, row, 3, item.get("os_version", ""), self._fmts["cell"])
                _safe_write(ws, row, 4, item.get("username", ""), self._fmts["cell"])
                _safe_write(ws, row, 5, _to_int(item.get("count", 0)), self._fmts["cell"])
                row += 1

    def _write_update_status(self) -> None:
        # jamf-cli pro report update-status --output json shape differs by version:
        #
        # v1.5 and earlier:
        #   {"summary": {"total_updates": N, "pending": N, "downloading": N,
        #                "installing": N, "installed": N, "errors": N},
        #    "ErrorDevices": [{"device_name":"...","serial":"...","os_version":"...",
        #                      "status":"...","product_key":"...","updated":"..."},...]}
        #
        # v1.6+:
        #   [{"total": N,
        #     "status_summary": [{"status": "...", "count": N}, ...],
        #     "plan_total": N,
        #     "plan_state_summary": [{"state": "...", "count": N}, ...]}]
        #   (error_devices and failed_plans only present with --scan-failures)
        #
        # Detect format by checking which summary key is present.
        raw = self._bridge.update_status()
        envelope = _extract_envelope(raw)
        if not envelope:
            raise RuntimeError("update-status returned no data")

        no_data_message = str(envelope.get("message", "") or "").strip()
        is_v16 = "status_summary" in envelope

        ws = self._wb.add_worksheet("Update Status")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Update Status"),
            f"Source: jamf-cli pro report update-status | Generated: {ts}",
            self._fmts,
            ncols=6,
        )
        ws.set_column(0, 0, 32)
        ws.set_column(1, 5, 20)

        if no_data_message and not is_v16 and not envelope.get("summary"):
            _safe_write(ws, row, 0, "Status", self._fmts["header"])
            _safe_write(ws, row, 1, "Details", self._fmts["header"])
            row += 1
            _safe_write(ws, row, 0, "No Data", self._fmts["yellow"])
            _safe_write(ws, row, 1, no_data_message, self._fmts["yellow"])
            return

        if is_v16:
            # v1.6+ format: status_summary list + plan_state_summary list
            status_summary = envelope.get("status_summary") or []
            plan_state_summary = envelope.get("plan_state_summary") or []
            total = _to_int(envelope.get("total", 0))
            plan_total = _to_int(envelope.get("plan_total", 0))

            _safe_write(ws, row, 0, f"Update Statuses ({total} total)", self._fmts["header"])
            _safe_write(ws, row, 1, "Count", self._fmts["header"])
            row += 1
            for item in sorted(status_summary, key=lambda x: -_to_int(x.get("count", 0))):
                _safe_write(ws, row, 0, item.get("status", ""), self._fmts["cell"])
                _safe_write(ws, row, 1, _to_int(item.get("count", 0)), self._fmts["cell"])
                row += 1

            if plan_state_summary:
                row += 1
                _safe_write(
                    ws, row, 0,
                    f"Update Plan States ({plan_total} total)",
                    self._fmts["header"],
                )
                _safe_write(ws, row, 1, "Count", self._fmts["header"])
                row += 1
                for item in sorted(
                    plan_state_summary, key=lambda x: -_to_int(x.get("count", 0))
                ):
                    _safe_write(ws, row, 0, item.get("state", ""), self._fmts["cell"])
                    _safe_write(ws, row, 1, _to_int(item.get("count", 0)), self._fmts["cell"])
                    row += 1
        else:
            # v1.5 and earlier: summary dict + ErrorDevices list
            summary = envelope.get("summary", {})
            error_devices = envelope.get("ErrorDevices", [])
            summary_fields = [
                ("total_updates", "Total Updates"),
                ("pending", "Pending"),
                ("downloading", "Downloading"),
                ("installing", "Installing"),
                ("installed", "Installed"),
                ("errors", "Errors"),
            ]
            for key, label in summary_fields:
                val = summary.get(key)
                if val is not None:
                    _safe_write(ws, row, 0, label, self._fmts["cell"])
                    _safe_write(ws, row, 1, val, self._fmts["cell"])
                    row += 1

            if error_devices:
                row += 1
                _safe_write(ws, row, 0, "Devices with Update Errors", self._fmts["header"])
                row += 1
                dev_headers = ["Device Name", "Serial", "OS Version", "Status",
                               "Product Key", "Updated"]
                for c, h in enumerate(dev_headers):
                    _safe_write(ws, row, c, h, self._fmts["header"])
                row += 1
                for item in error_devices[:200]:
                    _safe_write(ws, row, 0, item.get("device_name", ""), self._fmts["cell"])
                    _safe_write(ws, row, 1, item.get("serial", ""), self._fmts["cell"])
                    _safe_write(ws, row, 2, item.get("os_version", ""), self._fmts["cell"])
                    _safe_write(ws, row, 3, item.get("status", ""), self._fmts["cell"])
                    _safe_write(ws, row, 4, item.get("product_key", ""), self._fmts["cell"])
                    _safe_write(ws, row, 5, item.get("updated", ""), self._fmts["cell"])
                    row += 1

    def _write_update_failures(self) -> None:
        # jamf-cli pro report update-status --scan-failures --output json (v1.6+) returns:
        #   [{"total": N, "status_summary": [...],
        #     "error_devices": [{"name":"...","serial":"...","device_type":"...",
        #                        "os_version":"...","username":"...","status":"...",
        #                        "product_key":"...","updated":"..."}],
        #     "plan_total": N, "plan_state_summary": [...],
        #     "failed_plans": [{"name":"...","serial":"...","device_type":"...",
        #                       "os_version":"...","username":"...","state":"...",
        #                       "action":"...","version":"...","error":"...",
        #                       "last_event":"..."}]}]
        raw = self._bridge.update_device_failures()
        envelope = _extract_envelope(raw)
        if not envelope:
            raise RuntimeError("update-status --scan-failures returned no data")

        no_data_message = str(envelope.get("message", "") or "").strip()
        error_devices = envelope.get("error_devices") or []
        failed_plans = envelope.get("failed_plans") or []

        ws = self._wb.add_worksheet("Update Failures")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            self._t("Update Failures"),
            (
                "Source: jamf-cli pro report update-status --scan-failures"
                f" (v1.6.0+) | Generated: {ts}"
            ),
            self._fmts,
            ncols=8,
        )

        if no_data_message and not error_devices and not failed_plans:
            _safe_write(ws, row, 0, "No Data", self._fmts["yellow"])
            _safe_write(ws, row, 1, no_data_message, self._fmts["yellow"])
            return

        # --- Error Devices table ---
        ws.set_column(0, 0, 30)  # Device
        ws.set_column(1, 1, 16)  # Serial
        ws.set_column(2, 2, 14)  # Device Type
        ws.set_column(3, 3, 14)  # OS Version
        ws.set_column(4, 4, 22)  # Username
        ws.set_column(5, 5, 28)  # Status
        ws.set_column(6, 6, 24)  # Product Key
        ws.set_column(7, 7, 18)  # Updated

        _safe_write(
            ws, row, 0,
            f"Update Error Devices ({len(error_devices)})",
            self._fmts["header"],
        )
        row += 1
        if error_devices:
            err_headers = [
                "Device", "Serial", "Device Type", "OS Version",
                "Username", "Status", "Product Key", "Updated",
            ]
            for c, h in enumerate(err_headers):
                _safe_write(ws, row, c, h, self._fmts["header"])
            row += 1
            for item in error_devices:
                _safe_write(ws, row, 0, item.get("name", ""), self._fmts["cell"])
                _safe_write(ws, row, 1, item.get("serial", ""), self._fmts["cell"])
                _safe_write(ws, row, 2, item.get("device_type", ""), self._fmts["cell"])
                _safe_write(ws, row, 3, item.get("os_version", ""), self._fmts["cell"])
                _safe_write(ws, row, 4, item.get("username", ""), self._fmts["cell"])
                _safe_write(ws, row, 5, item.get("status", ""), self._fmts["red"])
                _safe_write(ws, row, 6, item.get("product_key", ""), self._fmts["cell"])
                _safe_write(ws, row, 7, item.get("updated", ""), self._fmts["cell"])
                row += 1
        else:
            _safe_write(ws, row, 0, "No update error devices found.", self._fmts["cell"])
            row += 1

        # --- Failed Plans table ---
        row += 1
        ws.set_column(5, 5, 16)  # Plan State (reuse col 5)
        ws.set_column(6, 6, 18)  # Action
        ws.set_column(7, 7, 14)  # Version

        _safe_write(
            ws, row, 0,
            f"Failed Update Plans ({len(failed_plans)})",
            self._fmts["header"],
        )
        row += 1
        if failed_plans:
            plan_headers = [
                "Device", "Serial", "Device Type", "OS Version",
                "Username", "Plan State", "Action", "Version",
                "Error Reasons", "Last Event",
            ]
            # Expand column count for extra columns
            ws.set_column(8, 8, 42)   # Error Reasons
            ws.set_column(9, 9, 28)   # Last Event
            for c, h in enumerate(plan_headers):
                _safe_write(ws, row, c, h, self._fmts["header"])
            row += 1
            for item in failed_plans:
                state = item.get("state", "")
                row_fmt = self._fmts["red"] if state == "PlanFailed" else self._fmts["yellow"]
                _safe_write(ws, row, 0, item.get("name", ""), self._fmts["cell"])
                _safe_write(ws, row, 1, item.get("serial", ""), self._fmts["cell"])
                _safe_write(ws, row, 2, item.get("device_type", ""), self._fmts["cell"])
                _safe_write(ws, row, 3, item.get("os_version", ""), self._fmts["cell"])
                _safe_write(ws, row, 4, item.get("username", ""), self._fmts["cell"])
                _safe_write(ws, row, 5, state, row_fmt)
                _safe_write(ws, row, 6, item.get("action", ""), self._fmts["cell"])
                _safe_write(ws, row, 7, item.get("version", ""), self._fmts["cell"])
                _safe_write(ws, row, 8, item.get("error", ""), self._fmts["cell"])
                _safe_write(ws, row, 9, item.get("last_event", ""), self._fmts["cell"])
                row += 1
        else:
            _safe_write(ws, row, 0, "No failed update plans found.", self._fmts["cell"])

    def _write_smart_groups(self) -> None:
        """Write a Smart Groups sheet from jamf-cli pro groups data.

        Columns: Group Name | Type | Smart | Member Count | Delta | Prior Count |
                 Scope Warning

        Delta = member_count - prior_count, derived from comparing the current
        groups JSON against the most-recent cached groups snapshot in the bridge's
        data_dir. A zero-member smart group that was non-zero in the prior run is
        highlighted in red (potential scope failure).

        Requires: JamfCLIBridge.groups() to be implemented (pending fixture
        validation). The sheet is automatically skipped via write_all's RuntimeError
        guard when the method raises NotImplementedError.

        JSON shape expected from jamf-cli pro groups --output json:
            [
              {
                "id": "123",
                "name": "All Managed Macs",
                "type": "computer",          # "computer" | "mobile_device"
                "is_smart": true,
                "member_count": 142
              },
              ...
            ]

        TODO: Implement after:
            1. Running `jamf-cli pro groups --output json` against a test instance
            2. Committing the result to tests/fixtures/jamf-cli-data/groups/groups.json
            3. Confirming the JSON shape matches the expected structure above
            4. Implementing JamfCLIBridge.groups()
            5. Adding a delta-comparison against the prior cached JSON

        Raises:
            RuntimeError: When groups data is unavailable or the bridge raises.
        """
        # Propagate NotImplementedError as RuntimeError so write_all skips gracefully.
        try:
            raw = self._bridge.groups()
        except NotImplementedError as exc:
            raise RuntimeError(str(exc)) from exc

        if not raw or not isinstance(raw, list):
            raise RuntimeError("groups returned no data")

        ws = self._wb.add_worksheet("Smart Groups")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row = _write_sheet_header(
            ws,
            "Smart Groups",
            f"Source: jamf-cli pro groups | Generated: {ts}",
            self._fmts,
            ncols=7,
        )
        headers = [
            "Group Name", "Type", "Smart Group", "Member Count",
            "Delta", "Prior Count", "Note",
        ]
        for col_i, h in enumerate(headers):
            _safe_write(ws, row, col_i, h, self._fmts["header"])
        row += 1

        # TODO: load prior cached groups JSON from data_dir for delta computation.
        # For now, delta and prior_count are blank.
        for item in sorted(raw, key=lambda g: g.get("name", "").casefold()):
            count = item.get("member_count", "")
            is_smart = item.get("is_smart", False)
            is_zero = isinstance(count, int) and count == 0
            row_fmt = self._fmts.get("row_red") or self._fmts["cell"] if is_zero else self._fmts["cell"]
            _safe_write(ws, row, 0, item.get("name", ""), row_fmt)
            _safe_write(ws, row, 1, item.get("type", ""), row_fmt)
            _safe_write(ws, row, 2, "Yes" if is_smart else "No", row_fmt)
            _safe_write(ws, row, 3, count, row_fmt)
            _safe_write(ws, row, 4, "", row_fmt)       # Delta — TODO
            _safe_write(ws, row, 5, "", row_fmt)       # Prior count — TODO
            _safe_write(ws, row, 6, "Zero members" if is_zero else "", row_fmt)
            row += 1

        ws.set_column(0, 0, 40)
        ws.set_column(1, 6, 16)


# ---------------------------------------------------------------------------
# CSVDashboard
# ---------------------------------------------------------------------------


class CSVDashboard:
    """Generates Excel sheets from a Jamf Pro CSV inventory export.

    Args:
        config: Loaded Config instance.
        csv_path: Path to the primary CSV file.
        workbook: Active xlsxwriter Workbook.
        fmts: Format dict from _build_formats.
        family_name: Report family driving CSV sheet generation.
        extra_csv_paths: Optional list of additional CSV paths to merge into the
            primary DataFrame.  Each extra CSV is loaded and concatenated; a
            "CSV Source" column is added to every row so downstream sheets can
            distinguish which export each record came from.  Missing columns are
            filled with empty strings.  Column names must match across files for
            the config mapping to work correctly — the primary CSV's schema is
            authoritative.
    """

    def __init__(
        self,
        config: Config,
        csv_path: str,
        workbook: xlsxwriter.Workbook,
        fmts: dict,
        family_name: str = "computers",
        extra_csv_paths: Optional[list[str]] = None,
    ) -> None:
        self._config = config
        self._family_name = family_name
        mapper_section = "mobile_columns" if family_name == "mobile" else "columns"
        self._mapper = ColumnMapper(config, mapper_section)
        self._wb = workbook
        self._fmts = fmts
        try:
            primary = pd.read_csv(csv_path, dtype=str, encoding="utf-8-sig").fillna("")
        except Exception as exc:
            raise SystemExit(
                f"Error: could not read CSV '{csv_path}': {exc}"
            ) from exc

        if extra_csv_paths:
            frames: list[Any] = []
            primary["CSV Source"] = Path(csv_path).name
            frames.append(primary)
            for extra_path in extra_csv_paths:
                try:
                    extra_df = pd.read_csv(
                        extra_path, dtype=str, encoding="utf-8-sig"
                    ).fillna("")
                    extra_df["CSV Source"] = Path(extra_path).name
                    frames.append(extra_df)
                    print(f"  Merged extra CSV: {Path(extra_path).name} ({len(extra_df)} rows)")
                except Exception as exc:
                    print(f"  [warn] Could not read extra CSV '{extra_path}': {exc} — skipping")
            all_col_sets = [set(f.columns) for f in frames]
            union_cols = set.union(*all_col_sets)
            for i, cols in enumerate(all_col_sets):
                missing = union_cols - cols
                if missing:
                    names = ", ".join(sorted(missing)[:5])
                    suffix = " ..." if len(missing) > 5 else ""
                    print(
                        f"  [warn] CSV {i + 1} is missing {len(missing)} column(s) present in"
                        f" other CSVs — those cells will be empty: {names}{suffix}"
                    )
            self._df = pd.concat(frames, ignore_index=True).fillna("")
            print(
                f"  Loaded {len(frames)} CSV(s): {len(self._df)} rows total,"
                f" {len(self._df.columns)} columns"
            )
        else:
            self._df = primary
            print(f"  Loaded CSV: {len(self._df)} rows, {len(self._df.columns)} columns")

    @property
    def _org_name(self) -> str:
        """Return the configured org name, or empty string."""
        return (self._config.get("branding", "org_name") or "").strip()

    def _t(self, base: str) -> str:
        """Return sheet title prefixed with org name when configured."""
        return _org_title(self._org_name, base)

    def write_all(self) -> list[str]:
        """Write all CSV-derived sheets. Returns list of sheet names written."""
        written = []
        if self._family_name == "mobile":
            sheets = [
                ("Mobile Device Inventory", self._write_mobile_inventory_csv),
                ("Mobile Stale Devices", self._write_mobile_stale),
            ]
        else:
            sheets = [
                ("Device Inventory", self._write_inventory),
                ("Stale Devices", self._write_stale),
                ("Security Controls", self._write_security_controls),
                ("Security Agents", self._write_security_agents),
                ("Compliance", self._write_compliance),
            ]
        for name, fn in sheets:
            try:
                fn()
                written.append(name)
                print(f"  [ok] {name}")
            except (KeyError, ValueError, RuntimeError) as exc:
                print(f"  [skip] {name}: {exc}")
            except Exception as exc:
                print(f"  [skip] {name}: unexpected error — {type(exc).__name__}: {exc}")

        for ea in self._config.custom_eas:
            ea_name = ea.get("name", "Custom EA")
            try:
                self._write_custom_ea(ea)
                written.append(ea_name)
                print(f"  [ok] {ea_name}")
            except (KeyError, ValueError, RuntimeError) as exc:
                print(f"  [skip] {ea_name}: {exc}")
            except Exception as exc:
                print(f"  [skip] {ea_name}: unexpected error — {type(exc).__name__}: {exc}")
        return written

    def _col(self, logical: str) -> Optional[str]:
        return self._mapper.get(logical)

    def _get(self, row: Any, logical: str) -> str:
        return self._mapper.extract(row, logical)

    def _device_name(self, row: Any) -> str:
        """Extract a device name from a DataFrame row using the configured column.

        Args:
            row: A pandas Series (DataFrame row) from self._df.iterrows().

        Returns:
            The device name string, or empty string if the column is not configured
            or not present in this row.
        """
        logical = "device_name" if self._family_name == "mobile" else "computer_name"
        col = self._col(logical)
        return str(row[col]) if col and col in row.index else ""

    def _write_inventory(self) -> None:
        stale_days = int(self._config.thresholds.get("stale_device_days", 30))
        name_col = self._col("computer_name")
        checkin_col = self._col("last_checkin")
        if not name_col:
            raise RuntimeError("computer_name column not configured")
        active_rows = []
        for _, row in self._df.iterrows():
            checkin = str(row.get(checkin_col, "")) if checkin_col else ""
            days = _days_since(checkin) if checkin else None
            if days is not None and 0 <= days <= stale_days:
                active_rows.append(row)

        ws = self._wb.add_worksheet("Device Inventory")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row_i = _write_sheet_header(
            ws,
            self._t("Device Inventory"),
            f"Active devices (checked in within {stale_days} days) | Generated: {ts}",
            self._fmts,
            ncols=6,
        )
        logical_cols = ["computer_name", "serial_number", "operating_system",
                        "last_checkin", "department", "model"]
        headers = [lc.replace("_", " ").title() for lc in logical_cols]
        col_names = [self._col(lc) for lc in logical_cols]
        for c, h in enumerate(headers):
            _safe_write(ws, row_i, c, h, self._fmts["header"])
        row_i += 1
        for row in active_rows:
            for c, cn in enumerate(col_names):
                val = str(row[cn]) if cn and cn in row.index else ""
                _safe_write(ws, row_i, c, val, self._fmts["cell"])
            row_i += 1
        ws.set_column(0, 0, 30)
        ws.set_column(1, 5, 22)

    def _write_stale(self) -> None:
        stale_days = int(self._config.thresholds.get("stale_device_days", 30))
        name_col = self._col("computer_name")
        checkin_col = self._col("last_checkin")
        if not name_col or not checkin_col:
            raise RuntimeError("computer_name or last_checkin column not configured")
        ws = self._wb.add_worksheet("Stale Devices")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row_i = _write_sheet_header(
            ws,
            self._t("Stale Devices"),
            f"Devices not checked in within {stale_days} days | Generated: {ts}",
            self._fmts,
            ncols=6,
        )
        headers = ["Computer Name", "Serial Number", "Days Stale", "OS", "Department", "Manager"]
        for c, h in enumerate(headers):
            _safe_write(ws, row_i, c, h, self._fmts["header"])
        row_i += 1
        for _, row in self._df.iterrows():
            checkin = str(row.get(checkin_col, ""))
            days = _days_since(checkin)
            if days is not None and days > stale_days:
                manager_raw = self._get(row, "manager")
                _safe_write(ws, row_i, 0, self._get(row, "computer_name"), self._fmts["cell"])
                _safe_write(ws, row_i, 1, self._get(row, "serial_number"), self._fmts["cell"])
                _safe_write(ws, row_i, 2, days, self._fmts["int"])
                _safe_write(ws, row_i, 3, self._get(row, "operating_system"), self._fmts["cell"])
                _safe_write(ws, row_i, 4, self._get(row, "department"), self._fmts["cell"])
                _safe_write(ws, row_i, 5, _parse_manager(manager_raw), self._fmts["cell"])
                row_i += 1
        ws.set_column(0, 0, 30)
        ws.set_column(1, 5, 22)

    def _write_mobile_inventory_csv(self) -> None:
        """Write a mobile CSV inventory sheet using mobile_columns mappings."""
        stale_days = int(self._config.thresholds.get("stale_device_days", 30))
        name_col = self._col("device_name")
        checkin_col = self._col("last_checkin")
        if not name_col:
            raise RuntimeError("mobile_columns.device_name is not configured")
        active_rows = []
        for _, row in self._df.iterrows():
            checkin = str(row.get(checkin_col, "")) if checkin_col else ""
            days = _days_since(checkin) if checkin else None
            if days is not None and 0 <= days <= stale_days:
                active_rows.append(row)

        ws = self._wb.add_worksheet("Mobile Device Inventory")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row_i = _write_sheet_header(
            ws,
            "Mobile Device Inventory",
            f"Active devices (checked in within {stale_days} days) | Generated: {ts}",
            self._fmts,
            ncols=9,
        )
        logical_cols = [
            "device_name",
            "serial_number",
            "operating_system",
            "last_checkin",
            "email",
            "model",
            "device_family",
            "managed",
            "supervised",
        ]
        headers = [
            "Device Name",
            "Serial Number",
            "OS Version",
            "Last Inventory Update",
            "Email Address",
            "Model",
            "Device Family",
            "Managed",
            "Supervised",
        ]
        col_names = [self._col(logical) for logical in logical_cols]
        for col_i, header in enumerate(headers):
            _safe_write(ws, row_i, col_i, header, self._fmts["header"])
        row_i += 1
        for row in active_rows:
            for col_i, col_name in enumerate(col_names):
                value = str(row[col_name]) if col_name and col_name in row.index else ""
                _safe_write(ws, row_i, col_i, value, self._fmts["cell"])
            row_i += 1
        ws.set_column(0, 0, 30)
        ws.set_column(1, 8, 22)

    def _write_mobile_stale(self) -> None:
        """Write stale mobile devices using mobile_columns mappings."""
        stale_days = int(self._config.thresholds.get("stale_device_days", 30))
        name_col = self._col("device_name")
        checkin_col = self._col("last_checkin")
        if not name_col or not checkin_col:
            raise RuntimeError(
                "mobile_columns.device_name or mobile_columns.last_checkin is not configured"
            )
        ws = self._wb.add_worksheet("Mobile Stale Devices")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row_i = _write_sheet_header(
            ws,
            "Mobile Stale Devices",
            f"Devices not checked in within {stale_days} days | Generated: {ts}",
            self._fmts,
            ncols=9,
        )
        headers = [
            "Device Name",
            "Serial Number",
            "Days Stale",
            "OS Version",
            "Email Address",
            "Model",
            "Device Family",
            "Managed",
            "Supervised",
        ]
        for col_i, header in enumerate(headers):
            _safe_write(ws, row_i, col_i, header, self._fmts["header"])
        row_i += 1
        for _, row in self._df.iterrows():
            checkin = str(row.get(checkin_col, ""))
            days = _days_since(checkin)
            if days is None or days <= stale_days:
                continue
            _safe_write(ws, row_i, 0, self._get(row, "device_name"), self._fmts["cell"])
            _safe_write(ws, row_i, 1, self._get(row, "serial_number"), self._fmts["cell"])
            _safe_write(ws, row_i, 2, days, self._fmts["int"])
            _safe_write(ws, row_i, 3, self._get(row, "operating_system"), self._fmts["cell"])
            _safe_write(ws, row_i, 4, self._get(row, "email"), self._fmts["cell"])
            _safe_write(ws, row_i, 5, self._get(row, "model"), self._fmts["cell"])
            _safe_write(ws, row_i, 6, self._get(row, "device_family"), self._fmts["cell"])
            _safe_write(ws, row_i, 7, self._get(row, "managed"), self._fmts["cell"])
            _safe_write(ws, row_i, 8, self._get(row, "supervised"), self._fmts["cell"])
            row_i += 1
        ws.set_column(0, 0, 30)
        ws.set_column(1, 8, 22)

    def _write_security_controls(self) -> None:
        control_fields = [
            "filevault",
            "sip",
            "firewall",
            "gatekeeper",
            "secure_boot",
            "bootstrap_token",
        ]
        configured = [(f, self._col(f)) for f in control_fields if self._col(f)]
        if not configured:
            raise RuntimeError("no security control columns configured")
        ws = self._wb.add_worksheet("Security Controls")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row_i = _write_sheet_header(
            ws, self._t("Security Controls"), f"Generated: {ts}", self._fmts, ncols=4
        )
        headers = ["Control", "Enabled/Compliant", "Not Compliant", "Compliance %"]
        for c, h in enumerate(headers):
            _safe_write(ws, row_i, c, h, self._fmts["header"])
        row_i += 1
        for logical, col in configured:
            total = len(self._df)
            compliant = self._df[col].apply(
                lambda value: _security_control_is_compliant(logical, value)
            ).sum()
            non_compliant = total - compliant
            pct = compliant / total if total > 0 else 0.0
            label = logical.replace("_", " ").title()
            _safe_write(ws, row_i, 0, label, self._fmts["cell"])
            _safe_write(ws, row_i, 1, int(compliant), self._fmts["cell"])
            _safe_write(ws, row_i, 2, int(non_compliant), self._fmts["cell"])
            _safe_write(ws, row_i, 3, pct, _pct_format(self._fmts, pct))
            row_i += 1
        ws.set_column(0, 0, 30)
        ws.set_column(1, 3, 22)

    def _write_security_agents(self) -> None:
        agents = self._config.security_agents
        if not agents:
            raise RuntimeError("no security agents configured")
        ws = self._wb.add_worksheet("Security Agents")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row_i = _write_sheet_header(
            ws, self._t("Security Agent Status"), f"Generated: {ts}", self._fmts, ncols=5
        )
        for agent in agents:
            col = agent.get("column", "")
            connected_val = _normalized_text(agent.get("connected_value", "connected"))
            if col not in self._df.columns:
                continue
            if not connected_val:
                print(
                    f"  [warn] Security agent {agent.get('name', col)!r}:"
                    " connected_value is empty — all non-empty cells will be counted as"
                    " connected. Set a specific value such as 'Installed' or 'Running'."
                )
                continue
            statuses = self._df[col].fillna("").astype(str)
            connected_mask = _contains_case_insensitive(statuses, connected_val)
            total = len(self._df)
            connected = int(connected_mask.sum())
            disconnected = total - connected
            pct = connected / total if total > 0 else 0.0
            _safe_write(ws, row_i, 0, agent.get("name", col), self._fmts["header"])
            row_i += 1
            for lbl, val in [("Connected", connected), ("Not Connected", disconnected),
                             ("Total", total)]:
                _safe_write(ws, row_i, 0, lbl, self._fmts["cell"])
                _safe_write(ws, row_i, 1, int(val), self._fmts["cell"])
                row_i += 1
            _safe_write(ws, row_i, 0, "Compliance %", self._fmts["cell"])
            _safe_write(ws, row_i, 1, pct, _pct_format(self._fmts, pct))
            row_i += 2
            non_connected_df = self._df[~connected_mask]
            if not non_connected_df.empty:
                _safe_write(ws, row_i, 0, "Non-Connected Devices", self._fmts["header"])
                row_i += 1
                for _, dr in non_connected_df.iterrows():
                    name = self._device_name(dr)
                    status = str(dr[col])
                    _safe_write(ws, row_i, 0, name, self._fmts["cell"])
                    _safe_write(ws, row_i, 1, status, self._fmts["cell"])
                    row_i += 1
            row_i += 1
        ws.set_column(0, 0, 35)
        ws.set_column(1, 1, 25)

    def _write_compliance(self) -> None:
        comp = self._config.compliance
        if not comp.get("enabled"):
            raise RuntimeError("compliance not enabled in config")
        count_col = comp.get("failures_count_column", "")
        list_col = comp.get("failures_list_column", "")
        label = _compliance_label(comp)
        if count_col and count_col not in self._df.columns:
            raise RuntimeError(f"failures_count_column '{count_col}' not in CSV")
        ws = self._wb.add_worksheet("Compliance")
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row_i = _write_sheet_header(
            ws, self._t(label), f"Generated: {ts}", self._fmts, ncols=4
        )
        if count_col:
            counts = pd.to_numeric(self._df[count_col], errors="coerce").fillna(0)
            compliant = int((counts == 0).sum())
            non_compliant = int((counts > 0).sum())
            total = len(self._df)
            pct = compliant / total if total > 0 else 0.0
            for lbl, val in [("Fully Compliant", compliant), ("Has Failures", non_compliant),
                             ("Total Devices", total)]:
                _safe_write(ws, row_i, 0, lbl, self._fmts["cell"])
                _safe_write(ws, row_i, 1, val, self._fmts["cell"])
                row_i += 1
            _safe_write(ws, row_i, 0, "Compliance %", self._fmts["cell"])
            _safe_write(ws, row_i, 1, pct, _pct_format(self._fmts, pct))
            row_i += 2

        if list_col and list_col in self._df.columns:
            rule_counts: dict[str, int] = {}
            for val in self._df[list_col]:
                for rule in _split_multi_value_cell(val):
                    rule_counts[rule] = rule_counts.get(rule, 0) + 1
            _safe_write(ws, row_i, 0, "Top Failing Rules", self._fmts["header"])
            _safe_write(ws, row_i, 1, "Device Count", self._fmts["header"])
            row_i += 1
            for rule, cnt in sorted(rule_counts.items(), key=lambda x: -x[1])[:30]:
                _safe_write(ws, row_i, 0, rule, self._fmts["cell"])
                _safe_write(ws, row_i, 1, cnt, self._fmts["cell"])
                row_i += 1
        ws.set_column(0, 0, 50)
        ws.set_column(1, 1, 20)

    def _write_custom_ea(self, ea: dict) -> None:
        """Write a single custom EA sheet based on its type configuration.

        Args:
            ea: Dict with keys: name, column, type, and type-specific options.
        """
        col = ea.get("column", "")
        ea_type = ea.get("type", "text")
        name = ea.get("name", col)
        if col not in self._df.columns:
            raise RuntimeError(f"column '{col}' not found in CSV")

        sheet_name = _excel_sheet_name(name, "", existing_names=set(self._wb.sheetnames))
        ws = self._wb.add_worksheet(sheet_name)
        ts = datetime.now().strftime("%Y-%m-%d %H:%M")
        row_i = _write_sheet_header(
            ws, name, f"Type: {ea_type} | Generated: {ts}", self._fmts, ncols=4
        )

        dispatch = {
            "boolean": self._ea_boolean,
            "percentage": self._ea_percentage,
            "version": self._ea_version,
            "text": self._ea_text,
            "date": self._ea_date,
        }
        handler = dispatch.get(ea_type, self._ea_text)
        handler(ws, row_i, col, ea)

    def _ea_boolean(self, ws: Any, row_i: int, col: str, ea: dict) -> None:
        true_val = ea.get("true_value", "Yes").lower()
        series = self._df[col].str.strip().str.lower()
        non_blank = series[series != ""]
        passed = int((non_blank == true_val).sum())
        unknown = int((series == "").sum())
        failed = len(non_blank) - passed
        reported = passed + failed
        pct = passed / reported if reported > 0 else 0.0
        for lbl, val in [("Pass", passed), ("Fail", failed), ("Total (reported)", reported)]:
            _safe_write(ws, row_i, 0, lbl, self._fmts["cell"])
            _safe_write(ws, row_i, 1, val, self._fmts["cell"])
            row_i += 1
        if unknown > 0:
            _safe_write(ws, row_i, 0, "Unknown / Not Reported", self._fmts["yellow"])
            _safe_write(ws, row_i, 1, unknown, self._fmts["yellow"])
            row_i += 1
        _safe_write(ws, row_i, 0, "Compliance %", self._fmts["cell"])
        _safe_write(ws, row_i, 1, pct, _pct_format(self._fmts, pct))
        row_i += 2
        failed_df = self._df[
            (self._df[col].str.strip().str.lower() != true_val)
            & (self._df[col].str.strip() != "")
        ]
        if not failed_df.empty:
            _safe_write(ws, row_i, 0, "Failing Devices", self._fmts["header"])
            _safe_write(ws, row_i, 1, "Value", self._fmts["header"])
            row_i += 1
            for _, dr in failed_df.iterrows():
                _safe_write(ws, row_i, 0, self._device_name(dr), self._fmts["cell"])
                _safe_write(ws, row_i, 1, str(dr[col]), self._fmts["cell"])
                row_i += 1

    def _ea_percentage(self, ws: Any, row_i: int, col: str, ea: dict) -> None:
        warn = float(
            ea.get("warning_threshold", self._config.thresholds.get("warning_disk_percent", 80))
        )
        crit = float(
            ea.get("critical_threshold", self._config.thresholds.get("critical_disk_percent", 90))
        )
        nums = pd.to_numeric(self._df[col].str.replace("%", "", regex=False), errors="coerce")
        critical_df = self._df[nums >= crit]
        warning_df = self._df[(nums >= warn) & (nums < crit)]
        for label, sub_df in [("Critical (>= {:.0f}%)".format(crit), critical_df),
                               ("Warning ({:.0f}% - {:.0f}%)".format(warn, crit), warning_df)]:
            _safe_write(ws, row_i, 0, label, self._fmts["header"])
            _safe_write(ws, row_i, 1, "Value", self._fmts["header"])
            row_i += 1
            for _, dr in sub_df.iterrows():
                _safe_write(ws, row_i, 0, self._device_name(dr), self._fmts["cell"])
                _safe_write(ws, row_i, 1, str(dr[col]), self._fmts["cell"])
                row_i += 1
            row_i += 1

    def _ea_version(self, ws: Any, row_i: int, col: str, ea: dict) -> None:
        current = [str(v).lower() for v in ea.get("current_versions", [])]
        dist = self._df[col].value_counts().to_dict()
        _safe_write(ws, row_i, 0, "Version", self._fmts["header"])
        _safe_write(ws, row_i, 1, "Count", self._fmts["header"])
        _safe_write(ws, row_i, 2, "Current?", self._fmts["header"])
        row_i += 1
        has_current_list = bool(current)
        for ver, cnt in sorted(dist.items(), key=lambda x: -x[1]):
            display_ver = _display_value(ver)
            is_current = (
                any(str(ver).lower().startswith(cv.lower()) for cv in current)
                if has_current_list else None
            )
            if not str(ver).strip():
                fmt = self._fmts["yellow"]
            elif is_current is True:
                fmt = self._fmts["green"]
            elif is_current is False:
                fmt = self._fmts["red"]
            else:
                fmt = self._fmts["cell"]
            _safe_write(ws, row_i, 0, display_ver, fmt)
            _safe_write(ws, row_i, 1, cnt, fmt)
            if is_current is True:
                current_label = "Yes"
            elif is_current is False:
                current_label = "No"
            else:
                current_label = ""
            _safe_write(ws, row_i, 2, current_label, fmt)
            row_i += 1

    def _ea_text(self, ws: Any, row_i: int, col: str, ea: dict) -> None:
        dist = self._df[col].value_counts().to_dict()
        _safe_write(ws, row_i, 0, "Value", self._fmts["header"])
        _safe_write(ws, row_i, 1, "Count", self._fmts["header"])
        row_i += 1
        for val, cnt in sorted(dist.items(), key=lambda x: -x[1]):
            fmt = self._fmts["yellow"] if not str(val).strip() else self._fmts["cell"]
            _safe_write(ws, row_i, 0, _display_value(val), fmt)
            _safe_write(ws, row_i, 1, cnt, fmt)
            row_i += 1

    def _ea_date(self, ws: Any, row_i: int, col: str, ea: dict) -> None:
        ea_name = ea.get("name", col)
        warn_days = int(ea.get("warning_days",
                                self._config.thresholds.get("cert_warning_days", 90)))
        _safe_write(ws, row_i, 0, "Device", self._fmts["header"])
        _safe_write(ws, row_i, 1, "Date Value", self._fmts["header"])
        _safe_write(ws, row_i, 2, "Days Until Expiry", self._fmts["header"])
        row_i += 1
        now = datetime.now(timezone.utc)
        rows_parsed = 0
        for _, dr in self._df.iterrows():
            raw = str(dr[col])
            parsed = pd.to_datetime(raw, errors="coerce", utc=True)
            if pd.isnull(parsed):
                continue
            rows_parsed += 1
            delta = parsed - now
            days_until = int(delta.total_seconds() / 86400)
            fmt = (self._fmts["red"] if days_until < 0 else
                   self._fmts["yellow"] if days_until < warn_days else
                   self._fmts["green"])
            _safe_write(ws, row_i, 0, self._device_name(dr), fmt)
            _safe_write(ws, row_i, 1, raw, fmt)
            _safe_write(ws, row_i, 2, days_until, fmt)
            row_i += 1
        if rows_parsed == 0:
            print(
                f"[WARN] EA '{ea_name}': no dates could be parsed from column '{col}'."
                " Check the date format in your CSV."
            )


# ---------------------------------------------------------------------------
# Chart generation
# ---------------------------------------------------------------------------

# Human-readable names for macOS major versions used in chart labels.
MACOS_NAMES: dict[str, str] = {
    "10": "Mac OS X 10",
    "11": "macOS Big Sur 11",
    "12": "macOS Monterey 12",
    "13": "macOS Ventura 13",
    "14": "macOS Sonoma 14",
    "15": "macOS Sequoia 15",
    "26": "macOS Tahoe 26",
}

# Colors per major version for the combined adoption timeline.
MAJOR_VERSION_COLORS: dict[str, str] = {
    "12": "#1f77b4",
    "13": "#2ca02c",
    "14": "#ff7f0e",
    "15": "#e377c2",
    "26": "#bcbd22",
}

INVENTORY_EXPORT_COLUMNS: list[str] = [
    "Jamf Pro ID",
    "Computer Name",
    "Serial Number",
    "Managed",
    "Operating System",
    "OS Build",
    "OS Rapid Security Response",
    "FileVault Status",
    "System Integrity Protection",
    "Firewall Enabled",
    "Bootstrap Token Escrowed",
    "Bootstrap Token Allowed",
    "Gatekeeper",
    "Model",
    "Asset Tag",
    "IP Address",
    "Last Check-in",
    "Last Report",
    "Last Enrollment",
    "Username",
    "Real Name",
    "Email Address",
    "Position",
    "Department",
    "Building",
    "Room",
    "Site",
    "UDID",
    "Management ID",
]

INVENTORY_SECURITY_DETAIL_COLUMNS: list[str] = [
    "FileVault Status",
    "System Integrity Protection",
    "Firewall Enabled",
    "Bootstrap Token Escrowed",
    "Bootstrap Token Allowed",
    "Gatekeeper",
]

INVENTORY_SECURITY_RESOURCE_MAP: dict[str, str] = {
    "FileVault": "FileVault Status",
    "SIP": "System Integrity Protection",
    "Firewall": "Firewall Enabled",
    "Bootstrap Token Escrowed": "Bootstrap Token Escrowed",
    "Bootstrap Token Allowed": "Bootstrap Token Allowed",
    "Gatekeeper": "Gatekeeper",
}


class ChartGenerator:
    """Generates charts from CSV snapshots and jamf-cli JSON snapshot history.

    Supported chart families:
    - macOS adoption timeline from historical CSV exports or jamf-cli
      inventory-summary snapshots.
    - Compliance trend from CSV snapshots containing a failures-count column.
    - Device state trend from jamf-cli device-compliance snapshots.

    Args:
        config: Loaded Config instance.
        csv_path: Path to the current CSV snapshot (optional).
        historical_dir: Directory of dated CSV snapshots for trend analysis.
        output_dir: Directory where generated PNG files are saved.
        workbook: Active xlsxwriter Workbook for embedding charts.
        jamf_cli_dir: Directory containing cached jamf-cli JSON snapshots.
        output_stem: Stem used as the prefix for PNG filenames.
    """

    def __init__(
        self,
        config: "Config",
        csv_path: Optional[str],
        historical_dir: Optional[str],
        output_dir: Path,
        workbook: xlsxwriter.Workbook,
        jamf_cli_dir: Optional[Path],
        output_stem: str,
    ) -> None:
        self._config = config
        self._csv_path = csv_path
        self._hist_dir = historical_dir
        self._out_dir = output_dir
        self._chart_dir = output_dir
        self._wb = workbook
        self._jamf_cli_dir = jamf_cli_dir.expanduser() if jamf_cli_dir else None
        self._chart_prefix = _filename_component(output_stem)

    def generate_all(self) -> tuple[list[str], list[str]]:
        """Generate enabled charts and return (png_paths, source_labels)."""
        if not _load_matplotlib():
            print("  [skip] Charts: matplotlib not installed (pip install matplotlib)")
            return [], []
        charts_cfg = self._config.get("charts") or {}
        if not charts_cfg.get("enabled", True):
            return [], []

        save_png = charts_cfg.get("save_png", True) is not False
        temp_chart_dir: Optional[Path] = None
        if save_png:
            self._chart_dir = self._out_dir
        else:
            temp_chart_dir = Path(tempfile.mkdtemp(prefix="jamf-reports-community-charts-"))
            self._chart_dir = temp_chart_dir

        try:
            png_paths: list[str] = []
            chart_sources: set[str] = set()
            csv_snapshots = self._load_snapshots(charts_cfg)

            if charts_cfg.get("os_adoption", {}).get("enabled", True):
                paths, source_label = self._generate_os_adoption(csv_snapshots, charts_cfg)
                png_paths.extend(paths)
                if source_label:
                    chart_sources.add(source_label)

            comp_cfg = charts_cfg.get("compliance_trend", {})
            if comp_cfg.get("enabled", True):
                fail_col = self._config.compliance.get("failures_count_column", "")
                if csv_snapshots and fail_col:
                    path = self._generate_compliance_trend(csv_snapshots, fail_col, comp_cfg)
                    if path:
                        png_paths.append(path)
                        chart_sources.add("CSV snapshots")
                else:
                    print("  [skip] Compliance trend: failures_count_column not configured")

            device_state_cfg = charts_cfg.get("device_state_trend", {})
            if device_state_cfg.get("enabled", True):
                path = self._generate_device_state_trend()
                if path:
                    png_paths.append(path)
                    chart_sources.add("jamf-cli snapshots")

            if not png_paths:
                print("  [skip] Charts: no CSV or jamf-cli snapshot history found")
                return [], []

            if png_paths and charts_cfg.get("embed_in_xlsx", True):
                self._embed_charts(png_paths)

            return png_paths, sorted(chart_sources)
        finally:
            if temp_chart_dir:
                shutil.rmtree(temp_chart_dir, ignore_errors=True)

    def _load_snapshots(self, charts_cfg: dict[str, Any]) -> list[tuple[datetime, Any]]:
        """Load all CSV snapshots sorted by date. Returns list of (date, DataFrame)."""
        loaded: list[dict[str, Any]] = []
        relevant_columns: set[str] = set()
        os_col = self._config.columns.get("operating_system", "")
        if charts_cfg.get("os_adoption", {}).get("enabled", True) and os_col:
            relevant_columns.add(os_col)
        fail_col = self._config.compliance.get("failures_count_column", "")
        if charts_cfg.get("compliance_trend", {}).get("enabled", True) and fail_col:
            relevant_columns.add(fail_col)

        if self._hist_dir and Path(self._hist_dir).is_dir():
            for f in sorted(
                p for p in Path(self._hist_dir).rglob("*.csv")
                if p.is_file() and not p.is_symlink()
            ):
                try:
                    header_df = pd.read_csv(f, nrows=0, encoding="utf-8-sig")
                except Exception as exc:
                    print(f"  [warn] Skipping unreadable CSV snapshot {f.name}: {exc}")
                    continue
                if relevant_columns and not relevant_columns.intersection(set(header_df.columns)):
                    continue
                dt = self._parse_date_from_path(f)
                try:
                    df = pd.read_csv(f, dtype=str, encoding="utf-8-sig").fillna("")
                except Exception as exc:
                    print(f"  [warn] Skipping CSV snapshot {f.name}: {exc}")
                    continue
                loaded.append(
                    {
                        "date": dt,
                        "df": df,
                        "path": f,
                        "schema": self._csv_snapshot_schema_key(df.columns),
                        "rows": len(df),
                    }
                )

        if self._csv_path and Path(self._csv_path).is_file():
            current_dt = datetime.now().replace(hour=0, minute=0, second=0, microsecond=0)
            try:
                df = pd.read_csv(
                    self._csv_path, dtype=str, encoding="utf-8-sig"
                ).fillna("")
                loaded.append(
                    {
                        "date": current_dt,
                        "df": df,
                        "path": Path(self._csv_path),
                        "schema": self._csv_snapshot_schema_key(df.columns),
                        "rows": len(df),
                    }
                )
            except Exception as exc:
                print(f"  [warn] Could not read current CSV: {exc}")

        snapshots = self._dedupe_csv_snapshots(loaded)
        snapshots.sort(key=lambda x: x[0])
        return snapshots

    def _csv_snapshot_schema_key(self, columns: Any) -> tuple[str, ...]:
        """Return a normalized schema key for a CSV snapshot."""
        return tuple(_normalized_text(column) for column in columns)

    def _dedupe_csv_snapshots(self, loaded: list[dict[str, Any]]) -> list[tuple[datetime, Any]]:
        """Keep one canonical CSV snapshot per day and schema.

        When multiple same-day exports share the same schema, prefer the one with
        the most rows and then the newest file mtime. This keeps an "All Devices"
        export ahead of smaller subset exports that landed in the same folder.
        """
        if not loaded:
            return []

        chosen: dict[tuple[Any, tuple[str, ...]], dict[str, Any]] = {}
        skipped = 0
        for item in loaded:
            key = (item["date"].date(), item["schema"])
            current = chosen.get(key)
            if current is None:
                chosen[key] = item
                continue
            item_rank = (item["rows"], item["path"].stat().st_mtime, str(item["path"]))
            current_rank = (
                current["rows"],
                current["path"].stat().st_mtime,
                str(current["path"]),
            )
            if item_rank > current_rank:
                chosen[key] = item
            skipped += 1

        if skipped:
            print(
                "  [note] Historical CSV dedupe kept the largest snapshot per day and schema;"
                f" ignored {skipped} same-day duplicate candidate(s)."
            )

        return [(item["date"], item["df"]) for item in chosen.values()]

    def _parse_date_from_path(self, path: Path) -> datetime:
        """Extract a timestamp from a filename, falling back to file mtime."""
        name = path.stem
        for pattern, fmt in (
            (r"(\d{4}-\d{2}-\d{2}[T_]\d{6})", "%Y-%m-%dT%H%M%S"),
            (r"(\d{4}-\d{2}-\d{2}[T_]\d{6})", "%Y-%m-%d_%H%M%S"),
            (r"(\d{8}[T_]\d{6})", "%Y%m%dT%H%M%S"),
            (r"(\d{8}[T_]\d{6})", "%Y%m%d_%H%M%S"),
        ):
            match = re.search(pattern, name)
            if not match:
                continue
            value = match.group(1)
            for candidate_fmt in {fmt, fmt.replace("T", "_"), fmt.replace("_", "T")}:
                try:
                    return datetime.strptime(value, candidate_fmt)
                except ValueError:
                    continue
        m = re.search(r"(\d{4}-\d{2}-\d{2})", name)
        if m:
            try:
                return datetime.strptime(m.group(1), "%Y-%m-%d")
            except ValueError:
                pass
        m = re.search(r"(\d{8})", name)
        if m:
            try:
                return datetime.strptime(m.group(1), "%Y%m%d")
            except ValueError:
                pass
        return datetime.fromtimestamp(path.stat().st_mtime).replace(
            hour=0, minute=0, second=0, microsecond=0
        )

    def _load_json_snapshots(self, report_names: list[str]) -> list[tuple[datetime, Any]]:
        """Load jamf-cli JSON snapshots sorted by timestamp."""
        if not self._jamf_cli_dir or not self._jamf_cli_dir.is_dir():
            return []

        candidates: dict[str, Path] = {}
        for report_name in report_names:
            report_dir = self._jamf_cli_dir / report_name
            if report_dir.is_dir():
                for path in report_dir.rglob("*.json"):
                    if ".partial" not in path.name:
                        candidates[str(path)] = path
                continue
            pattern = f"{report_name}_*.json"
            for path in self._jamf_cli_dir.rglob(pattern):
                if ".partial" not in path.name:
                    candidates[str(path)] = path

        snapshots: list[tuple[datetime, Any]] = []
        for path in sorted(candidates.values()):
            try:
                with open(path, encoding="utf-8") as fh:
                    data = json.load(fh)
            except (json.JSONDecodeError, OSError) as exc:
                print(f"  [warn] Skipping unreadable JSON snapshot {path.name}: {exc}")
                continue
            snapshots.append((self._parse_date_from_path(path), data))

        snapshots.sort(key=lambda item: item[0])
        return snapshots

    def _build_os_timeseries(
        self, snapshots: list[tuple[datetime, Any]], os_col: str
    ) -> Any:
        """Build a DataFrame of device counts per OS version over time.

        Returns a DataFrame with dates as index and OS version strings as columns.
        """
        records = []
        for dt, df in snapshots:
            if os_col not in df.columns:
                continue
            series = df[os_col].astype(str).str.strip()
            counts = series[series != ""].value_counts()
            if counts.empty:
                continue
            row: dict[str, Any] = {"date": dt}
            row.update(counts.to_dict())
            records.append(row)
        if not records:
            return pd.DataFrame()
        ts = pd.DataFrame(records).set_index("date").fillna(0).sort_index()
        return ts.astype(int)

    def _major_version(self, ver: str) -> str:
        """Extract the numeric major version from an OS version string.

        Handles both bare versions ("26.1.0" → "26") and full name strings
        from Jamf CSV exports ("macOS 26.1.0" → "26", "Mac OS X 10.15.7" → "10").
        """
        m = re.search(r"\b(\d+)\.\d", str(ver))
        return m.group(1) if m else str(ver).strip()

    def _build_inventory_summary_timeseries(self, snapshots: list[tuple[datetime, Any]]) -> Any:
        """Build a DataFrame of OS version counts from jamf-cli inventory-summary snapshots."""
        records = []
        for dt, data in snapshots:
            rows = data if isinstance(data, list) else []
            row: dict[str, Any] = {"date": dt}
            for item in rows:
                if not isinstance(item, dict):
                    continue
                version = str(item.get("os_version", "") or "").strip()
                if not version:
                    continue
                row[version] = row.get(version, 0) + _to_int(item.get("count", 0))
            if len(row) > 1:
                records.append(row)

        if not records:
            return pd.DataFrame()
        ts = pd.DataFrame(records).groupby("date", as_index=True).sum().sort_index().fillna(0)
        return ts.astype(int)

    def _generate_os_adoption(
        self,
        csv_snapshots: list[tuple[datetime, Any]],
        charts_cfg: dict,
    ) -> tuple[list[str], str]:
        """Generate OS adoption charts. Returns (PNG paths, source label)."""
        ts = pd.DataFrame()
        source_label = ""
        os_col = self._config.columns.get("operating_system", "")
        if csv_snapshots and os_col:
            ts = self._build_os_timeseries(csv_snapshots, os_col)
            if not ts.empty:
                source_label = "CSV snapshots"

        if ts.empty:
            inventory_snapshots = self._load_json_snapshots(
                ["inventory-summary", "inventory_summary"]
            )
            ts = self._build_inventory_summary_timeseries(inventory_snapshots)
            if not ts.empty:
                source_label = "jamf-cli snapshots"

        if ts.empty:
            print(
                "  [skip] OS adoption: no CSV operating_system history or"
                " jamf-cli inventory-summary snapshots found"
            )
            return [], ""

        png_paths = []
        combined_path = self._plot_adoption_combined(ts)
        if combined_path:
            png_paths.append(combined_path)
            print(f"  [chart] {Path(combined_path).name}")

        if charts_cfg.get("os_adoption", {}).get("per_major_charts", True):
            majors = sorted({self._major_version(v) for v in ts.columns})
            for major in majors:
                major_cols = [c for c in ts.columns if self._major_version(c) == major]
                if not major_cols:
                    continue
                path = self._plot_per_major(ts[major_cols], major)
                if path:
                    png_paths.append(path)
                    print(f"  [chart] {Path(path).name}")

        return png_paths, source_label

    def _chart_path(self, chart_name: str) -> str:
        """Return a chart PNG path using the current output stem as a prefix."""
        return str(self._chart_dir / f"{self._chart_prefix}_{chart_name}.png")

    def _plot_adoption_combined(self, ts: Any) -> Optional[str]:
        """Plot combined macOS adoption timeline grouped by major version."""
        major_ts: dict[str, Any] = {}
        for col in ts.columns:
            major = self._major_version(col)
            major_ts[major] = major_ts.get(major, 0) + ts[col]

        fig, ax = plt.subplots(figsize=(12, 6))
        for major, series in sorted(major_ts.items()):
            label = MACOS_NAMES.get(major, f"macOS {major}")
            color = MAJOR_VERSION_COLORS.get(major)
            kwargs: dict[str, Any] = {"label": label, "marker": "o", "markersize": 4}
            if color:
                kwargs["color"] = color
            ax.plot(series.index, series.values, **kwargs)

        ax.set_title("macOS Version Adoption Over Time", fontweight="bold")
        ax.set_xlabel("Date")
        ax.set_ylabel("Number of Devices")
        ax.legend(loc="best", fontsize=9)
        ax.grid(True, alpha=0.3)
        self._format_date_axis(ax, ts.index)
        fig.tight_layout()

        path = self._chart_path("adoption_timeline")
        fig.savefig(path, dpi=150, bbox_inches="tight")
        plt.close(fig)
        return path

    def _plot_per_major(self, ts: Any, major: str) -> Optional[str]:
        """Plot adoption over time for minor/patch versions within a single major."""
        fig, ax = plt.subplots(figsize=(12, 6))
        for col in ts.columns:
            ax.plot(ts.index, ts[col].values, label=col, marker="o", markersize=4)

        name = MACOS_NAMES.get(major, f"macOS {major}")
        ax.set_title(f"{name} — Version Adoption Over Time", fontweight="bold")
        ax.set_xlabel("Date")
        ax.set_ylabel("Number of Devices")
        ncol = max(1, len(ts.columns) // 5)
        ax.legend(loc="best", fontsize=8, ncol=ncol)
        ax.grid(True, alpha=0.3)
        self._format_date_axis(ax, ts.index)
        fig.tight_layout()

        path = self._chart_path(f"adoption_major_{major}")
        fig.savefig(path, dpi=150, bbox_inches="tight")
        plt.close(fig)
        return path

    def _generate_compliance_trend(
        self,
        snapshots: list[tuple[datetime, Any]],
        fail_col: str,
        comp_cfg: dict,
    ) -> Optional[str]:
        """Generate a stacked area compliance trend chart. Returns PNG path or None."""
        bands = comp_cfg.get("bands") or DEFAULT_CONFIG["charts"]["compliance_trend"]["bands"]

        records = []
        for dt, df in snapshots:
            if fail_col not in df.columns:
                continue
            counts_by_band: dict[str, int] = {b["label"]: 0 for b in bands}
            for val in df[fail_col]:
                try:
                    n = int(str(val).strip())
                except (ValueError, TypeError):
                    continue
                for band in bands:
                    if band["min_failures"] <= n <= band["max_failures"]:
                        counts_by_band[band["label"]] += 1
                        break
            row: dict[str, Any] = {"date": dt}
            row.update(counts_by_band)
            records.append(row)

        if not records:
            print(f"  [skip] Compliance trend: column '{fail_col}' not found in snapshots")
            return None

        ts = pd.DataFrame(records).set_index("date").sort_index()
        band_labels = [b["label"] for b in bands]
        colors = [b["color"] for b in bands]

        fig, ax = plt.subplots(figsize=(12, 6))
        ax.stackplot(
            ts.index,
            [ts[lbl].values for lbl in band_labels],
            labels=band_labels,
            colors=colors,
            alpha=0.85,
        )
        baseline_label = _compliance_label(self._config.compliance)
        ax.set_title(f"{baseline_label} Trend Over Time", fontweight="bold")
        ax.set_xlabel("Date")
        ax.set_ylabel("Number of Devices")
        ax.legend(loc="upper left", fontsize=9)
        ax.grid(True, alpha=0.2)
        self._format_date_axis(ax, ts.index)
        fig.tight_layout()

        path = self._chart_path("compliance_trend")
        fig.savefig(path, dpi=150, bbox_inches="tight")
        plt.close(fig)
        print(f"  [chart] {Path(path).name}")
        return path

    def _generate_device_state_trend(self) -> Optional[str]:
        """Generate a jamf-cli device state trend chart from device-compliance history."""
        snapshots = self._load_json_snapshots(["device-compliance", "device_compliance"])
        if not snapshots:
            print("  [skip] Device state trend: no jamf-cli device-compliance snapshots found")
            return None

        records = []
        for dt, data in snapshots:
            rows = data if isinstance(data, list) else []
            counts = {
                "Managed Current": 0,
                "Managed Stale": 0,
                "Unmanaged Current": 0,
                "Unmanaged Stale": 0,
            }
            for item in rows:
                if not isinstance(item, dict):
                    continue
                managed = _to_bool(item.get("managed"))
                stale = _to_bool(item.get("stale"))
                if managed and stale:
                    counts["Managed Stale"] += 1
                elif managed:
                    counts["Managed Current"] += 1
                elif stale:
                    counts["Unmanaged Stale"] += 1
                else:
                    counts["Unmanaged Current"] += 1
            if sum(counts.values()) > 0:
                row: dict[str, Any] = {"date": dt}
                row.update(counts)
                records.append(row)

        if not records:
            print("  [skip] Device state trend: device-compliance snapshots contained no rows")
            return None

        ts = pd.DataFrame(records).groupby("date", as_index=True).sum().sort_index()
        labels = ["Managed Current", "Managed Stale", "Unmanaged Current", "Unmanaged Stale"]
        colors = ["#2E86AB", "#F39C12", "#7FB069", "#C0392B"]

        fig, ax = plt.subplots(figsize=(12, 6))
        if len(ts.index) == 1:
            values = [int(ts.iloc[0][label]) for label in labels]
            ax.bar(labels, values, color=colors)
            plt.setp(ax.get_xticklabels(), rotation=20, ha="right")
            ax.set_xlabel("Device State")
        else:
            ax.stackplot(
                ts.index,
                [ts[label].values for label in labels],
                labels=labels,
                colors=colors,
                alpha=0.85,
            )
            ax.set_xlabel("Date")
            self._format_date_axis(ax, ts.index)

        ax.set_title("Device State Trend Over Time", fontweight="bold")
        ax.set_ylabel("Number of Devices")
        ax.legend(loc="upper left", fontsize=9)
        ax.grid(True, alpha=0.2)
        fig.tight_layout()

        path = self._chart_path("device_state_trend")
        fig.savefig(path, dpi=150, bbox_inches="tight")
        plt.close(fig)
        print(f"  [chart] {Path(path).name}")
        return path

    def _format_date_axis(self, ax: Any, dates: Any) -> None:
        """Apply readable date formatting to a matplotlib x-axis."""
        span = (dates.max() - dates.min()).days if len(dates) > 1 else 0
        if span <= 14:
            ax.xaxis.set_major_formatter(mdates.DateFormatter("%Y-%m-%d"))
            ax.xaxis.set_major_locator(mdates.DayLocator(interval=1))
        elif span <= 90:
            ax.xaxis.set_major_formatter(mdates.DateFormatter("%Y-%m-%d"))
            ax.xaxis.set_major_locator(mdates.WeekdayLocator())
        else:
            ax.xaxis.set_major_formatter(mdates.DateFormatter("%Y-%m"))
            ax.xaxis.set_major_locator(mdates.MonthLocator())
        plt.setp(ax.xaxis.get_majorticklabels(), rotation=45, ha="right")

    def _embed_charts(self, png_paths: list[str]) -> None:
        """Create a Charts sheet in the workbook with all PNGs embedded."""
        ws = self._wb.add_worksheet("Charts")
        row = 0
        for path in png_paths:
            if not Path(path).is_file():
                continue
            ws.insert_image(row, 0, path, {"x_scale": 0.85, "y_scale": 0.85})
            row += 33  # roughly 495px at ~15px/row — enough vertical space per chart


# ---------------------------------------------------------------------------
# Device deep-dive command
# ---------------------------------------------------------------------------


def _print_device_detail(data: Any, indent: int = 0) -> None:
    """Recursively format and print a jamf-cli device JSON response.

    Args:
        data: Parsed JSON (dict, list, or scalar) from jamf-cli pro device.
        indent: Current indentation level (spaces multiplied by 2).
    """
    pad = "  " * indent
    if isinstance(data, list):
        for item in data:
            _print_device_detail(item, indent)
            if isinstance(item, dict):
                print()
    elif isinstance(data, dict):
        for key, value in data.items():
            label = str(key).replace("_", " ").title()
            if isinstance(value, dict):
                print(f"{pad}{label}:")
                _print_device_detail(value, indent + 1)
            elif isinstance(value, list):
                print(f"{pad}{label} ({len(value)}):")
                _print_device_detail(value, indent + 1)
            else:
                display = "" if value is None else str(value)
                print(f"{pad}  {label:<28} {display}")
    else:
        print(f"{pad}{data}")


def cmd_device(config: Config, device_id: str) -> None:
    """Print a formatted device detail view using jamf-cli pro device.

    Args:
        config: Loaded Config instance.
        device_id: Device identifier (Jamf Pro computer ID or serial number).
    """
    if not _jamf_cli_enabled(config):
        raise SystemExit("Error: jamf-cli is disabled in config (jamf_cli.enabled: false).")
    bridge = _build_jamf_cli_bridge(config, save_output=False, use_cached_data=False)
    if not bridge.is_available():
        raise SystemExit(
            "Error: jamf-cli is not installed or not found.\n"
            "  Install via Homebrew: brew install jamf-cli\n"
            "  Or set JAMFCLI_PATH to the binary location."
        )
    print(f"Device: {device_id}")
    print("=" * (len(device_id) + 8))
    try:
        data = bridge.device_lookup(device_id)
    except RuntimeError as exc:
        raise SystemExit(f"Error: {exc}") from exc
    _print_device_detail(data)


# ---------------------------------------------------------------------------
# Scaffold command
# ---------------------------------------------------------------------------


_EA_SKIP_PATTERNS: frozenset[str] = frozenset({
    "udid", "managed", "supervised", "asset tag", "ip address", "mac address",
    "username", "real name", "email", "phone", "position", "department", "building",
    "room", "site", "notes", "comments", "po", "lease", "purchase", "vendor",
    "warranty", "model identifier", "processor", "ram", "storage", "battery",
    "serial", "jamf pro id", "management id",
})

_EA_TYPE_KEYWORDS: dict[str, list[str]] = {
    "version": ["version", "ver "],
    "date": ["date", "expir", "expires", "renewal"],
    "boolean": [
        "status", "enabled", "installed", "bound", "enrolled", "activation",
        "compliant", "connected", "running", "active", "detected",
    ],
}

_SCAFFOLD_COMPLIANCE_HINTS: dict[str, list[str]] = {
    "failures_count_column": [
        "failed mscp results count",
        "failed mscp result count",
        "failed results count",
        "failed rule count",
        "compliance failures count",
    ],
    "failures_list_column": [
        "failed mscp result list",
        "failed mscp results list",
        "failed results list",
        "failed rule list",
        "compliance failures list",
    ],
}


def _suggest_ea_type(header: str) -> str:
    """Return a suggested custom_ea type based on column name heuristics."""
    lower = header.lower()
    for ea_type, keywords in _EA_TYPE_KEYWORDS.items():
        if any(kw in lower for kw in keywords):
            return ea_type
    return "text"


def _looks_like_ea(header: str) -> bool:
    """Return True when a column is likely a custom Extension Attribute, not core inventory."""
    lower = header.lower()
    # Skip generic Jamf inventory column patterns
    if any(pat in lower for pat in _EA_SKIP_PATTERNS):
        return False
    # Columns with org-specific separators (" - ", ": ") are almost always EAs
    if " - " in header or ": " in header:
        return True
    # Multi-word columns ending in EA-type keywords are strong candidates
    if " " in header.strip() and any(
        lower.endswith(kw.strip()) for keywords in _EA_TYPE_KEYWORDS.values()
        for kw in keywords
    ):
        return True
    return False


def _suggest_custom_ea_candidates(unmatched_headers: list[str]) -> None:
    """Print custom_eas config suggestions for unmatched columns that look like EAs.

    Args:
        unmatched_headers: CSV column names that were not mapped to logical fields.
    """
    candidates = [(h, _suggest_ea_type(h)) for h in unmatched_headers if _looks_like_ea(h)]
    if not candidates:
        return

    print(f"\nCustom EA suggestions ({len(candidates)} columns):")
    print("  Add these to the custom_eas list in your config.yaml:")
    for header, ea_type in candidates:
        snippet = (
            f"  - name: {header!r}\n"
            f"    column: {header!r}\n"
            f"    type: {ea_type}"
        )
        if ea_type == "boolean":
            snippet += "\n    true_value: \"<compliant value>\"  # e.g. Installed, Enabled, Yes"
        elif ea_type == "version":
            snippet += "\n    current_versions: []  # e.g. [\"130.0\", \"131.0\"]"
        elif ea_type == "date":
            snippet += "\n    warning_days: 90"
        print(snippet)
    print()


def _yaml_scalar(value: Any) -> str:
    """Render a simple scalar value as YAML-compatible text."""
    if isinstance(value, bool):
        return "true" if value else "false"
    if value is None:
        return '""'
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        return str(value)
    return json.dumps(str(value))


def _render_scaffold_config(config_data: dict[str, Any], csv_path: Path) -> str:
    """Render scaffold output from config.example.yaml when available."""
    template_path = _seed_config_template_path()
    if not template_path.exists():
        return yaml.dump(config_data, default_flow_style=False, sort_keys=False)

    rendered = [
        "# Generated by jamf-reports-community.py scaffold",
        f"# Source CSV: {csv_path}",
        "# Review and adjust the remaining example sections before running generate.",
        "",
    ]
    current_section = ""
    with open(template_path, encoding="utf-8") as fh:
        for raw_line in fh.read().splitlines():
            top_level_match = re.match(r"^([A-Za-z_][A-Za-z0-9_]*):(?:\s*.*)?$", raw_line)
            if top_level_match:
                current_section = top_level_match.group(1)
                if current_section == "security_agents":
                    rendered.append("security_agents: []")
                    continue
                if current_section == "custom_eas":
                    rendered.append("custom_eas: []")
                    continue

            if current_section in {"security_agents", "custom_eas"}:
                if raw_line.strip() and not raw_line.lstrip().startswith("#"):
                    rendered.append(f"# {raw_line}")
                else:
                    rendered.append(raw_line)
                continue

            nested_match = re.match(r"^(\s{2})([A-Za-z_][A-Za-z0-9_]*):(?:\s*.*)?$", raw_line)
            if nested_match and current_section == "columns":
                key = nested_match.group(2)
                if key in config_data["columns"]:
                    rendered.append(f"  {key}: {_yaml_scalar(config_data['columns'][key])}")
                    continue
            if nested_match and current_section == "compliance":
                key = nested_match.group(2)
                if key in config_data["compliance"]:
                    rendered.append(f"  {key}: {_yaml_scalar(config_data['compliance'][key])}")
                    continue

            rendered.append(raw_line)

    return "\n".join(rendered) + "\n"


def _interactive_column_mapping(
    headers: list[str],
    matched: dict[str, str],
) -> dict[str, str]:
    """Walk the user through reviewing and correcting auto-matched column assignments.

    For each logical field, shows the auto-match (if any) alongside a numbered
    list of all available CSV columns so the user can accept, override, or skip.
    Falls back to returning the original matched dict when stdin is not a TTY.

    Args:
        headers: All column names from the CSV.
        matched: Auto-matched dict of logical_field -> csv_column (may be empty string).

    Returns:
        Updated mapping dict (logical -> csv_column or "").
    """
    if not sys.stdin.isatty():
        print("[interactive] stdin is not a terminal; using auto-matched results.")
        return matched

    logical_fields = list(DEFAULT_CONFIG["columns"].keys())
    result = dict(matched)

    print("\nInteractive column mapping")
    print("  Enter  — accept the auto-match (or skip if none)")
    print("  number — choose that column from the list")
    print("  s / 0  — leave this field blank\n")

    for idx, logical in enumerate(logical_fields, 1):
        auto = result.get(logical, "")
        # Available = all headers not already claimed by another field
        others_used = {v for k, v in result.items() if k != logical and v}
        available = [h for h in headers if h not in others_used]

        hints = ", ".join(COLUMN_HINTS.get(logical, []))
        print(f"[{idx}/{len(logical_fields)}] {logical}  (hints: {hints})")

        if auto:
            print(f"  Auto-match: {auto!r}  (press Enter to accept)")
        else:
            print("  No auto-match found")

        print("  0: <skip — leave blank>")
        for i, h in enumerate(available, 1):
            marker = "  <- auto-match" if h == auto else ""
            print(f"  {i}: {h}{marker}")

        prompt = f"  Choice [Enter={auto!r}]: " if auto else "  Choice [Enter=skip]: "
        while True:
            try:
                raw = input(prompt).strip()
            except (EOFError, KeyboardInterrupt):
                print("\nInteractive mapping aborted. Using current results.")
                return result

            if raw == "":
                break  # accept auto or skip
            if raw.lower() == "s" or raw == "0":
                result[logical] = ""
                break
            try:
                choice = int(raw)
                if 1 <= choice <= len(available):
                    result[logical] = available[choice - 1]
                    break
                print(f"  Please enter a number between 0 and {len(available)}.")
            except ValueError:
                print("  Invalid input — enter a number, 's' to skip, or Enter to accept.")

        print()

    return result


def cmd_scaffold(csv_path: str, out_path: str, interactive: bool = False) -> None:
    """Auto-generate a starter config.yaml from CSV headers.

    Args:
        csv_path: Path to the CSV file to inspect.
        out_path: Output path for generated config.yaml.
        interactive: If True, prompt the user to review each column mapping before writing.
    """
    csv_path_obj = _cli_path(csv_path)
    out_path_obj = _cli_path(out_path)
    if csv_path_obj is None or out_path_obj is None:
        raise SystemExit("Error: scaffold requires valid --csv and --out paths")
    if out_path_obj.exists():
        raise SystemExit(
            f"Error: scaffold output already exists: {out_path_obj}\n"
            "Refusing to overwrite the file. Choose a different --out path or move"
            " the existing config aside first."
        )

    try:
        df = pd.read_csv(csv_path_obj, nrows=0, encoding="utf-8-sig")
    except Exception as exc:
        raise SystemExit(f"Error: could not read CSV '{csv_path_obj}': {exc}") from exc
    headers = list(df.columns)
    print(f"Found {len(headers)} columns in CSV.")

    matched: dict[str, str] = {}
    used_headers: set[str] = set()
    for logical in DEFAULT_CONFIG["columns"]:
        best_header, best_score = _best_header_match(headers, logical, used_headers)
        if best_header and best_score > 0:
            matched[logical] = best_header
            used_headers.add(best_header)

    compliance_matches: dict[str, str] = {}
    for key, hints in _SCAFFOLD_COMPLIANCE_HINTS.items():
        best_header = _best_hint_match(headers, hints, used_headers)
        if best_header:
            compliance_matches[key] = best_header
            used_headers.add(best_header)

    unmatched = [header for header in headers if header not in used_headers]

    print(f"Auto-matched {len(matched)} logical fields:")
    for k, v in matched.items():
        print(f"  {k}: {v!r}")
    if compliance_matches:
        print("Auto-detected compliance fields:")
        for key, value in compliance_matches.items():
            print(f"  compliance.{key}: {value!r}")
        if len(compliance_matches) == 2:
            print("  compliance.enabled: true")
        else:
            print(
                "  compliance.enabled: false (complete the missing column before"
                " generating)"
            )
    if unmatched:
        print(f"Unmatched columns ({len(unmatched)}) — add manually to config if needed:")
        for h in unmatched:
            print(f"  - {h!r}")

    _suggest_custom_ea_candidates(unmatched)

    if interactive:
        matched = _interactive_column_mapping(headers, matched)

    config_data = copy.deepcopy(DEFAULT_CONFIG)
    for logical in config_data["columns"]:
        config_data["columns"][logical] = matched.get(logical, "")
    config_data["compliance"]["failures_count_column"] = compliance_matches.get(
        "failures_count_column", ""
    )
    config_data["compliance"]["failures_list_column"] = compliance_matches.get(
        "failures_list_column", ""
    )
    config_data["compliance"]["enabled"] = bool(
        config_data["compliance"]["failures_count_column"]
        and config_data["compliance"]["failures_list_column"]
    )

    config_str = _render_scaffold_config(config_data, csv_path_obj)
    out_path_obj.parent.mkdir(parents=True, exist_ok=True)
    with open(out_path_obj, "w", encoding="utf-8") as fh:
        fh.write(config_str)
    print(f"\nConfig written to: {out_path_obj}")


# ---------------------------------------------------------------------------
# Workspace bootstrap
# ---------------------------------------------------------------------------


def _resolve_workspace_profile_name(
    seed_config: Config,
    profile: Optional[str],
) -> str:
    """Return the profile name to use for a bootstrapped workspace."""
    explicit = str(profile or "").strip()
    if explicit:
        return explicit

    seeded = str(seed_config.jamf_cli.get("profile", "") or "").strip()
    if seeded:
        return seeded

    if sys.stdin.isatty():
        return _prompt_text("jamf-cli profile name")

    raise SystemExit(
        "Error: workspace-init requires --profile when the seed config does not"
        " already set jamf_cli.profile."
    )


def _resolve_workspace_root_dir(
    workspace_root: Optional[str],
    seed_config: Config,
) -> Path:
    """Return the parent directory under which a workspace will be created."""
    if workspace_root:
        return _expand_setup_path(workspace_root, Path.cwd())

    if not sys.stdin.isatty():
        return Path.cwd().resolve()

    default_root = str(seed_config.base_dir)
    return _expand_setup_path(
        _prompt_text("Workspace root directory", default_root),
        Path.cwd(),
    )


def _resolve_workspace_name(
    profile_name: str,
    workspace_name: Optional[str],
) -> str:
    """Return the filesystem directory name for a bootstrapped workspace."""
    explicit = str(workspace_name or "").strip()
    if explicit:
        return _filename_component(explicit)
    return _filename_component(profile_name)


def _workspace_seed_config_data(seed_config: Config, profile_name: str) -> dict[str, Any]:
    """Return a seed config adjusted for a per-profile workspace."""
    config_data = seed_config.to_dict()
    config_data.setdefault("jamf_cli", {})
    config_data.setdefault("output", {})
    config_data.setdefault("charts", {})
    config_data["jamf_cli"]["profile"] = profile_name
    config_data["jamf_cli"]["data_dir"] = "jamf-cli-data"
    config_data["output"]["output_dir"] = "Generated Reports"
    config_data["output"]["archive_dir"] = ""
    config_data["charts"]["historical_csv_dir"] = "snapshots"
    return config_data


def _write_config_yaml(path: Path, config_data: dict[str, Any], header_lines: list[str]) -> None:
    """Write a config YAML file with a small generated header."""
    config_str = yaml.dump(config_data, default_flow_style=False, sort_keys=False)
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w") as fh:
        for line in header_lines:
            fh.write(f"# {line}\n")
        fh.write("\n")
        fh.write(config_str)


def cmd_workspace_init(
    seed_config_path: Optional[str],
    profile: Optional[str],
    workspace_root: Optional[str],
    workspace_name: Optional[str],
    overwrite_config: bool,
) -> None:
    """Create a per-profile reporting workspace and seeded config.yaml."""
    seed_config, seed_source = _load_workspace_seed_config(seed_config_path)
    profile_name = _resolve_workspace_profile_name(seed_config, profile)
    root_dir = _resolve_workspace_root_dir(workspace_root, seed_config)
    directory_name = _resolve_workspace_name(profile_name, workspace_name)
    workspace_dir = (root_dir / directory_name).resolve()

    generated_paths = {
        "workspace": workspace_dir,
        "config": workspace_dir / "config.yaml",
        "jamf_cli_data": workspace_dir / "jamf-cli-data",
        "snapshots": workspace_dir / "snapshots",
        "output": workspace_dir / "Generated Reports",
        "csv_inbox": workspace_dir / "csv-inbox",
        "automation": workspace_dir / "automation",
        "logs": workspace_dir / "automation" / "logs",
    }

    for key, path in generated_paths.items():
        if key == "config":
            continue
        path.mkdir(parents=True, exist_ok=True)

    config_path = generated_paths["config"]
    wrote_config = False
    if config_path.exists() and not overwrite_config:
        print(f"Config already exists, leaving it in place: {config_path}")
    else:
        config_data = _workspace_seed_config_data(seed_config, profile_name)
        _write_config_yaml(
            config_path,
            config_data,
            [
                "Generated by jamf-reports-community.py workspace-init",
                f"Seed source: {seed_source}",
                f"jamf_cli.profile: {profile_name}",
                "This workspace is intended to isolate one Jamf tenant/profile.",
            ],
        )
        wrote_config = True

    print("\nWorkspace bootstrap summary")
    print(f"  profile: {profile_name}")
    print(f"  seed config: {seed_source}")
    print(f"  workspace: {workspace_dir}")
    print(f"  config: {config_path}")
    print(f"  jamf-cli data: {generated_paths['jamf_cli_data']}")
    print(f"  historical CSVs: {generated_paths['snapshots']}")
    print(f"  output: {generated_paths['output']}")
    print(f"  CSV inbox: {generated_paths['csv_inbox']}")
    print(f"  automation logs: {generated_paths['logs']}")
    if wrote_config:
        print("  config write: created or updated")
    else:
        print("  config write: skipped (use --overwrite-config to replace it)")

    print("\nNext steps")
    print(f"  1. Review {config_path}")
    print(f"  2. jamf-cli config validate -p {shlex.quote(profile_name)}")
    print(
        "  3. python3 jamf-reports-community.py check"
        f" --config {shlex.quote(str(config_path))}"
    )
    print(
        "  4. python3 jamf-reports-community.py launchagent-setup"
        f" --config {shlex.quote(str(config_path))}"
        f" --workspace-dir {shlex.quote(str(workspace_dir))}"
    )


# ---------------------------------------------------------------------------
# Check command
# ---------------------------------------------------------------------------

_VALID_EA_TYPES: frozenset[str] = frozenset({"boolean", "percentage", "version", "text", "date"})


def _validate_config_structure(config: Config) -> None:
    """Validate config keys for type correctness and internal consistency.

    Prints [ok] or [WARN] lines for each check. Does not raise — all issues
    are advisory so the user can fix them before running generate.

    Args:
        config: Loaded Config instance to validate.
    """
    issues: list[str] = []
    raw_config = config.to_dict()
    section_types: dict[str, type] = {
        "columns": dict,
        "mobile_columns": dict,
        "compliance": dict,
        "jamf_cli": dict,
        "protect": dict,
        "platform": dict,
        "thresholds": dict,
        "output": dict,
        "charts": dict,
        "report_families": dict,
        "security_agents": list,
        "custom_eas": list,
    }
    for section, expected_type in section_types.items():
        raw_value = raw_config.get(section)
        if raw_value is None:
            continue
        if not isinstance(raw_value, expected_type):
            issues.append(
                f"{section} is {type(raw_value).__name__} — expected"
                f" {expected_type.__name__}; defaults will be used for that section"
            )

    # Numeric threshold validation
    for key, label in [
        ("stale_device_days", "thresholds.stale_device_days"),
        ("critical_disk_percent", "thresholds.critical_disk_percent"),
        ("warning_disk_percent", "thresholds.warning_disk_percent"),
        ("cert_warning_days", "thresholds.cert_warning_days"),
    ]:
        val = config.thresholds.get(key)
        if val is not None:
            try:
                float(val)
            except (TypeError, ValueError):
                issues.append(f"{label} = {val!r} — must be a number, not {type(val).__name__}")

    jamf_cli_enabled = config.jamf_cli.get("enabled", True)
    if jamf_cli_enabled not in (True, False):
        issues.append(
            "jamf_cli.enabled must be true or false."
            f" Current value: {jamf_cli_enabled!r}"
        )

    # Compliance: if enabled, both columns must be set
    comp = config.compliance
    if comp.get("enabled"):
        for col_key in ("failures_count_column", "failures_list_column"):
            if not comp.get(col_key, ""):
                issues.append(
                    f"compliance.{col_key} is empty but compliance.enabled is true —"
                    " the Compliance sheet will fail at generate time"
                )

    for section_name, mapping in [
        ("columns", config.columns),
        ("mobile_columns", config.mobile_columns),
    ]:
        configured_columns: dict[str, list[str]] = {}
        for logical_field, raw_column in mapping.items():
            column_name = str(raw_column or "").strip()
            if not column_name:
                continue
            configured_columns.setdefault(column_name.casefold(), []).append(logical_field)
        duplicate_columns = {
            column_key: fields
            for column_key, fields in configured_columns.items()
            if len(fields) > 1
        }
        if duplicate_columns:
            for column_key, fields in duplicate_columns.items():
                canonical_name = next(
                    (str(mapping.get(field, "") or "").strip() for field in fields),
                    column_key,
                )
                issues.append(
                    f"{section_name} mapping reuses {canonical_name!r} for {', '.join(fields)}"
                    " — each logical field should point to a unique CSV column"
                )

    platform_cfg = config.platform
    if platform_cfg.get("enabled"):
        if not _platform_benchmark_titles(config):
            issues.append(
                "platform.enabled is true but platform.compliance_benchmarks is empty —"
                " benchmark-specific compliance sheets will be skipped"
            )

    report_families = config.report_families
    if not isinstance(report_families, dict):
        issues.append(
            "report_families must be a mapping when present."
            f" Current value: {type(report_families).__name__}"
        )
    else:
        for family_name in REPORT_FAMILY_NAMES:
            family = _report_family_config(config, family_name)
            enabled = family.get("enabled", False)
            if enabled not in (True, False):
                issues.append(
                    f"report_families.{family_name}.enabled must be true or false."
                    f" Current value: {enabled!r}"
                )
            if enabled is not True:
                continue
            if not str(family.get("current_dir", "") or "").strip():
                issues.append(
                    f"report_families.{family_name}.current_dir is empty while the family"
                    " is enabled"
                )
            if not str(family.get("historical_dir", "") or "").strip():
                issues.append(
                    f"report_families.{family_name}.historical_dir is empty while the family"
                    " is enabled"
                )
            include_globs = family.get("include_globs", [])
            if include_globs and not isinstance(include_globs, list):
                issues.append(
                    f"report_families.{family_name}.include_globs must be a list of glob"
                    f" strings, not {type(include_globs).__name__}"
                )
            exclude_globs = family.get("exclude_globs", [])
            if exclude_globs and not isinstance(exclude_globs, list):
                issues.append(
                    f"report_families.{family_name}.exclude_globs must be a list of glob"
                    f" strings, not {type(exclude_globs).__name__}"
                )
            preferred = family.get("prefer_name_contains", [])
            if preferred and not isinstance(preferred, list):
                issues.append(
                    f"report_families.{family_name}.prefer_name_contains must be a list of"
                    f" strings, not {type(preferred).__name__}"
                )

    # Custom EA type validation
    for ea in config.custom_eas:
        name = ea.get("name", "?")
        ea_type = ea.get("type", "")
        if ea_type not in _VALID_EA_TYPES:
            issues.append(
                f"custom_ea '{name}': type={ea_type!r} is not valid."
                f" Use one of: {', '.join(sorted(_VALID_EA_TYPES))}"
            )
        if ea_type == "boolean" and not ea.get("true_value", ""):
            issues.append(
                f"custom_ea '{name}': boolean EA has no true_value set."
                " Defaulting to 'Yes' — set true_value to the expected compliant value."
            )
        for num_key in ("warning_threshold", "critical_threshold", "warning_days"):
            val = ea.get(num_key)
            if val is not None:
                try:
                    float(str(val).rstrip("%"))
                except (TypeError, ValueError):
                    issues.append(
                        f"custom_ea '{name}': {num_key}={val!r} — must be numeric"
                    )

    for index, agent in enumerate(config.security_agents):
        if not isinstance(agent, dict):
            issues.append(
                f"security_agents[{index}] is {type(agent).__name__} —"
                " expected a mapping with name, column, and connected_value"
            )
            continue
        name = str(agent.get("name", "") or "").strip() or f"security_agents[{index}]"
        column = str(agent.get("column", "") or "").strip()
        if not column:
            issues.append(
                f"security_agents entry '{name}' has no column configured —"
                " the Security Agents sheet cannot evaluate it"
            )
        connected_value = str(agent.get("connected_value", "") or "").strip()
        if not connected_value:
            issues.append(
                f"security_agents entry '{name}' has an empty connected_value —"
                " any non-empty cell will count as connected"
            )

    # Charts: warn if enabled but matplotlib unavailable
    charts_enabled = config.get("charts", "enabled")
    if charts_enabled is None or charts_enabled:
        if not _load_matplotlib():
            issues.append(
                "charts.enabled is true but matplotlib is not installed."
                " Run: pip install matplotlib"
            )

    if issues:
        for issue in issues:
            print(f"  [WARN] {issue}")
    else:
        print("  [ok] config structure looks valid")


def cmd_check(config: Config, csv_path: Optional[str] = None) -> None:
    """Verify jamf-cli auth and config validity.

    When csv_path is provided, cross-references all configured column names
    against the actual CSV headers and reports any mismatches.

    Args:
        config: Loaded Config instance.
        csv_path: Optional path to CSV export for column validation.
    """
    print("--- Configuration check ---")
    print(f"  config base dir: {config.base_dir}")
    required = ["computer_name", "serial_number", "operating_system", "last_checkin"]
    any_configured = False
    for field in required:
        val = config.columns.get(field, "")
        status = "ok" if val else "not configured"
        print(f"  columns.{field}: {status}" + (f" -> {val!r}" if val else ""))
        if val:
            any_configured = True
    if not any_configured:
        print("  Note: No columns configured. Run scaffold or edit config.yaml to add mappings.")

    mobile_required = ["device_name", "serial_number", "operating_system", "last_checkin"]
    mobile_any_configured = False
    for field in mobile_required:
        val = config.mobile_columns.get(field, "")
        status = "ok" if val else "not configured"
        print(f"  mobile_columns.{field}: {status}" + (f" -> {val!r}" if val else ""))
        if val:
            mobile_any_configured = True
    if not mobile_any_configured:
        print("  Note: No mobile_columns configured. Mobile CSV validation will be limited.")

    print("\n--- Report Family Manifest ---")
    manifest_found = False
    for family_name in REPORT_FAMILY_NAMES:
        family = _report_family_config(config, family_name)
        if family.get("enabled") is not True:
            continue
        manifest_found = True
        current_dir = _report_family_current_dir(config, family_name)
        historical_dir = _report_family_historical_dir(config, family_name)
        latest_path, note = _latest_report_family_file(config, family_name)
        print(f"  {family_name}: enabled")
        if current_dir is not None:
            print(f"    current_dir: {current_dir}")
        if historical_dir is not None:
            print(f"    historical_dir: {historical_dir}")
        if latest_path is not None:
            print(f"    latest: {latest_path.name}")
        else:
            print(f"    latest: none")
        print(f"    note: {note}")
    if not manifest_found:
        print("  No enabled report_families entries.")

    selected_csv_family: Optional[str] = None
    if not csv_path:
        manifest_csv, selected_csv_family, manifest_note = _default_generate_csv(config)
        if manifest_csv is not None:
            csv_path = str(manifest_csv)
            print("\n  Using manifest-selected CSV for validation.")
            print(f"  {manifest_note}")

    if csv_path:
        print("\n--- CSV column validation ---")
        csv_path_obj = _resolve_cli_input_path(csv_path, config)
        try:
            if csv_path_obj is None:
                raise FileNotFoundError("missing CSV path")
            df = pd.read_csv(csv_path_obj, dtype=str, encoding="utf-8-sig").fillna("")
            headers = df.columns.tolist()
            csv_cols = set(headers)
        except Exception as exc:
            checked = _describe_cli_input_candidates(csv_path, config)
            print(f"  Could not read CSV {csv_path!r}: {exc}")
            if checked:
                print(f"  Checked: {checked}")
        else:
            if selected_csv_family is None and csv_path_obj is not None:
                selected_csv_family = _family_for_csv_path(config, csv_path_obj)
            if selected_csv_family is None:
                selected_csv_family = _guess_report_family_from_headers(config, headers)
            mapping = config.mobile_columns if selected_csv_family == "mobile" else config.columns
            mapping_label = "mobile_columns" if selected_csv_family == "mobile" else "columns"
            if selected_csv_family:
                print(f"  Detected CSV family: {selected_csv_family}")
            if selected_csv_family == "mobile" and not any(str(value or "").strip() for value in mapping.values()):
                print(
                    "  [WARN] Detected a mobile CSV, but mobile_columns is empty."
                    " Configure mobile_columns to validate or generate mobile CSV sheets."
                )
            mismatches = []
            for field, col in mapping.items():
                if not col:
                    continue
                if col in csv_cols:
                    print(f"  [ok] {mapping_label}.{field}: {col!r}")
                    suggested_col, suggested_score = _best_header_match(headers, field)
                    configured_score = _column_match_score(col, field)
                    if (
                        suggested_col
                        and suggested_col != col
                        and suggested_score > configured_score
                    ):
                        print(
                            f"  [SUGGEST] {mapping_label}.{field}: {suggested_col!r} looks like a"
                            f" better match than {col!r}"
                        )
                else:
                    print(f"  [MISSING] {mapping_label}.{field}: {col!r} — not found in CSV")
                    mismatches.append((f"{mapping_label}.{field}", col))
            for ea in config.custom_eas:
                col = ea.get("column", "")
                name = ea.get("name", "?")
                if not col:
                    continue
                if col in csv_cols:
                    print(f"  [ok] custom_ea '{name}': {col!r}")
                else:
                    print(f"  [MISSING] custom_ea '{name}': {col!r} — not found in CSV")
                    mismatches.append((f"custom_ea:{name}", col))
            if selected_csv_family != "mobile":
                for index, agent in enumerate(config.security_agents):
                    if not isinstance(agent, dict):
                        print(
                            f"  [MISSING] security_agents[{index}] is not a mapping —"
                            " cannot validate column settings"
                        )
                        mismatches.append((f"security_agents[{index}]", ""))
                        continue
                    name = str(agent.get("name", "") or "").strip() or f"security_agents[{index}]"
                    col = str(agent.get("column", "") or "").strip()
                    connected_value = str(agent.get("connected_value", "") or "").strip()
                    if not col:
                        print(
                            f"  [MISSING] security_agents '{name}': no column configured"
                        )
                        mismatches.append((f"security_agents:{name}", col))
                    elif col in csv_cols:
                        print(f"  [ok] security_agents '{name}': {col!r}")
                        if not connected_value:
                            print(
                                f"  [WARN] security_agents '{name}': connected_value is empty"
                                " and any non-empty cell will count as connected"
                            )
                    else:
                        print(
                            f"  [MISSING] security_agents '{name}': {col!r} — not found in CSV"
                        )
                        mismatches.append((f"security_agents:{name}", col))
                compliance_cols = [
                    ("failures_count_column", "compliance.failures_count_column"),
                    ("failures_list_column", "compliance.failures_list_column"),
                ]
                for key, label in compliance_cols:
                    col = str(config.compliance.get(key, "") or "").strip()
                    if not col:
                        continue
                    if col in csv_cols:
                        print(f"  [ok] {label}: {col!r}")
                    else:
                        print(f"  [MISSING] {label}: {col!r} — not found in CSV")
                        mismatches.append((label, col))
            if mismatches:
                print(
                    f"\n  {len(mismatches)} column(s) not found."
                    " Fix config.yaml or re-run scaffold."
                )
            else:
                print("\n  All configured columns found in CSV.")

            warnings = _semantic_warnings(config, df.head(50))
            if warnings:
                print("\n--- Semantic validation ---")
                for warning in warnings:
                    print(f"  [WARN] {warning}")

    print("\n--- Config validation ---")
    _validate_config_structure(config)
    isolation_guidance = _profile_isolation_guidance(config)
    if isolation_guidance:
        print("\n--- Profile isolation guidance ---")
        for item in isolation_guidance:
            print(f"  [NOTE] {item}")

    print("\n--- jamf-cli check ---")
    if not _jamf_cli_enabled(config):
        print("  jamf-cli: disabled in config (jamf_cli.enabled: false)")
        print("  Note: generate will skip live and cached jamf-cli sheets.")
        return
    jamf_cli_cfg = config.jamf_cli
    protect_enabled = config.get("protect", "enabled", default=False) is True
    platform_enabled = config.get("platform", "enabled", default=False) is True
    platform_benchmarks = _platform_benchmark_titles(config)
    jamf_cli_dir = config.resolve_path("jamf_cli", "data_dir", default="jamf-cli-data")
    jamf_cli_profile = str(jamf_cli_cfg.get("profile", "") or "").strip()
    live_overview_allowed = jamf_cli_cfg.get("allow_live_overview", True) is True
    bridge = _build_jamf_cli_bridge(config, save_output=False)
    print(f"  data dir: {bridge._data_dir}")
    if jamf_cli_profile:
        print(f"  profile: {jamf_cli_profile}")
    if live_overview_allowed:
        print("  live overview: enabled")
    else:
        print("  live overview: disabled (cached overview only)")
    print(f"  protect reporting: {'enabled' if protect_enabled else 'disabled'}")
    print(f"  platform reporting: {'enabled' if platform_enabled else 'disabled'}")
    if platform_enabled and platform_benchmarks:
        for bench in platform_benchmarks:
            print(f"  platform benchmark: {bench}")
    if bridge.is_available():
        print(f"  jamf-cli: found -> {bridge._binary}")
        report_commands = bridge._report_commands()
        if report_commands:
            current_core = {
                "security",
                "patch-status",
                "inventory-summary",
                "device-compliance",
                "ea-results",
                "software-installs",
            }
            future_optional = {"policy-status", "profile-status"}
            missing_current = sorted(current_core - report_commands)
            missing_optional = sorted(future_optional - report_commands)
            platform_required = {"blueprint-status", "ddm-status"} if platform_enabled else set()
            platform_benchmark_cmds = (
                {"compliance-rules", "compliance-devices"}
                if platform_enabled and platform_benchmarks else set()
            )
            missing_platform = sorted(
                (platform_required | platform_benchmark_cmds) - report_commands
            )
            print(f"  supported report commands: {', '.join(sorted(report_commands))}")
            if missing_current:
                print(f"  missing current core commands: {', '.join(missing_current)}")
            if missing_optional:
                print(f"  missing optional commands: {', '.join(missing_optional)}")
            if missing_platform:
                print(f"  missing platform commands: {', '.join(missing_platform)}")
        if protect_enabled:
            protect_commands = bridge._protect_commands()
            if protect_commands:
                print(f"  supported protect commands: {', '.join(sorted(protect_commands))}")
        try:
            live_bridge = _build_jamf_cli_bridge(
                config,
                save_output=False,
                use_cached_data=False,
            )
            probe_candidates: list[tuple[str, Any]] = []
            if not report_commands or "inventory-summary" in report_commands:
                probe_candidates.append(("inventory-summary", live_bridge.inventory_summary))
            if not report_commands or "security" in report_commands:
                probe_candidates.append(("security", live_bridge.security_report))
            if not report_commands or "patch-status" in report_commands:
                probe_candidates.append(("patch-status", live_bridge.patch_status))
            if live_overview_allowed:
                probe_candidates.append(("overview", live_bridge.overview))

            last_exc: Optional[RuntimeError] = None
            auth_ok = False
            for label, probe in probe_candidates:
                try:
                    probe()
                    print(f"  auth: ok ({label})")
                    auth_ok = True
                    break
                except RuntimeError as exc:
                    last_exc = exc

            if not auth_ok:
                if last_exc is not None:
                    print(f"  auth: failed — {last_exc}")
                else:
                    print("  auth: skipped — no safe jamf-cli report probe was available")
        except RuntimeError as exc:
            print(f"  auth: failed — {exc}")
        if protect_enabled:
            try:
                live_protect_bridge = _build_jamf_cli_bridge(
                    config,
                    save_output=False,
                    use_cached_data=False,
                )
                protect_overview = live_protect_bridge.protect_overview()
                if _protect_overview_has_data(protect_overview):
                    print("  protect auth: ok (overview returned live values)")
                else:
                    print("  protect auth: inconclusive (overview returned placeholder values)")
            except RuntimeError as exc:
                print(f"  protect auth: failed — {exc}")
        if platform_enabled:
            try:
                live_platform_bridge = _build_jamf_cli_bridge(
                    config,
                    save_output=False,
                    use_cached_data=False,
                )
                if not report_commands or "blueprint-status" in report_commands:
                    live_platform_bridge.blueprint_status()
                    print("  platform auth: ok (blueprint-status)")
                else:
                    print("  platform auth: skipped — blueprint-status not available")
                if platform_benchmarks:
                    if not report_commands or "compliance-rules" in report_commands:
                        live_platform_bridge.compliance_rules(platform_benchmarks[0])
                        print("  platform benchmark: ok (compliance-rules)")
                    else:
                        print("  platform benchmark: skipped — compliance-rules not available")
            except RuntimeError as exc:
                print(f"  platform auth: failed — {exc}")
        if bridge.has_cached_data(
            include_protect=protect_enabled,
            include_platform=platform_enabled,
            platform_benchmarks=platform_benchmarks,
        ):
            print("  cached snapshots: found")
    else:
        print("  jamf-cli: not found (set JAMFCLI_PATH or install via Homebrew)")
        if bridge.has_cached_data(
            include_protect=protect_enabled,
            include_platform=platform_enabled,
            platform_benchmarks=platform_benchmarks,
        ):
            print("  cached snapshots: found")
            print("  Note: core sheets can still render from cached jamf-cli JSON snapshots.")
        else:
            print("  Note: CSV-only reports will still work without jamf-cli.")


# ---------------------------------------------------------------------------
# Generate command
# ---------------------------------------------------------------------------


def _post_teams_notification(
    webhook_url: str,
    report_path: Path,
    sheets_written: int,
    generated_at: str,
) -> None:
    """Post an Adaptive Card summary to a Microsoft Teams incoming webhook.

    Args:
        webhook_url: Teams incoming webhook URL (from --notify or config).
        report_path: Path to the generated xlsx file.
        sheets_written: Total number of sheets written to the workbook.
        generated_at: Human-readable generation timestamp string.
    """
    payload = {
        "type": "message",
        "attachments": [
            {
                "contentType": "application/vnd.microsoft.card.adaptive",
                "contentUrl": None,
                "content": {
                    "$schema": "http://adaptivecards.io/schemas/adaptive-card.json",
                    "type": "AdaptiveCard",
                    "version": "1.4",
                    "body": [
                        {
                            "type": "TextBlock",
                            "size": "Large",
                            "weight": "Bolder",
                            "text": "Jamf Report Generated",
                        },
                        {
                            "type": "FactSet",
                            "facts": [
                                {"title": "File", "value": report_path.name},
                                {"title": "Sheets", "value": str(sheets_written)},
                                {"title": "Generated", "value": generated_at},
                            ],
                        },
                    ],
                },
            }
        ],
    }
    parsed_url = urllib.parse.urlparse(webhook_url)
    if parsed_url.scheme not in ("https", "http"):
        print("  [warn] Teams notification skipped: webhook URL must use https or http.")
        return
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        webhook_url,
        data=data,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=10) as resp:
            if resp.status not in (200, 202):
                print(f"  [warn] Teams webhook returned HTTP {resp.status}; notification may not appear.")
    except urllib.error.URLError as exc:
        print(f"  [warn] Teams webhook notification failed: {exc}")


def cmd_generate(
    config: Config,
    csv_path: Optional[str],
    out_file: Optional[str],
    historical_csv_dir: Optional[str] = None,
    notify_url: Optional[str] = None,
    csv_extra: Optional[list[str]] = None,
) -> Path:
    """Run all report generation and write the Excel file.

    Args:
        config: Loaded Config instance.
        csv_path: Optional path to CSV inventory export.
        out_file: Optional output file path override.
        historical_csv_dir: Optional directory of dated CSV snapshots for trend charts.
        notify_url: Optional Teams incoming webhook URL for post-generation notification.
        csv_extra: Optional list of additional CSV paths to merge with csv_path.
    """
    selected_family_name: Optional[str] = None
    selected_csv_origin = ""
    if not csv_path:
        manifest_csv, selected_family_name, manifest_note = _default_generate_csv(config)
        if manifest_csv is not None:
            csv_path = str(manifest_csv)
            selected_csv_origin = f"report_families.{selected_family_name}"
            print(f"Using manifest-selected {selected_family_name} CSV: {manifest_csv}")
        elif manifest_note:
            print(f"  [note] {manifest_note}")

    csv_path_obj = _resolve_cli_input_path(csv_path, config)
    if csv_path_obj is not None and selected_family_name is None:
        selected_family_name = _family_for_csv_path(config, csv_path_obj)
    if csv_path_obj is not None and selected_family_name is None:
        try:
            header_df = pd.read_csv(csv_path_obj, nrows=0, encoding="utf-8-sig")
            selected_family_name = _guess_report_family_from_headers(
                config, header_df.columns.tolist()
            )
        except Exception:
            selected_family_name = None
    if csv_path_obj is not None and not selected_csv_origin:
        if selected_family_name and _family_for_csv_path(config, csv_path_obj) == selected_family_name:
            selected_csv_origin = f"report_families.{selected_family_name}"
        else:
            selected_csv_origin = "--csv"
    csv_path_str = str(csv_path_obj) if csv_path_obj else None

    output_cfg = config.output
    out_dir = config.resolve_path("output", "output_dir", default="Generated Reports")
    if out_dir is None:
        out_dir = Path("Generated Reports")
    out_dir.mkdir(parents=True, exist_ok=True)
    timestamp_outputs = output_cfg.get("timestamp_outputs", True) is not False
    archive_enabled = output_cfg.get("archive_enabled", True) is not False
    keep_latest_runs = _to_int(output_cfg.get("keep_latest_runs", 10), 10)
    run_stamp = _file_stamp()

    jamf_cli_cfg = config.jamf_cli
    jamf_cli_enabled = _jamf_cli_enabled(config)
    protect_enabled = config.get("protect", "enabled", default=False) is True
    platform_enabled = config.get("platform", "enabled", default=False) is True
    platform_benchmarks = _platform_benchmark_titles(config)
    jamf_cli_dir = config.resolve_path("jamf_cli", "data_dir", default="jamf-cli-data")
    jamf_cli_profile = str(jamf_cli_cfg.get("profile", "") or "").strip()
    live_overview_allowed = jamf_cli_enabled and (
        jamf_cli_cfg.get("allow_live_overview", True) is True
    )
    bridge: Optional[JamfCLIBridge] = None
    jamf_cli_ready = False
    if jamf_cli_enabled:
        bridge = _build_jamf_cli_bridge(config, save_output=True)
        jamf_cli_ready = bridge.is_available() or bridge.has_cached_data(
            include_protect=protect_enabled,
            include_platform=platform_enabled,
            platform_benchmarks=platform_benchmarks,
        )

    if out_file:
        out_path = _timestamped_output_path(
            Path(out_file).expanduser(), run_stamp, timestamp_outputs
        )
    else:
        if csv_path_str and jamf_cli_ready:
            default_name = "jamf_report_csv_plus_jamf_cli.xlsx"
        elif csv_path_str:
            default_name = "jamf_report_csv_only.xlsx"
        else:
            default_name = "jamf_report_jamf_cli_only.xlsx"
        out_path = _timestamped_output_path(
            out_dir / default_name,
            run_stamp,
            timestamp_outputs,
        )
    out_path.parent.mkdir(parents=True, exist_ok=True)

    print(f"Output: {out_path}")
    print(f"  config base dir: {config.base_dir}")
    generated_at = datetime.now().strftime("%Y-%m-%d %H:%M")
    wb = xlsxwriter.Workbook(str(out_path))
    accent_color = (config.get("branding", "accent_color") or "#2D5EA2").strip()
    fmts = _build_formats(wb, accent_color)
    sheets_written = 0
    jamf_cli_written: list[str] = []
    csv_written: list[str] = []
    chart_source = ""
    source_details: list[dict[str, str]] = []
    archived_csv_path: Optional[Path] = None
    archived_csv_created: Optional[bool] = None
    if jamf_cli_enabled and bridge is not None:
        print(f"  jamf-cli data dir: {bridge._data_dir}")
        if jamf_cli_profile:
            print(f"  jamf-cli profile: {jamf_cli_profile}")
        if protect_enabled:
            print("  protect reporting: enabled (experimental)")
        if platform_enabled:
            print("  platform reporting: enabled (preview)")
            for bench in platform_benchmarks:
                print(f"  platform benchmark: {bench}")
        if not live_overview_allowed:
            print("  live overview disabled; Fleet Overview will use cache only.")
    else:
        print("  jamf-cli disabled in config; skipping live and cached jamf-cli sheets.")
    try:
        if jamf_cli_enabled and bridge is not None and jamf_cli_ready:

            print("\nGenerating jamf-cli sheets...")
            core = CoreDashboard(config, bridge, wb, fmts)
            jamf_cli_written = core.write_all()
            sheets_written += len(jamf_cli_written)
        elif jamf_cli_enabled:
            print("\nWarning: jamf-cli not available — skipping core dashboard sheets.")
            print(
                "  Set JAMFCLI_PATH, authenticate jamf-cli, or populate cached"
                " jamf-cli JSON snapshots."
            )
        else:
            print("\njamf-cli disabled in config — skipping core dashboard sheets.")

        if csv_path_str:
            print("\nGenerating CSV sheets...")
            try:
                csv_dash = CSVDashboard(
                    config,
                    csv_path_str,
                    wb,
                    fmts,
                    family_name=selected_family_name or "computers",
                    extra_csv_paths=csv_extra,
                )
            except (pd.errors.ParserError, UnicodeDecodeError, OSError) as exc:
                print(f"  [error] Cannot read CSV: {exc}")
                print("  Skipping CSV sheets. Verify the file is a valid UTF-8 CSV export.")
            else:
                csv_written = csv_dash.write_all()
                sheets_written += len(csv_written)
        else:
            print("\nNo CSV provided — skipping inventory sheets.")
            print("  Pass --csv path/to/export.csv to enable inventory analysis.")

        hist_dir_obj = _default_historical_dir(
            config,
            selected_family_name,
            historical_csv_dir,
        )
        hist_dir = str(hist_dir_obj) if hist_dir_obj else None
        if hist_dir_obj:
            print(f"  historical CSV dir: {hist_dir_obj}")
        if csv_path_str and hist_dir and config.get("charts", "archive_current_csv") is not False:
            try:
                archived_csv_path, archived_csv_created = _archive_csv_snapshot(csv_path_str, hist_dir)
                if archived_csv_path:
                    if archived_csv_created:
                        print(f"  Archived CSV snapshot: {archived_csv_path}")
                    else:
                        print(f"  Reusing existing identical CSV snapshot: {archived_csv_path}")
            except OSError as exc:
                print(f"  [warn] Could not archive CSV snapshot: {exc}")

        charts_enabled = config.get("charts", "enabled")
        if charts_enabled is None or charts_enabled:
            print("\nGenerating charts...")
            chart_csv_path = csv_path_str if selected_family_name != "mobile" else None
            chart_hist_dir = hist_dir if selected_family_name != "mobile" else None
            if selected_family_name == "mobile":
                print("  [note] Mobile CSV family selected; skipping CSV trend charts for this run.")
            chart_gen = ChartGenerator(
                config,
                chart_csv_path,
                chart_hist_dir,
                out_path.parent,
                wb,
                jamf_cli_dir if jamf_cli_enabled else None,
                out_path.stem,
            )
            png_paths, chart_sources = chart_gen.generate_all()
            if chart_sources:
                chart_source = " + ".join(chart_sources)
            if png_paths and config.get("charts", "save_png") is not False:
                print(f"  {len(png_paths)} PNG(s) saved to {out_path.parent}/")

        if sheets_written == 0:
            raise SystemExit(
                "Error: No data sources available. Enable jamf-cli in config and"
                " authenticate jamf-cli (pro and protect if configured), or provide"
                " --csv path/to/export.csv."
            )

        if bridge is not None:
            for report_type, info in sorted(bridge.all_source_info().items()):
                mode = str(info.get("mode", "") or "").strip()
                cached_path = info.get("cached_path")
                cached_obj = cached_path if isinstance(cached_path, Path) else None
                source_details.append(
                    {
                        "scope": report_type,
                        "kind": _source_kind_from_mode(mode),
                        "origin": f"jamf-cli {mode or 'cached'}".strip(),
                        "family": _source_family_for_report_type(report_type),
                        "path": str(cached_obj) if cached_obj else "",
                        "timestamp": _path_timestamp_label(cached_obj),
                        "age": _path_age_label(cached_obj),
                        "notes": "live fetch" if mode == "live" else "",
                    }
                )
        if csv_path_obj is not None:
            source_details.append(
                {
                    "scope": "csv_inventory",
                    "kind": (
                        "csv_manifest_selected"
                        if selected_csv_origin.startswith("report_families.")
                        else "csv_cli_explicit"
                    ),
                    "origin": selected_csv_origin or "--csv",
                    "family": selected_family_name or "",
                    "path": str(csv_path_obj),
                    "timestamp": _path_timestamp_label(csv_path_obj),
                    "age": _path_age_label(csv_path_obj),
                    "notes": "primary CSV input",
                }
            )
        if archived_csv_path is not None:
            archive_note = "archived current CSV snapshot"
            if archived_csv_created is False:
                archive_note = "reused existing identical snapshot"
            source_details.append(
                {
                    "scope": "csv_archive",
                    "kind": "csv_historical_archive",
                    "origin": hist_dir or "historical_csv_dir",
                    "family": selected_family_name or "",
                    "path": str(archived_csv_path),
                    "timestamp": _path_timestamp_label(archived_csv_path),
                    "age": _path_age_label(archived_csv_path),
                    "notes": archive_note,
                }
            )
        if selected_family_name == "mobile":
            source_details.append(
                {
                    "scope": "charts",
                    "kind": "jamf_cli_or_cache_only",
                    "origin": "mobile CSV workbook",
                    "family": "mobile",
                    "path": "",
                    "timestamp": "",
                    "age": "",
                    "notes": "CSV trend charts are skipped for mobile CSV families in this release.",
                }
            )

        _write_report_sources_sheet(
            wb,
            fmts,
            generated_at,
            config,
            csv_path_str,
            selected_family_name or "",
            selected_csv_origin,
            hist_dir,
            str(bridge._data_dir) if bridge is not None else "",
            jamf_cli_profile if jamf_cli_enabled else "",
            live_overview_allowed,
            bridge.overview_source_label() if "Fleet Overview" in jamf_cli_written else "",
            jamf_cli_written,
            csv_written,
            chart_source,
            org_name=config.get("branding", "org_name") or "",
            source_details=source_details,
        )

        wb.close()
    except SystemExit:
        wb.close()
        out_path.unlink(missing_ok=True)
        raise
    except Exception:
        try:
            wb.close()
        except Exception:
            pass
        out_path.unlink(missing_ok=True)
        raise
    if archive_enabled:
        archive_dir = config.resolve_path("output", "archive_dir")
        if archive_dir is None:
            archive_dir = out_path.parent / "archive"
        family_base = _strip_timestamp_suffix(out_path.stem)
        archived_paths = _archive_old_output_runs(
            out_path.parent,
            family_base,
            {".xlsx", ".png"},
            keep_latest_runs,
            archive_dir,
        )
        if archived_paths:
            print(f"  Archived {len(archived_paths)} older output file(s) to {archive_dir}")
    print(f"\nReport written: {out_path}")
    if notify_url:
        print("  Sending Teams notification...")
        _post_teams_notification(notify_url, out_path, sheets_written, generated_at)

    if config.output.get("export_pptx") is True:
        exporter = ReportExporter(
            out_path.parent,
            out_path.stem,
            generated_at,
            jamf_cli_written,
            csv_written,
            org_name=config.get("branding", "org_name") or "",
        )
        pptx_path = exporter.export_pptx()
        if pptx_path:
            print(f"  PPTX export: {pptx_path}")
    return out_path


# ---------------------------------------------------------------------------
# PPTX report exporter
# ---------------------------------------------------------------------------


class ReportExporter:
    """Generates a PPTX executive-summary deck alongside the xlsx report.

    The deck is intentionally minimal — title slide, sheet inventory, and data
    source summary.  It is designed to be opened in PowerPoint or Keynote as a
    starting point for presenting fleet data rather than as a complete document.

    Requires ``python-pptx`` (``pip install python-pptx``).  If the library is
    not installed the export is skipped with a warning.

    Args:
        output_dir: Directory where the xlsx was written.
        report_stem: Base filename (without extension) of the xlsx output.
        generated_at: Human-readable generation timestamp string.
        jamf_cli_sheets: Sheet names generated from jamf-cli data.
        csv_sheets: Sheet names generated from CSV data.
    """

    _TITLE_ACCENT = "#2D5EA2"      # blue matching the workbook header colour

    def __init__(
        self,
        output_dir: Path,
        report_stem: str,
        generated_at: str,
        jamf_cli_sheets: list[str],
        csv_sheets: list[str],
        org_name: str = "",
    ) -> None:
        self._output_dir = output_dir
        self._report_stem = report_stem
        self._generated_at = generated_at
        self._jamf_cli_sheets = jamf_cli_sheets
        self._csv_sheets = csv_sheets
        self._org_name = org_name

    def export_pptx(self) -> Optional[Path]:
        """Generate the PPTX summary file.

        Returns:
            Path to the ``.pptx`` file on success, or ``None`` if python-pptx
            is unavailable or an error occurs during generation.
        """
        if not _load_pptx():
            print("  [skip] PPTX export: python-pptx not installed (pip install python-pptx)")
            return None

        try:
            return self._build_pptx()
        except Exception as exc:
            print(f"  [warn] PPTX export failed: {exc}")
            return None

    def _build_pptx(self) -> Path:
        """Build and save the PPTX deck, returning the output path."""
        prs = pptx_Presentation()
        prs.slide_width = pptx_Inches(13.33)
        prs.slide_height = pptx_Inches(7.5)

        self._add_title_slide(prs)
        self._add_summary_slide(prs)
        if self._jamf_cli_sheets or self._csv_sheets:
            self._add_sheets_slide(prs)

        out_path = self._output_dir / f"{self._report_stem}.pptx"
        prs.save(str(out_path))
        return out_path

    def _add_title_slide(self, prs: Any) -> None:
        """Add the opening title slide."""
        layout = prs.slide_layouts[0]   # Title Slide layout
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = (
            f"{self._org_name} — Jamf Pro Fleet Report"
            if self._org_name
            else "Jamf Pro Fleet Report"
        )
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = self._generated_at

    def _add_summary_slide(self, prs: Any) -> None:
        """Add a slide summarising what data sources were included."""
        layout = prs.slide_layouts[1]   # Title and Content layout
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Report Summary"
        tf = slide.placeholders[1].text_frame
        tf.clear()

        total = len(self._jamf_cli_sheets) + len(self._csv_sheets)
        tf.paragraphs[0].text = f"Sheets generated: {total}"

        sources: list[str] = []
        if self._jamf_cli_sheets:
            sources.append("jamf-cli (live or cached snapshots)")
        if self._csv_sheets:
            sources.append("Jamf Pro CSV export")

        for source in sources:
            p = tf.add_paragraph()
            p.text = f"Data source: {source}"
            p.level = 1

        p = tf.add_paragraph()
        p.text = f"Generated: {self._generated_at}"
        p.level = 1

    def _add_sheets_slide(self, prs: Any) -> None:
        """Add a slide listing every sheet written to the workbook."""
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Included Sheets"
        tf = slide.placeholders[1].text_frame
        tf.clear()

        all_sheets = self._jamf_cli_sheets + self._csv_sheets
        tf.paragraphs[0].text = "jamf-cli sheets:" if self._jamf_cli_sheets else "Sheets:"
        for name in all_sheets:
            p = tf.add_paragraph()
            marker = "  •  " + name
            if name in self._csv_sheets and self._jamf_cli_sheets:
                marker = "  ○  " + name   # different bullet for CSV sheets
            p.text = marker
            p.level = 1


# ---------------------------------------------------------------------------
# HTML report
# ---------------------------------------------------------------------------
#
# Design credit: This HTML report is based on the work of DevliegereM
# (https://github.com/DevliegereM/). The layout, colour
# scheme, and section structure are adapted from that project into Python
# with permission under the spirit of open-source sharing.
# The original Bash/heredoc implementation inspired this Python port.
# ---------------------------------------------------------------------------


class HtmlReport:
    """Generates a self-contained, management-facing HTML status report."""
    _TREND_PALETTE = [
        "#0076B6", "#22c55e", "#8b5cf6", "#f59e0b", "#ef4444",
        "#06b6d4", "#f97316", "#84cc16", "#ec4899", "#64748b",
    ]

    def __init__(
        self,
        config: "Config",
        bridge: JamfCLIBridge,
        out_file: Path,
        no_open: bool = False,
    ) -> None:
        self._config = config
        self._bridge = bridge
        self._out_file = out_file
        self._no_open = no_open

    def generate(self) -> Path:
        """Fetch data and write the HTML report to disk."""
        print("\n--- HTML Report ---")
        data = self._fetch_all()
        html = self._render(data)
        self._out_file.parent.mkdir(parents=True, exist_ok=True)
        self._out_file.write_text(html, encoding="utf-8")
        print(f"  Written: {self._out_file}")
        if not self._no_open:
            subprocess.run(["open", str(self._out_file)], check=False)
        return self._out_file

    def _fetch_all(self) -> dict[str, Any]:
        """Fetch all data required for the report."""
        data: dict[str, Any] = {}
        fetch_status: dict[str, dict[str, Any]] = {}
        live_overview_allowed = self._config.jamf_cli.get("allow_live_overview", True) is True

        def _safe_fetch(key: str, fn: Any) -> None:
            try:
                data[key] = fn()
                count = len(data[key]) if isinstance(data[key], list) else 1
                source = self._bridge.source_info(key)
                fetch_status[key] = {
                    "status": "ok",
                    "records": count,
                    "source_mode": source.get("mode", ""),
                    "cached_path": str(source.get("cached_path", "") or ""),
                }
                print(f"  [ok]   {key} ({count} records)")
            except RuntimeError as exc:
                print(f"  [warn] {key}: {exc}")
                data[key] = []
                fetch_status[key] = {
                    "status": "warn",
                    "records": 0,
                    "source_mode": "",
                    "detail": str(exc),
                }

        print("  Fetching data from Jamf Pro...")
        _safe_fetch(
            "overview",
            lambda: self._bridge.overview(cached_only=not live_overview_allowed),
        )
        _safe_fetch("security", self._bridge.security_report)
        _safe_fetch("mobile_inventory", self._bridge.mobile_device_inventory_details)
        _safe_fetch("mobile_devices", self._bridge.mobile_devices_list)
        _safe_fetch("policies", self._bridge.classic_policies_list)
        _safe_fetch("macos_profiles", self._bridge.macos_profiles_list)
        _safe_fetch("ios_profiles", self._bridge.ios_profiles_list)
        _safe_fetch("smart_groups", self._bridge.smart_groups_list)
        _safe_fetch("scripts", self._bridge.scripts_list)
        _safe_fetch("packages", self._bridge.packages_list)
        _safe_fetch("categories", self._bridge.categories_list)
        _safe_fetch("device_enrollments", self._bridge.device_enrollments_list)
        _safe_fetch("sites", self._bridge.sites_list)
        _safe_fetch("buildings", self._bridge.buildings_list)
        _safe_fetch("departments", self._bridge.departments_list)
        data["_fetch_status"] = fetch_status
        return data

    @staticmethod
    def _ov(overview: Any, resource: str) -> str:
        """Extract a single value from the overview list by resource name."""
        if not isinstance(overview, list):
            return "N/A"
        for item in overview:
            if isinstance(item, dict) and item.get("resource") == resource:
                val = item.get("value", "")
                return str(val) if val not in (None, "") else "N/A"
        return "N/A"

    @staticmethod
    def _sec(security: Any, key: str) -> str:
        """Extract a value from the security report summary section."""
        if not isinstance(security, list):
            return "N/A"
        for item in security:
            if isinstance(item, dict) and item.get("section") == "summary":
                val = item.get("data", {}).get(key, "")
                return str(val) if val not in (None, "") else "N/A"
        return "N/A"

    @staticmethod
    def _to_float(value: Any) -> float:
        """Parse a percentage string or numeric string to float."""
        try:
            return float(str(value).replace("%", "").strip())
        except (ValueError, TypeError):
            return 0.0

    @staticmethod
    def _html_text(value: Any, default: str = "") -> str:
        """Escape a value for HTML text and attribute contexts."""
        from html import escape as _escape

        if value in (None, ""):
            return default
        return _escape(str(value), quote=True)

    @staticmethod
    def _safe_href(url: str) -> str:
        """Return url only if it uses https:// or http://; otherwise return '#'."""
        return url if url.startswith(("https://", "http://")) else "#"

    @classmethod
    def _safe_base_url(cls, url: Any) -> str:
        """Return a safe absolute console base URL, or an empty string."""
        candidate = str(url or "").strip().rstrip("/")
        safe = cls._safe_href(candidate)
        return "" if safe == "#" else safe

    @staticmethod
    def _json_text(value: Any) -> str:
        """Serialize data for an inline script without allowing script breakout."""
        text = json.dumps(value, ensure_ascii=False)
        return text.replace("&", "\\u0026").replace("<", "\\u003c").replace(">", "\\u003e")

    @staticmethod
    def _health_badge_class(status: Any) -> str:
        """Map health text to a badge class using exact positive matches."""
        normalized = " ".join(str(status).lower().replace("_", " ").split())
        if normalized in ("ok", "healthy", "operational", "online"):
            return "badge-ok"
        if "degraded" in normalized or "warning" in normalized:
            return "badge-warn"
        return "badge-err"

    @classmethod
    def _status_badge_html(cls, status: Any) -> str:
        """Render a status badge with escaped text."""
        text = cls._html_text(status)
        if not text:
            return ""
        status_lc = " ".join(str(status).lower().replace("_", " ").split())
        if status_lc in ("enabled", "ok", "active", "successful", "success"):
            badge_class = "badge-ok"
        elif status_lc in ("disabled", "inactive", "not required"):
            badge_class = "badge-dim"
        elif status_lc in ("warning", "degraded"):
            badge_class = "badge-warn"
        else:
            badge_class = "badge-blue"
        return f'<span class="badge {badge_class}">{text}</span>'

    @staticmethod
    def _list_names(items: Any, name_key: str = "name") -> list[str]:
        """Extract a sorted list of names from a list of dicts."""
        if not isinstance(items, list):
            return []
        return sorted(
            str(item.get(name_key, "Unnamed"))
            for item in items
            if isinstance(item, dict)
        )

    @staticmethod
    def _build_hierarchy(items: Any) -> list[dict[str, Any]]:
        """Group items by category, using naming-convention prefix as fallback."""
        if not isinstance(items, list):
            return []
        groups: dict[str, list[str]] = {}
        for item in items:
            if not isinstance(item, dict):
                continue
            name = str(item.get("name", "Unnamed"))
            cat_raw = item.get("category", None)
            if isinstance(cat_raw, dict):
                category = str(cat_raw.get("name", "") or "No Category")
            elif cat_raw and str(cat_raw).strip() not in ("", "None"):
                category = str(cat_raw).strip()
            else:
                parts = name.split(" - ")
                category = parts[0].strip() if len(parts) > 1 else "No Category"
            groups.setdefault(category, []).append(name)
        return [
            {"category": cat, "count": len(names), "items": sorted(names)}
            for cat, names in sorted(groups.items())
        ]

    @staticmethod
    def _overview_sections(overview: Any) -> dict[str, list[dict[str, str]]]:
        """Group overview rows by section for rendering."""
        if not isinstance(overview, list):
            return {}
        sections: dict[str, list[dict[str, str]]] = {}
        for item in overview:
            if not isinstance(item, dict):
                continue
            section = str(item.get("section", "General"))
            if section in ("Health & Alerts",):
                continue
            resource = str(item.get("resource", "") or "")
            if not resource:
                continue
            sections.setdefault(section, []).append(
                {
                    "resource": resource,
                    "value": str(item.get("value", "") or ""),
                    "status": str(item.get("status", "") or ""),
                }
            )
        return sections

    @classmethod
    def _resource_links(cls, console_url: str) -> dict[str, str]:
        """Return known Jamf console links for overview resources."""
        if not console_url:
            return {}
        base = console_url.rstrip("/")
        rel = {
            "Managed Computers": "/computers.html",
            "Unmanaged Computers": "/computers.html",
            "Managed Devices": "/mobileDevices.html",
            "Unmanaged Devices": "/mobileDevices.html",
            "Policies": "/policies.html",
            "macOS Config Profiles": "/OSXConfigurationProfiles.html",
            "iOS Config Profiles": "/mobileDeviceConfigurationProfiles.html",
            "Packages": "/view/settings/computer-management/packages",
            "Scripts": "/view/settings/computer-management/scripts",
            "App Installers": "/app-installers.html",
            "Patch Titles": "/patch.html",
            "eBooks": "/eBooks.html",
            "Webhooks": "/webhooks.html",
            "Sites": "/sites.html",
            "Buildings": "/view/settings/network-organization/buildings",
            "Departments": "/view/settings/network-organization/departments",
            "Categories": "/categories.html",
            "Computer Smart Groups": "/smartComputerGroups.html",
            "Computer Static Groups": "/staticComputerGroups.html",
            "Mobile Smart Groups": "/smartMobileDeviceGroups.html",
            "Mobile Static Groups": "/staticMobileDeviceGroups.html",
            "Static User Groups": "/staticUserGroups.html",
            "DEP Instances": "/deviceEnrollmentProgram.html",
            "ADE Instances": "/deviceEnrollmentProgram.html",
            "DEP Sync Status": "/deviceEnrollmentProgram.html",
            "Computer Prestages": "/computerPrestages.html",
            "Mobile Device Prestages": "/mobileDevicePrestages.html",
            "VPP Locations": "/volumePurchaseProgram.html",
            "Built-in CA Expires": "/view/settings/pki/certificate-authority",
            "Active Alerts": "/notifications.html",
            "Health Status": "/healthCheck.html",
            "MDM Auto Renew (Computers)": "/view/settings/global-management/mdm-profile-settings",
            "MDM Auto Renew (Mobile)": "/view/settings/global-management/mdm-profile-settings",
            "SSO (SAML)": "/view/settings/system-settings/sso",
            "LDAP/IdP Servers": "/ldapServers.html",
            "Patch Management": "/patch.html",
        }
        return {label: f"{base}{path}" for label, path in rel.items()}

    @classmethod
    def _quick_links(cls, console_url: str) -> list[tuple[str, str]]:
        """Return quick-link label/URL pairs."""
        if not console_url:
            return []
        base = console_url.rstrip("/")
        items = [
            ("Computers", "/computers.html"),
            ("Mobile Devices", "/mobileDevices.html"),
            ("Policies", "/policies.html"),
            ("macOS Profiles", "/OSXConfigurationProfiles.html"),
            ("iOS Profiles", "/mobileDeviceConfigurationProfiles.html"),
            ("Packages", "/view/settings/computer-management/packages"),
            ("Scripts", "/view/settings/computer-management/scripts"),
            ("Smart Groups", "/smartComputerGroups.html"),
            ("Categories", "/categories.html"),
            ("Patch Management", "/patch.html"),
            ("Notifications", "/notifications.html"),
            ("Check-In Settings", "/view/settings/computer-management/check-in"),
        ]
        return [(label, f"{base}{path}") for label, path in items]

    @staticmethod
    def _flagged_devices(security: Any) -> list[dict[str, Any]]:
        """Extract devices with at least one security concern from the security report."""
        if not isinstance(security, list):
            return []
        flagged = []
        for item in security:
            if not isinstance(item, dict) or item.get("section") != "device":
                continue
            fv = str(item.get("filevault", ""))
            gk = str(item.get("gatekeeper", ""))
            sip = str(item.get("sip", ""))
            fw = item.get("firewall", True)
            issues = (
                (fv not in ("ENCRYPTED", "") and fv)
                or gk.upper() == "DISABLED"
                or (sip not in ("ENABLED", "Enabled", "") and sip)
                or fw is False
            )
            if issues:
                flagged.append(
                    {
                        "name": item.get("name", ""),
                        "serial": item.get("serial", ""),
                        "os": item.get("os_version", ""),
                        "filevault": fv,
                        "gatekeeper": gk,
                        "sip": sip,
                        "firewall": "No" if fw is False else str(fw),
                    }
                )
        return flagged

    @staticmethod
    def _os_chart_data(security: Any) -> tuple[list[str], list[int]]:
        """Build OS version labels and counts for the donut chart."""
        if not isinstance(security, list):
            return [], []
        versions: dict[str, int] = {}
        for item in security:
            if not isinstance(item, dict) or item.get("section") != "os_version":
                continue
            ver = str(item.get("os_version", "")).removesuffix(".0") or "Unknown"
            count = int(item.get("count", 0))
            versions[ver] = versions.get(ver, 0) + count
        pairs = sorted(versions.items(), key=lambda x: x[0], reverse=True)
        return [pair[0] for pair in pairs], [pair[1] for pair in pairs]

    @staticmethod
    def _mobile_rows(items: Any) -> list[dict[str, Any]]:
        """Return normalized mobile inventory rows for the HTML report."""
        rows = [
            _normalize_mobile_inventory_row(item)
            for item in _extract_items(items)
            if isinstance(item, dict)
        ]
        return [
            row
            for row in rows
            if any(
                str(row.get(key, "")).strip()
                for key in ("Device Name", "Serial Number", "Model", "OS Version")
            )
        ]

    @staticmethod
    def _trend_labels(dates: list[datetime]) -> list[str]:
        """Format trend labels with time only when needed."""
        if len({dt.date() for dt in dates}) < len(dates):
            return [dt.strftime("%Y-%m-%d %H:%M") for dt in dates]
        return [dt.strftime("%Y-%m-%d") for dt in dates]

    def _chart_helper(self) -> "ChartGenerator":
        """Return a chart helper so HTML can reuse snapshot-loading logic."""
        hist_dir = self._config.resolve_path("charts", "historical_csv_dir")
        jamf_cli_dir = self._config.resolve_path(
            "jamf_cli", "data_dir", default="jamf-cli-data"
        )
        return ChartGenerator(
            self._config,
            None,
            str(hist_dir) if hist_dir else None,
            self._out_file.parent,
            None,
            jamf_cli_dir,
            self._out_file.stem,
        )

    def _adoption_trend_payload(self) -> dict[str, Any]:
        """Build trend payload for macOS adoption from existing snapshots."""
        charts_cfg = self._config.get("charts") or {}
        os_cfg = charts_cfg.get("os_adoption", {})
        if not charts_cfg.get("enabled", True) or not os_cfg.get("enabled", True):
            return {}

        helper = self._chart_helper()
        os_col = self._config.columns.get("operating_system", "")
        ts = pd.DataFrame()
        source = ""

        if os_col:
            csv_snapshots = helper._load_snapshots(charts_cfg)
            if csv_snapshots:
                ts = helper._build_os_timeseries(csv_snapshots, os_col)
                if not ts.empty:
                    source = "CSV snapshots"

        if ts.empty:
            inventory_snaps = helper._load_json_snapshots(
                ["inventory-summary", "inventory_summary"]
            )
            ts = helper._build_inventory_summary_timeseries(inventory_snaps)
            if not ts.empty:
                source = "jamf-cli snapshots"

        if ts.empty or len(ts.index) < 2:
            return {}

        major_series: dict[str, Any] = {}
        for column in ts.columns:
            major = helper._major_version(column)
            if major in major_series:
                major_series[major] = major_series[major] + ts[column]
            else:
                major_series[major] = ts[column]

        majors = sorted(
            major_series,
            key=lambda item: (not str(item).isdigit(), int(item) if str(item).isdigit() else item),
        )
        series = []
        for idx, major in enumerate(majors):
            color = MAJOR_VERSION_COLORS.get(str(major), self._TREND_PALETTE[idx % len(self._TREND_PALETTE)])
            series.append(
                {
                    "label": MACOS_NAMES.get(str(major), f"macOS {major}"),
                    "data": [int(value) for value in major_series[major].tolist()],
                    "borderColor": color,
                    "backgroundColor": f"{color}22",
                }
            )

        return {
            "labels": self._trend_labels(list(ts.index)),
            "series": series,
            "source": source,
        }

    def _security_trend_payload(self) -> dict[str, Any]:
        """Build trend payload for security posture from cached security snapshots."""
        charts_cfg = self._config.get("charts") or {}
        if not charts_cfg.get("enabled", True):
            return {}

        helper = self._chart_helper()
        snapshots = helper._load_json_snapshots(["security"])
        records: list[dict[str, Any]] = []

        for dt, payload in snapshots:
            if not isinstance(payload, list):
                continue
            total = self._sec(payload, "total_devices")
            if total in ("", "N/A"):
                continue
            fv = self._to_float(self._sec(payload, "filevault_encrypted_pct"))
            gk = self._to_float(self._sec(payload, "gatekeeper_enabled_pct"))
            sip = self._to_float(self._sec(payload, "sip_enabled_pct"))
            fw = self._to_float(self._sec(payload, "firewall_enabled_pct"))
            records.append(
                {
                    "date": dt,
                    "overall": round((fv + gk + sip + fw) / 4.0, 1),
                    "filevault": fv,
                    "gatekeeper": gk,
                    "sip": sip,
                    "firewall": fw,
                }
            )

        if len(records) < 2:
            return {}

        labels = self._trend_labels([record["date"] for record in records])
        return {
            "labels": labels,
            "series": [
                {
                    "label": "Overall",
                    "data": [record["overall"] for record in records],
                    "borderColor": "#004165",
                    "backgroundColor": "#00416522",
                    "fill": True,
                    "borderWidth": 2.5,
                },
                {
                    "label": "FileVault",
                    "data": [record["filevault"] for record in records],
                    "borderColor": "#22c55e",
                    "backgroundColor": "#22c55e22",
                    "fill": False,
                    "borderWidth": 1.6,
                },
                {
                    "label": "Gatekeeper",
                    "data": [record["gatekeeper"] for record in records],
                    "borderColor": "#0076B6",
                    "backgroundColor": "#0076B622",
                    "fill": False,
                    "borderWidth": 1.6,
                },
                {
                    "label": "SIP",
                    "data": [record["sip"] for record in records],
                    "borderColor": "#8b5cf6",
                    "backgroundColor": "#8b5cf622",
                    "fill": False,
                    "borderWidth": 1.6,
                },
                {
                    "label": "Firewall",
                    "data": [record["firewall"] for record in records],
                    "borderColor": "#f59e0b",
                    "backgroundColor": "#f59e0b22",
                    "fill": False,
                    "borderWidth": 1.6,
                },
            ],
            "source": "jamf-cli security snapshots",
        }

    def _trend_payload(self) -> dict[str, Any]:
        """Return all trend payloads that can be rendered in HTML."""
        return {
            "adoption": self._adoption_trend_payload(),
            "security": self._security_trend_payload(),
        }

    @staticmethod
    def _source_mode_label(mode: Any) -> str:
        """Return a concise human-readable label for a fetch source mode."""
        normalized = str(mode or "").strip().lower()
        if normalized == "live":
            return "Live"
        if normalized == "cached":
            return "Cached"
        if normalized == "cached-fallback":
            return "Cached fallback"
        return "Unavailable"

    @staticmethod
    def _source_badge_class(mode: Any, status: Any) -> str:
        """Map fetch status and source mode to a badge class."""
        if str(status or "").lower() == "warn":
            return "badge-warn"
        if str(mode or "").lower() == "live":
            return "badge-ok"
        if str(mode or "").lower() in {"cached", "cached-fallback"}:
            return "badge-blue"
        return "badge-dim"

    @staticmethod
    def _chart_step(values: list[float]) -> float:
        """Return a rounded chart ceiling based on the maximum value."""
        peak = max(values) if values else 0.0
        if peak <= 0:
            return 1.0
        magnitude = 10 ** max(len(str(int(peak))) - 1, 0)
        step = magnitude / 2 if peak / magnitude < 2 else magnitude
        return math.ceil(peak / step) * step

    @staticmethod
    def _polyline_points(
        values: list[float],
        width: float,
        height: float,
        left: float,
        top: float,
        y_max: float,
    ) -> list[tuple[float, float]]:
        """Project chart series values into SVG coordinates."""
        if not values:
            return []
        usable_width = max(width - left - 18.0, 1.0)
        usable_height = max(height - top - 34.0, 1.0)
        if len(values) == 1:
            x_positions = [left + usable_width / 2.0]
        else:
            x_positions = [
                left + (usable_width * idx / (len(values) - 1))
                for idx in range(len(values))
            ]
        max_value = y_max if y_max > 0 else 1.0
        points = []
        for x_pos, value in zip(x_positions, values):
            ratio = min(max(float(value) / max_value, 0.0), 1.0)
            y_pos = top + usable_height - (ratio * usable_height)
            points.append((round(x_pos, 1), round(y_pos, 1)))
        return points

    def _render_line_chart_svg(
        self,
        payload: dict[str, Any],
        percent_scale: bool = False,
    ) -> str:
        """Render a lightweight self-contained SVG line chart."""
        labels = payload.get("labels") or []
        series = payload.get("series") or []
        if len(labels) < 2 or not series:
            return '<p class="empty-note">Trend data unavailable.</p>'

        width = 760.0
        height = 260.0
        left = 44.0
        top = 18.0
        chart_bottom = height - 34.0
        y_max = 100.0 if percent_scale else self._chart_step(
            [float(value) for row in series for value in row.get("data", [])]
        )
        step_count = 4
        grid = []
        ticks = []
        for idx in range(step_count + 1):
            y = top + ((chart_bottom - top) * idx / step_count)
            value = y_max - ((y_max / step_count) * idx)
            label = f"{int(round(value))}% " if percent_scale else str(int(round(value)))
            grid.append(
                f'<line x1="{left}" y1="{y:.1f}" x2="{width - 18:.1f}" y2="{y:.1f}" '
                'stroke="var(--border)" stroke-width="1"/>'
            )
            ticks.append(
                f'<text x="{left - 8:.1f}" y="{y + 4:.1f}" text-anchor="end" '
                f'class="svg-axis">{self._html_text(label.strip())}</text>'
            )

        axis = (
            f'<line x1="{left}" y1="{chart_bottom:.1f}" x2="{width - 18:.1f}" y2="{chart_bottom:.1f}" '
            'stroke="var(--muted)" stroke-width="1.2"/>'
        )
        label_step = max(1, math.ceil(len(labels) / 6))
        x_labels = []
        x_points = self._polyline_points([0.0] * len(labels), width, height, left, top, y_max)
        for idx, (x_pos, _) in enumerate(x_points):
            if idx % label_step != 0 and idx != len(labels) - 1:
                continue
            x_labels.append(
                f'<text x="{x_pos:.1f}" y="{height - 10:.1f}" text-anchor="middle" class="svg-axis">'
                f"{self._html_text(labels[idx])}</text>"
            )

        paths = []
        for item in series:
            values = [float(value) for value in item.get("data", [])]
            points = self._polyline_points(values, width, height, left, top, y_max)
            if not points:
                continue
            path = " ".join(f"{x:.1f},{y:.1f}" for x, y in points)
            color = str(item.get("borderColor", "#0076B6"))
            paths.append(
                f'<polyline fill="none" stroke="{color}" stroke-width="{float(item.get("borderWidth", 2)):.1f}" '
                f'stroke-linejoin="round" stroke-linecap="round" points="{path}"/>'
            )
            for x_pos, y_pos in points:
                paths.append(
                    f'<circle cx="{x_pos:.1f}" cy="{y_pos:.1f}" r="3.2" fill="{color}" />'
                )

        legend = "".join(
            '<span class="chart-legend-item">'
            f'<span class="chart-legend-swatch" style="background:{self._html_text(item.get("borderColor", "#0076B6"))}"></span>'
            f'{self._html_text(item.get("label", "Series"))}</span>'
            for item in series
        )
        return (
            f'<svg class="trend-svg" viewBox="0 0 {int(width)} {int(height)}" '
            'role="img" aria-label="Trend chart">'
            f"{''.join(grid)}{axis}{''.join(ticks)}{''.join(x_labels)}{''.join(paths)}</svg>"
            f'<div class="chart-legend">{legend}</div>'
        )

    def _render_os_distribution_card(self, labels: list[str], counts: list[int]) -> str:
        """Render OS distribution as self-contained SVG plus a data table."""
        if not labels or not counts:
            return """<div class="chart-card">
  <div class="chart-title">macOS Version Distribution</div>
  <p class="empty-note">No macOS distribution data available.</p>
</div>"""

        colors = (self._TREND_PALETTE * ((len(labels) // len(self._TREND_PALETTE)) + 1))[:len(labels)]
        total = sum(counts) or 1
        max_count = max(counts) or 1
        bars = []
        rows = []
        width = 360.0
        height = max(140.0, float(len(labels) * 26 + 16))
        for idx, (label, count) in enumerate(zip(labels, counts)):
            y = 12.0 + (idx * 26.0)
            pct = (count / total) * 100.0
            bar_width = 90.0 + ((count / max_count) * 210.0)
            color = colors[idx]
            bars.append(
                f'<rect x="46" y="{y:.1f}" width="{bar_width:.1f}" height="14" rx="7" fill="{color}" />'
                f'<text x="12" y="{y + 11:.1f}" class="svg-axis">{idx + 1}</text>'
                f'<text x="{min(330.0, 52.0 + bar_width):.1f}" y="{y + 11:.1f}" class="svg-axis">{count}</text>'
            )
            rows.append(
                "<tr>"
                f'<td><span class="os-dot" style="background:{self._html_text(color)}"></span>{self._html_text(label)}</td>'
                f"<td>{count}</td>"
                f"<td>{pct:.1f}%</td>"
                "</tr>"
            )
        return f"""<div class="chart-card">
  <div class="chart-title">macOS Version Distribution</div>
  <svg class="trend-svg" viewBox="0 0 {int(width)} {int(height)}" role="img" aria-label="macOS version distribution">{''.join(bars)}</svg>
  <table class="os-table">
    <thead><tr><th>Version</th><th>Devices</th><th>%</th></tr></thead>
    <tbody>{''.join(rows)}</tbody>
  </table>
</div>"""

    def _render_source_status(self, fetch_status: dict[str, Any]) -> str:
        """Render source provenance and partial-failure visibility for HTML data."""
        if not fetch_status:
            return ""
        labels = {
            "overview": "Overview",
            "security": "Security",
            "mobile_inventory": "Mobile Inventory",
            "mobile_devices": "Mobile Devices",
            "policies": "Policies",
            "macos_profiles": "macOS Profiles",
            "ios_profiles": "iOS Profiles",
            "smart_groups": "Smart Groups",
            "scripts": "Scripts",
            "packages": "Packages",
            "categories": "Categories",
            "device_enrollments": "ADE Instances",
            "sites": "Sites",
            "buildings": "Buildings",
            "departments": "Departments",
        }
        rows = []
        for key, label in labels.items():
            info = fetch_status.get(key, {})
            badge_cls = self._source_badge_class(info.get("source_mode"), info.get("status"))
            source_label = self._source_mode_label(info.get("source_mode"))
            detail = info.get("detail", "")
            cached_name = Path(str(info.get("cached_path") or "")).name
            if cached_name:
                detail = cached_name
            if not detail:
                detail = f"{int(info.get('records', 0))} records"
            rows.append(
                "<tr>"
                f"<td>{self._html_text(label)}</td>"
                f"<td><span class='badge {badge_cls}'>{self._html_text(source_label)}</span></td>"
                f"<td>{self._html_text(detail)}</td>"
                "</tr>"
            )
        return f"""<div class="section-title">Report Sources</div>
<div class="card card-sm">
  <div class="table-note" style="margin-bottom:10px">
    Live and cached source state is shown here so partial-data runs are visible in the report itself.
  </div>
  <table class="data-table">
    <thead><tr><th>Dataset</th><th>Source</th><th>Detail</th></tr></thead>
    <tbody>{''.join(rows)}</tbody>
  </table>
</div>"""

    def _css(self) -> str:
        """Return the embedded CSS block.

        Colour palette and layout are adapted from DevliegereM's jamf-html-reports
        (github.com/DevliegereM/) with minor modifications for
        the Python port (no emojis; text-only status indicators).
        """
        accent = (self._config.get("branding", "accent_color") or "").strip()
        accent_dark = (self._config.get("branding", "accent_dark") or "").strip()
        overrides = ""
        if accent or accent_dark:
            parts = []
            if accent_dark:
                parts.append(f"    --blue-dark: {accent_dark};")
            if accent:
                parts.append(f"    --blue:      {accent};")
            overrides = ":root {\n" + "\n".join(parts) + "\n}\n"
        return overrides + """
:root {
    --blue-dark: #004165;
    --blue:      #0076B6;
    --blue-lt:   #e8f4fb;
    --green:     #22c55e;
    --amber:     #f59e0b;
    --red:       #ef4444;
    --purple:    #8b5cf6;
    --cyan:      #06b6d4;
    --bg:        #f0f4f8;
    --surface:   #ffffff;
    --surface-2: #f8fafc;
    --border:    #e2e8f0;
    --text:      #1e293b;
    --muted:     #64748b;
    --radius:    10px;
    --shadow:    0 1px 4px rgba(0,0,0,.08);
}
body.dark {
    --blue-dark: #5bb8e8;
    --blue: #7dcbf0;
    --blue-lt: #0f2d40;
    --green: #4ade80;
    --amber: #fbbf24;
    --red: #f87171;
    --purple: #a78bfa;
    --cyan: #22d3ee;
    --bg: #0f172a;
    --surface: #1e293b;
    --surface-2: #162032;
    --border: #334155;
    --text: #e2e8f0;
    --muted: #94a3b8;
    --shadow: 0 1px 4px rgba(0,0,0,.4);
}
body.dark .topbar { background: #0d1f30; border-bottom: 1px solid #334155; }
body.dark .data-table th { background: var(--surface-2); }
body.dark .data-table tr:hover td { background: #243447; }
body.dark .badge-ok  { background: #14532d; color: #4ade80; }
body.dark .badge-warn{ background: #451a03; color: #fbbf24; }
body.dark .badge-err { background: #450a0a; color: #f87171; }
body.dark .badge-dim { background: #1e293b; color: #94a3b8; }
body.dark .badge-blue{ background: #0f2d40; color: #7dcbf0; }
body.dark .sec-bar-track { background: #334155; }
body.dark .cat-toggle,
body.dark .sg-node,
body.dark .link-card { background: var(--surface-2); border-color: var(--border); }
body.dark .item-node:hover,
body.dark .link-card:hover { background: #1a3a52; }
body.dark .feat-on { background: #14532d; color: #4ade80; border-color: #166534; }
body.dark .feat-off { background: #1e293b; color: #94a3b8; border-color: #334155; }
body.dark .stat-value { color: var(--blue); }
body.dark .chart-title { color: var(--blue); }
body.dark .tree-tab.active { background: #10263a; }
body.dark .tree-search,
body.dark .table-action,
body.dark .table-sort { color: var(--text); }
body.dark .os-table th,
body.dark .os-table td { border-color: var(--border); }
*,*::before,*::after { box-sizing: border-box; margin: 0; padding: 0; }
body {
    font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
    background: var(--bg);
    color: var(--text);
    font-size: 14px;
    line-height: 1.5;
}
a { color: var(--blue); text-decoration: none; }
a:hover { text-decoration: underline; }
.topbar {
    background: var(--blue-dark);
    color: #fff;
    padding: 14px 24px;
    display: flex;
    align-items: center;
    gap: 16px;
    flex-wrap: wrap;
    justify-content: space-between;
    position: sticky;
    top: 0;
    z-index: 100;
}
.topbar-brand { font-size: 1.15rem; font-weight: 700; letter-spacing: -.3px; }
.topbar-meta { font-size: .8rem; opacity: .8; text-align: right; }
.topbar-meta strong { opacity: 1; font-size: .9rem; color: #fff; }
.dark-toggle {
    background: rgba(255,255,255,.12);
    border: 1px solid rgba(255,255,255,.25);
    color: #fff;
    border-radius: 20px;
    padding: 4px 12px;
    font-size: .78rem;
    font-weight: 600;
    cursor: pointer;
    transition: background .2s;
    white-space: nowrap;
}
.dark-toggle:hover { background: rgba(255,255,255,.22); }
.page { max-width: 1400px; margin: 0 auto; padding: 20px 20px 40px; }
.section-block { margin-top: 22px; padding-top: 8px; border-top: 1px solid var(--border); }
.section-block:first-of-type { border-top: none; padding-top: 0; }
.section-block-title {
    font-size: 1.4rem;
    font-weight: 700;
    letter-spacing: -.02em;
    color: var(--blue-dark);
    margin-bottom: 4px;
}
.section-block-subtitle { font-size: .88rem; color: var(--muted); margin-bottom: 16px; }
.section-title {
    font-size: .7rem;
    font-weight: 700;
    letter-spacing: .08em;
    text-transform: uppercase;
    color: var(--muted);
    margin: 28px 0 10px;
}
.grid { display: grid; gap: 14px; }
.grid-2 { grid-template-columns: repeat(2, 1fr); }
.grid-3 { grid-template-columns: repeat(3, 1fr); }
.grid-4 { grid-template-columns: repeat(4, 1fr); }
.grid-5 { grid-template-columns: repeat(5, 1fr); }
.grid-6 { grid-template-columns: repeat(6, 1fr); }
@media(max-width:1100px) {
    .grid-6,.grid-5 { grid-template-columns: repeat(3,1fr); }
    .grid-4 { grid-template-columns: repeat(2,1fr); }
}
@media(max-width:700px) {
    .grid-2,.grid-3,.grid-4,.grid-5,.grid-6 { grid-template-columns: 1fr; }
}
.card,
.chart-card {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: var(--radius);
    padding: 18px 20px;
    box-shadow: var(--shadow);
}
.card-sm { padding: 12px 14px; }
.stat-card { display: flex; flex-direction: column; gap: 4px; }
.stat-label {
    font-size: .73rem;
    font-weight: 600;
    text-transform: uppercase;
    letter-spacing: .05em;
    color: var(--muted);
}
.stat-value { font-size: 2rem; font-weight: 700; color: var(--blue-dark); line-height: 1.1; }
.stat-sub { font-size: .75rem; color: var(--muted); }
.stat-link { font-size: .75rem; color: var(--blue); margin-top: 4px; }
.badge {
    display: inline-block;
    font-size: .7rem;
    font-weight: 600;
    padding: 2px 8px;
    border-radius: 20px;
    letter-spacing: .03em;
    white-space: nowrap;
}
.badge-ok   { background: #dcfce7; color: #15803d; }
.badge-warn { background: #fef3c7; color: #92400e; }
.badge-err  { background: #fee2e2; color: #991b1b; }
.badge-dim  { background: #f1f5f9; color: var(--muted); }
.badge-blue { background: var(--blue-lt); color: var(--blue-dark); }
.health-strip {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: var(--radius);
    padding: 14px 20px;
    display: flex;
    flex-wrap: wrap;
    gap: 10px 24px;
    align-items: center;
}
.health-item { display: flex; align-items: center; gap: 6px; font-size: .8rem; }
.health-label { color: var(--muted); font-weight: 500; }
.chart-title { font-size: .85rem; font-weight: 700; color: var(--blue-dark); margin-bottom: 14px; }
.chart-sub { font-size: .73rem; color: var(--muted); margin-top: 8px; }
.chart-wrap { position: relative; height: 220px; }
.chart-wrap-lg { position: relative; height: 280px; }
.trend-svg { display: block; width: 100%; height: auto; overflow: visible; }
.svg-axis { font-size: 10px; fill: var(--muted); font-family: inherit; }
.chart-legend {
    display: flex;
    flex-wrap: wrap;
    gap: 8px 12px;
    margin-top: 10px;
}
.chart-legend-item {
    display: inline-flex;
    align-items: center;
    gap: 6px;
    font-size: .75rem;
    color: var(--muted);
}
.chart-legend-swatch {
    width: 10px;
    height: 10px;
    border-radius: 999px;
    display: inline-block;
}
.sec-bar-row { margin-bottom: 12px; }
.sec-bar-header { display: flex; justify-content: space-between; align-items: center; margin-bottom: 4px; }
.sec-bar-name { font-size: .8rem; font-weight: 600; color: var(--text); }
.sec-bar-pct  { font-size: .8rem; color: var(--muted); }
.sec-bar-track { background: #e2e8f0; border-radius: 4px; height: 10px; overflow: hidden; }
.sec-bar-fill { height: 100%; border-radius: 4px; transition: width .4s ease; }
.fill-fv  { background: var(--green); }
.fill-gk  { background: var(--blue); }
.fill-sip { background: var(--purple); }
.fill-fw  { background: var(--amber); }
.data-table { width: 100%; border-collapse: collapse; font-size: .8rem; }
.data-table th {
    text-align: left;
    padding: 7px 10px;
    background: var(--surface-2);
    border-bottom: 2px solid var(--border);
    color: var(--muted);
    font-weight: 600;
    text-transform: uppercase;
    font-size: .68rem;
    letter-spacing: .05em;
}
.data-table td { padding: 7px 10px; border-bottom: 1px solid var(--border); vertical-align: top; }
.data-table tr:last-child td { border-bottom: 0; }
.data-table tr:hover td { background: #fafbfc; }
.data-table .val { font-weight: 600; color: var(--blue-dark); text-align: right; }
.data-table .val-warn { color: var(--amber); }
.data-table .val-err  { color: var(--red); }
.data-table .val-ok   { color: var(--green); }
.overview-section-title {
    font-size: .72rem;
    font-weight: 700;
    text-transform: uppercase;
    letter-spacing: .07em;
    color: var(--muted);
    padding: 10px 10px 4px;
    border-top: 1px solid var(--border);
    margin-top: 6px;
}
.overview-section-title:first-child { border-top: none; margin-top: 0; }
.table-tools {
    display: flex;
    align-items: center;
    gap: 10px;
    flex-wrap: wrap;
    padding: 10px 14px;
    border-bottom: 1px solid var(--border);
    background: var(--surface-2);
}
.table-note { font-size: .75rem; color: var(--muted); }
.table-action,
.table-sort {
    background: transparent;
    border: 0;
    color: var(--blue-dark);
    cursor: pointer;
    font: inherit;
    padding: 0;
}
.table-action {
    border: 1px solid var(--border);
    border-radius: 6px;
    padding: 6px 10px;
    background: var(--surface);
    font-size: .76rem;
    font-weight: 600;
}
.table-sort:hover,
.table-action:hover { color: var(--blue); }
.table-wrap { overflow-x: auto; }
.tree-tabs { display: flex; flex-wrap: wrap; gap: 8px; margin-bottom: 14px; }
.tree-tab {
    border: 1px solid var(--border);
    background: var(--surface-2);
    color: var(--muted);
    border-radius: 999px;
    padding: 6px 12px;
    cursor: pointer;
    font-size: .79rem;
    font-weight: 600;
}
.tree-tab.active {
    color: var(--blue-dark);
    border-color: var(--blue);
    background: var(--blue-lt);
}
.tree-pane { display: none; }
.tree-pane.active { display: block; }
.tree-search {
    width: 100%;
    padding: 8px 10px;
    border: 1px solid var(--border);
    border-radius: 8px;
    font-size: .82rem;
    margin-bottom: 10px;
    background: var(--surface);
    color: var(--text);
}
.tree-summary { font-size: .75rem; color: var(--muted); margin-bottom: 10px; }
.cat-node { margin-bottom: 6px; }
.cat-toggle {
    cursor: pointer;
    background: var(--surface-2);
    border: 1px solid var(--border);
    border-radius: 8px;
    padding: 8px 12px;
    display: flex;
    justify-content: space-between;
    align-items: center;
    gap: 8px;
    font-weight: 600;
    font-size: .82rem;
    user-select: none;
}
.cat-toggle:hover { border-color: var(--blue); }
.cat-label { display: flex; align-items: center; gap: 8px; min-width: 0; }
.cat-caret { display: inline-block; width: 10px; color: var(--muted); transition: transform .18s; }
.cat-caret.open { transform: rotate(90deg); }
.cat-count {
    font-size: .72rem;
    color: var(--blue-dark);
    background: var(--blue-lt);
    border-radius: 999px;
    padding: 2px 8px;
    font-weight: 700;
}
.cat-items { display: none; padding: 6px 0 8px 24px; }
.cat-items.open { display: block; }
.item-node,
.sg-node {
    padding: 6px 10px;
    font-size: .79rem;
    border-radius: 6px;
    border: 1px solid var(--border);
    background: var(--surface);
    margin-bottom: 6px;
}
.item-node:hover,
.sg-node:hover { background: var(--blue-lt); }
.item-hidden { display: none !important; }
.os-table { width: 100%; border-collapse: collapse; font-size: .76rem; margin-top: 12px; }
.os-table th,.os-table td { padding: 5px 8px; border-bottom: 1px solid var(--border); }
.os-table th {
    color: var(--muted);
    font-weight: 600;
    text-transform: uppercase;
    font-size: .65rem;
    background: var(--surface-2);
}
.os-table td:last-child { text-align: right; font-weight: 600; color: var(--blue-dark); }
.os-dot { display: inline-block; width: 9px; height: 9px; border-radius: 50%; margin-right: 6px; }
.links-grid {
    display: grid;
    grid-template-columns: repeat(auto-fill, minmax(180px, 1fr));
    gap: 10px;
}
.link-card {
    display: flex;
    align-items: center;
    justify-content: center;
    text-align: center;
    min-height: 50px;
    padding: 10px 14px;
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 10px;
    color: var(--text);
    font-size: .8rem;
    font-weight: 600;
}
.link-card:hover {
    background: var(--blue-lt);
    border-color: var(--blue);
    text-decoration: none;
}
.empty-note { color: var(--muted); font-size: .82rem; }
	.footer {
	    text-align: center;
    font-size: .72rem;
    color: var(--muted);
    margin-top: 40px;
    padding-top: 16px;
    border-top: 1px solid var(--border);
}
	"""

    def _js(self) -> str:
        """Return the embedded JavaScript block for HTML interactivity."""
        return """
(() => {
  const key = 'jamfReportsHtmlDarkMode';
  const toggle = document.getElementById('darkToggle');
  const applyDark = (enabled) => {
    document.body.classList.toggle('dark', enabled);
    if (toggle) toggle.textContent = enabled ? 'Light mode' : 'Dark mode';
    try { localStorage.setItem(key, enabled ? '1' : '0'); } catch (err) {}
  };
  let saved = null;
  try { saved = localStorage.getItem(key); } catch (err) {}
  const prefersDark = window.matchMedia
    && window.matchMedia('(prefers-color-scheme: dark)').matches;
  applyDark(saved === null ? prefersDark : saved === '1');
  if (toggle) toggle.addEventListener('click', () => {
    applyDark(!document.body.classList.contains('dark'));
  });
})();

document.querySelectorAll('.cat-toggle').forEach((toggle) => {
  toggle.addEventListener('click', () => {
    const items = toggle.nextElementSibling;
    const caret = toggle.querySelector('.cat-caret');
    if (!items) return;
    const open = items.classList.toggle('open');
    if (caret) caret.classList.toggle('open', open);
  });
});

document.querySelectorAll('.tree-tab').forEach((tab) => {
  tab.addEventListener('click', () => {
    const targetId = tab.getAttribute('data-target');
    if (!targetId) return;
    document.querySelectorAll('.tree-tab').forEach((el) => el.classList.remove('active'));
    document.querySelectorAll('.tree-pane').forEach((el) => el.classList.remove('active'));
    tab.classList.add('active');
    const pane = document.getElementById(targetId);
    if (pane) pane.classList.add('active');
  });
});

const filterTreePane = (container, query) => {
  container.querySelectorAll('.cat-node').forEach((node) => {
    const items = Array.from(node.querySelectorAll('.item-node'));
    let visible = false;
    items.forEach((item) => {
      const match = !query || (item.dataset.name || '').includes(query);
      item.classList.toggle('item-hidden', !match);
      if (match) visible = true;
    });
    node.classList.toggle('item-hidden', !visible);
    const children = node.querySelector('.cat-items');
    const caret = node.querySelector('.cat-caret');
    if (children && caret) {
      children.classList.toggle('open', query ? visible : children.classList.contains('open'));
      caret.classList.toggle('open', children.classList.contains('open'));
    }
  });
};

const filterFlatPane = (container, query) => {
  container.querySelectorAll('.sg-node').forEach((node) => {
    const match = !query || (node.dataset.name || '').includes(query);
    node.classList.toggle('item-hidden', !match);
  });
};

document.querySelectorAll('.tree-search').forEach((input) => {
  input.addEventListener('input', () => {
    const targetId = input.getAttribute('data-filter-target');
    const kind = input.getAttribute('data-filter-kind');
    if (!targetId || !kind) return;
    const container = document.getElementById(targetId);
    if (!container) return;
    const query = input.value.toLowerCase().trim();
    if (kind === 'tree') filterTreePane(container, query);
    if (kind === 'flat') filterFlatPane(container, query);
  });
});

(() => {
  const tbody = document.getElementById('flaggedBody');
  const searchInput = document.getElementById('flaggedSearch');
  const exportBtn = document.getElementById('flaggedExport');
  if (!tbody || !searchInput || !exportBtn) return;

  const emptyRow = tbody.querySelector('.flagged-empty');
  let sortKey = 'name';
  let sortAsc = true;
  const safeCsvValue = (value) => {
    const cleaned = String(value ?? '').replace(/\r?\n/g, ' ').trim();
    return /^[=+-@]/.test(cleaned) ? "'" + cleaned : cleaned;
  };

  const deviceRows = () => Array.from(tbody.querySelectorAll('tr.flagged-device'));
  const applyFilter = () => {
    const query = searchInput.value.toLowerCase().trim();
    let visible = 0;
    deviceRows().forEach((row) => {
      const haystack = [row.dataset.name, row.dataset.serial, row.dataset.os].join(' ');
      const match = !query || haystack.includes(query);
      row.style.display = match ? '' : 'none';
      if (match) visible += 1;
    });
    if (emptyRow) emptyRow.hidden = visible !== 0;
  };

  document.querySelectorAll('[data-flagged-sort]').forEach((button) => {
    button.addEventListener('click', () => {
      const key = button.getAttribute('data-flagged-sort') || 'name';
      sortAsc = sortKey === key ? !sortAsc : true;
      sortKey = key;
      const rows = deviceRows().sort((a, b) => {
        const left = (a.dataset[sortKey] || '').toLowerCase();
        const right = (b.dataset[sortKey] || '').toLowerCase();
        return sortAsc ? left.localeCompare(right) : right.localeCompare(left);
      });
      rows.forEach((row) => tbody.appendChild(row));
      if (emptyRow) tbody.appendChild(emptyRow);
      applyFilter();
    });
  });

  searchInput.addEventListener('input', applyFilter);
  exportBtn.addEventListener('click', () => {
    const headers = ['Device', 'Serial', 'macOS', 'FileVault', 'Gatekeeper', 'SIP', 'Firewall'];
    const rows = [headers];
    deviceRows()
      .filter((row) => row.style.display !== 'none')
      .forEach((row) => {
        const cells = Array.from(row.querySelectorAll('td')).slice(0, 7);
        rows.push(cells.map((cell) => safeCsvValue(cell.textContent)));
      });
    const csv = rows
      .map((row) => row.map((value) => `"${String(value).replace(/"/g, '""')}"`).join(','))
      .join('\r\n');
    const blob = new Blob([csv], { type: 'text/csv;charset=utf-8' });
    const url = URL.createObjectURL(blob);
    const a = document.createElement('a');
    a.href = url;
    a.download = 'flagged-devices-' + new Date().toISOString().slice(0, 10) + '.csv';
    document.body.appendChild(a);
    a.click();
    document.body.removeChild(a);
    URL.revokeObjectURL(url);
  });

  applyFilter();
})();
"""

    def _logo_html(self) -> str:
        """Return an inline <img> tag with the logo base64-encoded, or empty string."""
        logo_path = self._config.resolve_path("branding", "logo_path")
        if not logo_path or not logo_path.exists():
            return ""
        try:
            import base64
            import mimetypes
            mime = mimetypes.guess_type(str(logo_path))[0] or "image/png"
            b64 = base64.b64encode(logo_path.read_bytes()).decode()
            return (
                f'<img src="data:{mime};base64,{b64}" alt="" '
                f'style="height:28px;margin-right:10px;vertical-align:middle;'
                f'border-radius:4px">'
            )
        except Exception:
            return ""

    def _render_stat_card(
        self,
        label: str,
        value: str,
        sub: str = "",
        link_url: str = "",
        link_text: str = "",
    ) -> str:
        """Render a single stat card HTML block."""
        link_html = ""
        safe_link = self._safe_href(link_url) if link_url else ""
        if safe_link and safe_link != "#" and link_text:
            link_html = (
                '<div class="stat-link"><a href="'
                f'{self._html_text(safe_link)}"'
                ' target="_blank" rel="noopener noreferrer">'
                f'{self._html_text(link_text)}</a></div>'
            )
        return f"""<div class="card stat-card">
  <div class="stat-label">{self._html_text(label)}</div>
  <div class="stat-value">{self._html_text(value, "N/A")}</div>
  {f'<div class="stat-sub">{self._html_text(sub)}</div>' if sub else ''}
  {link_html}
</div>"""

    def _render_sec_bar(self, name: str, pct: float, css_class: str) -> str:
        """Render a security posture progress bar."""
        return f"""<div class="sec-bar-row">
  <div class="sec-bar-header">
    <span class="sec-bar-name">{self._html_text(name)}</span>
    <span class="sec-bar-pct">{pct:.1f}%</span>
  </div>
  <div class="sec-bar-track">
    <div class="sec-bar-fill {css_class}" style="width:{min(pct, 100):.1f}%"></div>
  </div>
</div>"""

    def _render_hierarchy_tree_panel(
        self,
        pane_id: str,
        search_placeholder: str,
        groups: list[dict[str, Any]],
    ) -> str:
        """Render one grouped hierarchy pane."""
        summary = f"{len(groups)} categories · {sum(group['count'] for group in groups)} items"
        if not groups:
            body = '<p class="empty-note">No data available.</p>'
        else:
            chunks = []
            for group in groups:
                items = "".join(
                    '<div class="item-node" data-name="'
                    f'{self._html_text(name.lower())}">{self._html_text(name)}</div>'
                    for name in group["items"]
                )
                chunks.append(
                    '<div class="cat-node"><div class="cat-toggle"><span class="cat-label">'
                    '<span class="cat-caret">▶</span>'
                    f'<span>{self._html_text(group["category"])}</span>'
                    f'</span><span class="cat-count">{group["count"]}</span></div>'
                    f'<div class="cat-items">{items}</div></div>'
                )
            body = "".join(chunks)
        return f"""<div class="tree-pane" id="{self._html_text(pane_id)}">
  <input class="tree-search" type="search"
    placeholder="{self._html_text(search_placeholder)}"
    data-filter-target="{self._html_text(pane_id)}-list" data-filter-kind="tree">
  <div class="tree-summary">{self._html_text(summary)}</div>
  <div id="{self._html_text(pane_id)}-list">{body}</div>
</div>"""

    def _render_hierarchy_flat_panel(
        self,
        pane_id: str,
        search_placeholder: str,
        items: list[str],
        label: str,
    ) -> str:
        """Render one flat searchable hierarchy pane."""
        summary = f"{len(items)} {label}"
        if not items:
            body = '<p class="empty-note">No data available.</p>'
        else:
            body = "".join(
                '<div class="sg-node" data-name="'
                f'{self._html_text(item.lower())}">{self._html_text(item)}</div>'
                for item in items
            )
        return f"""<div class="tree-pane" id="{self._html_text(pane_id)}">
  <input class="tree-search" type="search"
    placeholder="{self._html_text(search_placeholder)}"
    data-filter-target="{self._html_text(pane_id)}-list" data-filter-kind="flat">
  <div class="tree-summary">{self._html_text(summary)}</div>
  <div id="{self._html_text(pane_id)}-list">{body}</div>
</div>"""

    def _render_hierarchy_tabs(
        self,
        pol_groups: list[dict[str, Any]],
        mcp_groups: list[dict[str, Any]],
        icp_groups: list[dict[str, Any]],
        scr_groups: list[dict[str, Any]],
        pkg_groups: list[dict[str, Any]],
        smart_groups: list[str],
    ) -> str:
        """Render the tabbed searchable deployment hierarchy."""
        tabs = [
            ("hier-policies", "Policies", sum(group["count"] for group in pol_groups)),
            ("hier-macos", "macOS Profiles", sum(group["count"] for group in mcp_groups)),
            ("hier-ios", "iOS Profiles", sum(group["count"] for group in icp_groups)),
            ("hier-scripts", "Scripts", sum(group["count"] for group in scr_groups)),
            ("hier-packages", "Packages", sum(group["count"] for group in pkg_groups)),
            ("hier-smart-groups", "Smart Groups", len(smart_groups)),
        ]
        tab_html = "".join(
            '<button type="button" class="tree-tab'
            f'{" active" if idx == 0 else ""}" data-target="{self._html_text(tab_id)}">'
            f'{self._html_text(label)} ({count})</button>'
            for idx, (tab_id, label, count) in enumerate(tabs)
        )
        panes = [
            self._render_hierarchy_tree_panel("hier-policies", "Search policies…", pol_groups),
            self._render_hierarchy_tree_panel("hier-macos", "Search macOS profiles…", mcp_groups),
            self._render_hierarchy_tree_panel("hier-ios", "Search iOS profiles…", icp_groups),
            self._render_hierarchy_tree_panel("hier-scripts", "Search scripts…", scr_groups),
            self._render_hierarchy_tree_panel("hier-packages", "Search packages…", pkg_groups),
            self._render_hierarchy_flat_panel(
                "hier-smart-groups",
                "Search smart groups…",
                smart_groups,
                "smart groups",
            ),
        ]
        panes[0] = panes[0].replace('class="tree-pane"', 'class="tree-pane active"', 1)
        return f"""<div class="section-title">Deployment Hierarchy</div>
<div class="card">
  <div class="tree-tabs">{tab_html}</div>
  {''.join(panes)}
</div>"""

    def _render_overview_table(self, overview: Any, console_url: str) -> str:
        """Render the full overview data as a grouped table with deep links."""
        sections = self._overview_sections(overview)
        if not sections:
            return ""
        links = self._resource_links(console_url)
        rows_html = ""
        for section_name, items in sections.items():
            rows_html += (
                f'<tr><td colspan="3" class="overview-section-title">'
                f'{self._html_text(section_name)}</td></tr>'
            )
            for item in items:
                badge = self._status_badge_html(item["status"])
                resource = self._html_text(item["resource"])
                link = links.get(item["resource"], "")
                if link:
                    resource = (
                        f'<a href="{self._html_text(self._safe_href(link))}" target="_blank" '
                        f'rel="noopener noreferrer">{resource}</a>'
                    )
                rows_html += (
                    f"<tr><td>{resource}</td>"
                    f"<td class='val'>{self._html_text(item['value'])}</td>"
                    f"<td>{badge}</td></tr>"
                )
        return f"""<div class="section-title">Full Overview</div>
<div class="card">
  <table class="data-table">
    <thead><tr>
      <th>Resource</th><th style="text-align:right">Value</th><th>Status</th>
    </tr></thead>
    <tbody>{rows_html}</tbody>
  </table>
</div>"""

    def _render_quick_links(self, console_url: str) -> str:
        """Render quick links back into the Jamf console."""
        links = self._quick_links(console_url)
        if not links:
            return ""
        cards = "".join(
            f'<a class="link-card" href="{self._html_text(self._safe_href(url))}" '
            f'target="_blank" rel="noopener noreferrer">{self._html_text(label)}</a>'
            for label, url in links
        )
        return f"""<div class="section-title">Quick Links</div>
<div class="links-grid">{cards}</div>"""

    def _render_flagged_table(self, flagged: list[dict[str, Any]], console_url: str) -> str:
        """Render the searchable and sortable flagged-devices table."""
        if not flagged:
            return """<div class="section-title">Devices with Security Issues</div>
<div class="card"><p class="empty-note">No devices with security issues found.</p></div>"""

        rows = []
        for dev in flagged:
            fv = self._html_text(dev["filevault"])
            gk = self._html_text(dev["gatekeeper"])
            sip = self._html_text(dev["sip"])
            fw = self._html_text(dev["firewall"])
            fv_cls = "val-ok" if fv.upper() == "ENCRYPTED" else "val-err"
            gk_cls = "val-ok" if gk.upper() not in ("DISABLED",) else "val-err"
            sip_cls = "val-ok" if sip.upper() in ("ENABLED",) else "val-err"
            fw_cls = "val-err" if fw.casefold() in {"no", "false", "off", "disabled"} else "val-ok"
            query_value = str(dev.get("serial") or dev.get("name") or "")
            open_link = ""
            if console_url and query_value:
                query = urllib.parse.quote(query_value, safe="")
                open_link = f"{console_url}/computers.html?query={query}&queryType=COMPUTERS&version="
            link_html = (
                f'<a href="{self._html_text(self._safe_href(open_link))}" target="_blank" '
                'rel="noopener noreferrer">Open</a>'
                if open_link
                else "—"
            )
            rows.append(
                "<tr class='flagged-device' data-name='"
                f"{self._html_text(str(dev['name']).lower())}' data-serial='"
                f"{self._html_text(str(dev['serial']).lower())}' data-os='"
                f"{self._html_text(str(dev['os']).lower())}'>"
                f"<td>{self._html_text(dev['name'])}</td>"
                f"<td>{self._html_text(dev['serial'])}</td>"
                f"<td>{self._html_text(dev['os'])}</td>"
                f"<td class='{fv_cls}'>{fv}</td>"
                f"<td class='{gk_cls}'>{gk}</td>"
                f"<td class='{sip_cls}'>{sip}</td>"
                f"<td class='{fw_cls}'>{fw}</td>"
                f"<td>{link_html}</td></tr>"
            )

        return f"""<div class="section-title">Devices with Security Issues ({len(flagged)})</div>
<div class="card" style="padding:0">
  <div class="table-tools">
    <input class="tree-search" id="flaggedSearch" type="search"
      placeholder="Filter by name, serial, or macOS…" style="margin:0;max-width:340px">
    <span class="table-note">Click a column header to sort.</span>
    <button type="button" class="table-action" id="flaggedExport">Export CSV</button>
  </div>
  <div class="table-wrap">
    <table class="data-table" id="flaggedTable">
      <thead><tr>
        <th><button type="button" class="table-sort" data-flagged-sort="name">Device ↕</button></th>
        <th><button type="button" class="table-sort" data-flagged-sort="serial">Serial ↕</button></th>
        <th><button type="button" class="table-sort" data-flagged-sort="os">macOS ↕</button></th>
        <th>FileVault</th><th>Gatekeeper</th><th>SIP</th><th>Firewall</th><th>Link</th>
      </tr></thead>
      <tbody id="flaggedBody">
        {''.join(rows)}
        <tr class="flagged-empty" hidden><td colspan="8" class="empty-note">No matching devices.</td></tr>
      </tbody>
    </table>
  </div>
</div>"""

    def _render_org_table(self, label: str, items: list[str]) -> str:
        """Render a small org-item name list as a card."""
        if not items:
            return ""
        rows = "".join(f"<tr><td>{self._html_text(name)}</td></tr>" for name in items)
        return f"""<div class="card card-sm" style="margin-bottom:10px">
  <div style="font-size:.75rem;font-weight:700;text-transform:uppercase;
      letter-spacing:.05em;color:var(--muted);margin-bottom:8px">{self._html_text(label)}</div>
  <table class="data-table"><tbody>{rows}</tbody></table>
</div>"""

    def _render_counter_table(
        self,
        title: str,
        left_header: str,
        counter: Counter,
        max_rows: int = 8,
    ) -> str:
        """Render a compact frequency table card."""
        if not counter:
            return f"""<div class="card card-sm">
  <div class="chart-title">{self._html_text(title)}</div>
  <p class="empty-note">No data available.</p>
</div>"""

        rows = "".join(
            f"<tr><td>{self._html_text(str(label))}</td><td class='val'>{count}</td></tr>"
            for label, count in counter.most_common(max_rows)
        )
        return f"""<div class="card card-sm">
  <div class="chart-title">{self._html_text(title)}</div>
  <table class="data-table">
    <thead><tr><th>{self._html_text(left_header)}</th><th style="text-align:right">Count</th></tr></thead>
    <tbody>{rows}</tbody>
  </table>
</div>"""

    def _render_mobile_inventory_table(
        self,
        rows: list[dict[str, Any]],
        stale_days: int,
    ) -> str:
        """Render a table of the mobile devices needing the most review."""
        if not rows:
            return """<div class="section-title">Mobile Inventory Review</div>
<div class="card"><p class="empty-note">No mobile inventory rows were available.</p></div>"""

        ranked = sorted(
            rows,
            key=lambda row: row.get("Days Since Inventory")
            if isinstance(row.get("Days Since Inventory"), int) else -1,
            reverse=True,
        )[:12]

        body = ""
        for row in ranked:
            days = row.get("Days Since Inventory")
            days_label = str(days) if isinstance(days, int) else "N/A"
            if isinstance(days, int) and days > stale_days:
                days_html = f"<td class='val-err'>{days_label}</td>"
            elif isinstance(days, int) and days > max(1, stale_days // 2):
                days_html = f"<td class='val-warn'>{days_label}</td>"
            else:
                days_html = f"<td>{days_label}</td>"
            user_label = row.get("Username") or row.get("Email") or "Unassigned"
            body += (
                f"<tr><td>{self._html_text(str(row.get('Device Name', '')))}</td>"
                f"<td>{self._html_text(str(row.get('Device Family', '')))}</td>"
                f"<td>{self._html_text(str(row.get('OS Version', '')))}</td>"
                f"<td>{self._html_text(user_label)}</td>"
                f"{days_html}"
                f"<td>{self._html_text(str(row.get('Managed', '')))}</td>"
                f"<td>{self._html_text(str(row.get('Supervised', '')))}</td></tr>"
            )
        return f"""<div class="section-title">Mobile Inventory Review</div>
<div class="card">
  <table class="data-table">
    <thead><tr>
      <th>Device</th><th>Family</th><th>OS</th><th>User</th>
      <th>Days Since Inventory</th><th>Managed</th><th>Supervised</th>
    </tr></thead>
    <tbody>{body}</tbody>
  </table>
</div>"""

    def _render_trends_section(self, trends: dict[str, Any]) -> str:
        """Render the HTML trend section when snapshot history is available."""
        cards = []
        adoption = trends.get("adoption") or {}
        security = trends.get("security") or {}
        if adoption.get("labels"):
            cards.append(
                f"""<div class="chart-card">
  <div class="chart-title">macOS Adoption Trend</div>
  {self._render_line_chart_svg(adoption)}
  <div class="chart-sub">Source: {self._html_text(adoption.get("source", ""))}</div>
</div>"""
            )
        if security.get("labels"):
            cards.append(
                f"""<div class="chart-card">
  <div class="chart-title">Security Posture Trend</div>
  {self._render_line_chart_svg(security, percent_scale=True)}
  <div class="chart-sub">Source: {self._html_text(security.get("source", ""))}</div>
</div>"""
            )
        if not cards:
            return ""
        return f"""<div class="section-title">Trends</div>
<div class="grid grid-2">{''.join(cards)}</div>"""

    def _render(self, data: dict[str, Any]) -> str:
        """Assemble the full HTML document from fetched data."""
        ov = data.get("overview", [])
        sec = data.get("security", [])
        mobile_items = data.get("mobile_inventory") or data.get("mobile_devices", [])
        mobile_rows = self._mobile_rows(mobile_items)
        mobile_summary = _summarize_mobile_inventory(mobile_rows)
        stale_days = int(self._config.thresholds.get("stale_device_days", 30))
        mobile_stale = sum(
            1
            for row in mobile_rows
            if isinstance(row.get("Days Since Inventory"), int)
            and row["Days Since Inventory"] > stale_days
        )

        instance_url = self._ov(ov, "Server URL")
        console_url = self._safe_base_url(instance_url)
        jamf_version = self._ov(ov, "Jamf Pro Version")
        health_status = self._ov(ov, "Health Status")
        active_alerts = self._ov(ov, "Active Alerts")
        managed_computers = self._ov(ov, "Managed Computers")
        unmanaged_computers = self._ov(ov, "Unmanaged Computers")
        managed_devices = self._ov(ov, "Managed Devices")
        checkin_freq = self._ov(ov, "Check-In Frequency")
        dep_token_exp = self._ov(ov, "DEP Token Expires")
        ca_expires = self._ov(ov, "Built-in CA Expires")
        ade_sync = self._ov(ov, "DEP Sync Status")
        vpp_locations = self._ov(ov, "VPP Locations")
        comp_prestages = self._ov(ov, "Computer Prestages")
        md_prestages = self._ov(ov, "Mobile Device Prestages")
        app_installers = self._ov(ov, "App Installers")
        webhooks = self._ov(ov, "Webhooks")
        jcds_files = self._ov(ov, "JCDS Files")
        patch_titles = self._ov(ov, "Patch Titles")
        ldap_servers = self._ov(ov, "LDAP/IdP Servers")
        mdm_renew_comp = self._ov(ov, "MDM Auto Renew (Computers)")
        mdm_renew_md = self._ov(ov, "MDM Auto Renew (Mobile)")
        sso = self._ov(ov, "SSO (SAML)")
        sso_jamf = self._ov(ov, "Jamf SSO")
        ade_instances = self._ov(ov, "DEP Instances")

        fv_pct = self._to_float(self._sec(sec, "filevault_encrypted_pct"))
        gk_pct = self._to_float(self._sec(sec, "gatekeeper_enabled_pct"))
        sip_pct = self._to_float(self._sec(sec, "sip_enabled_pct"))
        fw_pct = self._to_float(self._sec(sec, "firewall_enabled_pct"))
        total_scanned = self._sec(sec, "total_devices")
        os_labels, os_counts = self._os_chart_data(sec)
        flagged = self._flagged_devices(sec)
        trends = self._trend_payload()
        fetch_status = data.get("_fetch_status", {})

        pol_groups = self._build_hierarchy(data.get("policies", []))
        mcp_groups = self._build_hierarchy(data.get("macos_profiles", []))
        icp_groups = self._build_hierarchy(data.get("ios_profiles", []))
        scr_raw = data.get("scripts", [])
        scr_enriched = [
            dict(item, category={"name": item.get("categoryName", "No Category")})
            if isinstance(item, dict) else item
            for item in (scr_raw if isinstance(scr_raw, list) else [])
        ]
        scr_groups = self._build_hierarchy(scr_enriched)
        pkg_groups = self._build_hierarchy(data.get("packages", []))
        smart_group_names = self._list_names(data.get("smart_groups", []))

        site_names = self._list_names(data.get("sites", []))
        bldg_names = self._list_names(data.get("buildings", []))
        dept_names = self._list_names(data.get("departments", []))
        cat_names = self._list_names(data.get("categories", []))
        ade_names = self._list_names(data.get("device_enrollments", []))

        pol_count = len(data.get("policies", []))
        mcp_count = len(data.get("macos_profiles", []))
        icp_count = len(data.get("ios_profiles", []))
        scr_count = len(data.get("scripts", []))
        pkg_count = len(data.get("packages", []))
        sg_count = len(data.get("smart_groups", []))

        health_cls = self._health_badge_class(health_status)
        alert_cls = "badge-ok" if active_alerts in ("None", "N/A", "0") else "badge-err"
        report_date = datetime.now().strftime("%A %d %B %Y, %H:%M")

        def feat(label: str, value: str) -> str:
            value_lc = str(value).lower()
            enabled = "enabled" in value_lc or value_lc in ("yes", "true", "active")
            return (
                f'<span class="feat-pill {"feat-on" if enabled else "feat-off"}">'
                f'{self._html_text(label)} &mdash; {self._html_text("enabled" if enabled else "disabled")}'
                "</span>"
            )

        ade_preview = (
            ", ".join(ade_names[:3]) + ("..." if len(ade_names) > 3 else "")
            if ade_names else ""
        )
        feature_pills = (
            feat("MDM Auto Renew (Computers)", mdm_renew_comp)
            + feat("MDM Auto Renew (Mobile)", mdm_renew_md)
            + feat("SSO (SAML)", sso)
            + feat("Jamf SSO", sso_jamf)
            + feat(
                "Patch Management",
                "enabled" if patch_titles not in ("0", "N/A") else "disabled",
            )
        )

        css = self._css()
        js = self._js()
        org_name = (self._config.get("branding", "org_name") or "").strip()
        brand_label = (
            f"{org_name} \u2014 Jamf Pro Reporting Snapshot"
            if org_name
            else "Jamf Pro Reporting Snapshot"
        )
        page_title = (
            f"{org_name} \u2014 Jamf Pro Report \u2014 {instance_url}"
            if org_name
            else f"Jamf Pro Report \u2014 {instance_url}"
        )
        logo_html = self._logo_html()

        return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{page_title}</title>
<style>{css}</style>
</head>
<body>

<div class="topbar">
  <div class="topbar-brand">{logo_html}{brand_label}</div>
  <div style="display:flex;align-items:center;gap:16px;flex-wrap:wrap">
    <div class="topbar-meta">
      <strong>{self._html_text(instance_url, "N/A")}</strong><br>
      Version {self._html_text(jamf_version, "N/A")} &nbsp;&bull;&nbsp;
      Generated {self._html_text(report_date)}
      &nbsp;&bull;&nbsp; Check-in {self._html_text(checkin_freq, "N/A")}
    </div>
    <button class="dark-toggle" id="darkToggle">Dark mode</button>
  </div>
</div>

<div class="page">

  <div class="section-block">
    <div class="section-block-title">Overall Server Health</div>
    <div class="section-block-subtitle">
      Instance status, enrollment posture, organization footprint, and high-level services.
    </div>

    <div class="health-strip" style="margin-top:16px">
      <div class="health-item">
        <span class="health-label">Health</span>
        <span class="badge {health_cls}">{self._html_text(health_status)}</span>
      </div>
      <div class="health-item">
        <span class="health-label">Active Alerts</span>
        <span class="badge {alert_cls}">{self._html_text(active_alerts)}</span>
      </div>
      <div class="health-item">
        <span class="health-label">DEP Token Expires</span>
        <span class="badge badge-dim">{self._html_text(dep_token_exp)}</span>
      </div>
      <div class="health-item">
        <span class="health-label">Built-in CA Expires</span>
        <span class="badge badge-dim">{self._html_text(ca_expires)}</span>
      </div>
      <div class="health-item">
        <span class="health-label">DEP Sync</span>
        <span class="badge badge-dim">{self._html_text(ade_sync)}</span>
      </div>
    </div>

    <div class="section-title">Instance Summary</div>
    <div class="grid grid-4">
      {self._render_stat_card("Jamf Pro Version", jamf_version)}
      {self._render_stat_card("Health Status", health_status)}
      {self._render_stat_card("Active Alerts", active_alerts)}
      {self._render_stat_card("Check-In Frequency", checkin_freq)}
    </div>

    <div class="section-title">Enrollment &amp; Configuration</div>
    <div class="grid grid-5">
      {self._render_stat_card("ADE Instances", ade_instances, ade_preview)}
      {self._render_stat_card("VPP Locations", vpp_locations)}
      {self._render_stat_card("LDAP / IdP Servers", ldap_servers)}
      {self._render_stat_card("Computer Prestages", comp_prestages)}
      {self._render_stat_card("Mobile Prestages", md_prestages)}
    </div>
    <div class="grid grid-4" style="margin-top:14px">
      {self._render_stat_card("Webhooks", webhooks)}
      {self._render_stat_card("JCDS Files", jcds_files)}
      {self._render_stat_card("Patch Titles", patch_titles)}
      {self._render_stat_card("App Installers", app_installers)}
    </div>

    <div class="section-title">Organisation</div>
    <div class="grid grid-4">
      {self._render_org_table("Sites", site_names)}
      {self._render_org_table("Buildings", bldg_names)}
      {self._render_org_table("Departments", dept_names)}
      {self._render_org_table("Categories", cat_names)}
    </div>

    <div class="section-title">Enabled Features</div>
    <div class="card card-sm">{feature_pills}</div>

    {self._render_source_status(fetch_status)}
    {self._render_overview_table(ov, console_url)}
    {self._render_quick_links(console_url)}
  </div>

  <div class="section-block">
    <div class="section-block-title">macOS Fleet</div>
    <div class="section-block-subtitle">
      Computer inventory, security posture, and deployment coverage for macOS endpoints.
    </div>

    <div class="section-title">Computer Inventory</div>
    <div class="grid grid-6">
      {self._render_stat_card("Managed Computers", managed_computers, "", f"{console_url}/computers.html" if console_url else "", "Open in Jamf")}
      {self._render_stat_card("Unmanaged Computers", unmanaged_computers)}
      {self._render_stat_card("Policies", str(pol_count), "", f"{console_url}/policies.html" if console_url else "", "Open in Jamf")}
      {self._render_stat_card("macOS Profiles", str(mcp_count), "", f"{console_url}/OSXConfigurationProfiles.html" if console_url else "", "Open in Jamf")}
      {self._render_stat_card("Packages", str(pkg_count), "", f"{console_url}/view/settings/computer-management/packages" if console_url else "", "Open in Jamf")}
      {self._render_stat_card("Smart Groups", str(sg_count), "", f"{console_url}/smartComputerGroups.html" if console_url else "", "Open in Jamf")}
    </div>
    <div class="grid grid-3" style="margin-top:14px">
      {self._render_stat_card("Scripts", str(scr_count), "", f"{console_url}/view/settings/computer-management/scripts" if console_url else "", "Open in Jamf")}
      {self._render_stat_card("Patch Titles", patch_titles)}
      {self._render_stat_card("Security Rows Scanned", total_scanned, f"{len(flagged)} flagged devices")}
    </div>

    <div class="section-title">Security Posture &amp; OS Distribution</div>
    <div class="grid grid-2">
      <div class="chart-card">
        <div class="chart-title">Security Feature Compliance
          <span class="badge badge-dim" style="margin-left:6px">{self._html_text(total_scanned)} devices scanned</span>
          <span class="badge badge-warn" style="margin-left:6px">{len(flagged)} flagged</span>
        </div>
        {self._render_sec_bar("FileVault", fv_pct, "fill-fv")}
        {self._render_sec_bar("Gatekeeper", gk_pct, "fill-gk")}
        {self._render_sec_bar("SIP", sip_pct, "fill-sip")}
        {self._render_sec_bar("Firewall", fw_pct, "fill-fw")}
      </div>
      {self._render_os_distribution_card(os_labels, os_counts)}
    </div>

    {self._render_trends_section(trends)}
    {self._render_flagged_table(flagged, console_url)}
    {self._render_hierarchy_tabs(pol_groups, mcp_groups, icp_groups, scr_groups, pkg_groups, smart_group_names)}
  </div>

  <div class="section-block">
    <div class="section-block-title">Mobile Devices</div>
    <div class="section-block-subtitle">
      Mobile inventory, supervision posture, and iOS configuration profile coverage.
    </div>

    <div class="section-title">Mobile Inventory</div>
    <div class="grid grid-6">
      {self._render_stat_card("Total Mobile Devices", str(mobile_summary["total"]))}
      {self._render_stat_card("Managed", str(mobile_summary["managed"]))}
      {self._render_stat_card("Supervised", str(mobile_summary["supervised"]))}
      {self._render_stat_card("Shared iPad", str(mobile_summary["shared_ipad"]))}
      {self._render_stat_card("Assigned Users", str(mobile_summary["assigned"]))}
      {self._render_stat_card("iOS Profiles", str(icp_count), "", f"{console_url}/mobileDeviceConfigurationProfiles.html" if console_url else "", "Open in Jamf")}
    </div>
    <div class="grid grid-4" style="margin-top:14px">
      {self._render_stat_card("Activation Lock Enabled", str(mobile_summary["activation_lock"]))}
      {self._render_stat_card("Passcode Compliant", str(mobile_summary["passcode_compliant"]))}
      {self._render_stat_card("Inventory Age Known", str(mobile_summary["inventory_age_known"]))}
      {self._render_stat_card(f"Older Than {stale_days} Days", str(mobile_stale))}
    </div>

    <div class="section-title">Mobile Distribution</div>
    <div class="grid grid-3">
      {self._render_counter_table("Device Families", "Device Family", mobile_summary["families"])}
      {self._render_counter_table("OS Versions", "OS Version", mobile_summary["os_versions"])}
      {self._render_counter_table("Models", "Model", mobile_summary["models"])}
    </div>

    {self._render_mobile_inventory_table(mobile_rows, stale_days)}
  </div>

  <div class="footer">
    Generated by jamf-reports-community.
    HTML report design based on
    <a href="https://github.com/DevliegereM/" target="_blank" rel="noopener noreferrer">Github.com/DevliegereM</a>
  </div>

</div>

<script>{js}</script>
</body>
</html>"""


def cmd_html(
    config: Config,
    out_file: Optional[str],
    no_open: bool = False,
) -> None:
    """Generate a self-contained HTML status report from jamf-cli data.

    Requires jamf-cli to be installed and authenticated. Falls back to cached
    data when live calls fail, subject to the jamf_cli.use_cached_data setting.

    Args:
        config: Loaded Config instance.
        out_file: Destination file path. Defaults to the output_dir from config.
        no_open: When True, do not auto-open the file after writing.
    """
    if not _jamf_cli_enabled(config):
        raise SystemExit("Error: html requires jamf_cli.enabled: true in config.yaml.")
    bridge = _build_jamf_cli_bridge(config, save_output=True)
    if not bridge.is_available():
        print(
            "  [warn] jamf-cli not found — will attempt to use cached data only.\n"
            "         Install via: brew install Jamf-Concepts/tap/jamf-cli"
        )

    output_cfg = config.output
    timestamp_outputs = output_cfg.get("timestamp_outputs", True) is not False
    run_stamp = _file_stamp()
    if out_file:
        out_path = _timestamped_output_path(
            Path(out_file).expanduser(),
            run_stamp,
            timestamp_outputs,
        )
    else:
        out_dir = config.resolve_path("output", "output_dir") or Path("Generated Reports")
        out_path = _timestamped_output_path(
            out_dir / "JamfReport.html",
            run_stamp,
            timestamp_outputs,
        )

    report = HtmlReport(config, bridge, out_path, no_open=no_open)
    report.generate()


def _collect_snapshots(
    config: Config,
    csv_path: Optional[str] = None,
    historical_csv_dir: Optional[str] = None,
) -> tuple[int, bool]:
    """Collect live jamf-cli snapshots and optionally archive a CSV snapshot."""
    print("--- Collect snapshots ---")
    print(f"  config base dir: {config.base_dir}")

    collected = 0
    jamf_cli_enabled = _jamf_cli_enabled(config)
    jamf_cli_dir = config.resolve_path("jamf_cli", "data_dir", default="jamf-cli-data")
    jamf_cli_profile = str(config.jamf_cli.get("profile", "") or "").strip()
    protect_enabled = config.get("protect", "enabled", default=False) is True
    platform_enabled = config.get("platform", "enabled", default=False) is True
    platform_benchmarks = _platform_benchmark_titles(config)
    live_overview_allowed = jamf_cli_enabled and (
        config.jamf_cli.get("allow_live_overview", True) is True
    )
    bridge: Optional[JamfCLIBridge] = None
    if jamf_cli_enabled:
        bridge = _build_jamf_cli_bridge(config, save_output=True, use_cached_data=False)
        print(f"  jamf-cli data dir: {bridge._data_dir}")
        if jamf_cli_profile:
            print(f"  jamf-cli profile: {jamf_cli_profile}")
        if protect_enabled:
            print("  protect reporting: enabled (experimental)")
        if platform_enabled:
            print("  platform reporting: enabled (preview)")
            for bench in platform_benchmarks:
                print(f"  platform benchmark: {bench}")
    if jamf_cli_enabled and bridge is not None and bridge.is_available():

        stale_days = int(config.thresholds.get("stale_device_days", 30))
        commands = []
        if live_overview_allowed:
            commands.append(("Fleet Overview", bridge.overview))
        else:
            print("  [skip] Fleet Overview: live overview disabled in config")
        commands.extend(
            [
                ("Inventory Summary", bridge.inventory_summary),
                ("Mobile Inventory", bridge.mobile_device_inventory_details),
                ("Mobile Device List", bridge.mobile_devices_list),
                ("Mobile Config Profiles", bridge.ios_profiles_list),
                ("Security Posture", bridge.security_report),
                ("Device Compliance", lambda: bridge.device_compliance(stale_days)),
                ("EA Coverage", lambda: bridge.ea_results_report(include_all=True)),
                ("EA Definitions", bridge.computer_extension_attributes),
                ("Software Installs", bridge.software_installs),
                ("Patch Compliance", bridge.patch_status),
                ("Patch Failures", bridge.patch_device_failures),
                ("Policy Health", bridge.policy_status),
                ("Profile Status", bridge.profile_status),
            ]
        )
        if protect_enabled:
            commands.extend(
                [
                    ("Protect Overview", bridge.protect_overview),
                    ("Protect Computers", bridge.protect_computers_list),
                    ("Protect Analytics", bridge.protect_analytics),
                    ("Protect Plans", bridge.protect_plans),
                ]
            )
        if platform_enabled:
            commands.extend(
                [
                    ("Platform Blueprints", bridge.blueprint_status),
                    ("Platform DDM Status", bridge.ddm_status),
                ]
            )
            if platform_benchmarks:
                for bench in platform_benchmarks:
                    bench_label = bench[:40]
                    commands.extend(
                        [
                            (
                                f"Compliance Rules: {bench_label}",
                                lambda b=bench: bridge.compliance_rules(b),
                            ),
                            (
                                f"Compliance Devices: {bench_label}",
                                lambda b=bench: bridge.compliance_devices(b),
                            ),
                        ]
                    )
            else:
                print(
                    "  [skip] Platform Compliance: platform.compliance_benchmarks is empty"
                )
        for label, command in commands:
            try:
                command()
                collected += 1
                print(f"  [ok] {label}")
            except RuntimeError as exc:
                print(f"  [skip] {label}: {exc}")
    elif jamf_cli_enabled:
        print("  jamf-cli: not found; skipping live snapshot collection.")
    else:
        print("  jamf-cli disabled in config; skipping live snapshot collection.")

    archived = False
    csv_path_obj = _resolve_cli_input_path(csv_path, config)
    selected_family_name = _family_for_csv_path(config, csv_path_obj) if csv_path_obj else None
    hist_dir_obj = _default_historical_dir(config, selected_family_name, historical_csv_dir)
    if config.get("charts", "archive_current_csv") is not False:
        if csv_path_obj and hist_dir_obj:
            print(f"  historical CSV dir: {hist_dir_obj}")
            archived_path, created = _archive_csv_snapshot(str(csv_path_obj), str(hist_dir_obj))
            if archived_path:
                archived = True
                if created:
                    print(f"  [ok] Archived CSV snapshot: {archived_path}")
                else:
                    print(f"  [ok] Reusing existing identical CSV snapshot: {archived_path}")
        elif csv_path_obj is None:
            for family_name in REPORT_FAMILY_NAMES:
                family = _report_family_config(config, family_name)
                if family.get("enabled") is not True:
                    continue
                latest_path, note = _latest_report_family_file(config, family_name)
                hist_dir = _report_family_historical_dir(config, family_name)
                print(f"  {family_name}: {note}")
                if latest_path is None or hist_dir is None:
                    continue
                print(f"    historical dir: {hist_dir}")
                archived_path, created = _archive_csv_snapshot(str(latest_path), str(hist_dir))
                if archived_path:
                    archived = True
                    if created:
                        print(f"    [ok] Archived: {archived_path}")
                    else:
                        print(f"    [ok] Reusing identical snapshot: {archived_path}")

    return collected, archived


def cmd_collect(
    config: Config,
    csv_path: Optional[str] = None,
    historical_csv_dir: Optional[str] = None,
) -> None:
    """Collect live jamf-cli snapshots and optionally archive a CSV snapshot."""
    collected, archived = _collect_snapshots(config, csv_path, historical_csv_dir)
    if collected == 0 and not archived:
        if not _jamf_cli_enabled(config):
            raise SystemExit(
                "Error: No snapshots collected. jamf-cli is disabled in config, so"
                " pass --csv plus --historical-csv-dir or enable report_families to"
                " archive CSV history."
            )
        raise SystemExit(
            "Error: No snapshots collected. Authenticate jamf-cli for live data or"
            " pass --csv plus --historical-csv-dir or enable report_families to"
            " archive CSV history."
        )


def _default_inventory_csv_out_file(config: Config) -> Path:
    """Return the default inventory CSV destination path before timestamping."""
    out_dir = config.resolve_path("output", "output_dir", default="Generated Reports")
    if out_dir is None:
        out_dir = Path("Generated Reports")
    out_dir.mkdir(parents=True, exist_ok=True)
    jamf_cli_profile = str(config.jamf_cli.get("profile", "") or "").strip()
    if jamf_cli_profile:
        return out_dir / f"jamf_inventory_{_filename_component(jamf_cli_profile)}.csv"
    return out_dir / "jamf_inventory.csv"


def _automation_inventory_out_file(config: Config) -> Path:
    """Return the default automation inventory CSV destination path."""
    out_dir = config.resolve_path("output", "output_dir", default="Generated Reports")
    if out_dir is None:
        out_dir = Path("Generated Reports")
    out_dir.mkdir(parents=True, exist_ok=True)
    profile_name = str(config.jamf_cli.get("profile", "") or "").strip()
    slug_source = profile_name or config.path.stem or "default"
    stem = f"automation_inventory_{_filename_component(slug_source)}.csv"
    return out_dir / stem


def cmd_inventory_csv(config: Config, out_file: Optional[str]) -> Path:
    """Export a wide computer inventory CSV from jamf-cli inventory and EA data."""
    if not _jamf_cli_enabled(config):
        raise SystemExit("Error: inventory-csv requires jamf_cli.enabled: true in config.yaml.")
    output_cfg = config.output
    out_dir = config.resolve_path("output", "output_dir", default="Generated Reports")
    if out_dir is None:
        out_dir = Path("Generated Reports")
    out_dir.mkdir(parents=True, exist_ok=True)
    timestamp_outputs = output_cfg.get("timestamp_outputs", True) is not False
    archive_enabled = output_cfg.get("archive_enabled", True) is not False
    keep_latest_runs = _to_int(output_cfg.get("keep_latest_runs", 10), 10)

    run_stamp = _file_stamp()
    if out_file:
        out_path = _timestamped_output_path(
            Path(out_file).expanduser(), run_stamp, timestamp_outputs
        )
    else:
        out_path = _timestamped_output_path(
            _default_inventory_csv_out_file(config),
            run_stamp,
            timestamp_outputs,
        )
    out_path.parent.mkdir(parents=True, exist_ok=True)

    print(f"Output: {out_path}")
    print(f"  config base dir: {config.base_dir}")

    bridge = _build_jamf_cli_bridge(config, save_output=False, use_cached_data=False)
    if not bridge.is_available():
        raise SystemExit("Error: jamf-cli not found. Install it or set JAMFCLI_PATH.")

    computers_raw = bridge.computers_list()
    computers = computers_raw if isinstance(computers_raw, list) else []
    if not computers:
        raise SystemExit("Error: jamf-cli returned no computers.")

    counts_by_name: dict[str, int] = {}
    base_rows: list[dict[str, Any]] = []
    for item in computers:
        if not isinstance(item, dict):
            continue
        row = _inventory_export_row(item)
        name = row["Computer Name"]
        counts_by_name[name] = counts_by_name.get(name, 0) + 1
        base_rows.append(row)

    duplicate_names = sorted(
        name for name, count in counts_by_name.items() if count > 1 and name
    )
    if duplicate_names:
        names_preview = ", ".join(duplicate_names[:5])
        print(
            "  [warn] Duplicate computer names detected;"
            " matching will prefer Jamf Pro ID, serial number, UDID, and"
            f" management ID before name: {names_preview}"
        )

    row_index = _inventory_build_row_index(base_rows)
    detail_enriched, detail_failures, detail_unresolved = (
        _enrich_inventory_rows_with_security_details(
            bridge,
            computers,
            row_index,
        )
    )
    ea_columns: set[str] = set()
    unmatched_ea_rows = 0
    ea_raw = bridge.ea_results(include_all=True)
    ea_rows = ea_raw if isinstance(ea_raw, list) else []
    for item in ea_rows:
        if not isinstance(item, dict):
            continue
        device = str(item.get("device", "") or "").strip()
        ea_name = str(item.get("ea_name", "") or "").strip()
        if not device or not ea_name:
            continue
        row = _inventory_resolve_row(row_index, _inventory_ea_lookup_values(item))
        if row is None:
            unmatched_ea_rows += 1
            continue
        row[ea_name] = str(item.get("value", "") or "").strip()
        ea_columns.add(ea_name)

    ordered_columns = INVENTORY_EXPORT_COLUMNS + sorted(
        column for column in ea_columns if column not in INVENTORY_EXPORT_COLUMNS
    )
    df = pd.DataFrame(base_rows, columns=ordered_columns).fillna("")
    df.to_csv(out_path, index=False, encoding="utf-8-sig")

    print(f"  [ok] Exported {len(df)} computers")
    print(
        "  [ok] Added"
        f" {len(INVENTORY_SECURITY_DETAIL_COLUMNS)} generic security detail column(s)"
    )
    if detail_enriched:
        print(f"  [ok] Enriched security details for {detail_enriched} computer(s)")
    if detail_failures:
        print(
            "  [warn] Could not enrich device security details for"
            f" {detail_failures} computer(s)"
        )
    if detail_unresolved:
        print(
            "  [warn] Could not uniquely match security details for"
            f" {detail_unresolved} computer(s)"
        )
    print(f"  [ok] Included {len(ea_columns)} extension attribute columns")
    if unmatched_ea_rows:
        print(
            "  [warn] Could not uniquely match"
            f" {unmatched_ea_rows} EA row(s) to inventory rows"
        )
    print(
        "  Note: this export is built from jamf-cli computers list plus"
        " per-device security details and jamf-cli report ea-results."
    )
    if archive_enabled:
        archive_dir = config.resolve_path("output", "archive_dir")
        if archive_dir is None:
            archive_dir = out_path.parent / "archive"
        family_base = _strip_timestamp_suffix(out_path.stem)
        archived_paths = _archive_old_output_runs(
            out_path.parent,
            family_base,
            {".csv"},
            keep_latest_runs,
            archive_dir,
        )
        if archived_paths:
            print(f"  Archived {len(archived_paths)} older inventory export(s) to {archive_dir}")
    return out_path


def _resolve_time_choice(
    value: Optional[str],
    default: str = DEFAULT_AUTOMATION_TIME_OF_DAY,
) -> tuple[str, int, int]:
    """Return a validated HH:MM string and parsed hour/minute values."""
    candidate = value
    while True:
        raw = candidate if candidate is not None else _prompt_text(
            "Run time (24-hour HH:MM)", default,
        )
        try:
            hour, minute = _parse_time_of_day(raw)
            return raw, hour, minute
        except ValueError as exc:
            if candidate is not None:
                raise SystemExit(f"Error: {exc}") from None
            print(f"  {exc}")


def _resolve_weekday_choice(
    value: Optional[str],
    default: str = "Monday",
) -> tuple[int, str]:
    """Return a validated launchd weekday value and display name."""
    candidate = value
    while True:
        raw = candidate if candidate is not None else _prompt_text("Weekday", default)
        try:
            return _parse_weekday(raw)
        except ValueError as exc:
            if candidate is not None:
                raise SystemExit(f"Error: {exc}") from None
            print(f"  {exc}")


def _resolve_day_of_month_choice(value: Optional[int], default: int = 1) -> int:
    """Return a validated monthly day-of-month value."""
    candidate = value
    while True:
        raw: Any = candidate if candidate is not None else _prompt_text("Day of month (1-28)", str(default))
        try:
            return _parse_day_of_month(raw)
        except ValueError as exc:
            if candidate is not None:
                raise SystemExit(f"Error: {exc}") from None
            print(f"  {exc}")


def _resolve_csv_freshness_days(
    value: Any,
    default: int = DEFAULT_CSV_FRESHNESS_DAYS,
) -> int:
    """Return a validated positive CSV freshness window."""
    candidate = default if value is None else value
    while True:
        raw: Any = candidate if candidate is not None else _prompt_text(
            "CSV freshness window in days", str(default),
        )
        try:
            days = int(raw)
        except (TypeError, ValueError):
            if value is not None:
                raise SystemExit("Error: CSV freshness window must be a positive integer.") from None
            print("  CSV freshness window must be a positive integer.")
            candidate = None
            continue
        if days < 1:
            if value is not None:
                raise SystemExit("Error: CSV freshness window must be at least 1 day.") from None
            print("  CSV freshness window must be at least 1 day.")
            candidate = None
            continue
        return days


def _resolve_workspace_dir(value: Optional[str], config: Config) -> Path:
    """Return the automation workspace directory."""
    if value:
        return _expand_setup_path(value, config.base_dir)
    if not sys.stdin.isatty():
        return config.base_dir.resolve()
    return _expand_setup_path(
        _prompt_text("Workspace directory for automation files", str(config.base_dir)),
        config.base_dir,
    )


def _resolve_historical_csv_dir(
    value: Optional[str],
    workspace_dir: Path,
    mode: str,
) -> Optional[Path]:
    """Return the historical CSV directory to use for automation, if any."""
    if value:
        return _expand_setup_path(value, workspace_dir)
    if mode not in {"snapshot-only", "jamf-cli-full", "csv-assisted"}:
        return None
    if not sys.stdin.isatty():
        return (workspace_dir / "snapshots").resolve()
    enabled = _prompt_yes_no(
        "Create/use a historical CSV snapshots folder for trend reporting?",
        True,
    )
    if not enabled:
        return None
    return _expand_setup_path(
        _prompt_text("Historical CSV directory", str(workspace_dir / "snapshots")),
        workspace_dir,
    )


def _resolve_csv_inbox_settings(
    csv_inbox_dir: Optional[str],
    csv_freshness_days: Any,
    workspace_dir: Path,
    mode: str,
) -> tuple[Optional[Path], int]:
    """Return CSV inbox configuration for automation."""
    freshness_days = _resolve_csv_freshness_days(csv_freshness_days)
    if csv_inbox_dir:
        return _expand_setup_path(csv_inbox_dir, workspace_dir), freshness_days
    if mode not in {"snapshot-only", "csv-assisted"}:
        return None, freshness_days
    if not sys.stdin.isatty():
        if mode == "csv-assisted":
            return (workspace_dir / "csv-inbox").resolve(), freshness_days
        return None, freshness_days
    enabled = _prompt_yes_no(
        "Use a CSV inbox folder when Jamf emails or exports are available?",
        mode == "csv-assisted",
    )
    if not enabled:
        return None, freshness_days
    if csv_freshness_days is None:
        freshness_days = _resolve_csv_freshness_days(
            _prompt_text("CSV freshness window in days", str(freshness_days)),
        )
    inbox_dir = _expand_setup_path(
        _prompt_text("CSV inbox directory", str(workspace_dir / "csv-inbox")),
        workspace_dir,
    )
    return inbox_dir, freshness_days


def _launchagent_environment() -> dict[str, str]:
    """Return environment variables to persist into the LaunchAgent plist."""
    env = {
        "HOME": str(Path.home()),
        "PATH": DEFAULT_LAUNCHD_PATH,
        "PYTHONUNBUFFERED": "1",
    }
    xdg_config_home = os.environ.get("XDG_CONFIG_HOME", "").strip()
    if xdg_config_home:
        env["XDG_CONFIG_HOME"] = xdg_config_home
    jamf_cli_binary = _find_jamf_cli_binary()
    if jamf_cli_binary:
        env["JAMFCLI_PATH"] = jamf_cli_binary
    return env


def _launchagent_program_arguments(
    config_path: Path,
    mode: str,
    status_file: Path,
    historical_csv_dir: Optional[Path],
    csv_inbox_dir: Optional[Path],
    csv_freshness_days: int,
    notify_url: Optional[str],
) -> list[str]:
    """Return the ProgramArguments array for the generated LaunchAgent plist."""
    args = [
        str(Path(sys.executable).resolve()),
        str(Path(__file__).resolve()),
        "launchagent-run",
        "--config",
        str(config_path),
        "--mode",
        mode,
        "--status-file",
        str(status_file),
    ]
    if historical_csv_dir:
        args.extend(["--historical-csv-dir", str(historical_csv_dir)])
    if csv_inbox_dir:
        args.extend(
            [
                "--csv-inbox-dir",
                str(csv_inbox_dir),
                "--csv-freshness-days",
                str(csv_freshness_days),
            ]
        )
    if notify_url:
        args.extend(["--notify", notify_url])
    return args


def _write_launchagent_plist(
    plist_path: Path,
    label: str,
    program_arguments: list[str],
    working_directory: Path,
    schedule_items: list[dict[str, int]],
    stdout_path: Path,
    stderr_path: Path,
) -> None:
    """Write a LaunchAgent plist for the supplied automation plan."""
    payload = {
        "Label": label,
        "ProgramArguments": program_arguments,
        "WorkingDirectory": str(working_directory),
        "EnvironmentVariables": _launchagent_environment(),
        "StartCalendarInterval": schedule_items,
        "RunAtLoad": False,
        "StandardOutPath": str(stdout_path),
        "StandardErrorPath": str(stderr_path),
    }
    plist_path.parent.mkdir(parents=True, exist_ok=True)
    with open(plist_path, "wb") as fh:
        plistlib.dump(payload, fh, sort_keys=True)


def _load_launchagent(plist_path: Path, label: str, run_now: bool) -> str:
    """Bootstrap a LaunchAgent into the current GUI session."""
    target = f"gui/{os.getuid()}"
    label_target = f"{target}/{label}"
    subprocess.run(
        ["launchctl", "bootout", target, str(plist_path)],
        capture_output=True,
        text=True,
        check=False,
    )
    try:
        subprocess.run(
            ["launchctl", "bootstrap", target, str(plist_path)],
            capture_output=True,
            text=True,
            check=True,
        )
    except subprocess.CalledProcessError as exc:
        detail = (exc.stderr or exc.stdout).strip()
        raise SystemExit(f"Error: launchctl bootstrap failed: {detail}") from None
    subprocess.run(
        ["launchctl", "enable", label_target],
        capture_output=True,
        text=True,
        check=False,
    )
    if run_now:
        try:
            subprocess.run(
                ["launchctl", "kickstart", "-k", label_target],
                capture_output=True,
                text=True,
                check=True,
            )
        except subprocess.CalledProcessError as exc:
            detail = (exc.stderr or exc.stdout).strip()
            raise SystemExit(f"Error: launchctl kickstart failed: {detail}") from None
    return target


def cmd_launchagent_run(
    config: Config,
    mode: str,
    csv_inbox_dir: Optional[str],
    csv_freshness_days: int,
    historical_csv_dir: Optional[str],
    status_file: Optional[str],
    notify_url: Optional[str] = None,
) -> None:
    """Run a scheduled automation workflow from a generated LaunchAgent."""
    started_at = datetime.now(timezone.utc).isoformat()
    status: dict[str, Any] = {
        "config_path": str(config.path),
        "finished_at": None,
        "mode": mode,
        "profile": str(config.jamf_cli.get("profile", "") or "").strip(),
        "report_path": None,
        "selected_csv_family": None,
        "selected_csv_origin": None,
        "selected_csv": None,
        "started_at": started_at,
        "status_file": status_file,
        "success": False,
    }
    if historical_csv_dir:
        status["historical_csv_dir"] = str(Path(historical_csv_dir).expanduser())
    if csv_inbox_dir:
        status["csv_inbox_dir"] = str(Path(csv_inbox_dir).expanduser())

    try:
        selected_csv: Optional[Path] = None
        selected_family: Optional[str] = None
        selected_origin = ""
        if mode in {"snapshot-only", "csv-assisted"}:
            selected_csv, selected_family, selected_origin, selection_note = _select_automation_csv(
                config,
                csv_inbox_dir,
                csv_freshness_days,
            )
            print(selection_note)
            status["csv_selection_note"] = selection_note
            status["selected_csv_family"] = selected_family
            status["selected_csv_origin"] = selected_origin or None
            status["selected_csv"] = str(selected_csv) if selected_csv else None

        if mode == "snapshot-only":
            snapshot_csv = str(selected_csv) if selected_origin == "csv_inbox" and selected_csv else None
            collected, archived = _collect_snapshots(
                config,
                snapshot_csv,
                historical_csv_dir,
            )
            if collected == 0 and not archived:
                raise SystemExit(
                    "Error: snapshot-only automation had nothing to collect."
                    " Authenticate jamf-cli or configure report_families or a CSV inbox."
                )
            status["collected_snapshots"] = collected
            status["archived_csv"] = archived
        elif mode == "jamf-cli-only":
            report_path = cmd_generate(config, None, None, historical_csv_dir, notify_url)
            status["report_path"] = str(report_path)
        elif mode == "jamf-cli-full":
            inventory_path = cmd_inventory_csv(config, str(_automation_inventory_out_file(config)))
            status["inventory_csv_path"] = str(inventory_path)
            _collect_snapshots(config, None, historical_csv_dir)
            report_path = cmd_generate(
                config,
                str(inventory_path),
                None,
                historical_csv_dir,
                notify_url,
            )
            status["report_path"] = str(report_path)
        elif mode == "csv-assisted":
            _collect_snapshots(config, None, historical_csv_dir)
            report_path = cmd_generate(
                config,
                str(selected_csv) if selected_csv else None,
                None,
                historical_csv_dir,
                notify_url,
            )
            status["report_path"] = str(report_path)
        else:
            raise SystemExit(f"Error: unsupported automation mode: {mode}")
    except SystemExit as exc:
        status["error"] = str(exc)
        status["finished_at"] = datetime.now(timezone.utc).isoformat()
        _write_status_file(status_file, status)
        raise
    except Exception as exc:
        status["error"] = str(exc)
        status["finished_at"] = datetime.now(timezone.utc).isoformat()
        _write_status_file(status_file, status)
        raise

    status["success"] = True
    status["finished_at"] = datetime.now(timezone.utc).isoformat()
    _write_status_file(status_file, status)


def cmd_launchagent_setup(
    config: Config,
    config_path_value: str,
    label: Optional[str],
    mode: Optional[str],
    schedule: Optional[str],
    time_of_day: Optional[str],
    weekday: Optional[str],
    day_of_month: Optional[int],
    workspace_dir: Optional[str],
    launchagents_dir: Optional[str],
    csv_inbox_dir: Optional[str],
    csv_freshness_days: Any,
    historical_csv_dir: Optional[str],
    notify_url: Optional[str],
    skip_load: bool,
    run_now: bool,
) -> None:
    """Interactively create and optionally load a LaunchAgent automation job."""
    config_path = _require_existing_config_path(config_path_value)
    mode_options = [
        (key, f"{key}: {desc}") for key, desc in AUTOMATION_MODE_DESCRIPTIONS.items()
    ]
    schedule_options = [
        (key, f"{key}: {desc}") for key, desc in AUTOMATION_SCHEDULE_DESCRIPTIONS.items()
    ]
    selected_mode = mode or _prompt_choice(
        "Automation workflow",
        mode_options,
        DEFAULT_AUTOMATION_MODE,
    )
    selected_schedule = schedule or _prompt_choice(
        "Schedule type",
        schedule_options,
        DEFAULT_AUTOMATION_SCHEDULE,
    )
    _, hour, minute = _resolve_time_choice(time_of_day)
    weekday_value: Optional[int] = None
    weekday_name: Optional[str] = None
    if selected_schedule == "weekly":
        weekday_value, weekday_name = _resolve_weekday_choice(weekday)
    monthly_day: Optional[int] = None
    if selected_schedule == "monthly":
        monthly_day = _resolve_day_of_month_choice(day_of_month)

    automation_root = _resolve_workspace_dir(workspace_dir, config)
    csv_history_dir = _resolve_historical_csv_dir(historical_csv_dir, automation_root, selected_mode)
    csv_inbox_path, freshness_days = _resolve_csv_inbox_settings(
        csv_inbox_dir,
        csv_freshness_days,
        automation_root,
        selected_mode,
    )
    launchagents_root = _expand_setup_path(
        launchagents_dir or str(Path.home() / "Library" / "LaunchAgents"),
        automation_root,
    )
    job_label = label or _default_launchagent_label(config)
    job_slug = _filename_component(job_label)
    automation_dir = automation_root / "automation"
    logs_dir = automation_dir / "logs"
    status_path = automation_dir / f"{job_slug}_status.json"
    plist_path = launchagents_root / f"{job_label}.plist"
    stdout_path = logs_dir / f"{job_slug}.out.log"
    stderr_path = logs_dir / f"{job_slug}.err.log"

    output_dir = config.resolve_path("output", "output_dir", default="Generated Reports")
    jamf_cli_dir = config.resolve_path("jamf_cli", "data_dir", default="jamf-cli-data")
    for path in [automation_dir, logs_dir, output_dir, jamf_cli_dir, csv_history_dir, csv_inbox_path]:
        if path is not None:
            path.mkdir(parents=True, exist_ok=True)

    program_arguments = _launchagent_program_arguments(
        config_path,
        selected_mode,
        status_path,
        csv_history_dir,
        csv_inbox_path,
        freshness_days,
        notify_url,
    )
    schedule_items = _launchagent_schedule_items(
        selected_schedule,
        hour,
        minute,
        weekday_value,
        monthly_day,
    )
    _write_launchagent_plist(
        plist_path,
        job_label,
        program_arguments,
        config.base_dir,
        schedule_items,
        stdout_path,
        stderr_path,
    )

    print("\nLaunchAgent setup summary")
    print(f"  config: {config_path}")
    profile_name = str(config.jamf_cli.get("profile", "") or "").strip()
    print(f"  jamf-cli profile: {profile_name or 'default'}")
    if not profile_name:
        print(
            "  [warn] jamf_cli.profile is blank; this LaunchAgent will use jamf-cli's"
            " default profile."
        )
    print(f"  jamf-cli data dir: {jamf_cli_dir}")
    if not _find_jamf_cli_binary():
        print("  [warn] jamf-cli was not found during setup; runtime depends on PATH or JAMFCLI_PATH.")
    print(f"  workflow: {selected_mode}")
    print(f"  schedule: {_launchagent_schedule_summary(selected_schedule, hour, minute, weekday_name, monthly_day)}")
    print(f"  output dir: {output_dir}")
    if csv_history_dir:
        print(f"  historical CSV dir: {csv_history_dir}")
    if csv_inbox_path:
        print(f"  CSV inbox dir: {csv_inbox_path} (freshness: {freshness_days} day(s))")
    print(f"  plist: {plist_path}")
    print(f"  stdout log: {stdout_path}")
    print(f"  stderr log: {stderr_path}")
    print(f"  status file: {status_path}")
    print(f"  command: {shlex.join(program_arguments)}")
    isolation_guidance = _profile_isolation_guidance(config)
    if isolation_guidance:
        print("  profile isolation guidance:")
        for item in isolation_guidance:
            print(f"    - {item}")

    if skip_load:
        print("  LaunchAgent not loaded (--skip-load).")
        return

    target = _load_launchagent(plist_path, job_label, run_now)
    print(f"  Loaded into launchd target: {target}")
    if run_now:
        print("  Triggered one immediate run with launchctl kickstart.")


# ---------------------------------------------------------------------------
# Managed-state patching (requires jamf-cli v1.6.0+)
# ---------------------------------------------------------------------------


def _read_serials_file(path: str) -> list[str]:
    """Read serial numbers from a text file, one per line.

    Lines starting with '#' and blank lines are ignored.

    Args:
        path: Path to the serials file.

    Returns:
        List of stripped serial number strings.

    Raises:
        SystemExit: If the file cannot be opened.
    """
    serials: list[str] = []
    try:
        with open(path, encoding="utf-8") as fh:
            for line in fh:
                stripped = line.strip()
                if stripped and not stripped.startswith("#"):
                    serials.append(stripped)
    except OSError as exc:
        raise SystemExit(f"Error reading serials file {path}: {exc}") from exc
    return serials


def cmd_patch_managed(
    config: Config,
    managed_value: bool,
    dry_run: bool = False,
    serials_file: Optional[str] = None,
) -> None:
    """Set managed state on computers via jamf-cli computers-inventory patch.

    Requires jamf-cli v1.6.0+. Without --serials-file, queries device-compliance
    and patches all devices currently in the opposite managed state. With
    --serials-file, patches every listed serial to the target state.

    Args:
        config: Loaded Config object.
        managed_value: Target managed state (True = managed, False = unmanaged).
        dry_run: If True, print what would change without making API calls.
        serials_file: Optional path to a file with one serial number per line.
    """
    if not _jamf_cli_enabled(config):
        raise SystemExit("Error: patch-managed requires jamf_cli.enabled: true in config.yaml.")
    bridge = _build_jamf_cli_bridge(config, save_output=False, use_cached_data=False)

    if not bridge.is_available():
        raise SystemExit("Error: jamf-cli not found. Install it or set JAMFCLI_PATH.")

    target_str = "true" if managed_value else "false"
    print(f"Target managed state: {target_str}")

    if serials_file:
        serials = _read_serials_file(serials_file)
        if not serials:
            raise SystemExit(f"Error: no serial numbers found in {serials_file}")
        print(f"  {len(serials)} serial(s) from {serials_file}")
    else:
        stale_days = int(config.thresholds.get("stale_device_days", 30))
        raw = bridge.device_compliance(stale_days)
        rows = raw if isinstance(raw, list) else []
        if not rows:
            raise SystemExit("Error: device-compliance returned no data.")
        opposite = "unmanaged" if managed_value else "managed"
        serials = [
            str(item.get("serial", "") or "").strip()
            for item in rows
            if isinstance(item, dict)
            and _to_bool(item.get("managed")) is not managed_value
            and str(item.get("serial", "") or "").strip()
        ]
        print(f"  {len(serials)} currently-{opposite} device(s) from device-compliance")

    if not serials:
        print("  No devices to patch.")
        return

    if dry_run:
        print(f"\n[dry-run] Would set general.managed={target_str} on {len(serials)} device(s):")
        for serial in serials:
            print(f"  {serial}")
        return

    print(f"\nPatching {len(serials)} device(s)...")
    success = 0
    failed = 0
    for serial in serials:
        try:
            bridge.computers_inventory_patch(serial, {"general.managed": target_str})
            print(f"  [ok]   {serial}")
            success += 1
        except RuntimeError as exc:
            detail = str(exc)
            if "unknown command" in detail.lower() or "computers-inventory" in detail.lower():
                raise SystemExit(
                    "Error: 'computers-inventory patch' not available."
                    " Upgrade to jamf-cli v1.6.0+."
                ) from exc
            print(f"  [fail] {serial}: {exc}")
            failed += 1

    print(f"\nDone: {success} patched, {failed} failed.")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------


def main() -> None:
    """Parse arguments and dispatch to the appropriate command."""
    parser = argparse.ArgumentParser(
        description="Community Jamf Pro reporting tool.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Commands:\n"
            "  generate      Build the Excel report\n"
            "  html          Build a self-contained HTML status report for management\n"
            "  collect       Save jamf-cli snapshots and optional CSV history\n"
            "  inventory-csv Export a wide CSV from jamf-cli inventory plus EAs\n"
            "  workspace-init Create a per-profile reporting workspace skeleton\n"
            "  launchagent-setup Create a LaunchAgent for scheduled reporting\n"
            "  launchagent-run   Internal runner used by generated LaunchAgents\n"
            "  scaffold      Generate a starter config.yaml from a CSV\n"
            "  check         Verify jamf-cli auth and config\n"
            "  device        Print a device detail view from jamf-cli pro device\n"
            "  patch-managed Set managed state on computers (requires jamf-cli v1.6.0+)\n"
        ),
    )
    parser.add_argument(
        "command",
        choices=[
            "generate",
            "html",
            "collect",
            "inventory-csv",
            "workspace-init",
            "launchagent-setup",
            "launchagent-run",
            "scaffold",
            "check",
            "device",
            "patch-managed",
        ],
    )
    parser.add_argument("--config", default="config.yaml", help="Path to config.yaml")
    parser.add_argument("--csv", help="Path to Jamf Pro CSV export")
    parser.add_argument("--out-file", help="Output file path (generate or inventory-csv)")
    parser.add_argument("--out", default="config.yaml", help="Output config path (scaffold only)")
    parser.add_argument(
        "--historical-csv-dir",
        help="Directory of dated CSV snapshots for trend charts or collection",
    )
    parser.add_argument(
        "--notify",
        metavar="WEBHOOK_URL",
        help="Microsoft Teams incoming webhook URL; posts an Adaptive Card after generate",
    )
    parser.add_argument(
        "--interactive", "-i",
        action="store_true",
        help="Walk through each column mapping interactively after auto-matching (scaffold only)",
    )
    parser.add_argument(
        "--id",
        metavar="DEVICE_ID",
        help="Jamf Pro computer ID or serial number to look up (device command only)",
    )
    parser.add_argument(
        "--csv-extra",
        action="append",
        metavar="CSV_PATH",
        dest="csv_extra",
        help=(
            "Additional CSV export to merge with --csv (can be repeated). "
            "Useful for combining computers, mobile-device, and users exports. "
            "A 'CSV Source' column is added to identify which file each row came from."
        ),
    )
    parser.add_argument(
        "--managed",
        choices=["true", "false"],
        help="Target managed state: 'true' to manage, 'false' to unmanage (patch-managed only)",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Print what would be patched without making API calls (patch-managed only)",
    )
    parser.add_argument(
        "--serials-file",
        metavar="PATH",
        help=(
            "File with one serial number per line to patch"
            " (patch-managed only; default: auto-detect from device-compliance)"
        ),
    )
    parser.add_argument(
        "--mode",
        choices=sorted(AUTOMATION_MODE_DESCRIPTIONS),
        help="Automation workflow mode (launchagent commands only)",
    )
    parser.add_argument(
        "--profile",
        help="jamf-cli profile name (workspace-init only)",
    )
    parser.add_argument(
        "--schedule",
        choices=sorted(AUTOMATION_SCHEDULE_DESCRIPTIONS),
        help="Schedule preset for launchagent-setup",
    )
    parser.add_argument(
        "--time-of-day",
        help="Local run time in HH:MM 24-hour format (launchagent-setup only)",
    )
    parser.add_argument(
        "--weekday",
        help="Weekday name for weekly launchagent schedules (launchagent-setup only)",
    )
    parser.add_argument(
        "--day-of-month",
        type=int,
        help="Day of month 1-28 for monthly launchagent schedules",
    )
    parser.add_argument(
        "--label",
        help="LaunchAgent label override (launchagent-setup only)",
    )
    parser.add_argument(
        "--workspace-dir",
        help="Directory for automation logs, status, and helper folders",
    )
    parser.add_argument(
        "--launchagents-dir",
        help="LaunchAgents directory override (default: ~/Library/LaunchAgents)",
    )
    parser.add_argument(
        "--csv-inbox-dir",
        help="Folder containing emailed or exported Jamf CSVs for automation",
    )
    parser.add_argument(
        "--csv-freshness-days",
        type=int,
        help="Maximum CSV age in days before automation falls back",
    )
    parser.add_argument(
        "--status-file",
        help="Automation status JSON path (launchagent-run only)",
    )
    parser.add_argument(
        "--skip-load",
        action="store_true",
        help="Write the LaunchAgent plist without loading it",
    )
    parser.add_argument(
        "--run-now",
        action="store_true",
        help="Kickstart the LaunchAgent immediately after loading it",
    )
    parser.add_argument(
        "--seed-config",
        help="Seed config path for workspace-init; defaults to config.example.yaml when absent",
    )
    parser.add_argument(
        "--workspace-root",
        help="Parent directory under which workspace-init creates the profile workspace",
    )
    parser.add_argument(
        "--workspace-name",
        help="Directory name override for workspace-init",
    )
    parser.add_argument(
        "--overwrite-config",
        action="store_true",
        help="Allow workspace-init to replace an existing workspace config.yaml",
    )
    parser.add_argument(
        "--no-open",
        action="store_true",
        help="Do not auto-open the generated HTML file after writing (html command only)",
    )
    args = parser.parse_args()

    if args.command == "scaffold":
        if not args.csv:
            parser.error("scaffold requires --csv")
        cmd_scaffold(args.csv, args.out, interactive=args.interactive)
        return

    if args.command == "device":
        if not args.id:
            parser.error("device requires --id")
        config = Config(args.config)
        cmd_device(config, args.id)
        return

    if args.command == "workspace-init":
        cmd_workspace_init(
            args.seed_config,
            args.profile,
            args.workspace_root,
            args.workspace_name,
            args.overwrite_config,
        )
        return

    if args.command in {"launchagent-setup", "launchagent-run"}:
        _require_existing_config_path(args.config)
    config = Config(args.config)

    if args.command == "check":
        cmd_check(config, args.csv)
    elif args.command == "collect":
        cmd_collect(config, args.csv, args.historical_csv_dir)
    elif args.command == "inventory-csv":
        cmd_inventory_csv(config, args.out_file)
    elif args.command == "launchagent-setup":
        cmd_launchagent_setup(
            config,
            args.config,
            args.label,
            args.mode,
            args.schedule,
            args.time_of_day,
            args.weekday,
            args.day_of_month,
            args.workspace_dir,
            args.launchagents_dir,
            args.csv_inbox_dir,
            args.csv_freshness_days,
            args.historical_csv_dir,
            args.notify,
            args.skip_load,
            args.run_now,
        )
    elif args.command == "launchagent-run":
        if not args.mode:
            parser.error("launchagent-run requires --mode")
        cmd_launchagent_run(
            config,
            args.mode,
            args.csv_inbox_dir,
            _resolve_csv_freshness_days(args.csv_freshness_days),
            args.historical_csv_dir,
            args.status_file,
            args.notify,
        )
    elif args.command == "html":
        cmd_html(config, args.out_file, no_open=args.no_open)
    elif args.command == "generate":
        cmd_generate(
            config, args.csv, args.out_file, args.historical_csv_dir,
            args.notify, args.csv_extra,
        )
    elif args.command == "patch-managed":
        if not args.managed:
            parser.error("patch-managed requires --managed true|false")
        cmd_patch_managed(
            config,
            managed_value=args.managed == "true",
            dry_run=args.dry_run,
            serials_file=args.serials_file,
        )


if __name__ == "__main__":
    main()
